from __future__ import annotations

from dataclasses import dataclass
import hashlib
from html import escape
from io import BytesIO
from pathlib import Path
import re
from typing import Any, Callable

import numpy as np
import pandas as pd
import plotly.express as px
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt
import streamlit as st


st.set_page_config(
    page_title="국내 생산·포장 현황",
    layout="wide",
    initial_sidebar_state="expanded",
)


class DashboardConfigError(Exception):
    def __init__(self, messages: list[str]):
        super().__init__("\n".join(messages))
        self.messages = messages


@dataclass
class SourceFiles:
    request_file: Path
    packing_file: Path
    progress_file: Path | None = None
    inventory_file: Path | None = None
    daily_inventory_file: Path | None = None
    product_master_file: Path | None = None
    wip_file: Path | None = None


STATUS_ORDER = ["미착수", "진행중", "완료"]
UNIT_PACK = "PACK 기준"
UNIT_PCS = "PCS 기준"
UNIT_OPTIONS = [UNIT_PACK, UNIT_PCS]
DASHBOARD_TABS = ["제품 진도 현황", "일일 재고 대응", "생산코드 상세", "판매코드 상세"]
SIDEBAR_NAV_ITEMS = [
    ("제품 진도 현황", "product_progress"),
    ("일일 재고 대응", "daily_inventory"),
    ("생산코드 상세", "production_code"),
    ("판매코드 상세", "sales_code"),
]
SIDEBAR_NAV_KEY_TO_TAB = {key: tab for tab, key in SIDEBAR_NAV_ITEMS}
DAILY_INVENTORY_FILE_STANDARD = "클라렌사업본부 재고현황_YYMMDD.xlsx"
DAILY_INVENTORY_FILE_KEYWORDS = ["클라렌사업본부 재고현황", "재고현황_"]
SAMPLE_KEYWORDS = ["샘플"]
GROUP_ORDER = ["전체", "본품", "샘플", "PIA", "Clalen", "Toric", "1Day", "Color", "Monthly", "기타"]
FACTORY_GROUP_ORDER = ["전체", "A관", "C관", "S관"]
PERIOD_GROUP_ORDER = ["전체", "1-DAY", "FRP"]
FACTORY_GROUP_BY_CATEGORY = {
    "FRP_Sph": "A관",
    "FRP_Toric": "A관",
    "Si_FRP_Color_Sph": "A관",
    "Si_FRP_Sph": "A관",
    "Si_FRP_Toric": "A관",
    "1-Day_Color_Sph": "C관",
    "1-Day_Color_Toric": "C관",
    "1-Day_Sph": "S관",
    "Si_1-Day_Color_Sph": "S관",
    "Si_1-Day_Color_Toric": "S관",
    "Si_1-Day_Sph": "S관",
    "Si_1-Day_Toric": "S관",
}
DAILY_ITEM_STANDARD = {
    "S547": {"factory_group": "A관", "product_name": "Clalen O2O2 M_2P"},
    "S548": {"factory_group": "A관", "product_name": "Clalen O2O2 M_6P"},
    "S611": {"factory_group": "A관", "product_name": "Clalen O2O2 M_Natural Chocolat EX_2P"},
    "S129": {"factory_group": "C관", "product_name": "Iris SoulBrown_40팩"},
    "S172": {"factory_group": "C관", "product_name": "Iris Suzy Brown_30팩"},
    "S309": {"factory_group": "S관", "product_name": "Clalen O2O2 D_Micelia_Mute Brown_10P"},
    "S147": {"factory_group": "C관", "product_name": "Iris Toric Alicia Brown_30팩"},
    "S318": {"factory_group": "A관", "product_name": "Clalen O2O2 M_Micelia_Deep Black_45%_2P"},
    "S320": {"factory_group": "A관", "product_name": "Clalen O2O2 M_Micelia_Mute Brown_45%_2P"},
    "S524": {"factory_group": "S관", "product_name": "Clalen O2O2 D Toric_Ash Brown EX_30P"},
    "S154": {"factory_group": "C관", "product_name": "Iris Toric Halo Brown_30"},
    "S159": {"factory_group": "C관", "product_name": "Iris Toric Halo Gray_30"},
    "S162": {"factory_group": "C관", "product_name": "Iris BlueMoon_40팩"},
}
PRODUCTION_CODE_PACK_LABELS = ["1P", "2P", "5P", "6P", "10P", "30P", "40P", "80P", "90P"]
WIP_PROCESS_COLUMNS = ["검사접착", "누수규격검사"]
DATA_CACHE_VERSION = 37
REQUEST_DUE_MONTH = "2026-07"
REQUEST_DUE_MONTH_LABEL = "2026년 7월"
PRODUCTION_PROGRESS_DUE_MONTH = REQUEST_DUE_MONTH
PRODUCTION_PROGRESS_DUE_MONTH_LABEL = REQUEST_DUE_MONTH_LABEL
PACKING_RECEIPT_BASE_DATE_LABEL = "2026년 6월 24일"
DATA_BASIS_NOTE = (
    f"진도 기준: 생산지시물량 {REQUEST_DUE_MONTH_LABEL} 생산완료예상일 / "
    f"포장실적·용마입고량 기준: {PACKING_RECEIPT_BASE_DATE_LABEL}부터 / "
    "요청 대비 지시 수준은 3Q전체물량과 생산지시물량 비교"
)
MAIN_PRODUCT_FAMILY_ORDER = [
    "전체",
    "Clalen 1Day",
    "O2O2 1Day",
    "O2O2 D 컬러",
    "O2O2 D Micelia",
    "O2O2 D Toric",
    "O2O2 Monthly",
    "O2O2 M Micelia",
    "Clear",
    "PIA 1Day",
    "PIA Monthly",
    "Iris 컬러",
    "Iris Toric",
    "T38 Toric",
    "기타 Toric",
    "부자재/기타",
    "기타",
]
DETAIL_FAMILY_PLACEHOLDER = "본품분류 선택"
FAMILY_CARD_SECTION_ORDER = ["1DAY", "FRP", "기타"]
FAMILY_CARD_1DAY_NAMES = {
    "O2O2 D 컬러",
    "O2O2 D Micelia",
    "O2O2 D Toric",
    "Iris 컬러",
    "Iris Toric",
}
FAMILY_CARD_MISC_NAMES = {"PIA 1Day", "PIA Monthly"}
STANDARD_PACK_BUCKETS = ["5P", "10P", "30P", "80P", "90P"]
PRODUCT_QUERY_ALIASES = {
    "아이리스": ["Iris"],
    "클라렌": ["Clalen"],
    "오투오투": ["O2O2"],
    "원데이": ["1Day", "1-Day", "1D"],
    "토릭": ["Toric"],
    "딥블랙": ["Deep Black"],
    "레이크그레이": ["Lake Gray"],
    "뮤트브라운": ["Mute Brown"],
    "페일초코": ["Pale Choco"],
    "소울브라운": ["SoulBrown", "Soul Brown", "Iris SoulBrown", "Iris Soul Brown"],
    "수지그레이": ["Suzy Gray"],
    "수지브라운": ["Suzy Brown"],
    "알리샤브라운": ["Alicia Brown"],
    "알리사브라운": ["Alicia Brown"],
    "알리시아브라운": ["Alicia Brown"],
    "페즈브라운": ["Fez Brown"],
    "블루문": ["Blue Moon", "Bluemoon"],
    "미셀리아": ["Micelia"],
    "재즈블랙": ["JazzBlack", "Jazz Black"],
    "랩소디": ["Rhapsody"],
    "라틴": ["Latin"],
    "헤일로브라운": ["Halo Brown"],
    "헤일로그레이": ["Halo Gray"],
    "모카블랙": ["Mocha Black"],
    "퓨어초코": ["Pure Choco"],
    "로지브라운": ["Rosy Brown"],
    "샤이니브라운": ["Shiny Brown"],
    "세피아초코": ["Sepia Choco"],
    "애쉬브라운": ["Ash Brown"],
    "애쉬그레이": ["Ash Gray"],
    "블렌딩헤이즐": ["Blending Hazel"],
    "내추럴초콜릿": ["Natural Chocolat", "Natural Chocolate"],
    "클리어": ["Clear"],
    "프리덤": ["Freedom"],
    "스타플레어베이지": ["Starflare Beige"],
    "스타플레어그레이": ["Starflare Gray"],
    "데이글로우브라운": ["Dayglow Brown"],
    "데이글로우그레이": ["Dayglow Gray"],
    "아이코닉브라운": ["Iconic Brown"],
    "아포가토": ["Affogato"],
    "브륄레펄": ["Brulee Pearl"],
    "브릴레펄": ["Brulee Pearl"],
    "카푸치노": ["Cappuccino"],
    "체스트넛": ["Chestnut"],
    "커피젤리": ["Coffee Jelly"],
    "에스프레소": ["Espresso"],
    "우롱티": ["Oolong Tea"],
    "사쿠라무스": ["Sakura Mousse"],
    "쉬어블랙": ["Sheer Black"],
    "쉬어브라운": ["Sheer Brown"],
    "타르트타탱": ["Tarte Tatin"],
    "티라미수링": ["Tiramisu Ring"],
    "튤브라운": ["Tulle Brown"],
    "클린핏": ["CleanFit", "Clean Fit"],
    "리얼썸": ["Realsome", "Real Some"],
}
IRIS_PRODUCT_ALIAS_KEYS = {
    "소울브라운",
    "수지그레이",
    "수지브라운",
    "알리샤브라운",
    "알리사브라운",
    "알리시아브라운",
    "페즈브라운",
    "블루문",
    "재즈블랙",
    "랩소디",
    "라틴",
    "헤일로브라운",
    "헤일로그레이",
}
GENERIC_PRODUCT_ALIAS_KEYS = {"아이리스", "클라렌", "오투오투", "원데이", "토릭", "미셀리아"}
POWER_RE = re.compile(r"(-?\d+(?:\.\d+)?)\s*$")
CODE_KEY_RE = re.compile(r"[^0-9A-Za-z가-힣]+")
BASE_P_CODE_RE = re.compile(r"^(P\d+)")
PIA_TOKEN_RE = re.compile(r"\bPIA\b", re.IGNORECASE)
PACK_UNIT_RE = re.compile(r"(?:_|\b|\()(\d+(?:\.\d+)?)\s*(?:P|팩)?\)?$", re.IGNORECASE)
PACK_PREFIX_RE = re.compile(r"^(\d+(?:\.\d+)?)\s*(?:P|팩)_", re.IGNORECASE)
PACK_ANY_RE = re.compile(r"(?:^|[^0-9A-Za-z가-힣])(\d+(?:\.\d+)?)\s*(?:P|팩)(?:$|[^0-9A-Za-z가-힣])", re.IGNORECASE)
PACK_SUFFIX_RE = re.compile(
    r"(?:_샘플\(\d+(?:\.\d+)?P\)|_\d+(?:\.\d+)?\s*(?:P|팩)|_\d+(?:\.\d+)?)$",
    re.IGNORECASE,
)
PACK_PREFIX_SUFFIX_RE = re.compile(r"^\d+(?:\.\d+)?\s*(?:P|팩)_", re.IGNORECASE)

COLOR_BLUE = "#2563EB"
COLOR_TEAL = "#64748B"
COLOR_ORANGE = "#F97316"
COLOR_AMBER = "#7C3AED"
COLOR_DANGER = "#DC2626"
COLOR_ALERT_BG = "#FEF2F2"
COLOR_ALERT_BD = "#FECACA"
BG_PAGE = "#FFFFFF"
BG_CARD = "#FFFFFF"
BG_SECTION = "#FFFFFF"
TEXT_PRIMARY = "#111827"
TEXT_SECONDARY = "#6B7280"
TEXT_TERTIARY = "#9CA3AF"
BORDER_DEFAULT = "#E5E7EB"
BORDER_LIGHT = "#EEF0F3"

NAVY = COLOR_BLUE
SOFT_NAVY = COLOR_TEAL
WHITE = BG_CARD
LIGHT_GRAY = BG_PAGE
MID_GRAY = "#DEDED8"
TEXT_DARK = TEXT_PRIMARY
TEXT_MUTED = TEXT_SECONDARY
MUTED_ORANGE = COLOR_ORANGE
MUTED_RED = COLOR_ORANGE
PPT_FONT_NAME = "Pretendard"
REPORT_BG = "#FFFFFF"
REPORT_PANEL = "#FFFFFF"
REPORT_PANEL_LINE = "#E5E7EB"
REPORT_HEADER = "#111827"
REPORT_MUTED = "#6B7280"
REPORT_FAINT = "#E5E7EB"
REPORT_ROW_ALT = "#FFFFFF"
REPORT_ACCENT = COLOR_ORANGE
REPORT_NAVY = "#111827"
REPORT_TABLE_HEADER = "#FAFAFA"
REPORT_SOFT_BG = "#F8FAFC"
REPORT_BLUE_SOFT = "#EAF2F9"
REPORT_ACCENT_SOFT = "#FFF1E8"
REPORT_GREEN_SOFT = "#EAF7F1"

REQUEST_COLS = {
    "sales_code": ["판매코드", "판매 코드", "품목코드", "sales_code"],
    "product_name": ["제품명", "제품 명", "품명", "product_name"],
    "request_qty": [
        "수량(PACK)",
        "수량 (PACK)",
        "요청 PACK",
        "요청수량",
        "수량",
        "request_qty",
    ],
    "request_pcs": ["수량(PCS)", "수량 (PCS)", "요청 PCS", "요청수량(PCS)", "요청물량 PCS"],
    "units_per_pack": ["입수(낱개)", "입수", "팩당수량", "pack_size"],
    "due_date": [
        "계획일자",
        "계획 일자",
        "계획일",
        "생산계획일자",
        "생산 계획일자",
        "생산계획일",
        "생산 계획일",
        "납기일자",
        "납기 일자",
        "납기일",
        "due_date",
    ],
    "product_name_code": ["제품명코드", "제품명 코드", "제품규격", "품목코드"],
    "production_code": ["생산코드", "생산 코드", "제품코드", "제품 코드", "production_code"],
    "p_code": ["P코드(생산)", "P 코드(생산)", "P코드", "P 코드", "P code"],
    "q_code": ["Q코드(분리)", "Q 코드(분리)", "Q코드", "분리코드", "Q 코드"],
    "r_code": ["R코드(사출)", "R 코드(사출)", "R코드", "사출코드", "R 코드"],
    "market_type": ["국내/해외", "국내해외", "시장구분", "market_type"],
    "customer_name": ["거래처", "거래처명", "고객명", "고객 이름", "customer_name"],
    "category_summary": ["신규분류요약", "분류요약", "분류 요약", "category_summary"],
}

PRODUCT_CODE_MASTER_COLS = {
    "sales_code": ["품목코드", "판매코드", "판매 코드", "sales_code"],
    "product_name": ["품명", "제품명", "제품 명", "product_name"],
    "p_code": ["P 코드", "P코드", "P코드(생산)", "P code"],
    "production_code": ["제품코드", "제품 코드", "생산코드", "생산 코드", "production_code"],
    "q_code": ["분리코드", "Q코드", "Q코드(분리)", "Q 코드"],
    "r_code": ["사출코드", "R코드", "R코드(사출)", "R 코드"],
}

PACKING_COLS = {
    "sales_code": ["판매코드", "판매 코드", "품목코드", "sales_code"],
    "product_name": ["판매명", "생산명", "제품명", "제품 명", "품명", "product_name"],
    "lot_no": ["LOTNO", "LOT NO", "LOT", "lot_no"],
    "barcode_info": ["바코드정보", "바코드 정보", "barcode_info"],
    "packing_date": ["마킹일", "마킹시간", "포장일", "일자", "date"],
    "packing_qty": ["팩수량", "포장수량", "포장완료수량", "수량", "packing_qty"],
    "packing_pcs": ["낱개수량", "낱개 수량", "PCS수량", "PCS 수량", "pcs_qty", "packing_pcs"],
    "pack_unit": ["포장단위", "포장 단위", "입수", "pack_unit"],
}

YONGMA_COLS = {
    "sales_code": ["제품코드", "제품 코드", "판매코드", "판매 코드", "품목코드", "sales_code"],
    "product_name": ["품명", "제품명", "제품 명", "product_name"],
    "lot_no": ["LOTNO", "LOT NO", "LOT", "lot_no"],
    "receipt_qty": ["수량", "입고수량", "용마입고수량", "receipt_qty"],
}

SAMPLE_MOVEMENT_COLS = {
    "sales_code": ["품목코드", "제품코드", "제품 코드", "판매코드", "판매 코드", "sales_code"],
    "product_name": ["품명", "제품명", "제품 명", "product_name"],
    "lot_no": ["Lot No.", "LotNo", "LOTNO", "LOT NO", "LOT", "lot_no"],
    "movement_qty": ["이동수량", "이동 수량", "입고수량", "수량", "receipt_qty"],
    "status": ["상태", "처리상태", "status"],
    "inbound_warehouse": ["입고창고", "입고 창고", "입고처", "warehouse"],
}

SAMPLE_AVAILABLE_COLS = {
    "product_code": ["제품코드", "제품 코드", "품목코드", "생산코드", "production_code", "product_code"],
    "sample_available_qty": [
        "샘플 신청 가능 수량",
        "샘플신청가능수량",
        "샘플 신청가능수량",
        "샘플 신청 가능 수량 (당월)",
        "샘플신청가능수량(당월)",
        "sample_available_qty",
    ],
}
SAMPLE_AVAILABLE_QTY_COLUMN_INDEX = 9  # J열

WIP_COLS = {
    "product_code": ["제품 코드", "제품코드", "생산코드", "생산 코드", "production_code", "product_code"],
    "wip_qty": ["총 재공 수량", "총재공수량", "재공 수량", "재공수량", "수량", "wip_qty"],
    "process_name": ["WH_NAME", "공정명", "공정", "창고명", "재공 위치", "warehouse"],
}

WIP_PROCESS_ALIASES = {
    "검사접착": ["검사접착", "검사/접착", "검사 접착", "접착/멸균", "[55]접착/멸균"],
    "누수규격검사": ["누수규격검사", "누수/규격검사", "누수 규격검사", "누수 규격 검사", "[80]누수/규격검사"],
}

INVENTORY_STOCK_THRESHOLD_DEFAULT = 100

INVENTORY_COLS = {
    "sales_code": ["제품코드", "제품 코드", "판매코드", "판매 코드", "품목코드", "SKU", "sku"],
    "product_name": ["제품명", "제품 명", "품명", "product_name"],
    "available_stock_pack": [
        "실시간가용재고",
        "실시간 가용재고",
        "가용재고",
        "가용 재고",
        "재고수량",
        "재고 수량",
        "수량",
        "available_stock",
    ],
    "total_stock_pack": ["총수량", "총 수량", "총재고", "총 재고", "total_stock"],
    "product_spec": ["제품규격", "제품 규격", "규격", "product_spec"],
    "updated_at": ["전송일자", "전송 일자", "수집일자", "업데이트일자", "updated_at"],
}

PROCESS_STEPS = [
    {
        "id": "10",
        "header": "[10]사출조립",
        "label": "[10] 사출조립",
        "qty_col": "proc_10_qty",
        "due_col": "proc_10_due",
        "progress_pct": 20.0,
    },
    {
        "id": "20",
        "header": "[20]분리",
        "label": "[20] 분리",
        "qty_col": "proc_20_qty",
        "due_col": "proc_20_due",
        "progress_pct": 40.0,
    },
    {
        "id": "45",
        "header": "[45]하이드레이션/전면검사",
        "label": "[45] 하이드레이션/전면검사",
        "qty_col": "proc_45_qty",
        "due_col": "proc_45_due",
        "progress_pct": 60.0,
    },
    {
        "id": "55",
        "header": "[55]접착/멸균",
        "label": "[55] 접착/멸균",
        "qty_col": "proc_55_qty",
        "due_col": "proc_55_due",
        "progress_pct": 80.0,
    },
    {
        "id": "80",
        "header": "[80]누수/규격검사",
        "label": "[80] 누수/규격검사",
        "qty_col": "proc_80_qty",
        "due_col": "proc_80_due",
        "progress_pct": 100.0,
    },
]


def normalize_col(text: Any) -> str:
    return (
        str(text)
        .strip()
        .lower()
        .replace(" ", "")
        .replace("_", "")
        .replace("-", "")
        .replace("(", "")
        .replace(")", "")
    )


def clean_str(value: Any) -> str:
    if pd.isna(value):
        return ""
    return str(value).strip()


def factory_group_from_category(value: Any) -> str:
    text = clean_str(value)
    if not text:
        return "(미기재)"
    return FACTORY_GROUP_BY_CATEGORY.get(text, "(미기재)")


def factory_group_from_product_name(value: Any) -> str:
    text = clean_str(value)
    if not text:
        return ""
    upper = text.upper()
    if upper.startswith("PIA_KR_1D") or "FELIAMO" in upper or "EYESTAR" in upper:
        return "S관"
    if "REALSOME" in upper or "렌즈미" in text:
        return "S관"
    if "CLALEN O2O2 D" in upper:
        return "S관"
    if "CLALEN O2O2 M" in upper:
        return "A관"
    if "FRP" in upper:
        return "A관"
    if "IRIS" in upper:
        return "C관"
    return ""


def normalize_match_key(value: Any) -> str:
    return CODE_KEY_RE.sub("", clean_str(value)).upper()


def extract_pack_unit(value: Any) -> float:
    text = clean_str(value)
    if not text:
        return np.nan
    match = PACK_UNIT_RE.search(text)
    if not match:
        match = PACK_PREFIX_RE.search(text)
    if not match:
        return np.nan
    try:
        return float(match.group(1))
    except ValueError:
        return np.nan


def strip_pack_unit_suffix(value: Any) -> str:
    text = clean_str(value)
    stripped = PACK_PREFIX_SUFFIX_RE.sub("", text)
    stripped = PACK_SUFFIX_RE.sub("", stripped).strip("_ -")
    return stripped or text


def format_pack_unit_label(unit: Any, product_name: Any = "") -> str:
    num = pd.to_numeric(unit, errors="coerce")
    if pd.isna(num) or float(num) <= 0:
        return "(미기재)"
    value = f"{float(num):g}P"
    if "샘플" in clean_str(product_name):
        return f"{value} 샘플"
    return value


def extract_base_p_code_key(value: Any) -> str:
    match = BASE_P_CODE_RE.match(normalize_match_key(value))
    return match.group(1) if match else ""


def extract_code_measure_key(value: Any) -> str:
    values = extract_code_measure_values(value)
    if not values:
        return ""
    return "|".join(values[:2])


def extract_code_measure_values(value: Any) -> list[str]:
    text = clean_str(value).upper()
    if not text:
        return []
    return re.findall(r"[+-]?\d{1,2}\.\d{2}", text)


def normalize_measure_code(value: Any) -> str:
    return normalize_match_key(value)


def extract_production_family_key(value: Any) -> str:
    text = clean_str(value).upper()
    if not text:
        return ""
    base_p_code = extract_base_p_code_key(text)
    if not base_p_code:
        return ""
    without_measures = re.sub(r"[+-]?\d{1,2}\.\d{2}", "", text)
    key = normalize_match_key(without_measures)
    if re.fullmatch(r"P\d+[A-Z]?", key):
        return ""
    return key if key.startswith(base_p_code) else ""


def extract_production_code_template_parts(value: Any) -> dict[str, Any] | None:
    text = clean_str(value).upper()
    if not text:
        return None
    start = re.match(r"^P(\d+)([A-Z]?)", text)
    if not start:
        return None
    measures = list(re.finditer(r"[+-]?\d{1,2}\.\d{2}", text))
    if len(measures) < 2:
        return None
    tail = normalize_match_key(text[measures[1].end() :])
    if not tail:
        return None
    return {
        "number": int(start.group(1)),
        "width": len(start.group(1)),
        "letter": start.group(2),
        "tail": tail,
    }


def toric_cylinder_p_code_offset(value: Any) -> int | None:
    try:
        cylinder = abs(float(clean_str(value)))
    except ValueError:
        return None
    if np.isclose(cylinder, 0.75):
        return 0
    if np.isclose(cylinder, 1.25):
        return 1
    if np.isclose(cylinder, 1.75):
        return 2
    return None


def build_toric_progress_code_template_map(candidates: pd.DataFrame) -> dict[tuple[str, str], dict[str, Any]]:
    templates: dict[tuple[str, str], dict[str, Any]] = {}
    if candidates.empty:
        return templates
    for _, row in candidates.iterrows():
        sales_prefix = extract_sales_prefix(row.get("sales_code", ""))
        p_code_key = clean_str(row.get("p_code_key", ""))
        if not sales_prefix or not p_code_key:
            continue
        parts = extract_production_code_template_parts(row.get("production_code", ""))
        if not parts:
            continue
        templates.setdefault((sales_prefix, p_code_key), parts)
    return templates


def derive_toric_progress_code_key(row: pd.Series, templates: dict[tuple[str, str], dict[str, Any]]) -> str:
    sales_prefix = extract_sales_prefix(row.get("sales_code", ""))
    p_code_key = clean_str(row.get("p_code_key", ""))
    if not sales_prefix or not p_code_key:
        return ""
    template = templates.get((sales_prefix, p_code_key))
    if not template:
        return ""
    measures = extract_code_measure_values(row.get("sales_code", ""))
    if len(measures) < 2:
        return ""
    offset = toric_cylinder_p_code_offset(measures[1])
    if offset is None:
        return ""
    base_match = re.match(r"^P(\d+)$", p_code_key)
    if not base_match:
        return ""
    base_number = int(base_match.group(1))
    variant_number = base_number + offset
    return (
        f"P{variant_number:0{int(template['width'])}d}"
        f"{template['letter']}"
        f"{normalize_measure_code(measures[0])}"
        f"{normalize_measure_code(measures[1])}"
        f"{template['tail']}"
    )


def code_summary_production_family_keys(code_summary: pd.DataFrame) -> pd.Series:
    if code_summary.empty:
        return pd.Series(dtype=str)

    def column_values(col: str) -> pd.Series:
        if col in code_summary.columns:
            return code_summary[col]
        return pd.Series("", index=code_summary.index)

    return pd.Series(
        [
            first_nonempty(
                [
                    extract_production_family_key(production_code),
                    extract_production_family_key(p_code),
                ]
            )
            for production_code, p_code in zip(column_values("production_code"), column_values("p_code"))
        ],
        index=code_summary.index,
    )


def align_base_product_names_by_production_family(code_summary: pd.DataFrame) -> pd.DataFrame:
    out = code_summary.copy()
    if out.empty:
        return out
    if "base_product_name" not in out.columns:
        out["base_product_name"] = out.get("product_name", pd.Series("", index=out.index)).map(strip_pack_unit_suffix)

    family_key = code_summary_production_family_keys(out).map(clean_str)
    if family_key.empty or (family_key == "").all():
        return out

    rep_source = out.loc[family_key != "", ["base_product_name"]].copy()
    rep_source["_production_family_key"] = family_key.loc[family_key != ""]
    for col in ["request_pcs", "request_pack", "packing_pack", "yongma_in_pack"]:
        if col in out.columns:
            rep_source[col] = pd.to_numeric(out.loc[rep_source.index, col], errors="coerce").fillna(0.0)
        else:
            rep_source[col] = 0.0
    rep_source["base_product_name"] = rep_source["base_product_name"].map(clean_str)
    rep_source = rep_source[rep_source["base_product_name"] != ""].copy()
    if rep_source.empty:
        return out

    rep_source["_display_weight"] = rep_source["request_pcs"]
    fallback_request = rep_source["request_pack"]
    fallback_supply = rep_source["packing_pack"] + rep_source["yongma_in_pack"]
    rep_source["_display_weight"] = rep_source["_display_weight"].where(
        rep_source["_display_weight"] > 0,
        fallback_request,
    )
    rep_source["_display_weight"] = rep_source["_display_weight"].where(
        rep_source["_display_weight"] > 0,
        fallback_supply,
    )
    representatives = (
        rep_source.sort_values(
            ["_production_family_key", "_display_weight", "base_product_name"],
            ascending=[True, False, True],
        )
        .drop_duplicates("_production_family_key", keep="first")
        .set_index("_production_family_key")["base_product_name"]
    )
    representative_names = family_key.map(representatives).fillna("")
    replace_mask = representative_names.map(clean_str) != ""
    out.loc[replace_mask, "base_product_name"] = representative_names.loc[replace_mask]
    return out


def build_first_value_map(df: pd.DataFrame, key_col: str, value_col: str) -> dict[str, str]:
    if key_col not in df.columns or value_col not in df.columns:
        return {}
    pairs = df[[key_col, value_col]].copy()
    pairs[key_col] = pairs[key_col].map(clean_str)
    pairs[value_col] = pairs[value_col].map(clean_str)
    pairs = pairs[(pairs[key_col] != "") & (pairs[value_col] != "")]
    if pairs.empty:
        return {}
    return pairs.drop_duplicates(key_col, keep="first").set_index(key_col)[value_col].to_dict()


def min_datetime(series: pd.Series) -> pd.Timestamp:
    dates = pd.to_datetime(series, errors="coerce")
    dates = dates.dropna()
    if dates.empty:
        return pd.NaT
    return dates.min()


def max_datetime(series: pd.Series) -> pd.Timestamp:
    dates = pd.to_datetime(series, errors="coerce")
    dates = dates.dropna()
    if dates.empty:
        return pd.NaT
    return dates.max()


def sum_numeric_or_nan(series: pd.Series) -> float:
    numbers = pd.to_numeric(series, errors="coerce").dropna()
    if numbers.empty:
        return np.nan
    return float(numbers.sum())


def first_nonempty(series: pd.Series) -> str:
    for value in series:
        text = clean_str(value)
        if text:
            return text
    return ""


def join_unique(series: pd.Series, limit: int = 3) -> str:
    values = [clean_str(value) for value in series if clean_str(value)]
    unique = list(dict.fromkeys(values))
    if not unique:
        return ""
    if len(unique) <= limit:
        return ", ".join(unique)
    return f"{', '.join(unique[:limit])} 외 {len(unique) - limit}"


def factory_group_values(value: Any) -> list[str]:
    values = []
    for part in clean_str(value).split(","):
        text = clean_str(part)
        if text and text != "(미기재)":
            values.append(text)
    return values


def has_factory_group(value: Any) -> bool:
    return bool(factory_group_values(value))


def clean_factory_group_display(value: Any) -> str:
    return ", ".join(dict.fromkeys(factory_group_values(value)))


def join_factory_groups(series: pd.Series, limit: int = 3) -> str:
    values: list[str] = []
    for value in series:
        values.extend(factory_group_values(value))
    unique = list(dict.fromkeys(values))
    if not unique:
        return ""
    if len(unique) <= limit:
        return ", ".join(unique)
    return f"{', '.join(unique[:limit])} 외 {len(unique) - limit}"


def to_number(series: pd.Series) -> pd.Series:
    cleaned = (
        series.astype(str)
        .str.replace(",", "", regex=False)
        .str.replace(" ", "", regex=False)
        .str.strip()
    )
    return pd.to_numeric(cleaned, errors="coerce").fillna(0.0)


def to_number_value(value: Any) -> float:
    if value is None:
        return 0.0
    try:
        if pd.isna(value):
            return 0.0
    except (TypeError, ValueError):
        pass
    if isinstance(value, (int, float, np.integer, np.floating)):
        return float(value)
    text = clean_str(value).replace(",", "").replace(" ", "")
    if not text:
        return 0.0
    number = pd.to_numeric(pd.Series([text]), errors="coerce").iloc[0]
    return 0.0 if pd.isna(number) else float(number)


def find_column(df: pd.DataFrame, aliases: list[str]) -> str | None:
    normalized = {normalize_col(col): col for col in df.columns}
    for alias in aliases:
        key = normalize_col(alias)
        if key in normalized:
            return normalized[key]
    return None


def resolve_columns(
    df: pd.DataFrame,
    alias_map: dict[str, list[str]],
    required_keys: list[str],
    file_label: str,
) -> dict[str, str]:
    resolved: dict[str, str] = {}
    missing: list[str] = []
    for key, aliases in alias_map.items():
        col = find_column(df, aliases)
        if col is None:
            if key in required_keys:
                missing.append(f"{key} (후보: {', '.join(aliases)})")
        else:
            resolved[key] = col
    if missing:
        raise DashboardConfigError([f"[{file_label}] 필수 컬럼 누락: {'; '.join(missing)}"])
    return resolved


def list_excel_files(search_roots: list[Path]) -> list[Path]:
    seen: set[Path] = set()
    files: list[Path] = []
    for root in search_roots:
        if not root.exists():
            continue
        for pattern in ("*.xlsx", "*.xls"):
            for path in root.glob(pattern):
                if not path.is_file() or path.name.startswith("~$"):
                    continue
                real = path.resolve()
                if real in seen:
                    continue
                seen.add(real)
                files.append(path)
    files.sort(key=lambda p: p.stat().st_mtime, reverse=True)
    return files


def pick_latest_by_name(files: list[Path], keywords: list[str]) -> Path | None:
    for file in files:
        stem = file.stem.lower()
        if any(keyword.lower() in stem for keyword in keywords):
            return file
    return None


def has_alias(columns: list[str], aliases: list[str]) -> bool:
    normalized_cols = {normalize_col(col) for col in columns}
    return any(normalize_col(alias) in normalized_cols for alias in aliases)


def discover_source_files(base_dir: Path) -> SourceFiles:
    files = list_excel_files([base_dir / "data", base_dir])
    if not files:
        raise DashboardConfigError(
            [
                "엑셀 파일을 찾지 못했습니다.",
                f"- 검색 위치: {base_dir / 'data'}",
                f"- 검색 위치: {base_dir}",
            ]
        )

    request_file = pick_latest_by_name(files, ["생산요청등록", "국내", "요청"])
    packing_file = pick_latest_by_name(files, ["포장설비투입현황", "포장설비투입", "포장실적", "포장"])
    progress_file = pick_latest_by_name(files, ["수요정보", "전공정"])
    inventory_file = pick_latest_by_name(files, ["용마WMS재고현황", "WMS재고현황", "WMS"])
    daily_inventory_file = pick_latest_by_name(files, DAILY_INVENTORY_FILE_KEYWORDS)
    product_master_file = pick_latest_by_name(files, ["판매코드-제품코드 매칭 마스터", "제품코드 매칭 마스터", "매칭 마스터"])
    wip_file = pick_latest_by_name(files, ["ODV_WIP", "WIP", "재공"])

    if request_file is None or packing_file is None:
        for file in files:
            try:
                cols = pd.read_excel(file, nrows=0).columns.astype(str).tolist()
            except Exception:
                continue
            if request_file is None:
                if has_alias(cols, REQUEST_COLS["sales_code"]) and has_alias(cols, REQUEST_COLS["product_name"]) and has_alias(
                    cols, REQUEST_COLS["request_qty"]
                ):
                    request_file = file
            if packing_file is None:
                if has_alias(cols, PACKING_COLS["sales_code"]) and has_alias(cols, PACKING_COLS["packing_qty"]):
                    packing_file = file
            if inventory_file is None:
                if has_alias(cols, INVENTORY_COLS["sales_code"]) and has_alias(
                    cols,
                    INVENTORY_COLS["available_stock_pack"],
                ):
                    inventory_file = file

    messages: list[str] = []
    if request_file is None:
        messages.append("생산요청등록(국내) 파일을 찾지 못했습니다.")
    if packing_file is None:
        messages.append("포장설비투입현황 파일을 찾지 못했습니다.")
    if messages:
        raise DashboardConfigError(messages)

    return SourceFiles(
        request_file=request_file,
        packing_file=packing_file,
        progress_file=progress_file,
        inventory_file=inventory_file,
        daily_inventory_file=daily_inventory_file,
        product_master_file=product_master_file,
        wip_file=wip_file,
    )


def read_excel_preferred_sheet(path: Path, preferred_sheet: str) -> pd.DataFrame:
    try:
        xl = pd.ExcelFile(path)
    except Exception:
        return pd.read_excel(path)
    sheet_name = preferred_sheet if preferred_sheet in xl.sheet_names else xl.sheet_names[0]
    return xl.parse(sheet_name=sheet_name)


def select_total_request_sheet(sheet_names: list[str]) -> str | None:
    for sheet_name in sheet_names:
        if "전체물량" in clean_str(sheet_name):
            return sheet_name
    for sheet_name in sheet_names:
        name = clean_str(sheet_name)
        if "생산지시" not in name:
            return sheet_name
    return sheet_names[0] if sheet_names else None


def select_instruction_request_sheet(sheet_names: list[str]) -> str | None:
    for sheet_name in sheet_names:
        if "생산지시" in clean_str(sheet_name):
            return sheet_name
    if not any("전체물량" in clean_str(sheet_name) for sheet_name in sheet_names):
        return None
    return "Sheet1" if "Sheet1" in sheet_names else (sheet_names[0] if sheet_names else None)


def read_request_workbook_sheet(path: Path, preferred_sheet: str | None = None) -> pd.DataFrame:
    try:
        xl = pd.ExcelFile(path)
    except Exception:
        return pd.read_excel(path)
    sheet_name = preferred_sheet if preferred_sheet in xl.sheet_names else select_total_request_sheet(xl.sheet_names)
    return xl.parse(sheet_name=sheet_name)


def has_excel_sheet(path: Path, sheet_name: str) -> bool:
    try:
        return sheet_name in pd.ExcelFile(path).sheet_names
    except Exception:
        return False


def read_resolved_excel_sheet(
    xl: pd.ExcelFile,
    sheet_name: str,
    alias_map: dict[str, list[str]],
    required_keys: list[str],
    file_label: str,
) -> pd.DataFrame:
    header = xl.parse(sheet_name=sheet_name, nrows=0)
    cols = resolve_columns(
        header,
        alias_map,
        required_keys=required_keys,
        file_label=file_label,
    )
    usecols = list(dict.fromkeys(cols.values()))
    return xl.parse(sheet_name=sheet_name, usecols=usecols)


def normalize_product_code_master(path: Path | None) -> pd.DataFrame:
    columns = [
        "sales_code_key",
        "master_product_name",
        "master_p_code",
        "master_production_code",
        "master_q_code",
        "master_r_code",
    ]
    if path is None:
        return pd.DataFrame(columns=columns)

    try:
        raw = read_excel_preferred_sheet(path, "Sheet1")
        cols = resolve_columns(
            raw,
            PRODUCT_CODE_MASTER_COLS,
            required_keys=["sales_code", "p_code"],
            file_label=path.name,
        )
    except Exception:
        return pd.DataFrame(columns=columns)

    out = pd.DataFrame(
        {
            "sales_code_key": raw[cols["sales_code"]].map(normalize_match_key),
            "master_product_name": raw[cols["product_name"]].map(clean_str)
            if "product_name" in cols
            else "",
            "master_p_code": raw[cols["p_code"]].map(clean_str),
            "master_production_code": raw[cols["production_code"]].map(clean_str)
            if "production_code" in cols
            else "",
            "master_q_code": raw[cols["q_code"]].map(clean_str) if "q_code" in cols else "",
            "master_r_code": raw[cols["r_code"]].map(clean_str) if "r_code" in cols else "",
        }
    )
    out = out[(out["sales_code_key"] != "") & (out["master_p_code"].map(clean_str) != "")].copy()
    if out.empty:
        return pd.DataFrame(columns=columns)
    return out.drop_duplicates("sales_code_key", keep="first")[columns].copy()


def should_use_master_code(current_value: Any, master_value: Any, expected_prefix: str) -> bool:
    current = clean_str(current_value).upper()
    master = clean_str(master_value).upper()
    prefix = clean_str(expected_prefix).upper()
    if not prefix or not master.startswith(prefix):
        return False
    return current == "" or not current.startswith(prefix)


def enrich_request_from_product_master(request_df: pd.DataFrame, product_master_df: pd.DataFrame) -> pd.DataFrame:
    if request_df.empty or product_master_df.empty:
        return request_df.copy()

    out = request_df.copy()
    out["_sales_code_key_for_master"] = out["sales_code"].map(normalize_match_key)
    out = out.merge(
        product_master_df,
        left_on="_sales_code_key_for_master",
        right_on="sales_code_key",
        how="left",
        suffixes=("", "_master"),
    )
    fill_pairs = [
        ("p_code", "master_p_code", "P"),
        ("production_code", "master_production_code", "P"),
        ("q_code", "master_q_code", "Q"),
        ("r_code", "master_r_code", "R"),
    ]
    for target_col, master_col, expected_prefix in fill_pairs:
        if target_col not in out.columns:
            out[target_col] = ""
        if master_col not in out.columns:
            out[master_col] = ""
        use_master = [
            should_use_master_code(current, master, expected_prefix)
            for current, master in zip(out[target_col], out[master_col])
        ]
        out[target_col] = out[target_col].where(~pd.Series(use_master, index=out.index), out[master_col].fillna(""))
    if "master_product_name" in out.columns:
        master_product = out["master_product_name"].map(clean_str)
        out["product_name"] = out["master_product_name"].where(master_product != "", out["product_name"])
    return out.drop(
        columns=[
            "_sales_code_key_for_master",
            "sales_code_key",
            "master_product_name",
            "master_p_code",
            "master_production_code",
            "master_q_code",
            "master_r_code",
        ],
        errors="ignore",
    )


def normalize_request(
    path: Path,
    product_master_path: Path | None = None,
    preferred_sheet: str | None = None,
    product_master_df: pd.DataFrame | None = None,
) -> pd.DataFrame:
    raw = read_request_workbook_sheet(path, preferred_sheet)
    cols = resolve_columns(
        raw,
        REQUEST_COLS,
        required_keys=["sales_code", "product_name", "request_qty"],
        file_label=path.name,
    )
    out = pd.DataFrame(
        {
            "sales_code": raw[cols["sales_code"]].map(clean_str),
            "product_name": raw[cols["product_name"]].map(clean_str).replace("", "(제품명 미기재)"),
            "request_pack": to_number(raw[cols["request_qty"]]),
        }
    )
    request_pcs = to_number(raw[cols["request_pcs"]]) if "request_pcs" in cols else pd.Series(0.0, index=raw.index)
    raw_units_per_pack = (
        to_number(raw[cols["units_per_pack"]])
        if "units_per_pack" in cols
        else pd.Series(np.nan, index=raw.index)
    )
    name_units_per_pack = out["product_name"].map(extract_pack_unit)
    units_per_pack = raw_units_per_pack.where(raw_units_per_pack > 0, name_units_per_pack)
    out["pack_unit"] = units_per_pack.where(units_per_pack > 0, np.nan)
    out["pack_unit_label"] = [
        format_pack_unit_label(unit, name)
        for unit, name in zip(out["pack_unit"], out["product_name"])
    ]
    out["base_product_name"] = out["product_name"].map(strip_pack_unit_suffix)
    fallback_pcs = out["request_pack"] * units_per_pack.where(units_per_pack > 0, 1.0)
    out["request_pcs"] = request_pcs.where(request_pcs > 0, fallback_pcs)
    optional_text_cols = {
        "product_name_code": "product_name_code",
        "production_code": "production_code",
        "p_code": "p_code",
        "q_code": "q_code",
        "r_code": "r_code",
        "market_type": "market_type",
        "customer_name": "customer_name",
        "category_summary": "category_summary",
    }
    for source_key, output_col in optional_text_cols.items():
        if source_key in cols:
            out[output_col] = raw[cols[source_key]].map(clean_str)
        else:
            out[output_col] = "(미기재)" if output_col in {"customer_name", "category_summary"} else ""
    if "due_date" in cols:
        out["request_due_date"] = pd.to_datetime(raw[cols["due_date"]], errors="coerce")
    else:
        out["request_due_date"] = pd.NaT

    if "market_type" in out.columns:
        overseas_mask = out["market_type"].astype(str).str.contains("해외", case=False, na=False)
        out = out[~overseas_mask].copy()

    out = filter_request_for_due_month(out)
    master_df = product_master_df if product_master_df is not None else normalize_product_code_master(product_master_path)
    out = enrich_request_from_product_master(out, master_df)
    out["base_product_name"] = out["product_name"].map(strip_pack_unit_suffix)
    out["factory_group"] = out["category_summary"].map(factory_group_from_category)

    for col in ["sales_code", "product_name", "product_name_code", "production_code", "p_code", "q_code", "r_code"]:
        out[f"{col}_key"] = out[col].map(normalize_match_key)
    return out


def normalize_instruction_request(
    path: Path,
    product_master_path: Path | None = None,
    product_master_df: pd.DataFrame | None = None,
) -> pd.DataFrame:
    try:
        xl = pd.ExcelFile(path)
    except Exception:
        return pd.DataFrame()
    sheet_name = select_instruction_request_sheet(xl.sheet_names)
    if not sheet_name:
        return pd.DataFrame()
    return normalize_request(path, product_master_path, preferred_sheet=sheet_name, product_master_df=product_master_df)


def filter_request_for_due_month(
    request_df: pd.DataFrame,
    target_month: str = REQUEST_DUE_MONTH,
) -> pd.DataFrame:
    if request_df.empty or "request_due_date" not in request_df.columns:
        return request_df.copy()

    due_dates = pd.to_datetime(request_df["request_due_date"], errors="coerce")
    if due_dates.notna().sum() == 0:
        return request_df.copy()

    target_period = pd.Period(target_month, freq="M")
    return request_df.loc[due_dates.dt.to_period("M") == target_period].copy()


def normalize_packing_frame(raw: pd.DataFrame, file_label: str) -> pd.DataFrame:
    cols = resolve_columns(
        raw,
        PACKING_COLS,
        required_keys=["sales_code", "packing_qty"],
        file_label=file_label,
    )
    packing_pack = to_number(raw[cols["packing_qty"]])
    if "packing_pcs" in cols:
        packing_pcs = to_number(raw[cols["packing_pcs"]])
        if "pack_unit" in cols:
            pack_unit = to_number(raw[cols["pack_unit"]])
        elif "product_name" in cols:
            pack_unit = raw[cols["product_name"]].map(extract_pack_unit)
        else:
            pack_unit = pd.Series(0.0, index=raw.index)
        pack_unit = pd.to_numeric(pack_unit, errors="coerce").fillna(0.0)
        # Some exported pack counts are formatted as Excel dates; recover PACK from PCS when needed.
        derived_pack = (packing_pcs / pack_unit.where(pack_unit > 0, np.nan)).replace([np.inf, -np.inf], np.nan).fillna(0.0)
        packing_pack = packing_pack.where((packing_pack > 0) | (derived_pack <= 0), derived_pack)

    out = pd.DataFrame(
        {
            "sales_code": raw[cols["sales_code"]].map(clean_str),
            "packing_pack": packing_pack,
        }
    )
    out["sales_code_key"] = out["sales_code"].map(normalize_match_key)
    out["packing_product_name"] = raw[cols["product_name"]].map(clean_str) if "product_name" in cols else ""
    out["packing_lot"] = raw[cols["lot_no"]].map(clean_str) if "lot_no" in cols else ""
    out["packing_lot_key"] = out["packing_lot"].map(normalize_match_key)
    out["packing_barcode"] = raw[cols["barcode_info"]].map(clean_str) if "barcode_info" in cols else ""
    out["packing_barcode_key"] = out["packing_barcode"].map(normalize_match_key)
    out["packing_date"] = (
        parse_datetime_series(raw[cols["packing_date"]]) if "packing_date" in cols else pd.NaT
    )
    return out


def normalize_packing(path: Path) -> pd.DataFrame:
    try:
        xl = pd.ExcelFile(path)
        sheet_name = "포장실적" if "포장실적" in xl.sheet_names else xl.sheet_names[0]
        raw = read_resolved_excel_sheet(
            xl,
            sheet_name,
            PACKING_COLS,
            required_keys=["sales_code", "packing_qty"],
            file_label=path.name,
        )
    except DashboardConfigError:
        raise
    except Exception:
        raw = read_excel_preferred_sheet(path, "포장실적")
    return normalize_packing_frame(raw, path.name)


def empty_yongma_movement_df() -> pd.DataFrame:
    return pd.DataFrame(
        columns=[
            "sales_code",
            "sales_code_key",
            "yongma_product_name",
            "yongma_lot",
            "yongma_lot_key",
            "yongma_in_pack",
        ]
    )


def normalize_yongma_movement_frame(raw: pd.DataFrame, file_label: str) -> pd.DataFrame:
    cols = resolve_columns(
        raw,
        YONGMA_COLS,
        required_keys=["sales_code", "lot_no", "receipt_qty"],
        file_label=file_label,
    )
    out = pd.DataFrame(
        {
            "sales_code": raw[cols["sales_code"]].map(clean_str),
            "yongma_lot": raw[cols["lot_no"]].map(clean_str),
            "yongma_in_pack": to_number(raw[cols["receipt_qty"]]),
        }
    )
    out["sales_code_key"] = out["sales_code"].map(normalize_match_key)
    out["yongma_product_name"] = raw[cols["product_name"]].map(clean_str) if "product_name" in cols else ""
    out["yongma_lot_key"] = out["yongma_lot"].map(normalize_match_key)
    return out[(out["sales_code_key"] != "") & (out["yongma_in_pack"] > 0)].copy()


def normalize_yongma_movement(path: Path) -> pd.DataFrame:
    sheet_name = "용마이동현황"
    if not has_excel_sheet(path, sheet_name):
        return empty_yongma_movement_df()

    try:
        xl = pd.ExcelFile(path)
        raw = read_resolved_excel_sheet(
            xl,
            sheet_name,
            YONGMA_COLS,
            required_keys=["sales_code", "lot_no", "receipt_qty"],
            file_label=f"{path.name}:{sheet_name}",
        )
    except DashboardConfigError:
        raise
    except Exception:
        raw = pd.read_excel(path, sheet_name=sheet_name)
    return normalize_yongma_movement_frame(raw, f"{path.name}:{sheet_name}")


def normalize_sample_movement_frame(raw: pd.DataFrame, file_label: str) -> pd.DataFrame:
    cols = resolve_columns(
        raw,
        SAMPLE_MOVEMENT_COLS,
        required_keys=["sales_code", "movement_qty"],
        file_label=file_label,
    )
    out = pd.DataFrame(
        {
            "sales_code": raw[cols["sales_code"]].map(clean_str),
            "yongma_in_pack": to_number(raw[cols["movement_qty"]]),
        }
    )
    out["sales_code_key"] = out["sales_code"].map(normalize_match_key)
    out["yongma_product_name"] = raw[cols["product_name"]].map(clean_str) if "product_name" in cols else ""
    out["yongma_lot"] = raw[cols["lot_no"]].map(clean_str) if "lot_no" in cols else ""
    out["yongma_lot_key"] = out["yongma_lot"].map(normalize_match_key)

    keep = pd.Series(True, index=raw.index)
    if "status" in cols:
        status = raw[cols["status"]].map(clean_str)
        keep &= (status == "") | status.str.contains("확인|완료", regex=True, na=False)
    if "inbound_warehouse" in cols:
        inbound = raw[cols["inbound_warehouse"]].map(clean_str)
        if inbound.str.contains("샘플", regex=False, na=False).any():
            keep &= inbound.str.contains("샘플", regex=False, na=False)

    out = out[keep].copy()
    return out[(out["sales_code_key"] != "") & (out["yongma_in_pack"] > 0)].copy()


def normalize_sample_movement(path: Path) -> pd.DataFrame:
    sheet_name = "샘플이동"
    if not has_excel_sheet(path, sheet_name):
        return empty_yongma_movement_df()

    try:
        xl = pd.ExcelFile(path)
        raw = read_resolved_excel_sheet(
            xl,
            sheet_name,
            SAMPLE_MOVEMENT_COLS,
            required_keys=["sales_code", "movement_qty"],
            file_label=f"{path.name}:{sheet_name}",
        )
    except DashboardConfigError:
        raise
    except Exception:
        raw = pd.read_excel(path, sheet_name=sheet_name)
    return normalize_sample_movement_frame(raw, f"{path.name}:{sheet_name}")


def empty_sample_available_df() -> pd.DataFrame:
    return pd.DataFrame(
        columns=[
            "product_code",
            "production_code_key",
            "sample_available_pcs",
        ]
    )


def normalize_sample_available_frame(
    raw: pd.DataFrame,
    file_label: str,
    sample_available_col: str | None = None,
) -> pd.DataFrame:
    cols = resolve_columns(
        raw,
        SAMPLE_AVAILABLE_COLS,
        required_keys=["product_code"],
        file_label=file_label,
    )
    if sample_available_col is not None and sample_available_col in raw.columns:
        cols["sample_available_qty"] = sample_available_col
    elif "sample_available_qty" in cols:
        pass
    elif len(raw.columns) > SAMPLE_AVAILABLE_QTY_COLUMN_INDEX:
        cols["sample_available_qty"] = raw.columns[SAMPLE_AVAILABLE_QTY_COLUMN_INDEX]
    else:
        raise DashboardConfigError(
            [f"[{file_label}] 샘플신청가능수량 J열을 찾지 못했습니다."]
        )
    out = pd.DataFrame(
        {
            "product_code": raw[cols["product_code"]].map(clean_str),
            "sample_available_pcs": to_number(raw[cols["sample_available_qty"]]),
        }
    )
    out["production_code_key"] = out["product_code"].map(normalize_match_key)
    out = out[(out["production_code_key"] != "") & (out["sample_available_pcs"] > 0)].copy()
    if out.empty:
        return empty_sample_available_df()
    return (
        out.groupby("production_code_key", dropna=False)
        .agg(
            product_code=("product_code", first_nonempty),
            sample_available_pcs=("sample_available_pcs", "sum"),
        )
        .reset_index()[["product_code", "sample_available_pcs", "production_code_key"]]
        .copy()
    )


def read_sample_available_sheet(xl: pd.ExcelFile, sheet_name: str, file_label: str) -> pd.DataFrame:
    worksheet = xl.book[sheet_name]
    header_values = list(next(worksheet.iter_rows(min_row=1, max_row=1, values_only=True), []))
    normalized_headers = {normalize_col(value): idx for idx, value in enumerate(header_values)}
    product_col_idx = None
    for alias in SAMPLE_AVAILABLE_COLS["product_code"]:
        idx = normalized_headers.get(normalize_col(alias))
        if idx is not None:
            product_col_idx = idx
            break
    if product_col_idx is None:
        raise DashboardConfigError(
            [f"[{file_label}] 필수 컬럼 누락: product_code (후보: {', '.join(SAMPLE_AVAILABLE_COLS['product_code'])})"]
        )
    sample_col_idx = None
    for alias in SAMPLE_AVAILABLE_COLS["sample_available_qty"]:
        idx = normalized_headers.get(normalize_col(alias))
        if idx is not None:
            sample_col_idx = idx
            break
    if sample_col_idx is None:
        if len(header_values) <= SAMPLE_AVAILABLE_QTY_COLUMN_INDEX:
            raise DashboardConfigError([f"[{file_label}] 샘플신청가능수량 J열을 찾지 못했습니다."])
        sample_col_idx = SAMPLE_AVAILABLE_QTY_COLUMN_INDEX

    min_col_idx = min(product_col_idx, sample_col_idx)
    max_col_idx = max(product_col_idx, sample_col_idx)
    product_offset = product_col_idx - min_col_idx
    sample_offset = sample_col_idx - min_col_idx
    sample_by_key: dict[str, float] = {}
    product_by_key: dict[str, str] = {}
    for row in worksheet.iter_rows(
        min_row=2,
        min_col=min_col_idx + 1,
        max_col=max_col_idx + 1,
        values_only=True,
    ):
        product_code = clean_str(row[product_offset] if product_offset < len(row) else "")
        production_code_key = normalize_match_key(product_code)
        if not production_code_key:
            continue
        sample_qty = to_number_value(row[sample_offset] if sample_offset < len(row) else 0.0)
        if sample_qty <= 0:
            continue
        sample_by_key[production_code_key] = sample_by_key.get(production_code_key, 0.0) + sample_qty
        product_by_key.setdefault(production_code_key, product_code)

    if not sample_by_key:
        return empty_sample_available_df()
    return pd.DataFrame(
        {
            "product_code": [product_by_key[key] for key in sample_by_key],
            "sample_available_pcs": [sample_by_key[key] for key in sample_by_key],
            "production_code_key": list(sample_by_key),
        }
    )


def normalize_sample_available(path: Path) -> pd.DataFrame:
    sheet_name = "샘플신청가능수량"
    if not has_excel_sheet(path, sheet_name):
        return empty_sample_available_df()

    try:
        xl = pd.ExcelFile(path)
        return read_sample_available_sheet(xl, sheet_name, f"{path.name}:{sheet_name}")
    except DashboardConfigError:
        raise
    except Exception:
        raw = pd.read_excel(path, sheet_name=sheet_name)
    return normalize_sample_available_frame(raw, f"{path.name}:{sheet_name}")


def empty_wip_df() -> pd.DataFrame:
    return pd.DataFrame(
        columns=[
            "production_code_key",
            "production_code_display",
            "power_value",
            "POWER",
            *WIP_PROCESS_COLUMNS,
        ]
    )


def canonical_wip_process(value: Any) -> str:
    key = normalize_col(value)
    if not key:
        return ""
    mapping = getattr(canonical_wip_process, "_mapping", None)
    if mapping is None:
        mapping = {
            normalize_col(alias): canonical
            for canonical, aliases in WIP_PROCESS_ALIASES.items()
            for alias in aliases
        }
        setattr(canonical_wip_process, "_mapping", mapping)
    return mapping.get(key, "")


def normalize_wip(path: Path | None) -> pd.DataFrame:
    if path is None:
        return empty_wip_df()

    try:
        raw = read_excel_preferred_sheet(path, "Sheet1")
        cols = resolve_columns(
            raw,
            WIP_COLS,
            required_keys=["product_code", "wip_qty", "process_name"],
            file_label=path.name,
        )
    except DashboardConfigError:
        raise
    except Exception:
        return empty_wip_df()

    work = pd.DataFrame(
        {
            "production_code_display": raw[cols["product_code"]].map(clean_str),
            "production_code_key": raw[cols["product_code"]].map(normalize_match_key),
            "wip_qty": to_number(raw[cols["wip_qty"]]),
            "wip_process": raw[cols["process_name"]].map(canonical_wip_process),
        }
    )
    work = work[
        (work["production_code_key"].str.startswith("P"))
        & (work["wip_process"] != "")
        & (work["wip_qty"] > 0)
    ].copy()
    if work.empty:
        return empty_wip_df()

    work["power_value"] = work["production_code_display"].map(parse_power_from_sales_code)
    work["POWER"] = work["power_value"].map(format_power)
    grouped = (
        work.groupby(
            ["production_code_key", "production_code_display", "power_value", "POWER", "wip_process"],
            dropna=False,
        )["wip_qty"]
        .sum()
        .reset_index()
    )
    pivot = (
        grouped.pivot_table(
            index=["production_code_key", "production_code_display", "power_value", "POWER"],
            columns="wip_process",
            values="wip_qty",
            aggfunc="sum",
            fill_value=0.0,
        )
        .reset_index()
        .rename_axis(None, axis=1)
    )
    for col in WIP_PROCESS_COLUMNS:
        if col not in pivot.columns:
            pivot[col] = 0.0
        pivot[col] = pd.to_numeric(pivot[col], errors="coerce").fillna(0.0)
    return pivot[["production_code_key", "production_code_display", "power_value", "POWER", *WIP_PROCESS_COLUMNS]].copy()


def normalize_packing_workbook(path: Path) -> tuple[pd.DataFrame, pd.DataFrame, pd.DataFrame]:
    try:
        xl = pd.ExcelFile(path)
    except Exception:
        return (
            normalize_packing(path),
            pd.concat(
                [normalize_yongma_movement(path), normalize_sample_movement(path)],
                ignore_index=True,
            ),
            normalize_sample_available(path),
        )

    packing_sheet = "포장실적" if "포장실적" in xl.sheet_names else xl.sheet_names[0]
    packing_raw = read_resolved_excel_sheet(
        xl,
        packing_sheet,
        PACKING_COLS,
        required_keys=["sales_code", "packing_qty"],
        file_label=path.name,
    )
    packing_df = normalize_packing_frame(packing_raw, path.name)

    yongma_sheet = "용마이동현황"
    if yongma_sheet in xl.sheet_names:
        yongma_raw = read_resolved_excel_sheet(
            xl,
            yongma_sheet,
            YONGMA_COLS,
            required_keys=["sales_code", "lot_no", "receipt_qty"],
            file_label=f"{path.name}:{yongma_sheet}",
        )
        yongma_df = normalize_yongma_movement_frame(yongma_raw, f"{path.name}:{yongma_sheet}")
    else:
        yongma_df = empty_yongma_movement_df()

    sample_movement_sheet = "샘플이동"
    if sample_movement_sheet in xl.sheet_names:
        sample_movement_raw = read_resolved_excel_sheet(
            xl,
            sample_movement_sheet,
            SAMPLE_MOVEMENT_COLS,
            required_keys=["sales_code", "movement_qty"],
            file_label=f"{path.name}:{sample_movement_sheet}",
        )
        sample_movement_df = normalize_sample_movement_frame(
            sample_movement_raw,
            f"{path.name}:{sample_movement_sheet}",
        )
        if not sample_movement_df.empty:
            yongma_df = pd.concat([yongma_df, sample_movement_df], ignore_index=True)

    sample_sheet = "샘플신청가능수량"
    if sample_sheet in xl.sheet_names:
        sample_available_df = read_sample_available_sheet(xl, sample_sheet, f"{path.name}:{sample_sheet}")
    else:
        sample_available_df = empty_sample_available_df()

    return packing_df, yongma_df, sample_available_df


def empty_inventory_df() -> pd.DataFrame:
    return pd.DataFrame(
        columns=[
            "sales_code",
            "sales_code_key",
            "inventory_product_name",
            "available_stock_pack",
            "inventory_total_stock_pack",
            "inventory_product_spec",
            "inventory_updated_at",
        ]
    )


def normalize_inventory(path: Path | None) -> pd.DataFrame:
    if path is None:
        return empty_inventory_df()

    raw = pd.read_excel(path)
    cols = resolve_columns(
        raw,
        INVENTORY_COLS,
        required_keys=["sales_code", "available_stock_pack"],
        file_label=path.name,
    )
    out = pd.DataFrame(
        {
            "sales_code": raw[cols["sales_code"]].map(clean_str),
            "sales_code_key": raw[cols["sales_code"]].map(normalize_match_key),
            "available_stock_pack": to_number(raw[cols["available_stock_pack"]]),
        }
    )
    out["inventory_product_name"] = (
        raw[cols["product_name"]].map(clean_str) if "product_name" in cols else ""
    )
    out["inventory_total_stock_pack"] = (
        to_number(raw[cols["total_stock_pack"]]) if "total_stock_pack" in cols else np.nan
    )
    out["inventory_product_spec"] = (
        raw[cols["product_spec"]].map(clean_str) if "product_spec" in cols else ""
    )
    out["inventory_updated_at"] = (
        parse_datetime_series(raw[cols["updated_at"]]) if "updated_at" in cols else pd.NaT
    )
    return out[out["sales_code_key"] != ""].copy()


DAILY_INVENTORY_COLUMNS = [
    "제품명",
    "제품코드",
    "PACK",
    "POWER",
    "CP",
    "재고수량",
    "전일재고",
    "재고증감",
    "긴급요청",
    "대상품목",
]


def empty_daily_inventory_df() -> pd.DataFrame:
    return pd.DataFrame(columns=DAILY_INVENTORY_COLUMNS)


def numeric_scalar(value: Any, default: float = 0.0) -> float:
    number = pd.to_numeric(value, errors="coerce")
    if pd.isna(number):
        return default
    return float(number)


def daily_power_label(value: Any) -> str:
    text = clean_str(value).replace("−", "-").replace("–", "-").replace("—", "-")
    if not text:
        return ""
    if text.upper() == "PL":
        return format_power(0.0)
    number = pd.to_numeric(text, errors="coerce")
    if pd.isna(number):
        return ""
    return format_power(float(number))


def parse_daily_power_tokens(value: Any) -> list[str]:
    text = clean_str(value).replace("−", "-").replace("–", "-").replace("—", "-")
    if not text:
        return []
    tokens = re.findall(r"PL|[+-]?\d+(?:\.\d+)?", text, flags=re.IGNORECASE)
    powers = [daily_power_label(token) for token in tokens]
    return list(dict.fromkeys([power for power in powers if power]))


def normalize_toric_cp_label(value: Any) -> str:
    text = clean_str(value).replace("−", "-").replace("–", "-").replace("—", "-")
    if not text:
        return ""
    text = text.strip()
    compact_digits = re.sub(r"\D", "", text)
    if compact_digits in {"075", "125", "175"}:
        return f"{int(compact_digits) / 100:.2f}"

    number = pd.to_numeric(text, errors="coerce")
    if pd.notna(number) and abs(float(number)) in {0.75, 1.25, 1.75}:
        return f"{abs(float(number)):.2f}"
    return ""


def extract_toric_cp_label(value: Any) -> str:
    text = clean_str(value).replace("−", "-").replace("–", "-").replace("—", "-")
    if not text:
        return ""
    decimal_match = re.search(r"(?<!\d)(0?\.75|1\.25|1\.75)(?!\d)", text)
    if decimal_match:
        return normalize_toric_cp_label(decimal_match.group(1))
    compact_match = re.search(r"(?<!\d)(075|125|175)(?!\d)", text)
    if compact_match:
        return normalize_toric_cp_label(compact_match.group(1))
    return ""


def cp_label_from_sales_code(value: Any) -> str:
    text = clean_str(value).upper().replace("−", "-").replace("–", "-").replace("—", "-")
    if not text:
        return ""
    match = re.match(r"^[A-Z]\d+-\d{1,2}\.\d{2}[-_\s]+((?:0?\.75|1\.25|1\.75)|(?:075|125|175))", text)
    if match:
        return normalize_toric_cp_label(match.group(1))
    return ""


def extract_daily_pack_label(value: Any) -> str:
    text = clean_str(value)
    match = re.search(r"(\d+(?:\.\d+)?)\s*(?:P|p|팩|개입)", text)
    if match:
        return base_pack_label(float(match.group(1)))
    if extract_toric_cp_label(text):
        return ""
    unit = extract_pack_unit(text)
    if pd.notna(unit) and float(unit) > 0:
        return base_pack_label(unit)
    return ""


def daily_inventory_key(product_name: Any, pack_label: Any, power_label: Any, cp_label: Any = "") -> str:
    product = compact_query_text(product_name)
    return f"{product}|{clean_str(pack_label).upper()}|{clean_str(power_label)}|{clean_str(cp_label)}"


def normalize_daily_emergency_requests(xl: pd.ExcelFile) -> pd.DataFrame:
    sheet_name = "긴급요청"
    if sheet_name not in xl.sheet_names:
        return empty_daily_inventory_df()

    raw = xl.parse(sheet_name=sheet_name, header=None)
    header_idx = None
    for idx, row in raw.iterrows():
        values = [clean_str(value) for value in row.tolist()]
        if "제품명" in values and "제품코드" in values:
            header_idx = idx
            break
    if header_idx is None:
        return empty_daily_inventory_df()

    header_values = [clean_str(value) for value in raw.iloc[header_idx].tolist()]
    product_col = header_values.index("제품명")
    code_col = header_values.index("제품코드")
    target_start_col = code_col + 1

    rows: list[dict[str, Any]] = []
    for row_idx in range(header_idx + 1, len(raw)):
        row = raw.iloc[row_idx]
        product_name = clean_str(row.iloc[product_col] if product_col < len(row) else "")
        product_code = clean_str(row.iloc[code_col] if code_col < len(row) else "")
        if not product_name and not product_code:
            continue
        target_text = " ".join(clean_str(value) for value in row.iloc[target_start_col:].tolist() if clean_str(value))
        powers = parse_daily_power_tokens(target_text)
        pack_label = extract_daily_pack_label(product_name)
        cp_label = extract_toric_cp_label(product_name)
        for power in powers:
            rows.append(
                {
                    "제품명": product_name,
                    "제품코드": product_code,
                    "PACK": pack_label,
                    "POWER": power,
                    "CP": cp_label,
                    "재고수량": np.nan,
                    "전일재고": np.nan,
                    "재고증감": np.nan,
                    "긴급요청": True,
                    "대상품목": target_text,
                }
            )
    if not rows:
        return empty_daily_inventory_df()
    return pd.DataFrame(rows, columns=DAILY_INVENTORY_COLUMNS)


def normalize_daily_support_inventory(xl: pd.ExcelFile) -> pd.DataFrame:
    sheet_name = "지원파트 재고표"
    if sheet_name not in xl.sheet_names:
        return empty_daily_inventory_df()

    raw = xl.parse(sheet_name=sheet_name, header=None)
    rows: list[dict[str, Any]] = []
    current_product = ""
    current_pack = ""
    current_cp = ""

    for _, row in raw.iterrows():
        product_candidate = clean_str(row.iloc[1] if len(row) > 1 else "")
        pack_candidate = extract_daily_pack_label(product_candidate)
        cp_candidate = extract_toric_cp_label(product_candidate)
        if product_candidate and (pack_candidate or cp_candidate):
            current_product = product_candidate
            current_pack = pack_candidate
            current_cp = cp_candidate
        if not current_product or (not current_pack and not current_cp):
            continue

        for offset in range(18):
            power_col = 2 + offset
            current_col = 21 + offset
            previous_col = 39 + offset
            if power_col >= len(row):
                continue
            power = daily_power_label(row.iloc[power_col])
            if not power:
                continue
            current_stock = numeric_scalar(row.iloc[current_col] if current_col < len(row) else np.nan, np.nan)
            previous_stock = numeric_scalar(row.iloc[previous_col] if previous_col < len(row) else np.nan, np.nan)
            if pd.isna(current_stock) and pd.isna(previous_stock):
                continue
            rows.append(
                {
                    "제품명": current_product,
                    "제품코드": "",
                    "PACK": current_pack,
                    "POWER": power,
                    "CP": current_cp,
                    "재고수량": current_stock,
                    "전일재고": previous_stock,
                    "재고증감": current_stock - previous_stock if pd.notna(current_stock) and pd.notna(previous_stock) else np.nan,
                    "긴급요청": False,
                    "대상품목": "",
                }
            )

    if not rows:
        return empty_daily_inventory_df()
    return pd.DataFrame(rows, columns=DAILY_INVENTORY_COLUMNS)


def normalize_daily_inventory_file(path: Path | None) -> pd.DataFrame:
    if path is None:
        return empty_daily_inventory_df()
    try:
        xl = pd.ExcelFile(path)
    except Exception:
        return empty_daily_inventory_df()

    support = normalize_daily_support_inventory(xl)
    emergency = normalize_daily_emergency_requests(xl)
    if support.empty and emergency.empty:
        return empty_daily_inventory_df()

    support_work = support.copy()
    emergency_work = emergency.copy()
    support_work["_daily_key"] = [
        daily_inventory_key(product, pack, power, cp)
        for product, pack, power, cp in zip(support_work["제품명"], support_work["PACK"], support_work["POWER"], support_work["CP"])
    ]
    emergency_work["_daily_key"] = [
        daily_inventory_key(product, pack, power, cp)
        for product, pack, power, cp in zip(emergency_work["제품명"], emergency_work["PACK"], emergency_work["POWER"], emergency_work["CP"])
    ]

    merged = support_work.merge(
        emergency_work[["_daily_key", "제품명", "제품코드", "PACK", "POWER", "CP", "긴급요청", "대상품목"]].rename(
            columns={
                "제품명": "_emergency_product_name",
                "제품코드": "_emergency_product_code",
                "PACK": "_emergency_pack",
                "POWER": "_emergency_power",
                "CP": "_emergency_cp",
                "긴급요청": "_emergency_flag",
                "대상품목": "_emergency_target",
            }
        ),
        on="_daily_key",
        how="outer",
    )

    for col in DAILY_INVENTORY_COLUMNS:
        if col not in merged.columns:
            merged[col] = np.nan if col in {"재고수량", "전일재고", "재고증감"} else ""

    for base_col, emergency_col in [
        ("제품명", "_emergency_product_name"),
        ("제품코드", "_emergency_product_code"),
        ("PACK", "_emergency_pack"),
        ("POWER", "_emergency_power"),
        ("CP", "_emergency_cp"),
    ]:
        if emergency_col in merged.columns:
            merged[base_col] = merged[base_col].where(
                merged[base_col].map(clean_str) != "",
                merged[emergency_col],
            )
    if "_emergency_product_code" in merged.columns:
        merged["제품코드"] = merged["제품코드"].where(merged["제품코드"].map(clean_str) != "", merged["_emergency_product_code"])
    if "_emergency_flag" in merged.columns:
        base_flag = merged["긴급요청"].apply(lambda value: bool(value) if not pd.isna(value) else False)
        emergency_flag = merged["_emergency_flag"].apply(lambda value: bool(value) if not pd.isna(value) else False)
        merged["긴급요청"] = base_flag | emergency_flag
    if "_emergency_target" in merged.columns:
        merged["대상품목"] = merged["대상품목"].where(merged["대상품목"].map(clean_str) != "", merged["_emergency_target"])

    negative_stock = pd.to_numeric(merged["재고수량"], errors="coerce") < 0
    merged["긴급요청"] = merged["긴급요청"].apply(lambda value: bool(value) if not pd.isna(value) else False) | negative_stock
    merged.loc[negative_stock & (merged["대상품목"].map(clean_str) == ""), "대상품목"] = "재고표 음수 재고"

    merged["재고증감"] = pd.to_numeric(merged["재고증감"], errors="coerce")
    missing_delta = merged["재고증감"].isna()
    merged.loc[missing_delta, "재고증감"] = (
        pd.to_numeric(merged.loc[missing_delta, "재고수량"], errors="coerce")
        - pd.to_numeric(merged.loc[missing_delta, "전일재고"], errors="coerce")
    )
    return merged[DAILY_INVENTORY_COLUMNS].copy()


def clean_inventory_product_name(value: Any) -> str:
    text = clean_str(value)
    if not text:
        return ""
    text = re.sub(r"[/\\]+$", "", text).strip()
    text = re.sub(r"[/_ -]*[+-]?\d{1,2}\.\d{2}$", "", text).strip("_-/ ")
    return text or clean_str(value)


def inventory_power_label_from_sales_code(value: Any) -> str:
    power_value = parse_power_from_sales_code(value)
    if pd.isna(power_value):
        return ""
    return format_power(power_value)


def pack_label_from_inventory_name(product_name: Any, product_spec: Any = "") -> str:
    product_text = clean_str(product_name)
    inline_pack = re.search(
        r"(?:^|[_\s])(\d+(?:\.\d+)?)\s*(?:P|팩)(?=$|[_\s])",
        product_text,
        flags=re.IGNORECASE,
    )
    if inline_pack:
        return base_pack_label(float(inline_pack.group(1)))
    product_text = re.sub(r"[/\\]+$", "", product_text).strip()
    product_text = re.sub(r"[/_ -]*[+-]?\d{1,2}\.\d{2}$", "", product_text).strip("_-/ ")
    unit = extract_pack_unit(product_text)
    if pd.isna(unit):
        unit = extract_pack_unit(product_spec)
    return base_pack_label(unit) if pd.notna(unit) and float(unit) > 0 else ""


def build_daily_wms_catalog(inventory_df: pd.DataFrame) -> tuple[pd.DataFrame, pd.DataFrame]:
    exact_columns = ["제품코드", "POWER", "CP", "_wms_product_name", "_wms_pack", "_wms_stock"]
    prefix_columns = ["제품코드", "_wms_product_name", "_wms_pack"]
    if inventory_df is None or inventory_df.empty:
        return pd.DataFrame(columns=exact_columns), pd.DataFrame(columns=prefix_columns)

    work = inventory_df.copy()
    work["제품코드"] = work["sales_code"].map(extract_sales_prefix)
    work["POWER"] = work["sales_code"].map(inventory_power_label_from_sales_code)
    work["CP"] = work["sales_code"].map(cp_label_from_sales_code)
    work["_wms_product_name"] = work["inventory_product_name"].map(clean_inventory_product_name)
    work["_wms_pack"] = [
        pack_label_from_inventory_name(product_name, product_spec)
        for product_name, product_spec in zip(
            work.get("inventory_product_name", pd.Series("", index=work.index)),
            work.get("inventory_product_spec", pd.Series("", index=work.index)),
        )
    ]
    work["_wms_stock"] = pd.to_numeric(work.get("available_stock_pack", 0.0), errors="coerce")
    work = work[work["제품코드"].map(clean_str) != ""].copy()
    if work.empty:
        return pd.DataFrame(columns=exact_columns), pd.DataFrame(columns=prefix_columns)

    exact = (
        work[work["POWER"].map(clean_str) != ""]
        .groupby(["제품코드", "POWER", "CP"], dropna=False)
        .agg(
            _wms_product_name=("_wms_product_name", first_nonempty),
            _wms_pack=("_wms_pack", first_nonempty),
            _wms_stock=("_wms_stock", "sum"),
        )
        .reset_index()
    )
    prefix = (
        work.groupby("제품코드", dropna=False)
        .agg(
            _wms_product_name=("_wms_product_name", first_nonempty),
            _wms_pack=("_wms_pack", first_nonempty),
        )
        .reset_index()
    )
    return exact[exact_columns].copy(), prefix[prefix_columns].copy()


def fill_daily_product_code_from_wms(out: pd.DataFrame, inventory_df: pd.DataFrame) -> pd.DataFrame:
    if out.empty or inventory_df is None or inventory_df.empty:
        return out

    work = inventory_df.copy()
    work["제품코드"] = work["sales_code"].map(extract_sales_prefix)
    work["POWER"] = work["sales_code"].map(inventory_power_label_from_sales_code)
    work["CP"] = work["sales_code"].map(cp_label_from_sales_code)
    work["PACK"] = [
        pack_label_from_inventory_name(product_name, product_spec)
        for product_name, product_spec in zip(
            work.get("inventory_product_name", pd.Series("", index=work.index)),
            work.get("inventory_product_spec", pd.Series("", index=work.index)),
        )
    ]
    work["_wms_product_name"] = work["inventory_product_name"].map(clean_inventory_product_name)
    work = work[
        (work["제품코드"].map(clean_str) != "")
        & (work["POWER"].map(clean_str) != "")
        & (work["PACK"].map(clean_str) != "")
        & (work["_wms_product_name"].map(clean_str) != "")
    ].copy()
    if work.empty:
        return out

    catalog = (
        work.groupby(["제품코드", "PACK", "POWER", "CP"], dropna=False)
        .agg(_wms_product_name=("_wms_product_name", first_nonempty))
        .reset_index()
    )
    catalog_groups: dict[tuple[str, str, str], list[tuple[str, str, str]]] = {}
    for code, pack, power, cp, product_name in catalog[["제품코드", "PACK", "POWER", "CP", "_wms_product_name"]].itertuples(index=False, name=None):
        code_text = clean_str(code)
        pack_text = clean_str(pack)
        power_text = clean_str(power)
        cp_text = clean_str(cp)
        product_text = clean_str(product_name)
        catalog_groups.setdefault((pack_text, power_text, cp_text), []).append((code_text, product_text, product_text.lower()))

    filled = out.copy()
    needs_code = filled["제품코드"].map(clean_str) == ""
    for idx, row in filled[needs_code].iterrows():
        pack = clean_str(row.get("PACK", ""))
        power = clean_str(row.get("POWER", ""))
        cp = clean_str(row.get("CP", ""))
        terms = expand_product_query_terms(row.get("제품명", ""))
        if not pack or not power or not terms:
            continue
        candidates = catalog_groups.get((pack, power, cp), [])
        if not candidates:
            continue
        lowered_terms = [term.lower() for term in terms if clean_str(term)]
        matched = [
            (code, product_name)
            for code, product_name, product_name_lower in candidates
            if any(term in product_name_lower for term in lowered_terms)
        ]
        product_codes = list(dict.fromkeys([code for code, _product_name in matched if clean_str(code)]))
        if len(product_codes) != 1:
            continue
        filled.at[idx, "제품코드"] = product_codes[0]
        product_name = first_nonempty([name for _code, name in matched])
        if product_name and clean_str(filled.at[idx, "제품명"]) == "":
            filled.at[idx, "제품명"] = product_name
    return filled


def enrich_daily_inventory_from_wms(daily_inventory_df: pd.DataFrame, inventory_df: pd.DataFrame) -> pd.DataFrame:
    if daily_inventory_df is None or daily_inventory_df.empty or inventory_df is None or inventory_df.empty:
        return daily_inventory_df

    out = daily_inventory_df.copy()
    out["제품코드"] = out["제품코드"].map(clean_str).str.upper()
    out["POWER"] = out["POWER"].map(clean_str)
    out["CP"] = out["CP"].map(clean_str) if "CP" in out.columns else ""
    exact_catalog, prefix_catalog = build_daily_wms_catalog(inventory_df)

    if not exact_catalog.empty:
        out = out.merge(exact_catalog, on=["제품코드", "POWER", "CP"], how="left")
        out["제품명"] = out["제품명"].where(out["제품명"].map(clean_str) != "", out["_wms_product_name"])
        out["PACK"] = out["PACK"].where(out["PACK"].map(clean_str) != "", out["_wms_pack"])
        out["재고수량"] = pd.to_numeric(out["재고수량"], errors="coerce")
        out["재고수량"] = out["재고수량"].where(out["재고수량"].notna(), out["_wms_stock"])
        out = out.drop(columns=["_wms_product_name", "_wms_pack", "_wms_stock"], errors="ignore")

    if not prefix_catalog.empty:
        out = out.merge(prefix_catalog, on="제품코드", how="left")
        out["제품명"] = out["제품명"].where(out["제품명"].map(clean_str) != "", out["_wms_product_name"])
        out["PACK"] = out["PACK"].where(out["PACK"].map(clean_str) != "", out["_wms_pack"])
        out = out.drop(columns=["_wms_product_name", "_wms_pack"], errors="ignore")

    out = fill_daily_product_code_from_wms(out, inventory_df)
    return out[DAILY_INVENTORY_COLUMNS].copy()


def empty_progress_df() -> pd.DataFrame:
    columns = [
        "site_code",
        "customer_name",
        "order_no",
        "initial",
        "product_code",
        "demand_product_name",
        "demand_qty",
        "total_prod_qty",
        "total_due_date",
        "production_basis_qty",
        "product_code_key",
        "product_base_p_key",
        "demand_product_name_key",
        "linked_product_name",
        "match_source",
    ]
    for step in PROCESS_STEPS:
        columns.extend([step["qty_col"], step["due_col"]])
    return pd.DataFrame(columns=columns)

def find_progress_column_index(groups: pd.Series, fields: pd.Series, group_label: str, field_label: str) -> int | None:
    target_group = normalize_match_key(group_label)
    target_field = normalize_match_key(field_label)
    for idx in range(len(fields)):
        if normalize_match_key(groups.iloc[idx]) == target_group and normalize_match_key(fields.iloc[idx]) == target_field:
            return idx
    return None


def find_progress_base_column_index(fields: pd.Series, field_label: str) -> int | None:
    target_field = normalize_match_key(field_label)
    for idx in range(len(fields)):
        if normalize_match_key(fields.iloc[idx]) == target_field:
            return idx
    return None


def find_progress_due_column_index(groups: pd.Series, fields: pd.Series, group_label: str) -> int | None:
    for field_label in ["생산 계획일", "생산계획일", "계획일", "생산 계획일자", "생산계획일자", "계획일자", "납기일"]:
        idx = find_progress_column_index(groups, fields, group_label, field_label)
        if idx is not None:
            return idx
    return None


def normalize_progress(path: Path | None, request_df: pd.DataFrame) -> tuple[pd.DataFrame, dict[str, int]]:
    if path is None:
        return empty_progress_df(), {"total_rows": 0, "domestic_rows": 0, "code_rows": 0, "name_rows": 0}

    raw = pd.read_excel(path, sheet_name="Sheet1", header=None)
    if raw.shape[0] < 3:
        raise DashboardConfigError([f"[{path.name}] Sheet1에 처리할 데이터가 없습니다."])

    groups = raw.iloc[0]
    fields = raw.iloc[1]
    data = raw.iloc[2:].copy()

    base_indices = {
        "site_code": find_progress_base_column_index(fields, "설비 사이트 코드"),
        "customer_name": find_progress_base_column_index(fields, "고객 이름"),
        "order_no": find_progress_base_column_index(fields, "수주번호"),
        "initial": find_progress_base_column_index(fields, "이니셜"),
        "product_code": find_progress_base_column_index(fields, "제품 코드"),
        "demand_product_name": find_progress_base_column_index(fields, "수요 제품 이름"),
        "demand_qty": find_progress_base_column_index(fields, "수요 수량"),
    }
    missing = [name for name, idx in base_indices.items() if idx is None and name in {"product_code", "demand_product_name"}]
    if missing:
        raise DashboardConfigError([f"[{path.name}] 수요정보 필수 컬럼 누락: {', '.join(missing)}"])

    out = pd.DataFrame(index=data.index)
    for name, idx in base_indices.items():
        if idx is None:
            out[name] = "" if name != "demand_qty" else 0.0
            continue
        if name == "demand_qty":
            out[name] = to_number(data.iloc[:, idx])
        else:
            out[name] = data.iloc[:, idx].map(clean_str)

    for step in PROCESS_STEPS:
        qty_idx = find_progress_column_index(groups, fields, str(step["header"]), "생산 수량")
        due_idx = find_progress_due_column_index(groups, fields, str(step["header"]))
        out[step["qty_col"]] = to_number(data.iloc[:, qty_idx]) if qty_idx is not None else 0.0
        out[step["due_col"]] = pd.to_datetime(data.iloc[:, due_idx], errors="coerce") if due_idx is not None else pd.NaT

    total_qty_idx = find_progress_column_index(groups, fields, "총합계", "생산 수량")
    total_due_idx = find_progress_due_column_index(groups, fields, "총합계")
    out["total_prod_qty"] = to_number(data.iloc[:, total_qty_idx]) if total_qty_idx is not None else 0.0
    out["total_due_date"] = pd.to_datetime(data.iloc[:, total_due_idx], errors="coerce") if total_due_idx is not None else pd.NaT

    inspection_step = next(step for step in PROCESS_STEPS if step["id"] == "80")
    out["production_basis_qty"] = out[str(inspection_step["qty_col"])]

    out["product_code_key"] = out["product_code"].map(normalize_match_key)
    out["product_base_p_key"] = out["product_code"].map(extract_base_p_code_key)
    out["demand_product_name_key"] = out["demand_product_name"].map(normalize_match_key)

    request_production_keys = set(request_df["production_code_key"].map(clean_str)) - {""}
    request_p_keys = set(request_df["p_code_key"].map(clean_str)) - {""}
    request_name_keys = set(request_df["product_name_key"].map(clean_str)) - {""}

    exact_production_match = out["product_code_key"].isin(request_production_keys)
    exact_p_match = out["product_code_key"].isin(request_p_keys)
    base_p_match = out["product_base_p_key"].isin(request_p_keys)
    name_match = out["demand_product_name_key"].isin(request_name_keys)
    code_match = exact_production_match | exact_p_match | base_p_match
    domestic_match = code_match | name_match

    production_name_map = build_first_value_map(request_df, "production_code_key", "product_name")
    p_name_map = build_first_value_map(request_df, "p_code_key", "product_name")
    request_name_map = build_first_value_map(request_df, "product_name_key", "product_name")

    out["linked_product_name"] = ""
    out.loc[exact_production_match, "linked_product_name"] = out.loc[exact_production_match, "product_code_key"].map(
        production_name_map
    )
    p_link = exact_p_match & (out["linked_product_name"] == "")
    out.loc[p_link, "linked_product_name"] = out.loc[p_link, "product_code_key"].map(p_name_map)
    base_p_link = base_p_match & (out["linked_product_name"] == "")
    out.loc[base_p_link, "linked_product_name"] = out.loc[base_p_link, "product_base_p_key"].map(p_name_map)
    name_link = name_match & (out["linked_product_name"] == "")
    out.loc[name_link, "linked_product_name"] = out.loc[name_link, "demand_product_name_key"].map(request_name_map)
    out["linked_product_name"] = out["linked_product_name"].fillna("")
    empty_link = out["linked_product_name"] == ""
    out.loc[empty_link, "linked_product_name"] = out.loc[empty_link, "demand_product_name"]

    out["match_source"] = ""
    out.loc[name_match, "match_source"] = "제품명"
    out.loc[base_p_match, "match_source"] = "생산코드"
    out.loc[exact_p_match, "match_source"] = "생산코드"
    out.loc[exact_production_match, "match_source"] = "생산코드"

    filtered = out[domestic_match].copy()
    info = {
        "total_rows": int(len(out)),
        "domestic_rows": int(len(filtered)),
        "code_rows": int(code_match.sum()),
        "name_rows": int((~code_match & name_match).sum()),
    }
    return filtered, info


def format_date(value: Any) -> str:
    date = pd.to_datetime(value, errors="coerce")
    if pd.isna(date):
        return ""
    return date.strftime("%Y-%m-%d")


def parse_datetime_series(series: pd.Series) -> pd.Series:
    text = series.astype(str).str.strip()
    korean_ampm = text.str.replace("오전", "AM", regex=False).str.replace("오후", "PM", regex=False)
    parsed = pd.to_datetime(korean_ampm, format="%Y-%m-%d %p %I:%M:%S", errors="coerce")
    missing = parsed.isna()
    if missing.any():
        parsed.loc[missing] = pd.to_datetime(series.loc[missing], errors="coerce")
    return parsed


def summarize_progress(progress_df: pd.DataFrame, group_cols: list[str]) -> pd.DataFrame:
    if progress_df.empty:
        base_columns = group_cols + ["누수규격검사 생산수량"]
        return pd.DataFrame(columns=base_columns)

    agg_spec: dict[str, Any] = {
        "누수규격검사 생산수량": ("production_basis_qty", "sum"),
    }
    return progress_df.groupby(group_cols, dropna=False).agg(**agg_spec).reset_index()


def filter_progress_for_production_month(
    progress_df: pd.DataFrame,
    target_month: str = PRODUCTION_PROGRESS_DUE_MONTH,
) -> pd.DataFrame:
    if progress_df.empty:
        return progress_df.copy()

    inspection_step = next(step for step in PROCESS_STEPS if step["id"] == "80")
    inspection_due_col = str(inspection_step["due_col"])
    inspection_due_source = (
        progress_df[inspection_due_col]
        if inspection_due_col in progress_df.columns
        else pd.Series(pd.NaT, index=progress_df.index)
    )
    inspection_due = pd.to_datetime(
        inspection_due_source,
        errors="coerce",
    )
    production_due = inspection_due + pd.Timedelta(days=5)
    target_period = pd.Period(target_month, freq="M")
    return progress_df.loc[production_due.dt.to_period("M") == target_period].copy()


def classify_status(packing_pack: float, packing_progress_pct: float) -> str:
    if packing_progress_pct >= 100.0:
        return "완료"
    if packing_pack > 0:
        return "진행중"
    return "미착수"


def finalize_summary(summary: pd.DataFrame) -> pd.DataFrame:
    out = summary.copy()
    if "용마입고 PACK" not in out.columns:
        out["용마입고 PACK"] = 0.0
    out["포장부족수량"] = (out["요청 PACK"] - out["포장 PACK"]).clip(lower=0.0)
    out["미입고수량"] = (out["요청 PACK"] - out["용마입고 PACK"]).clip(lower=0.0)
    out["입고대기수량"] = (out["포장 PACK"] - out["용마입고 PACK"]).clip(lower=0.0)
    raw_progress = np.where(
        out["요청 PACK"] > 0,
        out["용마입고 PACK"] / out["요청 PACK"] * 100.0,
        0.0,
    )
    packing_progress = np.where(
        out["요청 PACK"] > 0,
        out["포장 PACK"] / out["요청 PACK"] * 100.0,
        0.0,
    )
    out["용마입고율"] = np.clip(raw_progress, 0.0, 100.0)
    out["포장진도율"] = np.clip(packing_progress, 0.0, 100.0)
    out["부족 PACK"] = out["미입고수량"]
    out["진도율(%)"] = out["용마입고율"]
    out["상태"] = [
        classify_status(float(packing), float(progress))
        for packing, progress in zip(out["포장 PACK"], out["포장진도율"])
    ]
    return out


def add_code_level_supply_basis(df: pd.DataFrame) -> pd.DataFrame:
    out = df.copy()
    index = out.index
    request = pd.to_numeric(out.get("request_pack", pd.Series(0.0, index=index)), errors="coerce").fillna(0.0)
    packing = pd.to_numeric(out.get("packing_pack", pd.Series(0.0, index=index)), errors="coerce").fillna(0.0)
    yongma = pd.to_numeric(out.get("yongma_in_pack", pd.Series(0.0, index=index)), errors="coerce").fillna(0.0)

    request = request.clip(lower=0.0)
    packing = packing.clip(lower=0.0)
    yongma = yongma.clip(lower=0.0)
    has_request = request > 0

    out["packing_recognized_pack"] = np.where(has_request, np.minimum(packing, request), 0.0)
    out["yongma_recognized_pack"] = np.where(has_request, np.minimum(yongma, request), 0.0)
    out["packing_over_pack"] = np.where(has_request, np.maximum(packing - request, 0.0), packing)
    out["yongma_over_pack"] = np.where(has_request, np.maximum(yongma - request, 0.0), yongma)
    out["code_packing_shortage_pack"] = np.where(
        has_request,
        np.maximum(request - out["packing_recognized_pack"], 0.0),
        0.0,
    )
    out["code_receipt_shortage_pack"] = np.where(
        has_request,
        np.maximum(request - out["yongma_recognized_pack"], 0.0),
        0.0,
    )
    out["code_receipt_wait_pack"] = np.maximum(
        pd.to_numeric(out["packing_recognized_pack"], errors="coerce").fillna(0.0)
        - pd.to_numeric(out["yongma_recognized_pack"], errors="coerce").fillna(0.0),
        0.0,
    )
    return out


def pcs_per_pack_for_rows(df: pd.DataFrame) -> pd.Series:
    index = df.index
    pack_unit = pd.to_numeric(df.get("pack_unit", pd.Series(np.nan, index=index)), errors="coerce")
    request_pack = pd.to_numeric(df.get("request_pack", pd.Series(0.0, index=index)), errors="coerce").fillna(0.0)
    request_pcs = pd.to_numeric(df.get("request_pcs", pd.Series(0.0, index=index)), errors="coerce").fillna(0.0)
    implied_unit = np.where(request_pack > 0, request_pcs / request_pack, np.nan)
    unit = pack_unit.where(pack_unit > 0, implied_unit)
    unit = pd.Series(unit, index=index).replace([np.inf, -np.inf], np.nan).fillna(1.0)
    return unit.where(unit > 0, 1.0)


def pack_quantity_to_pcs(df: pd.DataFrame, pack_qty: Any) -> pd.Series:
    index = df.index
    if isinstance(pack_qty, pd.Series):
        pack_source = pack_qty.reindex(index)
    else:
        pack_source = pd.Series(pack_qty, index=index)
    pack = pd.to_numeric(pack_source, errors="coerce").fillna(0.0)
    return (pack.clip(lower=0.0) * pcs_per_pack_for_rows(df)).clip(lower=0.0)


def recognized_packing_pcs(df: pd.DataFrame) -> pd.Series:
    pack = pd.to_numeric(
        df.get("packing_recognized_pack", df.get("packing_pack", pd.Series(0.0, index=df.index))),
        errors="coerce",
    ).fillna(0.0)
    request_pcs = pd.to_numeric(df.get("request_pcs", pd.Series(0.0, index=df.index)), errors="coerce").fillna(0.0)
    packing_pcs = pack_quantity_to_pcs(df, pack)
    return packing_pcs.where(request_pcs <= 0, packing_pcs.clip(upper=request_pcs))


def calc_production_progress_pct(request_qty: Any, production_shortage_qty: Any) -> Any:
    request = pd.to_numeric(request_qty, errors="coerce").fillna(0.0)
    shortage = pd.to_numeric(production_shortage_qty, errors="coerce").fillna(0.0)
    produced = (request - shortage).clip(lower=0.0)
    produced = produced.where(produced <= request, request)
    return np.where(request > 0, produced / request * 100.0, 0.0)


def build_summaries(
    request_df: pd.DataFrame,
    packing_df: pd.DataFrame,
    yongma_df: pd.DataFrame | None = None,
    product_master_df: pd.DataFrame | None = None,
) -> tuple[pd.DataFrame, float, pd.DataFrame]:
    request_work = request_df.copy()
    optional_cols = [
        "product_name_code",
        "production_code",
        "p_code",
        "q_code",
        "r_code",
        "request_due_date",
        "request_pcs",
        "pack_unit",
        "pack_unit_label",
        "base_product_name",
        "customer_name",
        "category_summary",
        "factory_group",
        "sales_code_key",
        "product_name_key",
        "product_name_code_key",
        "production_code_key",
        "p_code_key",
        "q_code_key",
        "r_code_key",
    ]
    for col in optional_cols:
        if col not in request_work.columns:
            if col == "request_due_date":
                request_work[col] = pd.NaT
            elif col == "request_pcs":
                request_work[col] = request_work["request_pack"]
            elif col == "pack_unit":
                request_work[col] = np.nan
            elif col == "pack_unit_label":
                request_work[col] = "(미기재)"
            elif col == "base_product_name":
                request_work[col] = request_work["product_name"].map(strip_pack_unit_suffix)
            elif col == "customer_name":
                request_work[col] = "(미기재)"
            elif col == "category_summary":
                request_work[col] = "(미기재)"
            elif col == "factory_group":
                request_work[col] = request_work.get("category_summary", pd.Series("", index=request_work.index)).map(factory_group_from_category)
            else:
                request_work[col] = ""

    request_work["_sales_prefix"] = request_work["sales_code"].map(extract_sales_prefix)
    prefix_meta = (
        request_work[request_work["_sales_prefix"].map(clean_str) != ""]
        .groupby("_sales_prefix", dropna=False)
        .agg(
            product_name=("product_name", first_nonempty),
            product_name_code=("product_name_code", first_nonempty),
            p_code=("p_code", first_nonempty),
            production_code=("production_code", first_nonempty),
            q_code=("q_code", first_nonempty),
            r_code=("r_code", first_nonempty),
            category_summary=("category_summary", first_nonempty),
            factory_group=("factory_group", join_factory_groups),
            sales_code=("sales_code", first_nonempty),
        )
        .to_dict(orient="index")
        if not request_work.empty
        else {}
    )
    product_master_by_key: dict[str, dict[str, Any]] = {}
    if product_master_df is not None and not product_master_df.empty and "sales_code_key" in product_master_df.columns:
        product_master_by_key = (
            product_master_df.drop_duplicates("sales_code_key", keep="first")
            .set_index("sales_code_key")
            .to_dict(orient="index")
        )

    group_cols = [
        "sales_code",
        "product_name",
        "product_name_code",
        "production_code",
        "p_code",
        "q_code",
        "r_code",
        "pack_unit",
        "pack_unit_label",
        "base_product_name",
        "customer_name",
        "category_summary",
        "factory_group",
        "sales_code_key",
        "product_name_key",
        "product_name_code_key",
        "production_code_key",
        "p_code_key",
        "q_code_key",
        "r_code_key",
    ]
    request_by_code = (
        request_work.groupby(group_cols, dropna=False)
        .agg(
            request_pack=("request_pack", "sum"),
            request_pcs=("request_pcs", "sum"),
            request_due_date=("request_due_date", min_datetime),
        )
        .reset_index()
    )
    if packing_df.empty:
        packing_by_key = pd.DataFrame(
            columns=["sales_code_key", "packing_sales_code", "packing_product_name", "packing_pack"]
        )
    else:
        packing_by_key = (
            packing_df.copy()
            .assign(sales_code_key=lambda df: df["sales_code_key"].map(clean_str))
            .groupby("sales_code_key", dropna=False)
            .agg(
                packing_sales_code=("sales_code", first_nonempty),
                packing_product_name=("packing_product_name", first_nonempty),
                packing_pack=("packing_pack", "sum"),
            )
            .reset_index()
        )
    if yongma_df is None or yongma_df.empty:
        yongma_by_key = pd.DataFrame(
            columns=["sales_code_key", "yongma_sales_code", "yongma_product_name", "yongma_in_pack"]
        )
    else:
        yongma_by_key = (
            yongma_df.copy()
            .assign(sales_code_key=lambda df: df["sales_code_key"].map(clean_str))
            .groupby("sales_code_key", dropna=False)
            .agg(
                yongma_sales_code=("sales_code", first_nonempty),
                yongma_product_name=("yongma_product_name", first_nonempty),
                yongma_in_pack=("yongma_in_pack", "sum"),
            )
            .reset_index()
        )

    matched_code_summary = request_by_code.merge(
        packing_by_key[["sales_code_key", "packing_pack"]],
        on="sales_code_key",
        how="left",
    )
    matched_code_summary["packing_pack"] = matched_code_summary["packing_pack"].fillna(0.0)
    matched_code_summary = matched_code_summary.merge(
        yongma_by_key[["sales_code_key", "yongma_in_pack"]],
        on="sales_code_key",
        how="left",
    )
    matched_code_summary["yongma_in_pack"] = matched_code_summary["yongma_in_pack"].fillna(0.0)

    request_keys = set(request_by_code["sales_code_key"].map(clean_str)) - {""}
    supply_by_key = packing_by_key.merge(yongma_by_key, on="sales_code_key", how="outer")
    for col in ["packing_pack", "yongma_in_pack"]:
        if col not in supply_by_key.columns:
            supply_by_key[col] = 0.0
        supply_by_key[col] = pd.to_numeric(supply_by_key[col], errors="coerce").fillna(0.0)
    unmatched_supply = supply_by_key[
        (supply_by_key["sales_code_key"].map(clean_str) != "")
        & ~supply_by_key["sales_code_key"].map(clean_str).isin(request_keys)
        & ((supply_by_key["packing_pack"] > 0) | (supply_by_key["yongma_in_pack"] > 0))
    ].copy()
    unmatched_packing_total = float(unmatched_supply["packing_pack"].sum()) if not unmatched_supply.empty else 0.0

    if not unmatched_supply.empty:
        unmatched_rows: list[dict[str, Any]] = []
        for _, row in unmatched_supply.iterrows():
            sales_code = clean_str(row.get("packing_sales_code", "")) or clean_str(row.get("yongma_sales_code", ""))
            sales_code_key = clean_str(row.get("sales_code_key", ""))
            sales_prefix = extract_sales_prefix(sales_code)
            prefix_values = prefix_meta.get(sales_prefix, {})
            master_values = product_master_by_key.get(sales_code_key, {})
            source_sales_code = clean_str(prefix_values.get("sales_code", ""))
            target_power = format_power(parse_power_from_sales_code(sales_code))
            source_power = format_power(parse_power_from_sales_code(source_sales_code))
            product_name = (
                clean_str(master_values.get("master_product_name", ""))
                or clean_str(prefix_values.get("product_name", ""))
                or clean_str(row.get("packing_product_name", ""))
                or clean_str(row.get("yongma_product_name", ""))
                or sales_code
            )
            production_code = clean_str(master_values.get("master_production_code", ""))
            if not production_code:
                production_code = replace_power_in_production_code(
                    prefix_values.get("production_code", ""),
                    source_power,
                    target_power,
                )
            q_code = clean_str(master_values.get("master_q_code", ""))
            if not q_code:
                q_code = replace_power_in_production_code(prefix_values.get("q_code", ""), source_power, target_power)
            r_code = clean_str(master_values.get("master_r_code", ""))
            if not r_code:
                r_code = replace_power_in_production_code(prefix_values.get("r_code", ""), source_power, target_power)
            p_code = clean_str(master_values.get("master_p_code", "")) or clean_str(prefix_values.get("p_code", ""))
            product_name_code = clean_str(prefix_values.get("product_name_code", "")) or product_name
            category_summary = clean_str(prefix_values.get("category_summary", "")) or "(포장실적)"
            factory_group = clean_factory_group_display(prefix_values.get("factory_group", ""))
            if not has_factory_group(factory_group):
                factory_group = clean_factory_group_display(factory_group_from_category(category_summary))
            if not has_factory_group(factory_group):
                factory_group = clean_factory_group_display(factory_group_from_product_name(product_name))
            pack_unit = extract_pack_unit(product_name)
            unmatched_rows.append(
                {
                    "sales_code": sales_code,
                    "product_name": product_name,
                    "product_name_code": product_name_code,
                    "production_code": production_code,
                    "p_code": p_code,
                    "q_code": q_code,
                    "r_code": r_code,
                    "pack_unit": pack_unit,
                    "pack_unit_label": format_pack_unit_label(pack_unit, product_name),
                    "base_product_name": strip_pack_unit_suffix(product_name),
                    "customer_name": "(포장실적)",
                    "category_summary": category_summary,
                    "factory_group": factory_group,
                    "sales_code_key": sales_code_key,
                    "product_name_key": normalize_match_key(product_name),
                    "product_name_code_key": normalize_match_key(product_name_code),
                    "production_code_key": normalize_match_key(production_code),
                    "p_code_key": normalize_match_key(p_code),
                    "q_code_key": normalize_match_key(q_code),
                    "r_code_key": normalize_match_key(r_code),
                    "request_pack": 0.0,
                    "request_pcs": 0.0,
                    "request_due_date": pd.NaT,
                    "packing_pack": float(row.get("packing_pack", 0.0)),
                    "yongma_in_pack": float(row.get("yongma_in_pack", 0.0)),
                }
            )
        matched_code_summary = pd.concat([matched_code_summary, pd.DataFrame(unmatched_rows)], ignore_index=True)

    matched_code_summary = align_base_product_names_by_production_family(matched_code_summary)
    matched_code_summary = add_code_level_supply_basis(matched_code_summary)

    product_summary = (
        matched_code_summary.groupby("base_product_name", dropna=False)[
            ["request_pack", "request_pcs", "packing_recognized_pack", "yongma_recognized_pack"]
        ]
        .sum()
        .reset_index()
        .rename(
            columns={
                "base_product_name": "제품명",
                "request_pack": "요청 PACK",
                "request_pcs": "요청 PCS",
                "packing_recognized_pack": "포장 PACK",
                "yongma_recognized_pack": "용마입고 PACK",
            }
        )
    )
    product_summary = finalize_summary(product_summary)
    product_summary["제품분류"] = product_summary["제품명"].map(classify_product_group)
    product_summary["본품분류"] = product_summary["제품명"].map(classify_main_product_family)
    product_summary = add_period_group_columns(product_summary)
    return product_summary, unmatched_packing_total, matched_code_summary


def attach_inventory_to_code_summary(code_summary: pd.DataFrame, inventory_df: pd.DataFrame) -> pd.DataFrame:
    out = code_summary.copy()
    if inventory_df.empty:
        out["available_stock_pack"] = np.nan
        out["inventory_total_stock_pack"] = np.nan
        out["inventory_product_name"] = ""
        out["inventory_product_spec"] = ""
        out["inventory_updated_at"] = pd.NaT
        out["inventory_matched"] = False
        return out

    inventory_by_code = (
        inventory_df.groupby("sales_code_key", dropna=False)
        .agg(
            available_stock_pack=("available_stock_pack", sum_numeric_or_nan),
            inventory_total_stock_pack=("inventory_total_stock_pack", sum_numeric_or_nan),
            inventory_product_name=("inventory_product_name", join_unique),
            inventory_product_spec=("inventory_product_spec", first_nonempty),
            inventory_updated_at=("inventory_updated_at", max_datetime),
        )
        .reset_index()
    )
    out = out.merge(inventory_by_code, on="sales_code_key", how="left")
    out["inventory_matched"] = out["available_stock_pack"].notna()
    return out


def attach_inventory_to_product_summary(product_summary: pd.DataFrame, code_summary: pd.DataFrame) -> pd.DataFrame:
    out = product_summary.copy()
    if code_summary.empty or "available_stock_pack" not in code_summary.columns:
        out["용마창고재고 (PACK)"] = np.nan
        out["재고매칭SKU수"] = 0
        return out

    work = with_operational_columns(code_summary)
    stock_by_product = (
        work.groupby("base_product_name", dropna=False)
        .agg(
            current_stock_pack=("available_stock_pack", sum_numeric_or_nan),
            inventory_matched_count=("inventory_matched", "sum"),
        )
        .reset_index()
        .rename(
            columns={
                "base_product_name": "제품명",
                "current_stock_pack": "용마창고재고 (PACK)",
                "inventory_matched_count": "재고매칭SKU수",
            }
        )
    )
    out = out.merge(stock_by_product, on="제품명", how="left")
    out["재고매칭SKU수"] = out["재고매칭SKU수"].fillna(0).astype(int)
    return out


def enrich_product_summary(product_summary: pd.DataFrame, progress_df: pd.DataFrame) -> pd.DataFrame:
    progress_work = progress_df.copy()
    if "linked_product_name" not in progress_work.columns:
        progress_work["linked_product_name"] = ""
    progress_work["linked_base_product_name"] = progress_work["linked_product_name"].map(strip_pack_unit_suffix)
    progress_by_product = summarize_progress(progress_work, ["linked_base_product_name"]).rename(
        columns={
            "linked_base_product_name": "제품명",
        }
    )
    out = product_summary.merge(
        progress_by_product[["제품명", "누수규격검사 생산수량"]],
        on="제품명",
        how="left",
    )
    out["누수규격검사 생산수량"] = out["누수규격검사 생산수량"].fillna(0.0)
    out["생산부족수량"] = out["누수규격검사 생산수량"].clip(lower=0.0)
    out["생산진도율"] = calc_production_progress_pct(out["요청 PCS"], out["생산부족수량"])
    return out


def enrich_product_summary_from_code_summary(product_summary: pd.DataFrame, code_summary: pd.DataFrame) -> pd.DataFrame:
    out = product_summary.copy()
    out = out.drop(columns=["누수규격검사 생산수량", "생산부족수량", "생산진도율"], errors="ignore")
    if code_summary.empty or "base_product_name" not in code_summary.columns:
        out["누수규격검사 생산수량"] = 0.0
        out["생산부족수량"] = 0.0
        out["생산진도율"] = calc_production_progress_pct(out["요청 PCS"], out["생산부족수량"])
        return out

    work = add_allocated_production_basis(code_summary.copy())
    if "production_basis_qty" not in work.columns:
        work["production_basis_qty"] = 0.0
    if "production_shortage_qty" not in work.columns:
        work["production_shortage_qty"] = work["production_basis_qty"]
    work["production_basis_qty"] = pd.to_numeric(work["production_basis_qty"], errors="coerce").fillna(0.0)
    work["production_shortage_qty"] = pd.to_numeric(work["production_shortage_qty"], errors="coerce").fillna(0.0)
    progress_by_product = (
        work.groupby("base_product_name", dropna=False)
        .agg(
            production_basis_qty=("production_basis_qty", "sum"),
            production_shortage_qty=("_allocated_production_shortage_qty", "sum"),
        )
        .reset_index()
        .rename(
            columns={
                "base_product_name": "제품명",
                "production_basis_qty": "누수규격검사 생산수량",
                "production_shortage_qty": "생산부족수량",
            }
        )
    )
    out = out.merge(progress_by_product, on="제품명", how="left")
    out["누수규격검사 생산수량"] = pd.to_numeric(
        out["누수규격검사 생산수량"],
        errors="coerce",
    ).fillna(0.0)
    out["생산부족수량"] = pd.to_numeric(out["생산부족수량"], errors="coerce").fillna(0.0).clip(lower=0.0)
    out["생산진도율"] = calc_production_progress_pct(out["요청 PCS"], out["생산부족수량"])
    return out


def progress_row_pack_unit(progress_row: pd.Series) -> float:
    for col in ["demand_product_name", "linked_product_name"]:
        unit = extract_pack_unit(progress_row.get(col, ""))
        if pd.notna(unit) and float(unit) > 0:
            return float(unit)
    return np.nan


def candidate_pack_units(candidates: pd.DataFrame) -> pd.Series:
    if "pack_unit" in candidates.columns:
        units = pd.to_numeric(candidates["pack_unit"], errors="coerce")
    else:
        units = pd.Series(np.nan, index=candidates.index)
    fallback = candidates.get("product_name", pd.Series("", index=candidates.index)).map(extract_pack_unit)
    return units.where(units.notna() & (units > 0), fallback)


def build_progress_source_key(progress_key: str, progress_index: Any) -> str:
    base = clean_str(progress_key) or "PROGRESS"
    row_key = normalize_match_key(progress_index)
    return f"{base}#ROW{row_key}" if row_key else base


def refine_progress_candidates_by_code_measure(candidates: pd.DataFrame, progress_row: pd.Series) -> pd.DataFrame:
    if candidates.empty or len(candidates) <= 1:
        return candidates

    narrowed = candidates
    progress_measure_key = extract_code_measure_key(progress_row.get("product_code", ""))
    if progress_measure_key and "_code_measure_key" in candidates.columns:
        narrowed_by_code = candidates[candidates["_code_measure_key"] == progress_measure_key].copy()
        if not narrowed_by_code.empty:
            narrowed = narrowed_by_code

    progress_pack_unit = progress_row_pack_unit(progress_row)
    if pd.notna(progress_pack_unit) and float(progress_pack_unit) > 0:
        pack_units = candidate_pack_units(narrowed)
        narrowed_by_pack = narrowed[np.isclose(pack_units, float(progress_pack_unit), rtol=0.0, atol=1e-6)].copy()
        if not narrowed_by_pack.empty:
            narrowed = narrowed_by_pack

    demand_qty = to_number_value(progress_row.get("demand_qty", np.nan))
    if demand_qty > 0 and "request_pack" in narrowed.columns:
        request_pack = pd.to_numeric(narrowed["request_pack"], errors="coerce").fillna(-1.0)
        narrowed_by_qty = narrowed[np.isclose(request_pack, demand_qty, rtol=0.0, atol=1e-6)].copy()
        if not narrowed_by_qty.empty:
            narrowed = narrowed_by_qty

    return narrowed


def build_progress_by_sales_code(code_summary: pd.DataFrame, progress_work: pd.DataFrame) -> pd.DataFrame:
    columns = [
        "sales_code_key",
        "production_basis_qty",
        "production_due_date",
        "production_plan_date",
        "production_complete_expected_date",
        "production_match_source",
        "production_progress_key",
    ]
    if code_summary.empty or progress_work.empty:
        return pd.DataFrame(columns=columns)

    candidates = code_summary.copy()
    required_cols = ["sales_code_key", "production_code_key", "p_code_key", "product_name_key", "product_name_code_key"]
    for col in required_cols:
        if col not in candidates.columns:
            candidates[col] = ""
        candidates[col] = candidates[col].map(clean_str)
    if "base_product_name" not in candidates.columns:
        candidates["base_product_name"] = candidates.get("product_name", pd.Series("", index=candidates.index)).map(
            strip_pack_unit_suffix
        )
    candidates["base_product_name_key"] = candidates["base_product_name"].map(normalize_match_key)
    candidates["_code_measure_key"] = [
        first_nonempty(
            [
                extract_code_measure_key(production_code),
                extract_code_measure_key(sales_code),
            ]
        )
        for production_code, sales_code in zip(
            candidates.get("production_code", pd.Series("", index=candidates.index)),
            candidates.get("sales_code", pd.Series("", index=candidates.index)),
        )
    ]
    if "request_pcs" not in candidates.columns:
        candidates["request_pcs"] = 0.0
    candidates["request_pcs"] = pd.to_numeric(candidates["request_pcs"], errors="coerce").fillna(0.0)
    candidates = candidates[(candidates["sales_code_key"] != "") & (candidates["request_pcs"] > 0)].copy()
    if candidates.empty:
        return pd.DataFrame(columns=columns)

    toric_templates = build_toric_progress_code_template_map(candidates)
    candidates["_derived_progress_code_key"] = [
        derive_toric_progress_code_key(row, toric_templates) for _, row in candidates.iterrows()
    ]

    by_production_code = {
        key: group.copy()
        for key, group in candidates[candidates["production_code_key"] != ""].groupby("production_code_key", dropna=False)
    }
    by_derived_progress_code = {
        key: group.copy()
        for key, group in candidates[candidates["_derived_progress_code_key"] != ""].groupby(
            "_derived_progress_code_key",
            dropna=False,
        )
    }
    by_sales_code = {
        key: group.copy()
        for key, group in candidates[candidates["sales_code_key"] != ""].groupby("sales_code_key", dropna=False)
    }
    by_p_code = {
        key: group.copy()
        for key, group in candidates[candidates["p_code_key"] != ""].groupby("p_code_key", dropna=False)
    }

    records: list[dict[str, Any]] = []
    for progress_index, progress_row in progress_work.iterrows():
        product_code_key = clean_str(progress_row.get("product_code_key", ""))
        product_base_p_key = clean_str(progress_row.get("product_base_p_key", ""))
        progress_key = first_nonempty([product_code_key, product_base_p_key])

        matched = by_sales_code.get(product_code_key, pd.DataFrame())
        match_source = "판매코드"
        if matched.empty:
            matched = by_production_code.get(product_code_key, pd.DataFrame())
            match_source = "생산코드"
        if matched.empty:
            matched = by_derived_progress_code.get(product_code_key, pd.DataFrame())
            match_source = "생산코드/판매코드 규칙"
        if matched.empty and product_code_key in by_p_code:
            matched = by_p_code[product_code_key]
            match_source = "생산코드"
        if matched.empty and product_base_p_key in by_p_code:
            matched = by_p_code[product_base_p_key]
            match_source = "P대표코드"
        if matched.empty:
            continue

        matched = refine_progress_candidates_by_code_measure(matched, progress_row)
        qty = to_number_value(progress_row.get("production_basis_qty", 0.0))
        if qty <= 0:
            continue
        source_progress_key = build_progress_source_key(progress_key, progress_index)

        for _, candidate_row in matched.iterrows():
            records.append(
                {
                    "sales_code_key": clean_str(candidate_row.get("sales_code_key", "")),
                    "production_basis_qty": qty,
                    "production_due_date": progress_row.get("production_due_date", pd.NaT),
                    "production_plan_date": progress_row.get("production_plan_date", pd.NaT),
                    "production_complete_expected_date": progress_row.get(
                        "production_complete_expected_date",
                        pd.NaT,
                    ),
                    "production_match_source": match_source,
                    "production_progress_key": source_progress_key,
                }
            )

    if not records:
        return pd.DataFrame(columns=columns)

    return (
        pd.DataFrame(records)
        .groupby("sales_code_key", dropna=False)
        .agg(
            production_basis_qty=("production_basis_qty", "sum"),
            production_due_date=("production_due_date", min_datetime),
            production_plan_date=("production_plan_date", min_datetime),
            production_complete_expected_date=("production_complete_expected_date", min_datetime),
            production_match_source=("production_match_source", join_unique),
            production_progress_key=("production_progress_key", join_unique),
        )
        .reset_index()[columns]
    )


def attach_progress_to_code_summary(code_summary: pd.DataFrame, progress_df: pd.DataFrame) -> pd.DataFrame:
    if progress_df.empty:
        progress_by_sales_code = pd.DataFrame(
            columns=[
                "sales_code_key",
                "production_basis_qty",
                "production_due_date",
                "production_plan_date",
                "production_complete_expected_date",
                "production_match_source",
                "production_progress_key",
            ]
        )
    else:
        progress_work = progress_df.copy()
        inspection_step = next(step for step in PROCESS_STEPS if step["id"] == "80")
        inspection_due = pd.to_datetime(progress_work.get(str(inspection_step["due_col"]), pd.NaT), errors="coerce")
        progress_work["production_plan_date"] = inspection_due
        progress_work["production_complete_expected_date"] = inspection_due + pd.Timedelta(days=5)
        progress_work["production_due_date"] = progress_work["production_complete_expected_date"]
        progress_by_sales_code = build_progress_by_sales_code(code_summary, progress_work)

    out = code_summary.merge(progress_by_sales_code, on="sales_code_key", how="left")
    out["production_basis_qty"] = out["production_basis_qty"].fillna(0.0)
    out["production_due_date"] = pd.to_datetime(out["production_due_date"], errors="coerce")
    for date_col in ["production_plan_date", "production_complete_expected_date"]:
        if date_col not in out.columns:
            out[date_col] = pd.NaT
        out[date_col] = pd.to_datetime(out[date_col], errors="coerce")
    request_pcs = pd.to_numeric(out.get("request_pcs", pd.Series(0.0, index=out.index)), errors="coerce").fillna(0.0)
    if "production_match_source" not in out.columns:
        out["production_match_source"] = ""
    out["production_match_source"] = out["production_match_source"].map(clean_str)
    if "production_progress_key" not in out.columns:
        out["production_progress_key"] = ""
    out["production_progress_key"] = out["production_progress_key"].map(clean_str)

    out["production_source_shortage_qty"] = out["production_basis_qty"].clip(lower=0.0)
    out["production_shortage_qty"] = out["production_source_shortage_qty"]
    out["production_progress_pct"] = calc_production_progress_pct(request_pcs, out["production_shortage_qty"])
    return out


def attach_sample_available_to_code_summary(code_summary: pd.DataFrame, sample_available_df: pd.DataFrame) -> pd.DataFrame:
    out = code_summary.copy()
    if "production_code_key" not in out.columns:
        out["production_code_key"] = ""
    out["sample_available_pcs"] = 0.0
    if sample_available_df.empty:
        return out

    sample_by_code = (
        sample_available_df.groupby("production_code_key", dropna=False)["sample_available_pcs"]
        .sum()
        .reset_index()
        .rename(columns={"sample_available_pcs": "_sample_available_pcs"})
    )
    out = out.merge(sample_by_code, on="production_code_key", how="left")
    out["sample_available_pcs"] = pd.to_numeric(out["_sample_available_pcs"], errors="coerce").fillna(0.0)
    return out.drop(columns=["_sample_available_pcs"])


def build_production_code_view(code_summary: pd.DataFrame) -> pd.DataFrame:
    work = add_code_level_supply_basis(code_summary)
    work["production_code"] = work["production_code"].replace("", "(생산코드 미기재)")
    if "production_shortage_qty" not in work.columns:
        work["production_shortage_qty"] = work.get("production_basis_qty", 0.0)
    grouped = (
        work.groupby("production_code", dropna=False)
        .agg(
            sales_code_count=("sales_code", "nunique"),
            product_name=("product_name", join_unique),
            request_pack=("request_pack", "sum"),
            request_pcs=("request_pcs", "sum"),
            packing_pack=("packing_recognized_pack", "sum"),
            production_basis_qty=("production_basis_qty", "max"),
            production_shortage_qty=("production_shortage_qty", "sum"),
            request_due_date=("request_due_date", min_datetime),
            production_due_date=("production_due_date", min_datetime),
        )
        .reset_index()
    )
    grouped["production_progress_pct"] = calc_production_progress_pct(
        grouped["request_pcs"],
        grouped["production_shortage_qty"],
    )
    grouped = grouped.rename(
        columns={
            "production_code": "생산코드",
            "sales_code_count": "연결 판매코드 수",
            "product_name": "제품명",
            "request_pack": "요청 PACK",
            "request_pcs": "요청 PCS",
            "packing_pack": "포장 PACK",
            "production_basis_qty": "누수규격검사 생산수량",
            "production_shortage_qty": "생산부족수량",
            "production_progress_pct": "생산진도율",
            "request_due_date": "납기일",
            "production_due_date": "생산완료예상일",
        }
    )
    grouped = finalize_summary(grouped)
    return grouped


def is_sample_name(name: str) -> bool:
    text = str(name)
    return any(keyword in text for keyword in SAMPLE_KEYWORDS)


def split_main_sample(df: pd.DataFrame) -> tuple[pd.DataFrame, pd.DataFrame]:
    sample_mask = df["제품명"].astype(str).map(is_sample_name)
    sample_df = df[sample_mask].copy()
    main_df = df[~sample_mask].copy()
    return main_df, sample_df


def parse_power_from_sales_code(value: Any) -> float:
    text = clean_str(value)
    if not text:
        return np.nan
    tail = text.split("-", 1)[1] if "-" in text else text
    match = re.search(r"(-?\d+(?:\.\d+)?)", tail)
    if match is None:
        match = POWER_RE.search(text)
    if not match:
        return np.nan
    try:
        number = float(match.group(1))
    except ValueError:
        return np.nan
    # Sales code values are typically encoded as 00.50, 01.00 ... and represent minus diopters.
    return round(-abs(number), 2)


def format_power(value: Any) -> str:
    num = pd.to_numeric(value, errors="coerce")
    if pd.isna(num):
        return "(미기재)"
    return f"-{abs(float(num)):05.2f}"


def sort_power_detail_default(
    df: pd.DataFrame,
    extra_cols: list[str] | None = None,
    extra_ascending: list[bool] | None = None,
) -> pd.DataFrame:
    if df.empty:
        return df.copy()

    out = df.copy()
    temp_cols = [
        "_power_default_missing",
        "_power_default_abs",
        "_power_default_value",
        "_power_default_label",
    ]
    if "POWER" in out.columns:
        power_label = out["POWER"].map(clean_str)
    elif "power_value" in out.columns:
        power_label = out["power_value"].map(format_power)
    elif "_power_sort" in out.columns:
        power_label = out["_power_sort"].map(format_power)
    else:
        return out

    power_number = pd.to_numeric(
        power_label.astype(str).str.replace("-00.00", "0", regex=False),
        errors="coerce",
    )
    out["_power_default_missing"] = power_number.isna().astype(int)
    out["_power_default_abs"] = power_number.abs().fillna(999999.0)
    out["_power_default_value"] = power_number.fillna(999999.0)
    out["_power_default_label"] = power_label

    sort_cols = temp_cols.copy()
    ascending = [True, True, True, True]
    for idx, col in enumerate(extra_cols or []):
        if col in out.columns and col not in sort_cols:
            sort_cols.append(col)
            if extra_ascending and idx < len(extra_ascending):
                ascending.append(extra_ascending[idx])
            else:
                ascending.append(True)

    out = out.sort_values(
        sort_cols,
        ascending=ascending,
        na_position="last",
        kind="stable",
    )
    return out.drop(columns=temp_cols, errors="ignore")


def classify_product_group(product_name: str) -> str:
    text = clean_str(product_name)
    upper = text.upper()
    if is_sample_name(text):
        return "샘플"
    if upper.startswith("PIA_") or upper.startswith("PIA ") or PIA_TOKEN_RE.search(upper):
        return "PIA"
    if "CLALEN" in upper:
        return "Clalen"
    if "TORIC" in upper or "사축" in text:
        return "Toric"
    if "1DAY" in upper or "원데이" in text:
        return "1Day"
    if "COLOR" in upper or "컬러" in text:
        return "Color"
    if "M_" in upper or " M " in upper or "먼슬리" in text or "MONTHLY" in upper:
        return "Monthly"
    return "기타"


def classify_main_product_family(product_name: str) -> str:
    text = clean_str(product_name)
    upper = text.upper()
    if is_sample_name(text):
        return "샘플"
    if upper.startswith("PIA_") or upper.startswith("PIA ") or PIA_TOKEN_RE.search(upper):
        if "_1D" in upper or " 1D" in upper:
            return "PIA 1Day"
        if "_1M" in upper or " 1M" in upper:
            return "PIA Monthly"
        return "PIA 기타"
    if "O2O2 D TORIC" in upper or "O2O2 D_TORIC" in upper or "O2O2 D TORIC_" in upper:
        return "O2O2 D Toric"
    if "IRIS TORIC" in upper:
        return "Iris Toric"
    if upper.startswith("T38") or "T38" in upper or "사축" in text or "정축" in text:
        return "T38 Toric"
    if "TORIC" in upper:
        return "기타 Toric"
    if "O2O2 1DAY" in upper:
        return "O2O2 1Day"
    if "O2O2 D_MICELIA" in upper or "O2O2 D MICELIA" in upper:
        return "O2O2 D Micelia"
    if "O2O2 D_" in upper or "O2O2 D " in upper:
        return "O2O2 D 컬러"
    if "O2O2 M_MICELIA" in upper or "O2O2 M MICELIA" in upper:
        return "O2O2 M Micelia"
    if "O2O2 M" in upper:
        return "O2O2 Monthly"
    if "CLALEN 1DAY" in upper or "CLALEN1DAY" in upper.replace(" ", ""):
        return "Clalen 1Day"
    if "CLEAR" in upper:
        return "Clear"
    if "IRIS" in upper:
        return "Iris 컬러"
    if upper.startswith("S38") or upper.startswith("S45") or upper.startswith("US38") or "BANDAGE" in upper:
        return "부자재/기타"
    return "기타"


def build_power_detail(code_summary: pd.DataFrame) -> pd.DataFrame:
    work = with_operational_columns(code_summary)
    work = work[work["power_value"].notna()].copy()
    if work.empty:
        return pd.DataFrame(
            columns=[
                "제품분류",
                "기간구분",
                "제품명",
                "POWER",
                "요청수량",
                "요청PCS",
                "포장수량",
                "부족수량",
                "진도율",
                "상태",
                "생산부족수량",
                "생산진도율",
                "power_value",
            ]
        )

    work["제품분류"] = work["product_name"].map(classify_product_group)
    if "factory_group" not in work.columns:
        work = with_operational_columns(work)
    work["factory_group"] = work["factory_group"].map(clean_str).replace("", "(미기재)").fillna("(미기재)")
    work["POWER"] = work["power_value"].map(format_power)

    grouped = (
        work.groupby(["제품분류", "product_name", "power_value", "POWER"], dropna=False)
        .agg(
            factory_group=("factory_group", join_unique),
            period_group=("period_group", first_nonempty),
            request_pack=("request_pack", "sum"),
            request_pcs=("request_pcs", "sum"),
            packing_pack=("packing_recognized_pack", "sum"),
        )
        .reset_index()
        .rename(
            columns={
                "product_name": "제품명",
                "factory_group": "공장구분",
                "period_group": "기간구분",
                "request_pack": "요청수량",
                "request_pcs": "요청PCS",
                "packing_pack": "포장수량",
            }
        )
    )
    grouped["부족수량"] = (grouped["요청수량"] - grouped["포장수량"]).clip(lower=0.0)
    grouped["진도율"] = np.where(
        grouped["요청수량"] > 0,
        grouped["포장수량"] / grouped["요청수량"] * 100.0,
        0.0,
    )
    grouped["진도율"] = np.clip(grouped["진도율"], 0.0, 100.0)
    grouped["상태"] = [
        classify_status(float(packing), float(progress))
        for packing, progress in zip(grouped["포장수량"], grouped["진도율"])
    ]

    progress_source = work.copy()
    if "production_shortage_qty" not in progress_source.columns:
        progress_source["production_shortage_qty"] = progress_source.get("production_basis_qty", 0.0)
    progress_source["_progress_dedupe_key"] = np.where(
        progress_source["production_code_key"].map(clean_str) != "",
        progress_source["production_code_key"],
        progress_source["sales_code_key"],
    )
    progress_source = progress_source.drop_duplicates(["product_name", "power_value", "_progress_dedupe_key"])
    progress_grouped = (
        progress_source.groupby(["product_name", "power_value"], dropna=False)
        .agg(
            production_basis_qty=("production_basis_qty", "sum"),
            production_shortage_qty=("production_shortage_qty", "sum"),
        )
        .reset_index()
        .rename(
            columns={
                "product_name": "제품명",
                "production_basis_qty": "누수규격검사 생산수량",
                "production_shortage_qty": "생산부족수량",
            }
        )
    )
    grouped = grouped.merge(progress_grouped, on=["제품명", "power_value"], how="left")
    grouped["누수규격검사 생산수량"] = grouped["누수규격검사 생산수량"].fillna(0.0)
    grouped["생산부족수량"] = pd.to_numeric(grouped["생산부족수량"], errors="coerce").fillna(0.0).clip(lower=0.0)
    grouped["생산진도율"] = calc_production_progress_pct(grouped["요청PCS"], grouped["생산부족수량"])
    return grouped


def get_group_options(power_df: pd.DataFrame) -> list[str]:
    options = GROUP_ORDER.copy()
    if power_df.empty:
        return options
    available = set(power_df["제품분류"].astype(str))
    extras = sorted(available - set(options))
    return options + extras


def filter_power_detail(
    power_df: pd.DataFrame,
    group_name: str,
    product_name: str,
    high_power_only: bool,
    shortage_only: bool,
    not_started_only: bool,
) -> pd.DataFrame:
    out = power_df.copy()
    if group_name == "본품":
        out = out[~out["제품명"].astype(str).map(is_sample_name)]
    elif group_name == "샘플":
        out = out[out["제품명"].astype(str).map(is_sample_name)]
    elif group_name != "전체":
        out = out[out["제품분류"] == group_name]
    if product_name != "전체":
        out = out[out["제품명"] == product_name]
    if high_power_only:
        out = out[out["power_value"] <= -5.0]
    if shortage_only:
        out = out[out["부족수량"] > 0]
    if not_started_only:
        out = out[out["상태"] == "미착수"]
    return out


def calc_power_ops_kpi(power_df: pd.DataFrame) -> dict[str, float]:
    if power_df.empty:
        return {
            "rows": 0,
            "shortage_rows": 0,
            "not_started_rows": 0,
            "high_power_shortage_rows": 0,
            "shortage_qty": 0.0,
        }
    return {
        "rows": int(len(power_df)),
        "shortage_rows": int((power_df["부족수량"] > 0).sum()),
        "not_started_rows": int((power_df["상태"] == "미착수").sum()),
        "high_power_shortage_rows": int(((power_df["power_value"] <= -5.0) & (power_df["부족수량"] > 0)).sum()),
        "shortage_qty": float(power_df["부족수량"].sum()),
    }


def render_power_ops_table(power_df: pd.DataFrame, max_rows: int = 2000) -> None:
    if power_df.empty:
        st.warning("조건에 맞는 POWER 상세 데이터가 없습니다.")
        return

    ordered = power_df.sort_values(["power_value", "부족수량"], ascending=[True, False], kind="stable").head(max_rows).copy()
    rows: list[str] = []
    for _, row in ordered.iterrows():
        power_value = float(row["power_value"])
        power_label = escape(str(row["POWER"]))
        req = format_int(float(row["요청수량"]))
        packed = format_int(float(row["포장수량"]))
        shortage = float(row["부족수량"])
        shortage_txt = format_int(shortage)
        progress = float(row["진도율"])
        prod_progress = float(row.get("생산진도율", 0.0))

        power_class = "power-cell high" if power_value <= -5.0 else "power-cell"
        shortage_class = "num shortage" if shortage > 0 else "num"
        progress_html = progress_cell_html(progress, "포장")
        prod_progress_html = progress_cell_html(prod_progress, "생산")

        rows.append(
            "<tr>"
            f"<td class='{power_class}'>{power_label}</td>"
            f"<td class='num'>{req}</td>"
            f"<td class='num'>{packed}</td>"
            f"<td class='{shortage_class}'>{shortage_txt}</td>"
            f"<td>{prod_progress_html}</td>"
            f"<td>{progress_html}</td>"
            "</tr>"
        )

    header = (
        "<tr>"
        "<th>POWER</th>"
        "<th class='num'>요청수량</th>"
        "<th class='num'>포장수량</th>"
        "<th class='num'>부족수량</th>"
        "<th>생산진도율</th>"
        "<th>포장진도율</th>"
        "</tr>"
    )
    table_html = (
        "<div class='table-wrap'>"
        "<table class='ops-table'>"
        f"<thead>{header}</thead>"
        f"<tbody>{''.join(rows)}</tbody>"
        "</table>"
        "</div>"
    )
    st.markdown(table_html, unsafe_allow_html=True)


def build_power_heatmap(power_df: pd.DataFrame) -> px.imshow | None:
    if power_df.empty:
        return None

    power_order = (
        power_df[["POWER", "power_value"]]
        .drop_duplicates()
        .sort_values("power_value", ascending=True)["POWER"]
        .tolist()
    )
    product_order = (
        power_df.groupby("제품명", dropna=False)["부족수량"]
        .sum()
        .sort_values(ascending=False)
        .index.tolist()
    )
    matrix = power_df.pivot_table(index="제품명", columns="POWER", values="진도율", aggfunc="mean")
    matrix = matrix.reindex(index=product_order, columns=power_order)
    if matrix.shape[0] > 35:
        matrix = matrix.iloc[:35]

    fig = px.imshow(
        matrix,
        aspect="auto",
        zmin=0,
        zmax=100,
        color_continuous_scale=[(0.0, BG_PAGE), (0.5, "#D9E7F3"), (1.0, COLOR_BLUE)],
        labels={"x": "POWER", "y": "제품명", "color": "포장진도율(%)"},
        title="제품/POWER 포장진도율 Heatmap",
    )
    fig.update_layout(
        paper_bgcolor=WHITE,
        plot_bgcolor=WHITE,
        margin=dict(l=8, r=8, t=52, b=8),
    )
    fig.update_traces(
        hovertemplate="제품명: %{y}<br>POWER: %{x}<br>포장진도율: %{z:.1f}%<extra></extra>"
    )
    fig.update_xaxes(type="category")
    return fig


def apply_filters(df: pd.DataFrame, query: str, statuses: list[str]) -> pd.DataFrame:
    out = df.copy()
    q = query.strip()
    if q:
        out = out[out["제품명"].astype(str).str.contains(q, case=False, na=False)]
    if statuses:
        out = out[out["상태"].isin(statuses)]
    else:
        out = out.iloc[0:0]
    return out


def product_scope_options(df: pd.DataFrame) -> list[str]:
    base = ["전체", "본품", "샘플"]
    if "제품분류" not in df.columns:
        return base
    extras = [value for value in GROUP_ORDER if value not in base and value in set(df["제품분류"].astype(str))]
    remaining = sorted(set(df["제품분류"].astype(str)) - set(base) - set(extras))
    return base + extras + remaining


def product_family_options(df: pd.DataFrame) -> list[str]:
    base = ["전체"]
    if "본품분류" not in df.columns:
        return base
    available = set(df["본품분류"].dropna().astype(str))
    ordered = [value for value in MAIN_PRODUCT_FAMILY_ORDER if value not in base and value in available]
    remaining = sorted(available - set(base) - set(ordered))
    return base + ordered + remaining


def apply_product_scope_filter(df: pd.DataFrame, scope: str) -> pd.DataFrame:
    if scope == "전체" or "제품분류" not in df.columns:
        return df.copy()
    if scope == "본품":
        return df[~df["제품명"].astype(str).map(is_sample_name)].copy()
    if scope == "샘플":
        return df[df["제품명"].astype(str).map(is_sample_name)].copy()
    return df[df["제품분류"].astype(str) == scope].copy()


def apply_product_family_filter(df: pd.DataFrame, family: str) -> pd.DataFrame:
    if family == "전체" or "본품분류" not in df.columns:
        return df.copy()
    return df[df["본품분류"].astype(str) == family].copy()


def calc_kpi(df: pd.DataFrame) -> dict[str, float]:
    request_pack = float(df["요청 PACK"].sum()) if not df.empty else 0.0
    packing_pack = float(df["포장 PACK"].sum()) if not df.empty else 0.0
    yongma_in_pack = (
        float(df["용마입고 PACK"].sum()) if "용마입고 PACK" in df.columns and not df.empty else 0.0
    )
    shortage_pack = float(df["미입고수량"].sum()) if "미입고수량" in df.columns and not df.empty else max(0.0, request_pack - yongma_in_pack)
    progress = (yongma_in_pack / request_pack * 100.0) if request_pack > 0 else 0.0
    packing_progress = (packing_pack / request_pack * 100.0) if request_pack > 0 else 0.0
    request_pcs = float(df["요청 PCS"].sum()) if "요청 PCS" in df.columns and not df.empty else 0.0
    if "미입고 PCS" in df.columns and not df.empty:
        shortage_pcs = float(pd.to_numeric(df["미입고 PCS"], errors="coerce").fillna(0.0).sum())
    elif request_pack > 0 and request_pcs > 0:
        shortage_pcs = max(0.0, request_pcs * shortage_pack / request_pack)
    else:
        shortage_pcs = 0.0
    production_shortage_qty = (
        float(df["생산부족수량"].sum()) if "생산부족수량" in df.columns and not df.empty else 0.0
    )
    production_progress = (
        (request_pcs - production_shortage_qty) / request_pcs * 100.0
        if request_pcs > 0
        else 0.0
    )
    return {
        "request_pack": request_pack,
        "request_pcs": request_pcs,
        "packing_pack": packing_pack,
        "yongma_in_pack": yongma_in_pack,
        "shortage_pack": shortage_pack,
        "shortage_pcs": shortage_pcs,
        "production_shortage_pcs": production_shortage_qty,
        "progress_pct": min(100.0, max(0.0, progress)),
        "packing_progress_pct": min(100.0, max(0.0, packing_progress)),
        "production_progress_pct": min(100.0, max(0.0, production_progress)),
        "production_shortage_products": int((df["생산부족수량"] > 0).sum()) if "생산부족수량" in df.columns else 0,
        "packing_shortage_products": int((df["포장부족수량"] > 0).sum()) if "포장부족수량" in df.columns else 0,
        "not_started_products": int((df["상태"] == "미착수").sum()) if "상태" in df.columns else 0,
        "completed_products": int((df["상태"] == "완료").sum()) if "상태" in df.columns else 0,
    }


def code_summary_for_products(code_summary: pd.DataFrame, product_names: pd.Series) -> pd.DataFrame:
    if code_summary.empty:
        return code_summary.copy()
    names = set(product_names.dropna().astype(str))
    if not names:
        return code_summary.iloc[0:0].copy()

    work = code_summary.copy()
    if "base_product_name" in work.columns:
        base_names = work["base_product_name"].astype(str)
    else:
        base_names = work["product_name"].map(strip_pack_unit_suffix).astype(str)
    return work[base_names.isin(names)].copy()


def add_allocated_production_basis(code_summary: pd.DataFrame) -> pd.DataFrame:
    work = code_summary.copy()
    if work.empty:
        work["_allocated_production_shortage_qty"] = 0.0
        work["_allocated_production_source_shortage_qty"] = 0.0
        work["_basis_difference_pcs"] = 0.0
        work["_allocated_sample_available_pcs"] = 0.0
        return work

    if "production_basis_qty" not in work.columns:
        work["production_basis_qty"] = 0.0
    if "production_shortage_qty" not in work.columns:
        work["production_shortage_qty"] = work["production_basis_qty"]
    if "production_source_shortage_qty" not in work.columns:
        work["production_source_shortage_qty"] = work["production_shortage_qty"]
    if "sample_available_pcs" not in work.columns:
        work["sample_available_pcs"] = 0.0
    work["production_basis_qty"] = pd.to_numeric(work["production_basis_qty"], errors="coerce").fillna(0.0)
    work["production_shortage_qty"] = pd.to_numeric(work["production_shortage_qty"], errors="coerce").fillna(0.0)
    work["production_source_shortage_qty"] = pd.to_numeric(
        work["production_source_shortage_qty"],
        errors="coerce",
    ).fillna(0.0)
    work["sample_available_pcs"] = pd.to_numeric(work["sample_available_pcs"], errors="coerce").fillna(0.0)

    progress_key = work.get("production_progress_key", pd.Series("", index=work.index)).map(clean_str)
    production_key = work.get("production_code_key", pd.Series("", index=work.index)).map(clean_str)
    sales_key = work.get("sales_code_key", pd.Series("", index=work.index)).map(clean_str)
    fallback_key = work.get("sales_code", pd.Series("", index=work.index)).map(clean_str)
    work["_production_alloc_key"] = progress_key.where(progress_key != "", production_key)
    work["_production_alloc_key"] = work["_production_alloc_key"].where(work["_production_alloc_key"] != "", sales_key)
    work["_production_alloc_key"] = work["_production_alloc_key"].where(work["_production_alloc_key"] != "", fallback_key)

    raw_shortage = work["production_shortage_qty"].clip(lower=0.0)
    request_pcs = pd.to_numeric(work.get("request_pcs", pd.Series(0.0, index=work.index)), errors="coerce").fillna(0.0)
    source_shortage = work["production_source_shortage_qty"].clip(lower=0.0)
    source_produced_pcs = (request_pcs - source_shortage).clip(lower=0.0)
    packing_pcs = recognized_packing_pcs(work)
    basis_difference_pcs = (packing_pcs - source_produced_pcs).clip(lower=0.0)
    duplicated_progress = (
        (work["_production_alloc_key"].map(clean_str) != "")
        & (raw_shortage > 0)
        & work["_production_alloc_key"].duplicated(keep="first")
    )
    duplicated_source = (
        (work["_production_alloc_key"].map(clean_str) != "")
        & (source_shortage > 0)
        & work["_production_alloc_key"].duplicated(keep="first")
    )
    work["_allocated_production_shortage_qty"] = raw_shortage.where(~duplicated_progress, 0.0)
    work["_allocated_production_source_shortage_qty"] = source_shortage.where(~duplicated_source, 0.0)
    work["_basis_difference_pcs"] = basis_difference_pcs.where(~duplicated_source, 0.0)
    work["_allocated_sample_available_pcs"] = work["sample_available_pcs"].clip(lower=0.0).round(0).astype("int64")
    return work


def calc_kpi_from_code_summary(code_summary: pd.DataFrame) -> dict[str, float]:
    if code_summary.empty:
        return {
            "request_pack": 0.0,
            "request_pcs": 0.0,
            "packing_pack": 0.0,
            "yongma_in_pack": 0.0,
            "shortage_pack": 0.0,
            "shortage_pcs": 0.0,
            "production_shortage_pcs": 0.0,
            "packable_pcs": 0.0,
            "progress_pct": 0.0,
            "packing_progress_pct": 0.0,
            "production_progress_pct": 0.0,
        }

    work = add_code_level_supply_basis(code_summary)
    request_pack = float(pd.to_numeric(work.get("request_pack", pd.Series(dtype=float)), errors="coerce").fillna(0.0).sum())
    packing_pack = float(
        pd.to_numeric(work.get("packing_recognized_pack", work.get("packing_pack", pd.Series(dtype=float))), errors="coerce")
        .fillna(0.0)
        .sum()
    )
    yongma_in_pack = float(
        pd.to_numeric(work.get("yongma_recognized_pack", work.get("yongma_in_pack", pd.Series(dtype=float))), errors="coerce")
        .fillna(0.0)
        .sum()
    )
    shortage_pack = max(0.0, request_pack - yongma_in_pack)
    yongma_in_pcs = float(
        pack_quantity_to_pcs(
            work,
            pd.to_numeric(
                work.get("yongma_recognized_pack", work.get("yongma_in_pack", pd.Series(0.0, index=work.index))),
                errors="coerce",
            ).fillna(0.0),
        ).sum()
    )
    receipt_progress = (yongma_in_pack / request_pack * 100.0) if request_pack > 0 else 0.0
    packing_progress = (packing_pack / request_pack * 100.0) if request_pack > 0 else 0.0

    work = add_allocated_production_basis(work)
    request_pcs = float(work["request_pcs"].sum())
    receipt_shortage_pcs = max(0.0, request_pcs - yongma_in_pcs)
    production_shortage_pcs = float(work["_allocated_production_shortage_qty"].sum())
    packable_pcs = max(0.0, request_pcs - production_shortage_pcs)
    production_progress = ((request_pcs - production_shortage_pcs) / request_pcs * 100.0) if request_pcs > 0 else 0.0

    return {
        "request_pack": request_pack,
        "request_pcs": request_pcs,
        "packing_pack": packing_pack,
        "yongma_in_pack": yongma_in_pack,
        "shortage_pack": shortage_pack,
        "shortage_pcs": receipt_shortage_pcs,
        "production_shortage_pcs": production_shortage_pcs,
        "packable_pcs": packable_pcs,
        "progress_pct": min(100.0, max(0.0, receipt_progress)),
        "packing_progress_pct": min(100.0, max(0.0, packing_progress)),
        "production_progress_pct": min(100.0, max(0.0, production_progress)),
    }


def format_int(value: float) -> str:
    return f"{value:,.0f}"


def progress_tone(progress: float) -> str:
    if progress >= 100:
        return "done"
    if progress >= 80:
        return "active"
    if progress >= 50:
        return "warn"
    return "risk"


def status_class(status: str) -> str:
    if status in {"완료", "입고완료"}:
        return "done"
    if status in {"진행중"}:
        return "active"
    if status in {"부족"}:
        return "warn"
    if status in {"입고대기"}:
        return "waiting"
    return "risk"


def progress_cell_html(progress: float, label: str = "", show_label: bool = True) -> str:
    width = max(0.0, min(100.0, float(progress)))
    semantic = " packing"
    if label in {"생산"}:
        semantic = " production"
    elif label in {"포장"}:
        semantic = " packing"
    elif label in {"입고", "용마입고"}:
        semantic = " receipt"
    prefix = f"<span class='progress-name'>{escape(label)}</span>" if label and show_label else ""
    return (
        "<div class='progress-cell'>"
        f"{prefix}"
        "<div class='progress-track'>"
        f"<div class='progress-fill{semantic}' style='width:{width:.1f}%'></div>"
        "</div>"
        f"<span class='progress-text'>{progress:.1f}%</span>"
        "</div>"
    )


def render_ops_table(
    df: pd.DataFrame,
    compact: bool = False,
    max_rows: int = 500,
    show_family: bool = False,
) -> None:
    if df.empty:
        st.warning("조건에 맞는 데이터가 없습니다.")
        return

    source = df.copy()
    if "용마입고율" not in source.columns:
        source["용마입고율"] = 0.0
    receipt_shortage_col = "미입고수량" if "미입고수량" in source.columns else "포장부족수량"
    receipt_progress_col = "용마입고율"

    ordered = source.sort_values(
        [receipt_shortage_col, "생산부족수량", "요청 PACK"],
        ascending=[False, False, False],
        kind="stable",
    ).head(max_rows).copy()

    rows: list[str] = []
    for _, row in ordered.iterrows():
        product = escape(str(row["제품명"]))
        family = escape(str(row.get("본품분류", ""))) if show_family else ""
        req = format_int(float(row["요청 PACK"]))
        receipt_shortage = float(row.get(receipt_shortage_col, 0.0))
        receipt_shortage_txt = format_int(receipt_shortage)
        receipt_progress = float(row.get(receipt_progress_col, 0.0))
        production_shortage = float(row.get("생산부족수량", 0.0))
        production_shortage_txt = format_int(production_shortage)
        prod_progress = float(row.get("생산진도율", 0.0))
        status = escape(str(row["상태"]))
        badge = f"<span class='status-badge {status_class(str(row['상태']))}'>{status}</span>"

        receipt_shortage_class = "num shortage" if receipt_shortage > 0 else "num"
        production_shortage_class = "num shortage" if production_shortage > 0 else "num"
        receipt_progress_html = progress_cell_html(receipt_progress, "입고")
        production_progress_html = progress_cell_html(prod_progress, "생산")

        if compact:
            rows.append(
                "<tr>"
                f"<td class='left'>{product}</td>"
                f"<td>{production_progress_html}</td>"
                f"<td>{receipt_progress_html}</td>"
                f"<td class='{receipt_shortage_class}'>{receipt_shortage_txt}</td>"
                f"<td>{badge}</td>"
                "</tr>"
            )
        else:
            rows.append(
                "<tr>"
                f"<td class='left'>{product}</td>"
                f"{f'<td>{family}</td>' if show_family else ''}"
                f"<td class='num'>{req}</td>"
                f"<td>{production_progress_html}</td>"
                f"<td>{receipt_progress_html}</td>"
                f"<td class='{production_shortage_class}'>{production_shortage_txt}</td>"
                f"<td class='{receipt_shortage_class}'>{receipt_shortage_txt}</td>"
                f"<td>{badge}</td>"
                "</tr>"
            )

    header = (
        "<tr>"
        "<th class='left'>제품명</th>"
        "<th>생산진도율</th>"
        "<th>용마입고율</th>"
        "<th class='num'>미입고수량</th>"
        "<th>상태</th>"
        "</tr>"
        if compact
        else "<tr>"
        "<th class='left'>제품명</th>"
        f"{'<th>본품분류</th>' if show_family else ''}"
        "<th class='num'>요청 PACK</th>"
        "<th>생산진도율</th>"
        "<th>용마입고율</th>"
        "<th class='num'>생산부족수량</th>"
        "<th class='num'>미입고수량</th>"
        "<th>상태</th>"
        "</tr>"
    )

    table_html = (
        "<div class='table-wrap'>"
        "<table class='ops-table'>"
        f"<thead>{header}</thead>"
        f"<tbody>{''.join(rows)}</tbody>"
        "</table>"
        "</div>"
    )
    st.markdown(table_html, unsafe_allow_html=True)


def build_family_progress_view(product_df: pd.DataFrame) -> pd.DataFrame:
    if product_df.empty or "본품분류" not in product_df.columns:
        return pd.DataFrame(
            columns=[
                "본품분류",
                "요청 PACK",
                "요청 PCS",
                "포장 PACK",
                "용마입고 PACK",
                "생산부족수량",
                "포장부족수량",
                "미입고수량",
                "생산진도율",
                "포장진도율",
                "용마입고율",
            ]
        )

    grouped = (
        product_df.groupby("본품분류", dropna=False)
        .agg(
            request_pack=("요청 PACK", "sum"),
            request_pcs=("요청 PCS", "sum"),
            packing_pack=("포장 PACK", "sum"),
            yongma_in_pack=("용마입고 PACK", "sum"),
            production_shortage_qty=("생산부족수량", "sum"),
            packing_shortage_qty=("포장부족수량", "sum"),
            receipt_shortage_qty=("미입고수량", "sum"),
        )
        .reset_index()
        .rename(
            columns={
                "request_pack": "요청 PACK",
                "request_pcs": "요청 PCS",
                "packing_pack": "포장 PACK",
                "yongma_in_pack": "용마입고 PACK",
                "production_shortage_qty": "생산부족수량",
                "packing_shortage_qty": "포장부족수량",
                "receipt_shortage_qty": "미입고수량",
            }
        )
    )
    request_pack = pd.to_numeric(grouped["요청 PACK"], errors="coerce").fillna(0.0)
    request_pcs = pd.to_numeric(grouped["요청 PCS"], errors="coerce").fillna(0.0)
    grouped = grouped[(request_pack > 0) | (request_pcs > 0)].copy()
    grouped["생산진도율"] = calc_production_progress_pct(grouped["요청 PCS"], grouped["생산부족수량"])
    grouped["용마입고율"] = np.where(
        grouped["요청 PACK"] > 0,
        grouped["용마입고 PACK"] / grouped["요청 PACK"] * 100.0,
        0.0,
    )
    grouped["용마입고율"] = np.clip(grouped["용마입고율"], 0.0, 100.0)
    grouped["포장진도율"] = np.where(
        grouped["요청 PACK"] > 0,
        grouped["포장 PACK"] / grouped["요청 PACK"] * 100.0,
        0.0,
    )
    grouped["포장진도율"] = np.clip(grouped["포장진도율"], 0.0, 100.0)
    grouped["_order"] = grouped["본품분류"].map(
        {name: idx for idx, name in enumerate(MAIN_PRODUCT_FAMILY_ORDER)}
    ).fillna(999)
    return grouped.sort_values(
        ["_order", "포장부족수량", "요청 PACK"],
        ascending=[True, False, False],
        kind="stable",
    ).drop(columns=["_order"])


def family_card_section(family: Any) -> str:
    text = clean_str(family)
    upper = text.upper()
    if text in FAMILY_CARD_MISC_NAMES:
        return "기타"
    if text in FAMILY_CARD_1DAY_NAMES:
        return "1DAY"
    if "1DAY" in upper or "1 DAY" in upper:
        return "1DAY"
    if text in {"부자재/기타", "기타", "샘플"}:
        return "기타"
    return "FRP"


def period_group_from_category(value: Any) -> str:
    normalized = normalize_col(value)
    if "1day" in normalized:
        return "1-DAY"
    if "frp" in normalized:
        return "FRP"
    return ""


def period_group_from_family(value: Any) -> str:
    return "1-DAY" if family_card_section(value) == "1DAY" else "FRP"


def period_group_from_product_name(value: Any) -> str:
    category_period = period_group_from_category(value)
    if category_period:
        return category_period
    return period_group_from_family(classify_main_product_family(clean_str(value)))


def add_period_group_columns(df: pd.DataFrame) -> pd.DataFrame:
    if df is None or df.empty:
        return df.copy() if df is not None else pd.DataFrame()
    out = df.copy()
    period = pd.Series("", index=out.index, dtype="object")
    for col in ["category_summary", "신규분류요약"]:
        if col in out.columns:
            candidate = out[col].map(period_group_from_category)
            period = period.where(period.map(clean_str) != "", candidate)
    for col in ["본품분류"]:
        if col in out.columns:
            candidate = out[col].map(period_group_from_family)
            period = period.where(period.map(clean_str) != "", candidate)
    for col in ["base_product_name", "product_name", "제품명", "대표 제품명", "요청제품명", "마스터제품명", "재고표 제품명"]:
        if col in out.columns:
            candidate = out[col].map(period_group_from_product_name)
            period = period.where(period.map(clean_str) != "", candidate)
    period = period.map(clean_str).replace("", "FRP").fillna("FRP")
    out["period_group"] = period
    out["기간구분"] = period
    return out


def filter_by_period_group(df: pd.DataFrame, period_group: str) -> pd.DataFrame:
    if df is None:
        return pd.DataFrame()
    if period_group == "전체" or df.empty:
        return df.copy()
    work = add_period_group_columns(df)
    return work[work["period_group"] == period_group].copy()


def family_progress_bar_html(value: float, tone: str) -> str:
    width = max(0.0, min(100.0, float(value)))
    return (
        "<div class='family-progress-metric'>"
        "<div class='family-progress-track'>"
        f"<div class='family-progress-fill {escape(tone)}' style='width:{width:.1f}%'></div>"
        "</div>"
        f"<span>{float(value):.1f}%</span>"
        "</div>"
    )


def family_progress_row_html(row: pd.Series, idx: int) -> str:
    family = escape(clean_str(row["본품분류"]))
    request_pcs = format_int(float(row.get("요청 PCS", 0.0)))
    request_pack = format_int(float(row["요청 PACK"]))
    production_progress = float(row["생산진도율"])
    packing_progress = float(row.get("포장진도율", 0.0))
    receipt_progress = float(row.get("용마입고율", 0.0))
    production_shortage_value = float(row["생산부족수량"])
    production_shortage = format_int(production_shortage_value)
    shortage_class = "danger" if production_shortage_value > 0 else "normal"
    dot_class = f"dot-{idx % 8}"
    return (
        "<div class='family-table-row'>"
        f"<div class='family-name'><span class='family-dot {dot_class}'></span><b>{family}</b></div>"
        "<div class='family-request'>"
        f"<strong>{request_pcs}</strong>"
        f"<span>({request_pack} PACK)</span>"
        "</div>"
        f"<div>{family_progress_bar_html(production_progress, 'production')}</div>"
        f"<div>{family_progress_bar_html(packing_progress, 'packing')}</div>"
        f"<div>{family_progress_bar_html(receipt_progress, 'receipt')}</div>"
        f"<div class='family-num shortage {shortage_class}'>{production_shortage}</div>"
        "</div>"
    )


def render_family_progress_cards(family_df: pd.DataFrame, max_rows: int = 14) -> None:
    if family_df.empty:
        st.warning("본품 분류별 진도현황을 표시할 데이터가 없습니다.")
        return

    view = family_df.copy()
    view["_section"] = view["본품분류"].map(family_card_section)
    available_sections = [section for section in FAMILY_CARD_SECTION_ORDER if not view[view["_section"] == section].empty]
    if not available_sections:
        st.warning("본품 분류별 진도현황을 표시할 데이터가 없습니다.")
        return
    filter_key = "family_progress_section_filter"
    if st.session_state.get(filter_key) not in available_sections:
        st.session_state[filter_key] = available_sections[0]
    selected_section = st.segmented_control(
        "제품 분류 탭",
        options=available_sections,
        default=available_sections[0],
        label_visibility="collapsed",
        key=filter_key,
    )
    selected_section = str(selected_section or available_sections[0])
    scoped = view[view["_section"] == selected_section].head(max_rows).copy()
    rows = "".join(family_progress_row_html(row, idx) for idx, (_, row) in enumerate(scoped.iterrows()))
    st.markdown(
        "<div class='panel-box dashboard-card family-progress-panel'>"
        "<div class='family-table'>"
        "<div class='family-table-row family-table-head'>"
        "<div>제품 분류</div>"
        "<div class='family-request-head'>요청 PCS</div>"
        "<div>생산진도율</div>"
        "<div>포장진도율</div>"
        "<div>용마입고율</div>"
        "<div class='family-num'>생산부족 PCS</div>"
        "</div>"
        f"{rows}"
        "</div>"
        "</div>",
        unsafe_allow_html=True,
    )


def build_category_request_summary_view(
    code_summary: pd.DataFrame,
    instruction_df: pd.DataFrame | None = None,
) -> pd.DataFrame:
    request_only_columns = ["신규분류요약", "요청 PACK", "요청 PCS"]
    instruction_columns = [
        "신규분류요약",
        "요청 PACK",
        "지시 PACK",
        "PACK 지시율",
        "요청 PCS",
        "지시 PCS",
        "PCS 지시율",
        "미지시 PCS",
    ]
    if code_summary.empty:
        return pd.DataFrame(columns=instruction_columns if instruction_df is not None and not instruction_df.empty else request_only_columns)

    work = code_summary.copy()
    if "category_summary" not in work.columns:
        work["category_summary"] = "(미기재)"
    work["category_summary"] = work["category_summary"].map(clean_str).replace("", "(미기재)")
    for col in ["request_pack", "request_pcs"]:
        if col not in work.columns:
            work[col] = 0.0
        work[col] = pd.to_numeric(work[col], errors="coerce").fillna(0.0)
    work = work[(work["request_pack"] > 0) | (work["request_pcs"] > 0)].copy()
    if work.empty:
        return pd.DataFrame(columns=instruction_columns if instruction_df is not None and not instruction_df.empty else request_only_columns)

    request_grouped = (
        work.groupby("category_summary", dropna=False)
        .agg(
            request_pack=("request_pack", "sum"),
            request_pcs=("request_pcs", "sum"),
        )
        .reset_index()
        .rename(
            columns={
                "category_summary": "신규분류요약",
                "request_pack": "요청 PACK",
                "request_pcs": "요청 PCS",
            }
        )
    )

    if instruction_df is None or instruction_df.empty:
        grouped = request_grouped.sort_values(["요청 PACK", "요청 PCS"], ascending=[False, False], kind="stable")
        total = pd.DataFrame(
            [
                {
                    "신규분류요약": "합계",
                    "요청 PACK": float(grouped["요청 PACK"].sum()),
                    "요청 PCS": float(grouped["요청 PCS"].sum()),
                }
            ]
        )
        return pd.concat([grouped, total], ignore_index=True)[request_only_columns]

    instruction_work = instruction_df.copy()
    if "category_summary" not in instruction_work.columns:
        instruction_work["category_summary"] = "(미기재)"
    instruction_work["category_summary"] = instruction_work["category_summary"].map(clean_str).replace("", "(미기재)")
    for col in ["request_pack", "request_pcs"]:
        if col not in instruction_work.columns:
            instruction_work[col] = 0.0
        instruction_work[col] = pd.to_numeric(instruction_work[col], errors="coerce").fillna(0.0)
    instruction_grouped = (
        instruction_work.groupby("category_summary", dropna=False)
        .agg(
            instruction_pack=("request_pack", "sum"),
            instruction_pcs=("request_pcs", "sum"),
        )
        .reset_index()
        .rename(
            columns={
                "category_summary": "신규분류요약",
                "instruction_pack": "지시 PACK",
                "instruction_pcs": "지시 PCS",
            }
        )
    )
    grouped = request_grouped.merge(instruction_grouped, on="신규분류요약", how="outer").fillna(0.0)
    for col in ["요청 PACK", "요청 PCS", "지시 PACK", "지시 PCS"]:
        grouped[col] = pd.to_numeric(grouped[col], errors="coerce").fillna(0.0)
    grouped["PACK 지시율"] = np.where(grouped["요청 PACK"] > 0, grouped["지시 PACK"] / grouped["요청 PACK"] * 100.0, 0.0)
    grouped["PCS 지시율"] = np.where(grouped["요청 PCS"] > 0, grouped["지시 PCS"] / grouped["요청 PCS"] * 100.0, 0.0)
    grouped["미지시 PCS"] = (grouped["요청 PCS"] - grouped["지시 PCS"]).clip(lower=0.0)
    grouped = grouped.sort_values(["요청 PCS", "요청 PACK"], ascending=[False, False], kind="stable")
    total = pd.DataFrame(
        [
            {
                "신규분류요약": "합계",
                "요청 PACK": float(grouped["요청 PACK"].sum()),
                "지시 PACK": float(grouped["지시 PACK"].sum()),
                "PACK 지시율": (
                    float(grouped["지시 PACK"].sum()) / float(grouped["요청 PACK"].sum()) * 100.0
                    if float(grouped["요청 PACK"].sum()) > 0
                    else 0.0
                ),
                "요청 PCS": float(grouped["요청 PCS"].sum()),
                "지시 PCS": float(grouped["지시 PCS"].sum()),
                "PCS 지시율": (
                    float(grouped["지시 PCS"].sum()) / float(grouped["요청 PCS"].sum()) * 100.0
                    if float(grouped["요청 PCS"].sum()) > 0
                    else 0.0
                ),
                "미지시 PCS": float(grouped["미지시 PCS"].sum()),
            }
        ]
    )
    return pd.concat([grouped, total], ignore_index=True)[instruction_columns]


def render_request_instruction_level_cards(summary_view: pd.DataFrame) -> None:
    if summary_view.empty or "PCS 지시율" not in summary_view.columns:
        return
    total = summary_view[summary_view["신규분류요약"].map(clean_str) == "합계"]
    if total.empty:
        total = summary_view.tail(1)
    row = total.iloc[0]
    request_pcs = float(pd.to_numeric(pd.Series([row.get("요청 PCS", 0.0)]), errors="coerce").fillna(0.0).iloc[0])
    instruction_pcs = float(pd.to_numeric(pd.Series([row.get("지시 PCS", 0.0)]), errors="coerce").fillna(0.0).iloc[0])
    missing_pcs = float(pd.to_numeric(pd.Series([row.get("미지시 PCS", 0.0)]), errors="coerce").fillna(0.0).iloc[0])
    instruction_rate = (instruction_pcs / request_pcs * 100.0) if request_pcs > 0 else 0.0
    render_metric_card_grid(
        [
            ("요청 PCS", format_int(request_pcs), "normal"),
            ("지시 PCS", format_int(instruction_pcs), "normal"),
            ("요청 대비 지시율", f"{instruction_rate:.1f}%", "warn" if instruction_rate < 90.0 else "normal"),
            ("미지시 PCS", format_int(missing_pcs), "warn" if missing_pcs > 0 else "normal"),
        ]
    )


def render_category_request_summary_table(summary_view: pd.DataFrame) -> None:
    if summary_view.empty:
        st.warning("신규분류요약별 요청량을 표시할 데이터가 없습니다.")
        return
    st.dataframe(
        summary_view,
        hide_index=True,
        height=dataframe_auto_height(len(summary_view), 360),
        width="stretch",
        column_config=drilldown_column_config(),
    )


def build_top_shortage_view(product_df: pd.DataFrame, top_n: int = 10) -> pd.DataFrame:
    columns = ["순위", "제품명", "미입고 PACK", "생산진도율", "포장진도율", "용마입고율"]
    if product_df.empty:
        return pd.DataFrame(columns=columns)

    view = product_df.copy()
    for col in ["미입고수량", "생산진도율", "포장진도율", "용마입고율"]:
        if col not in view.columns:
            view[col] = 0.0
        view[col] = pd.to_numeric(view[col], errors="coerce").fillna(0.0)
    view = view[view["미입고수량"] > 0].sort_values("미입고수량", ascending=False, kind="stable").head(top_n).copy()
    if view.empty:
        return pd.DataFrame(columns=columns)
    view["미입고 PACK"] = view["미입고수량"]
    view["순위"] = range(1, len(view) + 1)
    return view[columns].copy()


def render_top_shortage_list(top_df: pd.DataFrame) -> None:
    if top_df.empty:
        st.warning("미입고 제품이 없습니다.")
        return
    rows: list[str] = []
    for _, row in top_df.iterrows():
        rows.append(
            "<tr>"
            f"<td class='num muted'>{format_int(float(row.get('순위', 0.0)))}</td>"
            f"<td class='left'>{escape(str(row.get('제품명', '')))}</td>"
            f"<td class='num shortage'>{format_int(float(row.get('미입고 PACK', 0.0)))}</td>"
            f"<td>{progress_cell_html(float(row.get('생산진도율', 0.0)), '생산', show_label=False)}</td>"
            f"<td>{progress_cell_html(float(row.get('포장진도율', 0.0)), '포장', show_label=False)}</td>"
            f"<td>{progress_cell_html(float(row.get('용마입고율', 0.0)), '용마입고', show_label=False)}</td>"
            "</tr>"
        )
    st.markdown(
        "<div class='table-wrap compact-table'>"
        "<table class='ops-table progress-summary-table main-summary-table top-shortage-summary-table'>"
        "<colgroup>"
        "<col class='summary-rank-col'>"
        "<col class='summary-product-col'>"
        "<col class='summary-number-col'>"
        "<col class='summary-progress-col'>"
        "<col class='summary-progress-col'>"
        "<col class='summary-progress-col'>"
        "</colgroup>"
        "<thead><tr>"
        "<th class='num'>순위</th>"
        "<th class='left'>제품명</th>"
        "<th class='num'>미입고 PACK</th>"
        "<th>생산진도율</th>"
        "<th>포장진도율</th>"
        "<th>용마입고율</th>"
        "</tr></thead>"
        f"<tbody>{''.join(rows)}</tbody>"
        "</table></div>",
        unsafe_allow_html=True,
    )


def render_top_shortage_compact(top_df: pd.DataFrame) -> None:
    if top_df.empty:
        st.warning("미입고 제품이 없습니다.")
        return
    max_shortage = float(pd.to_numeric(top_df["미입고 PACK"], errors="coerce").fillna(0.0).max())
    rows: list[str] = []
    for _, row in top_df.iterrows():
        shortage = float(pd.to_numeric(pd.Series([row.get("미입고 PACK", 0.0)]), errors="coerce").fillna(0.0).iloc[0])
        width = (shortage / max_shortage * 100.0) if max_shortage > 0 else 0.0
        rows.append(
            "<div class='rank-list-row'>"
            f"<span class='rank-num'>{format_int(float(row.get('순위', 0.0)))}</span>"
            f"<span class='rank-name'>{escape(clean_str(row.get('제품명', '')))}</span>"
            "<span class='rank-bar'><i style='width:"
            f"{max(4.0, min(100.0, width)):.1f}%'></i></span>"
            f"<span class='rank-value'>{format_int(shortage)}</span>"
            "</div>"
        )
    st.markdown(
        f"<div class='panel-box dashboard-card'><div class='rank-list compact-rank-list'>{''.join(rows)}</div></div>",
        unsafe_allow_html=True,
    )


def build_gap_top_view(product_df: pd.DataFrame, top_n: int = 10) -> pd.DataFrame:
    columns = ["순위", "제품명", "생산진도율", "용마입고율", "GAP"]
    if product_df.empty:
        return pd.DataFrame(columns=columns)

    source = product_df.copy()
    if "생산진도율" not in source.columns:
        source["생산진도율"] = 0.0
    if "용마입고율" not in source.columns:
        source["용마입고율"] = 0.0
    if "미입고수량" not in source.columns:
        source["미입고수량"] = 0.0
    if "요청 PACK" not in source.columns:
        source["요청 PACK"] = 0.0
    source["생산진도율"] = pd.to_numeric(source["생산진도율"], errors="coerce").fillna(0.0)
    source["용마입고율"] = pd.to_numeric(source["용마입고율"], errors="coerce").fillna(0.0)
    source["GAP"] = source["생산진도율"] - source["용마입고율"]
    source = source[source["GAP"] > 0].copy()
    if source.empty:
        return pd.DataFrame(columns=columns)

    out = (
        source.sort_values(
            ["GAP", "미입고수량", "요청 PACK"],
            ascending=[False, False, False],
            kind="stable",
        )
        .head(top_n)
        .copy()
    )
    out["순위"] = range(1, len(out) + 1)
    return out[columns].copy()


def render_gap_top_list(gap_df: pd.DataFrame) -> None:
    if gap_df.empty:
        st.warning("생산진도율이 용마입고율보다 높은 GAP 제품이 없습니다.")
        return
    rows: list[str] = []
    for _, row in gap_df.iterrows():
        rows.append(
            "<tr>"
            f"<td class='num muted'>{format_int(float(row.get('순위', 0.0)))}</td>"
            f"<td class='left'>{escape(str(row.get('제품명', '')))}</td>"
            f"<td>{progress_cell_html(float(row.get('생산진도율', 0.0)), '생산', show_label=False)}</td>"
            f"<td>{progress_cell_html(float(row.get('용마입고율', 0.0)), '용마입고', show_label=False)}</td>"
            f"<td class='num shortage'>+{float(row.get('GAP', 0.0)):.1f}</td>"
            "</tr>"
        )
    st.markdown(
        "<div class='table-wrap compact-table'>"
        "<table class='ops-table progress-summary-table main-summary-table gap-summary-table'>"
        "<colgroup>"
        "<col class='summary-rank-col'>"
        "<col class='summary-product-col'>"
        "<col class='summary-progress-col'>"
        "<col class='summary-progress-col'>"
        "<col class='summary-number-col'>"
        "</colgroup>"
        "<thead><tr>"
        "<th class='num'>순위</th>"
        "<th class='left'>제품명</th>"
        "<th>생산진도율</th>"
        "<th>용마입고율</th>"
        "<th class='num'>GAP</th>"
        "</tr></thead>"
        f"<tbody>{''.join(rows)}</tbody>"
        "</table></div>",
        unsafe_allow_html=True,
    )


def kpi_metric_item_html(label: str, value: str, tone: str = "normal") -> str:
    return (
        "<div class='kpi-metric'>"
        f"<div class='metric-label'>{escape(label)}</div>"
        f"<div class='metric-value {escape(tone)}'>{escape(value)}</div>"
        "</div>"
    )


def kpi_quantity_item_html(label: str, pcs_value: float, pack_value: float, tone: str = "normal") -> str:
    return (
        "<div class='kpi-metric quantity-metric'>"
        f"<div class='metric-label'>{escape(label)}</div>"
        f"<div class='metric-value quantity {escape(tone)}'>"
        f"<span class='metric-number'>{escape(format_int(pcs_value))}</span>"
        "<span class='metric-unit'>PCS</span>"
        "</div>"
        f"<div class='metric-subvalue'>({escape(format_int(pack_value))} PACK)</div>"
        "</div>"
    )


def kpi_progress_line_html(label: str, value: float, tone: str) -> str:
    width = max(0.0, min(100.0, float(value)))
    return (
        "<div class='kpi-progress-row'>"
        f"<span>{escape(label)}</span>"
        "<div class='kpi-progress-track'>"
        f"<div class='kpi-progress-fill {escape(tone)}' style='width:{width:.1f}%'></div>"
        "</div>"
        f"<b>{float(value):.1f}%</b>"
        "</div>"
    )


def render_kpi_panel(title: str, kpi: dict[str, float], unit_mode: str = UNIT_PACK) -> None:
    progress = float(kpi["progress_pct"])
    production_progress = float(kpi.get("production_progress_pct", 0.0))
    packing_progress = float(kpi.get("packing_progress_pct", 0.0))
    shortage_tone = "danger" if kpi["shortage_pack"] > 0 else "normal"
    production_shortage_tone = "danger" if kpi.get("production_shortage_pcs", 0.0) > 0 else "normal"
    scope_class = "sample-kpi" if "샘플" in title else "main-kpi" if "본품" in title else ""

    if unit_mode == UNIT_PCS:
        metrics = [
            ("요청수량", format_int(kpi.get("request_pcs", 0.0)), "normal"),
            ("생산부족 PCS", format_int(kpi.get("production_shortage_pcs", 0.0)), production_shortage_tone),
            ("생산진도율", f"{production_progress:.1f}%", "primary"),
            ("용마입고 PACK", format_int(kpi.get("yongma_in_pack", 0.0)), "normal"),
            ("용마입고율", f"{progress:.1f}%", "purple"),
        ]
        metric_html = "".join(kpi_metric_item_html(label, value, tone) for label, value, tone in metrics)
    else:
        metric_html = "".join(
            [
                kpi_quantity_item_html(
                    "요청수량",
                    float(kpi.get("request_pcs", 0.0)),
                    float(kpi.get("request_pack", 0.0)),
                    "normal",
                ),
                kpi_metric_item_html("생산진도율", f"{production_progress:.1f}%", "primary"),
                kpi_metric_item_html("포장진도율", f"{packing_progress:.1f}%", "warning"),
                kpi_metric_item_html("용마입고율", f"{progress:.1f}%", "purple"),
                kpi_quantity_item_html(
                    "미입고수량",
                    float(kpi.get("shortage_pcs", 0.0)),
                    float(kpi.get("shortage_pack", 0.0)),
                    shortage_tone,
                ),
                kpi_metric_item_html(
                    "생산부족 PCS",
                    format_int(kpi.get("production_shortage_pcs", 0.0)),
                    production_shortage_tone,
                ),
            ]
        )
    panel_html = f"""
    <div class='kpi-panel scope-kpi {scope_class}'>
      <div class='kpi-panel-head'>
        <div class='kpi-title'>{escape(title)}</div>
      </div>
      <div class='kpi-divider-grid'>{metric_html}</div>
      <div class='kpi-progress-stack'>
        {kpi_progress_line_html("생산", production_progress, "production")}
        {kpi_progress_line_html("포장", packing_progress, "packing")}
        {kpi_progress_line_html("용마입고", progress, "receipt")}
      </div>
    </div>
    """
    st.markdown(panel_html, unsafe_allow_html=True)


def render_kpi_scope_panels(
    code_summary: pd.DataFrame,
    product_names: pd.Series | None = None,
    unit_mode: str = UNIT_PACK,
) -> None:
    work = add_allocated_production_basis(code_summary)
    if product_names is not None:
        work = code_summary_for_products(work, product_names)

    scope_kpis = [
        (f"{name} KPI", kpi)
        for name, kpi in build_scope_kpis(work)
        if name in {"본품", "샘플"}
    ]
    kpi_cols = st.columns(len(scope_kpis), gap="small")
    for col, (title, kpi) in zip(kpi_cols, scope_kpis):
        with col:
            render_kpi_panel(title, kpi, unit_mode=unit_mode)


def render_product_scope_kpi_panels(product_summary: pd.DataFrame, unit_mode: str = UNIT_PACK) -> None:
    main_products, sample_products = split_main_sample(product_summary)
    scopes = [
        ("본품 KPI", main_products),
        ("샘플 KPI", sample_products),
    ]
    kpi_cols = st.columns(2, gap="small")
    for col, (title, scope_df) in zip(kpi_cols, scopes):
        with col:
            render_kpi_panel(title, calc_kpi(scope_df), unit_mode=unit_mode)


def metric_progress_tone(progress: float) -> str:
    if progress >= 80:
        return "good"
    if progress >= 50:
        return "mid"
    if progress > 0:
        return "warn"
    return "muted"


def render_metric_card_grid(items: list[tuple[str, str, str] | tuple[str, str, str, str]]) -> None:
    def metric_html(item: tuple[str, str, str] | tuple[str, str, str, str]) -> str:
        label, value, tone = item[:3]
        note = item[3] if len(item) > 3 else ""
        note_html = f"<div class='metric-note'>{escape(note)}</div>" if note else ""
        return (
            "<div class='kpi-metric'>"
            f"<div class='metric-label'>{escape(label)}</div>"
            f"<div class='metric-value {tone}'>{escape(value)}</div>"
            f"{note_html}"
            "</div>"
        )

    metrics = "".join(
        metric_html(item)
        for item in items
    )
    st.markdown(
        f"<div class='kpi-panel metric-strip'><div class='kpi-divider-grid'>{metrics}</div></div>",
        unsafe_allow_html=True,
    )


def render_status_board(
    product_summary: pd.DataFrame,
    code_summary: pd.DataFrame,
    daily_inventory_df: pd.DataFrame | None,
    sample_available_df: pd.DataFrame | None,
    stock_threshold_pack: float,
    exception_kpis: dict[str, float] | None = None,
) -> None:
    kpi = calc_operation_kpis(product_summary, code_summary, stock_threshold_pack)
    if exception_kpis is None:
        exception_kpis, _exception_detail = build_daily_exception_report_view(
            daily_inventory_df,
            code_summary,
            sample_available_df,
        )
    request_pack = float(kpi.get("request_pack", 0.0))
    request_pcs = float(kpi.get("request_pcs", 0.0))
    yongma_in_pack = float(kpi.get("yongma_in_pack", 0.0))
    missing_pack = float(kpi.get("packing_shortage_pack", 0.0))
    missing_pcs = float(kpi.get("packing_shortage_pcs", kpi.get("receipt_shortage_pcs", 0.0)))
    production_shortage = float(kpi.get("production_shortage_pcs", 0.0))
    packing_progress = float(kpi.get("packing_progress_pct", 0.0))
    receipt_progress = float(kpi.get("receipt_progress_pct", 0.0))
    production_progress = float(kpi.get("production_progress_pct", 0.0))
    packing_todo_pack = float(kpi.get("packing_todo_pack", 0.0))
    receipt_wait_pack = float(kpi.get("receipt_wait_pack", 0.0))
    priority_products = int(kpi.get("priority_products", 0.0))
    request_out_count = int(exception_kpis.get("request_out_count", 0.0))

    receipt_width = max(0.0, min(100.0, receipt_progress))
    missing_width = max(0.0, min(100.0 - receipt_width, 100.0))
    emergency_count = request_out_count
    if emergency_count > 0 or priority_products > 0:
        board_tone = "risk"
    elif missing_pack > 0 or production_shortage > 0:
        board_tone = "warn"
    else:
        board_tone = "good"

    metric_html = "".join(
        [
            kpi_quantity_item_html("요청수량", request_pcs, request_pack, "normal"),
            kpi_metric_item_html("생산진도", f"{production_progress:.1f}%", "primary"),
            kpi_metric_item_html("포장진도", f"{packing_progress:.1f}%", "warning"),
            kpi_metric_item_html("용마입고율", f"{receipt_progress:.1f}%", "purple"),
            kpi_quantity_item_html("미입고수량", missing_pcs, missing_pack, "danger" if missing_pack > 0 else "normal"),
        ]
    )

    board_html = f"""
    <div class='kpi-dashboard-block status-board {board_tone}'>
      <div class='overall-kpi-card'>
        <div class='overall-kpi-title'>
          <span>전체 KPI</span>
          <b>요청 대비 진행 현황</b>
        </div>
        <div class='kpi-divider-grid overall'>{metric_html}</div>
        <div class='kpi-progress-stack overall'>
          {kpi_progress_line_html("생산", production_progress, "production")}
          {kpi_progress_line_html("포장", packing_progress, "packing")}
          {kpi_progress_line_html("용마입고", receipt_progress, "receipt")}
        </div>
        <div class='overall-kpi-foot'>
          <span>용마입고 {format_int(yongma_in_pack)} PACK</span>
          <span>포장대기 {format_int(receipt_wait_pack)} PACK</span>
          <span>포장필요 {format_int(packing_todo_pack)} PACK</span>
          <span>생산부족 {format_int(production_shortage)} PCS</span>
        </div>
      </div>
    </div>
    """
    st.markdown(board_html, unsafe_allow_html=True)


def build_scope_kpis(code_summary: pd.DataFrame) -> list[tuple[str, dict[str, float]]]:
    sample_mask = (
        code_summary["product_name"].astype(str).map(is_sample_name)
        if "product_name" in code_summary.columns
        else pd.Series(False, index=code_summary.index)
    )
    return [
        ("전체", calc_kpi_from_code_summary(code_summary)),
        ("본품", calc_kpi_from_code_summary(code_summary[~sample_mask].copy())),
        ("샘플", calc_kpi_from_code_summary(code_summary[sample_mask].copy())),
    ]


def base_pack_label(value: Any) -> str:
    num = pd.to_numeric(value, errors="coerce")
    if pd.isna(num) or float(num) <= 0:
        return "(미기재)"
    return f"{float(num):g}P"


def extract_pack_label_from_text(value: Any) -> str:
    unit = extract_pack_unit(value)
    if pd.notna(unit) and float(unit) > 0:
        return base_pack_label(unit)
    text = clean_str(value)
    if not text:
        return "(미기재)"
    match = PACK_ANY_RE.search(text)
    if not match:
        return "(미기재)"
    try:
        return base_pack_label(float(match.group(1)))
    except ValueError:
        return "(미기재)"


def standard_pack_bucket(label: Any) -> str:
    text = clean_str(label)
    return text if text in STANDARD_PACK_BUCKETS else "기타팩"


def row_pack_bucket(row: pd.Series) -> str:
    for value in [row.get("_pack_label", ""), row.get("product_name", ""), row.get("sales_code", "")]:
        label = extract_pack_label_from_text(value) if not str(value).endswith("P") else clean_str(value)
        bucket = standard_pack_bucket(label)
        if bucket != "기타팩":
            return bucket
    return "기타팩"


def pack_sort_key(label: Any) -> tuple[int, float, str]:
    text = clean_str(label)
    if text == "기타팩":
        return (1, 999998.0, text)
    match = re.match(r"^(\d+(?:\.\d+)?)P$", text)
    if match:
        return (0, float(match.group(1)), text)
    if text == "(미기재)":
        return (2, 999999.0, text)
    return (1, 999999.0, text)


def pack_sort_rank(label: Any) -> float:
    group, value, _ = pack_sort_key(label)
    return group * 1000000.0 + value


def sorted_pack_labels(labels: list[str]) -> list[str]:
    unique = list(dict.fromkeys([clean_str(label) for label in labels if clean_str(label)]))
    return sorted(unique, key=pack_sort_key)


def pack_pcs_label(label: Any) -> str:
    return f"{clean_str(label)}(PCS)"


def with_operational_columns(code_summary: pd.DataFrame) -> pd.DataFrame:
    work = add_code_level_supply_basis(code_summary)
    if "base_product_name" not in work.columns:
        work["base_product_name"] = work["product_name"].map(strip_pack_unit_suffix)
    if "pack_unit" not in work.columns:
        work["pack_unit"] = work["product_name"].map(extract_pack_unit)
    if "_pack_label" not in work.columns:
        work["_pack_label"] = work["pack_unit"].map(base_pack_label)
    if "_pack_sort" not in work.columns:
        work["_pack_sort"] = work["_pack_label"].map(pack_sort_rank)
    if "제품분류" not in work.columns:
        work["제품분류"] = work["base_product_name"].map(classify_product_group)
    if "본품분류" not in work.columns:
        work["본품분류"] = work["base_product_name"].map(classify_main_product_family)
    work = add_period_group_columns(work)
    if "본품/샘플" not in work.columns:
        work["본품/샘플"] = np.where(work["base_product_name"].astype(str).map(is_sample_name), "샘플", "본품")
    if "factory_group" not in work.columns:
        category_source = work.get("category_summary", pd.Series("", index=work.index))
        work["factory_group"] = category_source.map(factory_group_from_category)
    work["factory_group"] = work["factory_group"].map(clean_str)
    category_source = work.get("category_summary", pd.Series("", index=work.index))
    category_factory = category_source.map(factory_group_from_category)
    product_source = work.get("product_name", work["base_product_name"])
    product_factory = product_source.map(factory_group_from_product_name)
    missing_factory = ~work["factory_group"].map(has_factory_group)
    work.loc[missing_factory & category_factory.map(has_factory_group), "factory_group"] = category_factory
    missing_factory = ~work["factory_group"].map(has_factory_group)
    work.loc[missing_factory & product_factory.map(has_factory_group), "factory_group"] = product_factory
    work["factory_group"] = work["factory_group"].map(clean_factory_group_display).replace("", "(미기재)").fillna("(미기재)")
    work["공장구분"] = work["factory_group"]
    if "customer_name" not in work.columns:
        work["customer_name"] = "(미기재)"
    work["customer_name"] = work["customer_name"].replace("", "(미기재)").fillna("(미기재)")
    if "power_value" not in work.columns:
        sales_power = work["sales_code"].map(parse_power_from_sales_code)
        production_power = work["production_code"].map(parse_power_from_sales_code)
        product_power = work["product_name"].map(parse_power_from_sales_code)
        work["power_value"] = sales_power.fillna(production_power).fillna(product_power)
    if "POWER" not in work.columns:
        work["POWER"] = work["power_value"].map(format_power)
    if "production_code_display" not in work.columns:
        work["production_code_display"] = work["production_code"].map(clean_str)
    if "_pack_bucket" not in work.columns:
        work["_pack_bucket"] = work.apply(row_pack_bucket, axis=1)
    if "_pack_bucket_sort" not in work.columns:
        work["_pack_bucket_sort"] = work["_pack_bucket"].map(pack_sort_rank)
    return work


def available_pack_options(code_summary: pd.DataFrame) -> list[str]:
    work = with_operational_columns(code_summary)
    labels = sorted_pack_labels(work["_pack_label"].dropna().astype(str).tolist())
    return ["전체"] + labels


def available_product_group_options(code_summary: pd.DataFrame) -> list[str]:
    work = with_operational_columns(code_summary)
    available = set(work["제품분류"].dropna().astype(str))
    ordered = [value for value in GROUP_ORDER if value not in {"본품", "샘플"} and value in available]
    remaining = sorted(available - set(ordered))
    return ["전체"] + ordered + remaining


def ordered_factory_group_options(values: Any) -> list[str]:
    series = pd.Series(values, dtype="object") if values is not None else pd.Series(dtype="object")
    available = {clean_str(value) for value in series.dropna().tolist()}
    available = {value for value in available if value and value != "전체"}
    ordered = [value for value in FACTORY_GROUP_ORDER if value != "전체"]
    remaining = sorted(available - set(ordered))
    return ["전체"] + ordered + remaining


def available_factory_group_options(code_summary: pd.DataFrame) -> list[str]:
    work = add_allocated_production_basis(with_operational_columns(code_summary))
    return ordered_factory_group_options(work["factory_group"])


def available_power_options(code_summary: pd.DataFrame) -> list[str]:
    work = with_operational_columns(code_summary)
    source = work[work["power_value"].notna()][["POWER", "power_value"]].drop_duplicates()
    source = source.sort_values("power_value", ascending=True, kind="stable")
    return ["전체"] + source["POWER"].astype(str).tolist()


def available_customer_options(code_summary: pd.DataFrame) -> list[str]:
    work = with_operational_columns(code_summary)
    values = sorted(work["customer_name"].dropna().astype(str).unique().tolist())
    return ["전체"] + values


def filter_operational_code_summary(
    code_summary: pd.DataFrame,
    product_query: str = "",
    production_query: str = "",
    sales_query: str = "",
    pack_label: str = "전체",
    product_group: str = "전체",
    sample_scope: str = "전체",
    power_label: str = "전체",
    customer_name: str = "전체",
    factory_group: str = "전체",
    period_group: str = "전체",
) -> pd.DataFrame:
    out = with_operational_columns(code_summary)
    product_q = product_query.strip()
    if product_q:
        name_match = out["product_name"].astype(str).str.contains(product_q, case=False, na=False)
        base_match = out["base_product_name"].astype(str).str.contains(product_q, case=False, na=False)
        out = out[name_match | base_match]
    production_q = production_query.strip()
    if production_q:
        out = out[out["production_code_display"].astype(str).str.contains(production_q, case=False, na=False)]
    sales_q = sales_query.strip()
    if sales_q:
        out = out[out["sales_code"].astype(str).str.contains(sales_q, case=False, na=False)]
    if pack_label != "전체":
        out = out[out["_pack_label"] == pack_label]
    if product_group != "전체":
        out = out[out["제품분류"] == product_group]
    if sample_scope == "본품":
        out = out[out["본품/샘플"] == "본품"]
    elif sample_scope == "샘플":
        out = out[out["본품/샘플"] == "샘플"]
    if power_label != "전체":
        out = out[out["POWER"] == power_label]
    if customer_name != "전체":
        out = out[out["customer_name"] == customer_name]
    if factory_group != "전체":
        out = out[out["factory_group"] == factory_group]
    if period_group != "전체":
        out = out[out["period_group"] == period_group]
    return out.copy()


def build_pack_pivot(
    code_summary: pd.DataFrame,
    index_cols: list[str],
    pack_labels: list[str],
) -> pd.DataFrame:
    if code_summary.empty:
        return pd.DataFrame(columns=index_cols + pack_labels)
    work = with_operational_columns(code_summary)
    work["_pivot_request_pack"] = pd.to_numeric(work["request_pack"], errors="coerce").fillna(0.0).astype(float)
    pivot = (
        work.pivot_table(
            index=index_cols,
            columns="_pack_label",
            values="_pivot_request_pack",
            aggfunc="sum",
            dropna=False,
        )
        .fillna(0.0)
        .reset_index()
        .rename_axis(None, axis=1)
    )
    for label in pack_labels:
        if label not in pivot.columns:
            pivot[label] = 0.0
    return pivot[index_cols + pack_labels]


def display_date_or_dash(value: Any) -> str:
    text = format_date(value)
    return text if text else "-"


def pct_from_parts(done: Any, total: Any) -> float:
    done_num = float(pd.to_numeric(done, errors="coerce") or 0.0)
    total_num = float(pd.to_numeric(total, errors="coerce") or 0.0)
    if total_num <= 0:
        return 0.0
    return min(100.0, max(0.0, done_num / total_num * 100.0))


def d_day_number(value: Any) -> float:
    due = pd.to_datetime(value, errors="coerce")
    if pd.isna(due):
        return np.nan
    today = pd.Timestamp.now(tz="Asia/Seoul").tz_localize(None).normalize()
    return float((due.normalize() - today).days)


def d_day_text(value: Any) -> str:
    days = d_day_number(value)
    if pd.isna(days):
        return "-"
    day_count = int(days)
    if day_count < 0:
        return f"D+{abs(day_count)}"
    if day_count == 0:
        return "D-Day"
    return f"D-{day_count}"


def priority_grade(shortage_pack: Any, due_date: Any, current_stock_pack: Any, stock_threshold_pack: float) -> str:
    shortage = float(pd.to_numeric(shortage_pack, errors="coerce") or 0.0)
    if shortage <= 0:
        return "완료"

    days = d_day_number(due_date)
    stock = pd.to_numeric(current_stock_pack, errors="coerce")
    stock_low = pd.notna(stock) and float(stock) <= float(stock_threshold_pack)
    due_over = pd.notna(days) and days <= 0
    due_very_close = pd.notna(days) and days <= 3
    due_close = pd.notna(days) and days <= 7

    if due_over or (due_very_close and stock_low):
        return "A 긴급"
    if due_close or stock_low:
        return "B 주의"
    return "C 일반"


def priority_sort_value(value: Any) -> int:
    order = {"A 긴급": 0, "B 주의": 1, "C 일반": 2, "완료": 3}
    return order.get(str(value), 9)


def add_priority_columns(
    df: pd.DataFrame,
    stock_threshold_pack: float,
    shortage_col: str = "포장부족수량",
    due_col: str = "request_due_date",
    stock_col: str = "용마창고재고 (PACK)",
    request_col: str | None = None,
) -> pd.DataFrame:
    out = df.copy()
    if due_col not in out.columns:
        out[due_col] = pd.NaT
    if stock_col not in out.columns:
        out[stock_col] = np.nan
    stock = pd.to_numeric(out[stock_col], errors="coerce")
    out["재고기준(PACK)"] = float(stock_threshold_pack)
    if request_col and request_col in out.columns:
        request = pd.to_numeric(out[request_col], errors="coerce").fillna(0.0)
        out["재고부족(PACK)"] = np.where(stock.notna(), (request - stock).clip(lower=0.0), np.nan)
    else:
        out["재고부족(PACK)"] = np.where(stock.notna(), (float(stock_threshold_pack) - stock).clip(lower=0.0), np.nan)
    out["D-Day"] = out[due_col].map(d_day_text)
    out["우선등급"] = [
        priority_grade(shortage, due, current_stock, stock_threshold_pack)
        for shortage, due, current_stock in zip(out[shortage_col], out[due_col], out[stock_col])
    ]
    out["_priority_sort"] = out["우선등급"].map(priority_sort_value)
    out["_request_due_date_sort"] = pd.to_datetime(out[due_col], errors="coerce")
    return out


def calc_operation_kpis(
    product_summary: pd.DataFrame,
    code_summary: pd.DataFrame,
    stock_threshold_pack: float,
) -> dict[str, float]:
    today = pd.Timestamp.now(tz="Asia/Seoul").tz_localize(None).normalize()
    due_limit = today + pd.Timedelta(days=7)
    work = add_allocated_production_basis(with_operational_columns(code_summary))
    work["request_due_date"] = pd.to_datetime(work["request_due_date"], errors="coerce")
    if "available_stock_pack" not in work.columns:
        work["available_stock_pack"] = np.nan
    if "yongma_in_pack" not in work.columns:
        work["yongma_in_pack"] = 0.0
    work["_packing_shortage_pack"] = pd.to_numeric(
        work.get("code_packing_shortage_pack", pd.Series(0.0, index=work.index)),
        errors="coerce",
    ).fillna(0.0)
    work["_receipt_shortage_pack"] = pd.to_numeric(
        work.get("code_receipt_shortage_pack", pd.Series(0.0, index=work.index)),
        errors="coerce",
    ).fillna(0.0)
    product_priority = (
        work.groupby("base_product_name", dropna=False)
        .agg(
            packing_shortage_pack=("_receipt_shortage_pack", "sum"),
            request_due_date=("request_due_date", min_datetime),
            current_stock_pack=("available_stock_pack", sum_numeric_or_nan),
        )
        .reset_index()
    )
    product_priority["우선등급"] = [
        priority_grade(shortage, due, stock, stock_threshold_pack)
        for shortage, due, stock in zip(
            product_priority["packing_shortage_pack"],
            product_priority["request_due_date"],
            product_priority["current_stock_pack"],
        )
    ]
    shortage_mask = product_priority["packing_shortage_pack"] > 0
    due_mask = product_priority["request_due_date"].notna() & (product_priority["request_due_date"] <= due_limit)
    stock = pd.to_numeric(product_priority["current_stock_pack"], errors="coerce")
    stock_shortage_mask = stock.notna() & (stock <= float(stock_threshold_pack))
    priority_mask = product_priority["우선등급"].isin(["A 긴급", "B 주의"])
    # 전체/본품/샘플 KPI는 모두 같은 생산코드 요약 기준을 사용해 진도율 흔들림을 막는다.
    progress_kpi = calc_kpi_from_code_summary(code_summary)
    request_pack = float(progress_kpi.get("request_pack", 0.0))
    request_pcs = float(progress_kpi.get("request_pcs", 0.0))
    packing_pack = float(progress_kpi.get("packing_pack", 0.0))
    yongma_in_pack = float(progress_kpi.get("yongma_in_pack", 0.0))
    code_packing_shortage_pack = float(pd.to_numeric(work["_packing_shortage_pack"], errors="coerce").fillna(0.0).sum())
    code_receipt_wait_pack = float(
        pd.to_numeric(work.get("code_receipt_wait_pack", pd.Series(dtype=float)), errors="coerce").fillna(0.0).sum()
    )
    packing_shortage_pack = code_packing_shortage_pack or (
        float(product_summary["포장부족수량"].sum())
        if "포장부족수량" in product_summary.columns and not product_summary.empty
        else max(0.0, request_pack - packing_pack)
    )
    receipt_shortage_pack = float(progress_kpi.get("shortage_pack", max(0.0, request_pack - yongma_in_pack)))
    receipt_shortage_pcs = float(progress_kpi.get("shortage_pcs", 0.0))
    receipt_wait_pack = code_receipt_wait_pack or (
        float(product_summary["입고대기수량"].sum())
        if "입고대기수량" in product_summary.columns and not product_summary.empty
        else max(0.0, packing_pack - yongma_in_pack)
    )
    production_shortage_pcs = float(progress_kpi.get("production_shortage_pcs", 0.0))
    packable_pcs = float(progress_kpi.get("packable_pcs", max(0.0, request_pcs - production_shortage_pcs)))
    packing_progress = float(progress_kpi.get("packing_progress_pct", 0.0))
    receipt_progress = float(progress_kpi.get("progress_pct", 0.0))
    production_progress = float(progress_kpi.get("production_progress_pct", 0.0))
    return {
        "priority_products": float((shortage_mask & priority_mask).sum()),
        "urgent_products": float((shortage_mask & due_mask).sum()),
        "stock_shortage_products": float((shortage_mask & stock_shortage_mask).sum()),
        "request_pack": request_pack,
        "request_pcs": request_pcs,
        "yongma_in_pack": yongma_in_pack,
        "packing_done_pack": packing_pack,
        "packing_todo_pack": packing_shortage_pack,
        "receipt_shortage_pack": receipt_shortage_pack,
        "receipt_shortage_pcs": receipt_shortage_pcs,
        "receipt_wait_pack": receipt_wait_pack,
        "packing_shortage_pack": receipt_shortage_pack,
        "packing_shortage_pcs": receipt_shortage_pcs,
        "production_shortage_pcs": production_shortage_pcs,
        "packable_pcs": packable_pcs,
        "packing_progress_pct": min(100.0, max(0.0, packing_progress)),
        "receipt_progress_pct": min(100.0, max(0.0, receipt_progress)),
        "production_progress_pct": min(100.0, max(0.0, production_progress)),
    }


def render_operation_kpis(
    product_summary: pd.DataFrame,
    code_summary: pd.DataFrame,
    stock_threshold_pack: float,
    unit_mode: str = UNIT_PACK,
) -> None:
    kpi = calc_operation_kpis(product_summary, code_summary, stock_threshold_pack)
    if unit_mode == UNIT_PCS:
        render_metric_card_grid(
            [
                ("요청 PCS", format_int(kpi["request_pcs"]), "normal"),
                ("생산부족 PCS", format_int(kpi["production_shortage_pcs"]), "warn"),
                ("생산진도율", f"{kpi['production_progress_pct']:.1f}%", metric_progress_tone(kpi["production_progress_pct"])),
                ("포장진도율", f"{kpi['packing_progress_pct']:.1f}%", metric_progress_tone(kpi["packing_progress_pct"])),
                ("용마입고율", f"{kpi['receipt_progress_pct']:.1f}%", metric_progress_tone(kpi["receipt_progress_pct"])),
            ]
        )
        return
    render_metric_card_grid(
        [
            ("긴급 대응 품목 수", f"{int(kpi['priority_products']):,}", "normal"),
            ("재고부족 품목 수", f"{int(kpi['stock_shortage_products']):,}", "normal"),
            ("미입고 PACK", format_int(kpi["packing_shortage_pack"]), "warn"),
            ("생산부족 PCS", format_int(kpi["production_shortage_pcs"]), "warn"),
        ]
    )


def render_unit_selector(key: str) -> str:
    unit_mode = st.radio(
        "조회 단위 선택",
        UNIT_OPTIONS,
        index=0,
        horizontal=True,
        key=key,
    )
    if unit_mode == UNIT_PCS:
        st.caption("포장가능재고·생산부족 기준 조회")
    else:
        st.caption("용마입고·포장부족 기준 조회")
    return unit_mode


def visible_columns(df: pd.DataFrame, columns: list[str]) -> list[str]:
    return [col for col in columns if col in df.columns]


def dataframe_for_streamlit(df: pd.DataFrame) -> pd.DataFrame:
    if df is None or df.empty:
        return pd.DataFrame()
    return df.loc[:, ~pd.Index(df.columns).duplicated()].copy()


def product_progress_column_order(df: pd.DataFrame, pack_labels: list[str], unit_mode: str) -> list[str]:
    if unit_mode == UNIT_PCS:
        columns = [
            "우선등급",
            "D-Day",
            "제품명",
            "요청합계(PCS)",
            "생산필요수량(PCS)",
            "생산부족수량(PCS)",
            "생산진도율",
            "용마입고율",
        ]
    else:
        columns = [
            "우선등급",
            "D-Day",
            "제품명",
            *pack_labels,
            "요청합계(PACK)",
            "용마입고 PACK",
            "재고부족(PACK)",
            "미입고(PACK)",
            "생산부족수량(PCS)",
            "용마입고율",
            "생산진도율",
        ]
    return visible_columns(df, columns)


def production_progress_column_order(df: pd.DataFrame, pack_labels: list[str], unit_mode: str) -> list[str]:
    columns = [
        "생산코드",
        "기간구분",
        "대표 제품명",
        "요청합계(PACK)",
        "포장부족(PACK)",
        "포장가능재고(PCS)",
        "생산부족수량(PCS)",
        "기준차이",
        "생산진도율",
        "포장진도율",
        "생산완료예상일",
    ]
    return visible_columns(df, columns)


def production_power_detail_column_order(df: pd.DataFrame, pack_labels: list[str]) -> list[str]:
    columns = [
        "생산코드 전체",
        "기간구분",
        "대표 제품명",
        "POWER",
        *pack_labels,
        "요청합계(PACK)",
        "포장부족(PACK)",
        "포장가능재고(PCS)",
        *WIP_PROCESS_COLUMNS,
        "생산부족수량(PCS)",
        "기준차이",
        "생산진도율",
        "포장진도율",
        "생산완료예상일",
    ]
    return visible_columns(df, columns)


def sales_progress_column_order(df: pd.DataFrame, unit_mode: str) -> list[str]:
    if unit_mode == UNIT_PCS:
        columns = [
            "우선등급",
            "기간구분",
            "판매코드",
            "생산코드",
            "제품명",
            "POWER",
            "PACK",
            "생산요청물량(PCS)",
            "용마입고수량(PCS)",
            "용마입고대기수량(PCS)",
            "포장가능재고(PCS)",
            "포장부족(PCS)",
            "생산부족(PCS)",
            "생산진도율",
            "생산완료예상일",
            "상태",
        ]
    else:
        columns = [
            "우선등급",
            "기간구분",
            "판매코드",
            "생산코드",
            "제품명",
            "POWER",
            "PACK",
            "생산요청물량(PACK)",
            "용마입고수량(PACK)",
            "용마입고대기수량(PACK)",
            "포장부족(PACK)",
            "용마입고율",
            "생산완료예상일",
            "상태",
        ]
    return visible_columns(df, columns)


def sales_group_column_order(df: pd.DataFrame, unit_mode: str) -> list[str]:
    if unit_mode == UNIT_PCS:
        columns = [
            "우선등급",
            "기간구분",
            "판매코드",
            "대표 제품명",
            "POWER 수",
            "PACK",
            "생산요청물량(PCS)",
            "용마입고수량(PCS)",
            "용마입고대기수량(PCS)",
            "포장가능재고(PCS)",
            "포장부족(PCS)",
            "생산부족(PCS)",
            "생산진도율",
            "생산완료예상일",
            "상태",
        ]
    else:
        columns = [
            "우선등급",
            "기간구분",
            "판매코드",
            "대표 제품명",
            "POWER 수",
            "PACK",
            "생산요청물량(PACK)",
            "용마입고수량(PACK)",
            "용마입고대기수량(PACK)",
            "포장부족(PACK)",
            "용마입고율",
            "생산완료예상일",
            "상태",
        ]
    return visible_columns(df, columns)


def power_progress_column_order(df: pd.DataFrame, unit_mode: str) -> list[str]:
    if unit_mode == UNIT_PCS:
        columns = [
            "POWER",
            "기간구분",
            "요청합계(PCS)",
            "생산필요수량(PCS)",
            "생산부족수량(PCS)",
            "생산진도율",
            "포장진도율",
        ]
    else:
        columns = [
            "POWER",
            "기간구분",
            "요청합계(PACK)",
            "포장 PACK",
            "포장부족(PACK)",
            "생산부족수량(PCS)",
            "포장진도율",
            "생산진도율",
        ]
    return visible_columns(df, columns)


def build_product_progress_main_view(
    product_summary: pd.DataFrame,
    code_summary: pd.DataFrame,
    pack_labels: list[str],
    stock_threshold_pack: float = INVENTORY_STOCK_THRESHOLD_DEFAULT,
) -> pd.DataFrame:
    work = with_operational_columns(code_summary)
    due_by_product = (
        work.groupby("base_product_name", dropna=False)
        .agg(
            production_due_date=("production_due_date", min_datetime),
            request_due_date=("request_due_date", min_datetime),
        )
        .reset_index()
        .rename(columns={"base_product_name": "제품명"})
    )
    pivot = build_pack_pivot(work, ["base_product_name"], pack_labels).rename(columns={"base_product_name": "제품명"})

    out = product_summary.merge(pivot, on="제품명", how="left").merge(due_by_product, on="제품명", how="left")
    for label in pack_labels:
        out[label] = out[label].fillna(0.0)
    if "용마창고재고 (PACK)" not in out.columns:
        out["용마창고재고 (PACK)"] = np.nan
    if "재고매칭SKU수" not in out.columns:
        out["재고매칭SKU수"] = 0
    out["제품필요수량"] = out.get("생산부족수량", 0.0)
    out["생산필요수량(PCS)"] = out["제품필요수량"]
    out["생산부족수량(PCS)"] = out.get("생산부족수량", 0.0)
    out["진도율"] = out.get("생산진도율", 0.0)
    out["전체진도율"] = out["용마입고율"]
    out = add_priority_columns(out, stock_threshold_pack, shortage_col="미입고수량", request_col="요청 PACK")
    out = out.rename(columns={"요청 PACK": "요청합계(PACK)", "요청 PCS": "요청합계(PCS)"})
    out["포장부족(PACK)"] = out["포장부족수량"]
    out["미입고(PACK)"] = out["미입고수량"]
    ordered = [
        "우선등급",
        "D-Day",
        "제품명",
        *pack_labels,
        "요청합계(PACK)",
        "요청합계(PCS)",
        "용마입고 PACK",
        "용마창고재고 (PACK)",
        "재고부족(PACK)",
        "미입고(PACK)",
        "포장부족(PACK)",
        "생산필요수량(PCS)",
        "생산부족수량(PCS)",
        "제품필요수량",
        "진도율",
        "생산진도율",
        "포장부족수량",
        "미입고수량",
        "포장진도율",
        "용마입고율",
        "전체진도율",
        "상태",
        "_priority_sort",
        "_request_due_date_sort",
    ]
    return out[ordered].copy()


def build_product_sku_detail_view(code_summary: pd.DataFrame, product_name: str) -> pd.DataFrame:
    work = with_operational_columns(code_summary)
    scope = work[work["base_product_name"] == product_name].copy()
    if scope.empty:
        scope = work[work["product_name"] == product_name].copy()
    if scope.empty:
        return pd.DataFrame(columns=["SKU", "생산코드", "판매코드 수", "요청 PACK", "용마입고 PACK", "미입고 PACK", "용마입고율"])
    grouped = (
        scope.groupby(["product_name", "production_code_display"], dropna=False)
        .agg(
            sales_code_count=("sales_code", "nunique"),
            request_pack=("request_pack", "sum"),
            packing_pack=("packing_recognized_pack", "sum"),
            yongma_in_pack=("yongma_recognized_pack", "sum"),
        )
        .reset_index()
        .rename(
            columns={
                "product_name": "SKU",
                "production_code_display": "생산코드",
                "sales_code_count": "판매코드 수",
                "request_pack": "요청 PACK",
                "yongma_in_pack": "용마입고 PACK",
            }
        )
    )
    grouped["미입고 PACK"] = (grouped["요청 PACK"] - grouped["용마입고 PACK"]).clip(lower=0.0)
    grouped["용마입고율"] = np.where(grouped["요청 PACK"] > 0, grouped["용마입고 PACK"] / grouped["요청 PACK"] * 100.0, 0.0)
    grouped["용마입고율"] = np.clip(grouped["용마입고율"], 0.0, 100.0)
    return grouped[["SKU", "생산코드", "판매코드 수", "요청 PACK", "용마입고 PACK", "미입고 PACK", "용마입고율"]].sort_values(
        ["미입고 PACK", "요청 PACK"], ascending=[False, False], kind="stable"
    )


def compact_query_text(value: Any) -> str:
    return re.sub(r"\s+", "", clean_str(value)).lower()


def compact_search_text(value: Any) -> str:
    return re.sub(r"[\W_]+", "", clean_str(value), flags=re.UNICODE).lower()


def expand_product_query_terms(query: str) -> list[str]:
    text = clean_str(query)
    if not text:
        return []
    terms = [text]
    compact = compact_search_text(text)
    for alias, values in PRODUCT_QUERY_ALIASES.items():
        alias_compact = compact_search_text(alias)
        if alias_compact and (alias_compact in compact or compact in alias_compact):
            terms.extend(values)
    return list(dict.fromkeys([term for term in terms if clean_str(term)]))


def contains_any_query_term(series: pd.Series, terms: list[str]) -> pd.Series:
    if not terms:
        return pd.Series(False, index=series.index)
    text = series.astype(str).str.lower()
    compact_text = series.astype(str).map(compact_search_text)
    mask = pd.Series(False, index=series.index)
    for term in terms:
        raw_term = clean_str(term).lower()
        compact_term = compact_search_text(term)
        if raw_term:
            mask = mask | text.str.contains(raw_term, na=False, regex=False)
        if compact_term:
            mask = mask | compact_text.str.contains(compact_term, na=False, regex=False)
    return mask


def split_lookup_query_terms(query: str) -> list[str]:
    text = clean_str(query)
    if not text:
        return []
    tokens = [clean_str(token) for token in re.split(r"[,，]+|\s+", text)]
    return [token for token in tokens if token]


def parse_quick_lookup_direct_query(
    query: str,
    pack_options: list[str],
    power_options: list[str],
) -> tuple[str, str, list[str]]:
    text = clean_str(query)
    if not text:
        return "", "전체", []

    if "," in text or "，" in text:
        tokens = [clean_str(token) for token in re.split(r"[,，]+", text)]
    else:
        tokens = [clean_str(token) for token in text.split()]
    tokens = [token for token in tokens if token]

    available_packs = set(pack_options)
    available_powers = set(power_options)
    product_terms: list[str] = []
    pack_label = "전체"
    power_labels: list[str] = []

    for token in tokens:
        normalized = token.replace("−", "-").replace("–", "-").replace("—", "-")
        pack_match = re.fullmatch(r"(\d+(?:\.\d+)?)\s*(?:P|팩|개입)", normalized, flags=re.IGNORECASE)
        if pack_match:
            candidate_pack = base_pack_label(float(pack_match.group(1)))
            pack_label = candidate_pack if candidate_pack in available_packs else candidate_pack
            continue

        power_match = re.fullmatch(r"[+-]?\d+(?:\.\d+)?", normalized)
        if power_match and (normalized.startswith(("+", "-")) or "." in normalized):
            candidate_power = format_power(float(normalized))
            power_labels.append(candidate_power if candidate_power in available_powers else candidate_power)
            continue

        product_terms.append(token)

    return " ".join(product_terms).strip(), pack_label, list(dict.fromkeys(power_labels))


def build_product_pack_power_quick_view(
    code_summary: pd.DataFrame,
    product_query: str,
    pack_label: str,
    power_labels: list[str],
) -> pd.DataFrame:
    columns = [
        "제품명",
        "SKU",
        "PACK",
        "POWER",
        "판매코드 수",
        "요청 PACK",
        "포장 PACK",
        "용마입고 PACK",
        "입고대기 PACK",
        "미입고 PACK",
        "요청 PCS",
        "생산부족 PCS",
        "용마입고율",
        "생산진도율",
        "생산완료예상일",
    ]
    if code_summary.empty:
        return pd.DataFrame(columns=columns)

    work = add_allocated_production_basis(with_operational_columns(code_summary))
    query = product_query.strip()
    if query:
        terms = expand_product_query_terms(query)
        mask = (
            contains_any_query_term(work["base_product_name"], terms)
            | contains_any_query_term(work["product_name"], terms)
            | contains_any_query_term(work["sales_code"], terms)
            | contains_any_query_term(work["production_code_display"], terms)
        )
        work = work[mask].copy()
    if pack_label != "전체":
        work = work[work["_pack_label"] == pack_label].copy()
    if power_labels:
        work = work[work["POWER"].isin(power_labels)].copy()
    if work.empty:
        return pd.DataFrame(columns=columns)
    work = refresh_scope_production_shortage(work)
    if "yongma_in_pack" not in work.columns:
        work["yongma_in_pack"] = 0.0

    grouped = (
        work.groupby(["base_product_name", "product_name", "_pack_label", "POWER", "power_value"], dropna=False)
        .agg(
            sales_code_count=("sales_code", "nunique"),
            request_pack=("request_pack", "sum"),
            packing_pack=("packing_recognized_pack", "sum"),
            yongma_in_pack=("yongma_recognized_pack", "sum"),
            request_pcs=("request_pcs", "sum"),
            production_shortage_pcs=("_allocated_production_shortage_qty", "sum"),
            production_due_date=("production_due_date", max_datetime),
        )
        .reset_index()
        .rename(
            columns={
                "base_product_name": "제품명",
                "product_name": "SKU",
                "_pack_label": "PACK",
                "sales_code_count": "판매코드 수",
                "request_pack": "요청 PACK",
                "packing_pack": "포장 PACK",
                "yongma_in_pack": "용마입고 PACK",
                "request_pcs": "요청 PCS",
                "production_shortage_pcs": "생산부족 PCS",
            }
        )
    )
    grouped["입고대기 PACK"] = (grouped["포장 PACK"] - grouped["용마입고 PACK"]).clip(lower=0.0)
    grouped["미입고 PACK"] = (grouped["요청 PACK"] - grouped["용마입고 PACK"]).clip(lower=0.0)
    grouped["용마입고율"] = np.where(
        grouped["요청 PACK"] > 0,
        grouped["용마입고 PACK"] / grouped["요청 PACK"] * 100.0,
        0.0,
    )
    grouped["용마입고율"] = np.clip(grouped["용마입고율"], 0.0, 100.0)
    grouped["생산부족 PCS"] = pd.to_numeric(grouped["생산부족 PCS"], errors="coerce").fillna(0.0).round(0)
    grouped["생산진도율"] = calc_production_progress_pct(grouped["요청 PCS"], grouped["생산부족 PCS"])
    grouped["생산완료예상일"] = grouped["production_due_date"].map(display_date_or_dash)
    grouped["_pack_sort"] = grouped["PACK"].map(pack_sort_key)
    return grouped.sort_values(
        ["power_value", "_pack_sort", "미입고 PACK", "요청 PACK"],
        ascending=[True, True, False, False],
        na_position="last",
        kind="stable",
    )[columns].copy()


def render_product_pack_power_quick_lookup(code_summary: pd.DataFrame) -> None:
    st.markdown("<div style='height:10px'></div>", unsafe_allow_html=True)
    render_panel_title(
        "제품·PACK·POWER 간편 조회",
        "제품명 일부, PACK, POWER 조합으로 포장·용마입고·생산 상태를 확인합니다.",
    )
    pack_options = available_pack_options(code_summary)
    power_options = available_power_options(code_summary)
    direct_query = st.text_input(
        "직접 검색",
        value="",
        placeholder="예: 소울브라운, 40P, -06.50",
        key="quick_lookup_direct_query",
    )
    q1, q2, q3 = st.columns([2.3, 1.1, 2.0], gap="small")
    with q1:
        product_query = st.text_input(
            "제품명/SKU/코드 검색",
            value="",
            placeholder="예: 소울브라운",
            key="quick_lookup_product_query",
        )
    with q2:
        pack_label = st.selectbox(
            "PACK 선택",
            options=pack_options,
            index=0,
            key="quick_lookup_pack",
        )
    with q3:
        power_labels = st.multiselect(
            "POWER 선택",
            options=power_options[1:],
            default=[],
            key="quick_lookup_power",
        )
    if direct_query.strip():
        product_query, pack_label, power_labels = parse_quick_lookup_direct_query(
            direct_query,
            pack_options=pack_options,
            power_options=power_options[1:],
        )

    if not product_query.strip() and pack_label == "전체" and not power_labels:
        return

    quick_view = build_product_pack_power_quick_view(code_summary, product_query, pack_label, power_labels)
    total_request = float(quick_view["요청 PACK"].sum()) if not quick_view.empty else 0.0
    total_yongma = float(quick_view["용마입고 PACK"].sum()) if not quick_view.empty else 0.0
    total_shortage = float(quick_view["미입고 PACK"].sum()) if not quick_view.empty else 0.0
    total_request_pcs = float(quick_view["요청 PCS"].sum()) if not quick_view.empty else 0.0
    total_production_shortage = float(quick_view["생산부족 PCS"].sum()) if not quick_view.empty else 0.0
    receipt_progress = (total_yongma / total_request * 100.0) if total_request > 0 else 0.0
    production_progress = (
        (total_request_pcs - total_production_shortage) / total_request_pcs * 100.0
        if total_request_pcs > 0
        else 0.0
    )
    render_metric_card_grid(
        [
            ("요청 PACK", format_int(total_request), "normal"),
            ("용마입고 PACK", format_int(total_yongma), "normal"),
            ("미입고 PACK", format_int(total_shortage), "warn" if total_shortage > 0 else "normal"),
            ("생산부족 PCS", format_int(total_production_shortage), "warn" if total_production_shortage > 0 else "normal"),
            ("용마입고율", f"{min(100.0, max(0.0, receipt_progress)):.1f}%", metric_progress_tone(receipt_progress)),
            ("생산진도율", f"{min(100.0, max(0.0, production_progress)):.1f}%", metric_progress_tone(production_progress)),
        ]
    )
    st.dataframe(
        quick_view,
        hide_index=True,
        height=dataframe_auto_height(len(quick_view), 360),
        width="stretch",
        column_config=drilldown_column_config(),
    )


def extract_sales_prefix(value: Any) -> str:
    text = clean_str(value).upper()
    match = re.match(r"^([A-Z]\d+)", text)
    return match.group(1) if match else ""


def build_daily_request_match_view(code_summary: pd.DataFrame) -> pd.DataFrame:
    columns = [
        "제품코드",
        "기간구분",
        "PACK",
        "POWER",
        "CP",
        "요청 PACK",
        "포장 PACK",
        "용마입고 PACK",
        "미입고 PACK",
        "요청 PCS",
        "생산부족 PCS",
        "용마입고대기 PACK",
        "포장가능재고(PCS)",
        "샘플신청가능수량",
        "생산진도율",
        "생산완료예상일",
        "요청제품명",
        "판매코드 수",
    ]
    if code_summary.empty:
        return pd.DataFrame(columns=columns)

    work = add_allocated_production_basis(with_operational_columns(code_summary))
    work["_sales_prefix"] = work["sales_code"].map(extract_sales_prefix)
    work["CP"] = work["sales_code"].map(cp_label_from_sales_code)
    work = work[work["_sales_prefix"] != ""].copy()
    if work.empty:
        return pd.DataFrame(columns=columns)
    if "yongma_in_pack" not in work.columns:
        work["yongma_in_pack"] = 0.0

    grouped = (
        work.groupby(["_sales_prefix", "_pack_label", "POWER", "CP"], dropna=False)
        .agg(
            request_pack=("request_pack", "sum"),
            packing_pack=("packing_recognized_pack", "sum"),
            yongma_in_pack=("yongma_recognized_pack", "sum"),
            factory_group=("factory_group", join_factory_groups),
            period_group=("period_group", first_nonempty),
            request_pcs=("request_pcs", "sum"),
            production_shortage_pcs=("_allocated_production_shortage_qty", "sum"),
            sample_available_pcs=("_allocated_sample_available_pcs", "sum"),
            production_due_date=("production_due_date", max_datetime),
            product_name=("product_name", first_nonempty),
            sales_code_count=("sales_code", "nunique"),
        )
        .reset_index()
        .rename(
            columns={
                "_sales_prefix": "제품코드",
                "_pack_label": "PACK",
                "request_pack": "요청 PACK",
                "packing_pack": "포장 PACK",
                "yongma_in_pack": "용마입고 PACK",
                "factory_group": "공장구분",
                "period_group": "기간구분",
                "request_pcs": "요청 PCS",
                "production_shortage_pcs": "생산부족 PCS",
                "sample_available_pcs": "샘플신청가능수량",
                "product_name": "요청제품명",
                "sales_code_count": "판매코드 수",
            }
        )
    )
    grouped["미입고 PACK"] = (grouped["요청 PACK"] - grouped["용마입고 PACK"]).clip(lower=0.0)
    grouped["용마입고대기 PACK"] = (grouped["포장 PACK"] - grouped["용마입고 PACK"]).clip(lower=0.0)
    grouped["포장가능재고(PCS)"] = (
        grouped["요청 PCS"] - grouped["생산부족 PCS"] + grouped["샘플신청가능수량"]
    ).clip(lower=0.0)
    grouped["생산진도율"] = calc_production_progress_pct(grouped["요청 PCS"], grouped["생산부족 PCS"])
    grouped["생산완료예상일"] = grouped["production_due_date"].map(display_date_or_dash)
    return grouped[columns].copy()


def build_daily_production_power_catalog(code_summary: pd.DataFrame) -> pd.DataFrame:
    columns = [
        "제품코드",
        "POWER",
        "CP",
        "_production_request_pack",
        "_production_request_pcs",
        "_production_shortage_pcs",
        "_production_sample_available_pcs",
        "_production_available_stock_pcs",
        "_production_progress_pct",
    ]
    if code_summary.empty:
        return pd.DataFrame(columns=columns)

    work = add_allocated_production_basis(with_operational_columns(code_summary))
    work["_sales_prefix"] = work["sales_code"].map(extract_sales_prefix)
    work["CP"] = work["sales_code"].map(cp_label_from_sales_code)
    work["_production_key"] = work.get("production_code_key", pd.Series("", index=work.index)).map(clean_str)
    fallback_key = work.get("sales_code_key", pd.Series("", index=work.index)).map(clean_str)
    work["_production_key"] = work["_production_key"].where(work["_production_key"] != "", fallback_key)
    work = work[
        (work["_sales_prefix"].map(clean_str) != "")
        & (work["POWER"].map(clean_str) != "")
        & (work["_production_key"].map(clean_str) != "")
    ].copy()
    if work.empty:
        return pd.DataFrame(columns=columns)

    for col in ["request_pack", "request_pcs", "_allocated_production_shortage_qty", "sample_available_pcs"]:
        if col not in work.columns:
            work[col] = 0.0
        work[col] = pd.to_numeric(work[col], errors="coerce").fillna(0.0)

    by_production = (
        work.groupby(["_sales_prefix", "POWER", "CP", "_production_key"], dropna=False)
        .agg(
            request_pack=("request_pack", "sum"),
            request_pcs=("request_pcs", "sum"),
            production_shortage_pcs=("_allocated_production_shortage_qty", "sum"),
            sample_available_pcs=("sample_available_pcs", "max"),
        )
        .reset_index()
    )
    grouped = (
        by_production.groupby(["_sales_prefix", "POWER", "CP"], dropna=False)
        .agg(
            _production_request_pack=("request_pack", "sum"),
            _production_request_pcs=("request_pcs", "sum"),
            _production_shortage_pcs=("production_shortage_pcs", "sum"),
            _production_sample_available_pcs=("sample_available_pcs", "sum"),
        )
        .reset_index()
        .rename(columns={"_sales_prefix": "제품코드"})
    )
    grouped["_production_available_stock_pcs"] = (
        grouped["_production_request_pcs"]
        - grouped["_production_shortage_pcs"]
        + grouped["_production_sample_available_pcs"]
    ).clip(lower=0.0)
    grouped["_production_progress_pct"] = calc_production_progress_pct(
        grouped["_production_request_pcs"],
        grouped["_production_shortage_pcs"],
    )
    for col in [
        "_production_request_pack",
        "_production_request_pcs",
        "_production_shortage_pcs",
        "_production_sample_available_pcs",
        "_production_available_stock_pcs",
    ]:
        grouped[col] = pd.to_numeric(grouped[col], errors="coerce").fillna(0.0).round(0)
    return grouped[columns].copy()


def build_daily_base_power_production_catalog(code_summary: pd.DataFrame) -> pd.DataFrame:
    columns = [
        "_daily_base_product_name",
        "POWER",
        "CP",
        "_base_production_request_pack",
        "_base_production_request_pcs",
        "_base_production_shortage_pcs",
        "_base_production_sample_available_pcs",
        "_base_production_available_stock_pcs",
        "_base_production_progress_pct",
    ]
    if code_summary.empty:
        return pd.DataFrame(columns=columns)

    work = add_allocated_production_basis(with_operational_columns(code_summary))
    work["CP"] = work["sales_code"].map(cp_label_from_sales_code)
    if "base_product_name" in work.columns:
        work["_daily_base_product_name"] = work["base_product_name"].map(clean_str)
    else:
        work["_daily_base_product_name"] = work["product_name"].map(strip_pack_unit_suffix).map(clean_str)
    work["_production_key"] = work.get("production_code_key", pd.Series("", index=work.index)).map(clean_str)
    fallback_key = work.get("sales_code_key", pd.Series("", index=work.index)).map(clean_str)
    work["_production_key"] = work["_production_key"].where(work["_production_key"] != "", fallback_key)
    work = work[
        (work["_daily_base_product_name"].map(clean_str) != "")
        & (work["POWER"].map(clean_str) != "")
        & (work["_production_key"].map(clean_str) != "")
    ].copy()
    if work.empty:
        return pd.DataFrame(columns=columns)

    for col in ["request_pack", "request_pcs", "_allocated_production_shortage_qty", "sample_available_pcs"]:
        if col not in work.columns:
            work[col] = 0.0
        work[col] = pd.to_numeric(work[col], errors="coerce").fillna(0.0)

    by_production = (
        work.groupby(["_daily_base_product_name", "POWER", "CP", "_production_key"], dropna=False)
        .agg(
            request_pack=("request_pack", "sum"),
            request_pcs=("request_pcs", "sum"),
            production_shortage_pcs=("_allocated_production_shortage_qty", "sum"),
            sample_available_pcs=("sample_available_pcs", "max"),
        )
        .reset_index()
    )
    grouped = (
        by_production.groupby(["_daily_base_product_name", "POWER", "CP"], dropna=False)
        .agg(
            _base_production_request_pack=("request_pack", "sum"),
            _base_production_request_pcs=("request_pcs", "sum"),
            _base_production_shortage_pcs=("production_shortage_pcs", "sum"),
            _base_production_sample_available_pcs=("sample_available_pcs", "sum"),
        )
        .reset_index()
    )
    grouped["_base_production_available_stock_pcs"] = (
        grouped["_base_production_request_pcs"]
        - grouped["_base_production_shortage_pcs"]
        + grouped["_base_production_sample_available_pcs"]
    ).clip(lower=0.0)
    grouped["_base_production_progress_pct"] = calc_production_progress_pct(
        grouped["_base_production_request_pcs"],
        grouped["_base_production_shortage_pcs"],
    )
    for col in [
        "_base_production_request_pack",
        "_base_production_request_pcs",
        "_base_production_shortage_pcs",
        "_base_production_sample_available_pcs",
        "_base_production_available_stock_pcs",
    ]:
        grouped[col] = pd.to_numeric(grouped[col], errors="coerce").fillna(0.0).round(0)
    return grouped[columns].copy()


def pack_unit_from_label(value: Any) -> float:
    match = re.search(r"(\d+(?:\.\d+)?)", clean_str(value))
    if not match:
        return 1.0
    try:
        number = float(match.group(1))
    except ValueError:
        return 1.0
    return number if number > 0 else 1.0


def build_item_code_from_prefix_power(product_code: Any, power_label: Any, cp_label: Any = "") -> str:
    prefix = clean_str(product_code).upper()
    power = clean_str(power_label)
    if not prefix or not power:
        return ""
    cp = clean_str(cp_label)
    return f"{prefix}{power}-{cp}" if cp else f"{prefix}{power}"


def replace_power_in_production_code(template: Any, source_power: Any, target_power: Any) -> str:
    text = clean_str(template)
    target = clean_str(target_power)
    source = clean_str(source_power)
    if not text or not target:
        return ""
    if source and source in text:
        return text.replace(source, target, 1)
    match = re.search(r"[+-]\d{1,2}\.\d{2}", text)
    if not match:
        return ""
    return f"{text[:match.start()]}{target}{text[match.end():]}"


def build_sample_available_lookup(sample_available_df: pd.DataFrame | None) -> dict[str, float]:
    if sample_available_df is None or sample_available_df.empty:
        return {}
    if "production_code_key" not in sample_available_df.columns or "sample_available_pcs" not in sample_available_df.columns:
        return {}
    grouped = (
        sample_available_df.copy()
        .assign(sample_available_pcs=lambda df: pd.to_numeric(df["sample_available_pcs"], errors="coerce").fillna(0.0))
        .groupby("production_code_key", dropna=False)["sample_available_pcs"]
        .sum()
    )
    return {clean_str(key): float(value) for key, value in grouped.items() if clean_str(key)}


def build_daily_product_catalog(code_summary: pd.DataFrame) -> pd.DataFrame:
    columns = ["제품코드", "PACK", "마스터제품명", "마스터공장구분", "생산코드템플릿", "생산코드POWER"]
    if code_summary.empty:
        return pd.DataFrame(columns=columns)

    work = add_allocated_production_basis(with_operational_columns(code_summary))
    work["_sales_prefix"] = work["sales_code"].map(extract_sales_prefix)
    work = work[(work["_sales_prefix"] != "") & (work["_pack_label"].map(clean_str) != "")].copy()
    if work.empty:
        return pd.DataFrame(columns=columns)

    grouped = (
        work.groupby(["_sales_prefix", "_pack_label"], dropna=False)
        .agg(
            product_name=("product_name", first_nonempty),
            factory_group=("factory_group", join_factory_groups),
            production_code=("production_code", first_nonempty),
            power=("POWER", first_nonempty),
        )
        .reset_index()
        .rename(
            columns={
                "_sales_prefix": "제품코드",
                "_pack_label": "PACK",
                "product_name": "마스터제품명",
                "factory_group": "마스터공장구분",
                "production_code": "생산코드템플릿",
                "power": "생산코드POWER",
            }
        )
    )
    return grouped[columns].copy()


def build_daily_code_power_catalog(code_summary: pd.DataFrame) -> pd.DataFrame:
    columns = ["제품코드", "POWER", "CP", "_code_pack", "_code_product_name"]
    if code_summary.empty:
        return pd.DataFrame(columns=columns)

    work = with_operational_columns(code_summary)
    work["_sales_prefix"] = work["sales_code"].map(extract_sales_prefix)
    work["CP"] = work["sales_code"].map(cp_label_from_sales_code)
    work = work[(work["_sales_prefix"] != "") & (work["POWER"].map(clean_str) != "")].copy()
    if work.empty:
        return pd.DataFrame(columns=columns)

    grouped = (
        work.groupby(["_sales_prefix", "POWER", "CP"], dropna=False)
        .agg(
            pack=("_pack_label", first_nonempty),
            product_name=("product_name", first_nonempty),
        )
        .reset_index()
        .rename(
            columns={
                "_sales_prefix": "제품코드",
                "pack": "_code_pack",
                "product_name": "_code_product_name",
            }
        )
    )
    return grouped[columns].copy()


def build_daily_code_catalog(code_summary: pd.DataFrame) -> pd.DataFrame:
    columns = ["제품코드", "_code_default_pack", "_code_default_product_name"]
    if code_summary.empty:
        return pd.DataFrame(columns=columns)

    work = with_operational_columns(code_summary)
    work["_sales_prefix"] = work["sales_code"].map(extract_sales_prefix)
    work = work[work["_sales_prefix"] != ""].copy()
    if work.empty:
        return pd.DataFrame(columns=columns)

    grouped = (
        work.groupby("_sales_prefix", dropna=False)
        .agg(
            pack=("_pack_label", first_nonempty),
            product_name=("product_name", first_nonempty),
        )
        .reset_index()
        .rename(
            columns={
                "_sales_prefix": "제품코드",
                "pack": "_code_default_pack",
                "product_name": "_code_default_product_name",
            }
        )
    )
    return grouped[columns].copy()


def fill_daily_product_code_from_code_summary(out: pd.DataFrame, code_summary: pd.DataFrame) -> pd.DataFrame:
    if out.empty or code_summary.empty:
        return out

    work = with_operational_columns(code_summary)
    work["_sales_prefix"] = work["sales_code"].map(extract_sales_prefix)
    work["CP"] = work["sales_code"].map(cp_label_from_sales_code)
    work = work[
        (work["_sales_prefix"].map(clean_str) != "")
        & (work["POWER"].map(clean_str) != "")
        & (work["CP"].map(clean_str) != "")
        & (work["product_name"].map(clean_str) != "")
    ].copy()
    if work.empty:
        return out

    catalog = (
        work.groupby(["_sales_prefix", "_pack_label", "POWER", "CP"], dropna=False)
        .agg(product_name=("product_name", first_nonempty))
        .reset_index()
        .rename(columns={"_sales_prefix": "제품코드", "_pack_label": "PACK"})
    )
    catalog_groups: dict[tuple[str, str], list[tuple[str, str, str, str]]] = {}
    for code, pack, power, cp, product_name in catalog[["제품코드", "PACK", "POWER", "CP", "product_name"]].itertuples(index=False, name=None):
        product_text = clean_str(product_name)
        product_compact = compact_search_text(product_text)
        catalog_groups.setdefault((clean_str(power), clean_str(cp)), []).append(
            (clean_str(code), clean_str(pack), product_text, product_compact)
        )

    filled = out.copy()
    if "CP" not in filled.columns:
        filled["CP"] = ""
    needs_code = filled["제품코드"].map(clean_str) == ""
    for idx, row in filled[needs_code].iterrows():
        power = clean_str(row.get("POWER", ""))
        cp = clean_str(row.get("CP", ""))
        if not power or not cp:
            continue
        candidates = catalog_groups.get((power, cp), [])
        if not candidates:
            continue
        terms = expand_product_query_terms(row.get("제품명", ""))
        term_compacts = [compact_search_text(term) for term in terms if clean_str(term)]
        matched = [
            (code, pack, product_name)
            for code, pack, product_name, product_compact in candidates
            if any(term and (term in product_compact or product_compact in term) for term in term_compacts)
        ]
        product_codes = list(dict.fromkeys([code for code, _pack, _product_name in matched if clean_str(code)]))
        if len(product_codes) != 1:
            continue
        filled.at[idx, "제품코드"] = product_codes[0]
        pack = first_nonempty([pack for _code, pack, _product_name in matched])
        if pack and clean_str(filled.at[idx, "PACK"]) == "":
            filled.at[idx, "PACK"] = pack
        product_name = first_nonempty([product_name for _code, _pack, product_name in matched])
        if product_name and clean_str(filled.at[idx, "제품명"]) == "":
            filled.at[idx, "제품명"] = product_name
    return filled


def enrich_daily_inventory_from_code_summary(daily_inventory_df: pd.DataFrame, code_summary: pd.DataFrame) -> pd.DataFrame:
    if daily_inventory_df is None or daily_inventory_df.empty or code_summary.empty:
        return daily_inventory_df

    exact_catalog = build_daily_code_power_catalog(code_summary)
    code_catalog = build_daily_code_catalog(code_summary)
    if exact_catalog.empty and code_catalog.empty:
        return daily_inventory_df

    out = daily_inventory_df.copy()
    out["제품코드"] = out["제품코드"].map(clean_str).str.upper()
    out["POWER"] = out["POWER"].map(clean_str)
    out["CP"] = out["CP"].map(clean_str) if "CP" in out.columns else ""
    out = fill_daily_product_code_from_code_summary(out, code_summary)
    if not exact_catalog.empty:
        out = out.merge(exact_catalog, on=["제품코드", "POWER", "CP"], how="left")
        catalog_pack = out["_code_pack"].map(clean_str)
        catalog_product = out["_code_product_name"].map(clean_str)
        out["PACK"] = out["_code_pack"].where(catalog_pack != "", out["PACK"])
        out["제품명"] = out["_code_product_name"].where(catalog_product != "", out["제품명"])
        out = out.drop(columns=["_code_pack", "_code_product_name"], errors="ignore")
    if not code_catalog.empty:
        out = out.merge(code_catalog, on="제품코드", how="left")
        current_pack = out["PACK"].map(clean_str)
        current_product = out["제품명"].map(clean_str)
        out["PACK"] = out["PACK"].where(current_pack != "", out["_code_default_pack"])
        out["제품명"] = out["제품명"].where(current_product != "", out["_code_default_product_name"])
        out = out.drop(columns=["_code_default_pack", "_code_default_product_name"], errors="ignore")
    return out[DAILY_INVENTORY_COLUMNS].copy()


def classify_daily_inventory_status(row: pd.Series) -> str:
    urgent = bool(row.get("긴급요청", False))
    stock = pd.to_numeric(row.get("재고수량", np.nan), errors="coerce")
    request_pack = pd.to_numeric(row.get("요청 PACK", 0.0), errors="coerce")
    has_request = pd.notna(request_pack) and float(request_pack) > 0
    stock_negative = pd.notna(stock) and float(stock) < 0
    if urgent and has_request:
        return "요청내 긴급"
    if urgent:
        return "요청외 긴급"
    if stock_negative and has_request:
        return "요청내 재고부족"
    if stock_negative:
        return "재고 음수"
    if has_request:
        return "요청내 재고확인"
    return "재고 모니터링"


def complete_daily_response_mask(df: pd.DataFrame) -> pd.Series:
    required_cols = ["품목코드", "제품명", "제품코드", "PACK", "POWER"]
    mask = pd.Series(True, index=df.index)
    for col in required_cols:
        if col not in df.columns:
            return pd.Series(False, index=df.index)
        mask &= df[col].map(clean_str) != ""
    return mask


def daily_item_code_base(value: Any) -> str:
    match = re.match(r"^(S\d{3})", clean_str(value).upper())
    return match.group(1) if match else ""


def canonical_product_base_from_alias(value: Any) -> str:
    text = clean_str(value)
    compact = compact_search_text(text)
    if not compact:
        return ""

    for alias, values in PRODUCT_QUERY_ALIASES.items():
        if alias in GENERIC_PRODUCT_ALIAS_KEYS:
            continue
        alias_compact = compact_search_text(alias)
        candidates = [candidate for candidate in values if clean_str(candidate)]
        candidate_compacts = [compact_search_text(candidate) for candidate in candidates]
        matched = (
            alias_compact and alias_compact in compact
        ) or any(candidate and (candidate in compact or compact in candidate) for candidate in candidate_compacts)
        if not matched:
            continue

        iris_candidate = next(
            (candidate for candidate in candidates if clean_str(candidate).lower().startswith("iris ")),
            "",
        )
        if iris_candidate:
            return clean_str(iris_candidate)
        if alias in IRIS_PRODUCT_ALIAS_KEYS and candidates:
            return f"Iris {clean_str(candidates[0])}"
        if candidates:
            return clean_str(candidates[0])

    return ""


def daily_product_name_needs_standardization(value: Any) -> bool:
    text = clean_str(value)
    if not text:
        return False
    if PACK_PREFIX_SUFFIX_RE.match(text):
        return True
    hangul_marker_text = re.sub(r"(팩|개입)", "", text, flags=re.IGNORECASE)
    return bool(re.search(r"[가-힣]", hangul_marker_text))


def daily_standard_product_name(product_name: Any, pack_label: Any = "", item_code: Any = "") -> str:
    item_base = daily_item_code_base(item_code)
    standard_name = clean_str(DAILY_ITEM_STANDARD.get(item_base, {}).get("product_name", ""))
    if standard_name:
        return standard_name

    raw_name = clean_str(product_name)
    if not daily_product_name_needs_standardization(raw_name):
        return raw_name

    standard_base = canonical_product_base_from_alias(raw_name)
    if not standard_base:
        return raw_name

    pack_unit = pack_unit_from_label(pack_label) if clean_str(pack_label) else extract_pack_unit(raw_name)
    if pd.notna(pack_unit) and float(pack_unit) > 0:
        return f"{standard_base}_{float(pack_unit):g}팩"
    return standard_base


def apply_daily_item_standard(view: pd.DataFrame, product_col: str = "대표 제품명") -> pd.DataFrame:
    if view.empty or "품목코드" not in view.columns:
        return view
    out = view.copy()
    item_codes = out["품목코드"].map(daily_item_code_base)
    if "공장구분" in out.columns:
        standard_factory = item_codes.map(lambda code: DAILY_ITEM_STANDARD.get(code, {}).get("factory_group", ""))
        current_factory = out["공장구분"].map(clean_factory_group_display)
        standard_factory = standard_factory.map(clean_factory_group_display)
        factory_mask = (standard_factory != "") & ~current_factory.map(has_factory_group)
        out["공장구분"] = current_factory
        out.loc[factory_mask, "공장구분"] = standard_factory[factory_mask]
    if product_col in out.columns:
        pack_values = out["PACK"] if "PACK" in out.columns else pd.Series("", index=out.index)
        out[product_col] = [
            daily_standard_product_name(product, pack, item_code)
            for product, pack, item_code in zip(out[product_col], pack_values, out["품목코드"])
        ]
    return out


def daily_inventory_status_rank(value: Any) -> int:
    ranks = {
        "요청외 긴급": 0,
        "요청내 긴급": 1,
        "요청내 재고부족": 2,
        "재고 음수": 3,
        "요청내 재고확인": 4,
        "재고 모니터링": 5,
    }
    return ranks.get(clean_str(value), 99)


def first_daily_inventory_status(series: pd.Series) -> str:
    values = [clean_str(value) for value in series if clean_str(value)]
    if not values:
        return ""
    return min(values, key=daily_inventory_status_rank)


def build_daily_lot_wait_lookup(lot_status_df: pd.DataFrame | None) -> dict[str, float]:
    if lot_status_df is None or lot_status_df.empty:
        return {}
    required_cols = ["제품코드", "입고대기수량"]
    if any(col not in lot_status_df.columns for col in required_cols):
        return {}

    work = lot_status_df[required_cols].copy()
    work["_lot_item_key"] = work["제품코드"].map(normalize_match_key)
    work["입고대기수량"] = pd.to_numeric(work["입고대기수량"], errors="coerce").fillna(0.0)
    work = work[work["_lot_item_key"] != ""].copy()
    if work.empty:
        return {}

    grouped = work.groupby("_lot_item_key", dropna=False)["입고대기수량"].sum()
    return {clean_str(key): float(value) for key, value in grouped.items() if clean_str(key)}


@st.cache_data(show_spinner=False, max_entries=16)
def build_daily_inventory_response_view(
    daily_inventory_df: pd.DataFrame,
    code_summary: pd.DataFrame,
    sample_available_df: pd.DataFrame | None = None,
    lot_status_df: pd.DataFrame | None = None,
) -> pd.DataFrame:
    columns = [
        "대응상태",
        "품목코드",
        "기간구분",
        "제품명",
        "재고표 제품명",
        "제품코드",
        "PACK",
        "POWER",
        "CP",
        "재고수량",
        "전일재고",
        "재고증감",
        "재고부족수량",
        "긴급요청",
        "요청 PACK",
        "용마입고 PACK",
        "미입고 PACK",
        "포장 PACK",
        "요청 PCS",
        "생산부족 PCS",
        "용마입고대기 PACK",
        "포장부족(재고 PCS)",
        "포장가능재고(PCS)",
        "생산진도율",
        "생산완료예상일",
        "요청제품명",
        "판매코드 수",
        "대상품목",
    ]
    if daily_inventory_df.empty:
        return pd.DataFrame(columns=columns)

    daily = daily_inventory_df.copy()
    daily = enrich_daily_inventory_from_code_summary(daily, code_summary)
    daily["제품코드"] = daily["제품코드"].map(clean_str).str.upper()
    daily["PACK"] = daily["PACK"].map(clean_str)
    daily["POWER"] = daily["POWER"].map(clean_str)
    daily["CP"] = daily["CP"].map(clean_str) if "CP" in daily.columns else ""
    daily["_daily_base_product_name"] = daily["제품명"].map(strip_pack_unit_suffix).map(clean_str)
    daily["재고수량"] = pd.to_numeric(daily["재고수량"], errors="coerce")
    daily["전일재고"] = pd.to_numeric(daily["전일재고"], errors="coerce")
    daily["재고증감"] = pd.to_numeric(daily["재고증감"], errors="coerce")
    daily["긴급요청"] = daily["긴급요청"].apply(lambda value: bool(value) if not pd.isna(value) else False)

    request_match = build_daily_request_match_view(code_summary)
    out = daily.merge(
        request_match,
        on=["제품코드", "PACK", "POWER", "CP"],
        how="left",
    )
    production_catalog = build_daily_production_power_catalog(code_summary)
    out = out.merge(production_catalog, on=["제품코드", "POWER", "CP"], how="left")
    base_production_catalog = build_daily_base_power_production_catalog(code_summary)
    out = out.merge(base_production_catalog, on=["_daily_base_product_name", "POWER", "CP"], how="left")
    product_catalog = build_daily_product_catalog(code_summary)
    out = out.merge(product_catalog, on=["제품코드", "PACK"], how="left")
    out["재고표 제품명"] = out["제품명"].map(clean_str)
    out["품목코드"] = [
        build_item_code_from_prefix_power(product_code, power, cp)
        for product_code, power, cp in zip(out["제품코드"], out["POWER"], out["CP"])
    ]
    out["제품명"] = out["요청제품명"].where(out["요청제품명"].map(clean_str) != "", out["마스터제품명"])
    out["제품명"] = out["제품명"].where(out["제품명"].map(clean_str) != "", out["재고표 제품명"])
    out = add_period_group_columns(out)
    for col in ["생산완료예상일", "요청제품명", "대상품목", "마스터제품명", "마스터공장구분", "재고표 제품명", "품목코드", "CP"]:
        if col in out.columns:
            out[col] = out[col].fillna("")
    if "공장구분" not in out.columns:
        out["공장구분"] = ""
    out["공장구분"] = out["공장구분"].map(clean_factory_group_display)
    if "마스터공장구분" in out.columns:
        out["마스터공장구분"] = out["마스터공장구분"].map(clean_factory_group_display)
        out["공장구분"] = out["공장구분"].where(
            out["공장구분"].map(has_factory_group),
            out["마스터공장구분"],
        )
    out = out[complete_daily_response_mask(out)].copy()
    if out.empty:
        return pd.DataFrame(columns=columns)
    sample_lookup = build_sample_available_lookup(sample_available_df)
    inferred_production_code = [
        replace_power_in_production_code(template, source_power, power)
        for template, source_power, power in zip(
            out.get("생산코드템플릿", pd.Series("", index=out.index)),
            out.get("생산코드POWER", pd.Series("", index=out.index)),
            out["POWER"],
        )
    ]
    inferred_sample_available = [
        sample_lookup.get(normalize_match_key(code), 0.0)
        for code in inferred_production_code
    ]
    out["_추정샘플신청가능수량"] = inferred_sample_available
    numeric_cols = [
        "요청 PACK",
        "포장 PACK",
        "용마입고 PACK",
        "미입고 PACK",
        "요청 PCS",
        "생산부족 PCS",
        "용마입고대기 PACK",
        "포장가능재고(PCS)",
        "샘플신청가능수량",
        "생산진도율",
        "판매코드 수",
        "_production_request_pack",
        "_production_request_pcs",
        "_production_shortage_pcs",
        "_production_sample_available_pcs",
        "_production_available_stock_pcs",
        "_production_progress_pct",
        "_base_production_request_pack",
        "_base_production_request_pcs",
        "_base_production_shortage_pcs",
        "_base_production_sample_available_pcs",
        "_base_production_available_stock_pcs",
        "_base_production_progress_pct",
    ]
    for col in numeric_cols:
        if col in out.columns:
            out[col] = pd.to_numeric(out[col], errors="coerce").fillna(0.0)
    for col in [
        "용마입고대기 PACK",
        "포장가능재고(PCS)",
        "샘플신청가능수량",
        "_production_request_pcs",
        "_production_shortage_pcs",
        "_production_sample_available_pcs",
        "_production_available_stock_pcs",
        "_production_progress_pct",
        "_base_production_request_pcs",
        "_base_production_shortage_pcs",
        "_base_production_sample_available_pcs",
        "_base_production_available_stock_pcs",
        "_base_production_progress_pct",
    ]:
        if col not in out.columns:
            out[col] = 0.0
    has_request = out["요청 PACK"] > 0
    has_code_production_context = (
        (out["_production_request_pcs"] > 0)
        | (out["_production_shortage_pcs"] > 0)
    )
    has_base_production_context = (
        (out["_base_production_request_pcs"] > 0)
        | (out["_base_production_shortage_pcs"] > 0)
    )
    use_base_production = ~has_code_production_context & has_base_production_context
    matched_production_shortage_pcs = np.where(
        use_base_production,
        out["_base_production_shortage_pcs"],
        out["_production_shortage_pcs"],
    )
    matched_production_sample_pcs = np.where(
        use_base_production,
        out["_base_production_sample_available_pcs"],
        out["_production_sample_available_pcs"],
    )
    matched_production_available_pcs = np.where(
        use_base_production,
        out["_base_production_available_stock_pcs"],
        out["_production_available_stock_pcs"],
    )
    matched_production_progress_pct = np.where(
        use_base_production,
        out["_base_production_progress_pct"],
        out["_production_progress_pct"],
    )
    has_production_context = has_code_production_context | has_base_production_context
    out["샘플신청가능수량"] = np.where(
        has_production_context,
        matched_production_sample_pcs,
        out["샘플신청가능수량"],
    )
    out["샘플신청가능수량"] = out["샘플신청가능수량"].where(
        has_request | has_production_context,
        out["_추정샘플신청가능수량"],
    )
    lot_wait_lookup = build_daily_lot_wait_lookup(lot_status_df)
    if lot_wait_lookup:
        item_keys = out["품목코드"].map(normalize_match_key)
        out["용마입고대기 PACK"] = item_keys.map(lambda key: lot_wait_lookup.get(key, 0.0))
    else:
        out["용마입고대기 PACK"] = 0.0
    out["생산부족 PCS"] = np.where(
        has_production_context,
        matched_production_shortage_pcs,
        out["생산부족 PCS"],
    )
    out["생산진도율"] = np.where(
        has_production_context,
        matched_production_progress_pct,
        out["생산진도율"],
    )
    out["생산부족 PCS"] = pd.to_numeric(out["생산부족 PCS"], errors="coerce").fillna(0.0).round(0)
    exact_available_stock_pcs = np.where(
        has_request,
        (out["요청 PCS"] - out["생산부족 PCS"] + out["샘플신청가능수량"]).clip(lower=0.0),
        0.0,
    )
    available_stock_pcs = np.where(
        has_production_context,
        matched_production_available_pcs,
        exact_available_stock_pcs,
    )
    out["재고부족수량"] = (-out["재고수량"]).clip(lower=0.0).fillna(0.0)
    pack_units = out["PACK"].map(pack_unit_from_label)
    stock_shortage_pcs = (out["재고부족수량"] * pack_units).clip(lower=0.0)
    stock_available_pcs = (out["재고수량"].clip(lower=0.0).fillna(0.0) * pack_units).clip(lower=0.0)
    sample_available_pcs = pd.to_numeric(out["샘플신청가능수량"], errors="coerce").fillna(0.0)
    fallback_available_stock_pcs = sample_available_pcs.where(
        sample_available_pcs >= stock_available_pcs,
        stock_available_pcs,
    )
    has_supply_context = has_request | has_production_context
    out["포장부족(재고 PCS)"] = np.where(
        has_supply_context,
        (stock_shortage_pcs - available_stock_pcs).clip(lower=0.0),
        (stock_shortage_pcs - fallback_available_stock_pcs).clip(lower=0.0),
    )
    out["포장가능재고(PCS)"] = np.where(
        has_supply_context,
        available_stock_pcs,
        fallback_available_stock_pcs,
    )
    out["포장가능재고(PCS)"] = pd.to_numeric(out["포장가능재고(PCS)"], errors="coerce").fillna(0.0).round(0)
    out["대응상태"] = out.apply(classify_daily_inventory_status, axis=1)
    out = apply_daily_item_standard(out, product_col="제품명")
    out["_urgent_sort"] = out["긴급요청"].astype(int)
    out["_negative_sort"] = (out["재고수량"] < 0).astype(int)
    out["_request_sort"] = (out["요청 PACK"] > 0).astype(int)
    out["_pack_sort"] = out["PACK"].map(pack_sort_key)
    out["_power_sort"] = pd.to_numeric(out["POWER"].str.replace("-00.00", "0", regex=False), errors="coerce").fillna(0.0)
    out = out.sort_values(
        ["_urgent_sort", "_negative_sort", "_request_sort", "재고부족수량", "제품명", "_pack_sort", "_power_sort"],
        ascending=[False, False, False, False, True, True, True],
        kind="stable",
    )
    return out[columns].copy()


def build_daily_inventory_main_view(response_view: pd.DataFrame) -> pd.DataFrame:
    visible_columns = [
        "대응상태",
        "품목코드",
        "기간구분",
        "대표 제품명",
        "긴급요청 수",
        "재고수량",
        "요청 PACK",
        "용마입고 PACK",
        "용마입고대기 PACK",
        "포장가능재고(PCS)",
        "생산부족 PCS",
        "생산진도율",
        "생산완료예상일",
    ]
    if response_view.empty:
        return pd.DataFrame(columns=visible_columns + ["_daily_item_code_base", "_daily_expected_date_sort"])

    work = response_view.copy()
    work["_daily_item_code_base"] = work["품목코드"].map(daily_item_code_base)
    work = work[work["_daily_item_code_base"] != ""].copy()
    if work.empty:
        return pd.DataFrame(columns=visible_columns + ["_daily_item_code_base", "_daily_expected_date_sort"])

    numeric_cols = [
        "재고수량",
        "재고부족수량",
        "요청 PACK",
        "용마입고 PACK",
        "미입고 PACK",
        "포장 PACK",
        "요청 PCS",
        "생산부족 PCS",
        "용마입고대기 PACK",
        "포장가능재고(PCS)",
    ]
    for col in numeric_cols:
        if col not in work.columns:
            work[col] = 0.0
        work[col] = pd.to_numeric(work[col], errors="coerce").fillna(0.0)
    work["_expected_date_sort"] = pd.to_datetime(work.get("생산완료예상일", pd.NaT), errors="coerce")

    grouped = (
        work.groupby("_daily_item_code_base", dropna=False)
        .agg(
            status=("대응상태", first_daily_inventory_status),
            period_group=("기간구분", first_nonempty),
            product_name=("제품명", first_nonempty),
            detail_count=("품목코드", "count"),
            urgent_count=("긴급요청", "sum"),
            stock_qty=("재고수량", "sum"),
            stock_shortage_qty=("재고부족수량", "sum"),
            request_pack=("요청 PACK", "sum"),
            yongma_in_pack=("용마입고 PACK", "sum"),
            yongma_wait_pack=("용마입고대기 PACK", "sum"),
            packable_stock_pcs=("포장가능재고(PCS)", "sum"),
            production_shortage_pcs=("생산부족 PCS", "sum"),
            request_pcs=("요청 PCS", "sum"),
            expected_date=("_expected_date_sort", max_datetime),
        )
        .reset_index()
    )
    grouped["생산진도율"] = calc_production_progress_pct(grouped["request_pcs"], grouped["production_shortage_pcs"])
    grouped["_daily_status_sort"] = grouped["status"].map(daily_inventory_status_rank)
    grouped["_daily_negative_sort"] = (grouped["stock_qty"] < 0).astype(int)
    grouped["_daily_expected_date_sort"] = pd.to_datetime(grouped["expected_date"], errors="coerce")
    grouped["생산완료예상일"] = grouped["expected_date"].map(display_date_or_dash)

    out = grouped.rename(
        columns={
            "_daily_item_code_base": "품목코드",
            "status": "대응상태",
            "period_group": "기간구분",
            "product_name": "대표 제품명",
            "detail_count": "상세 건수",
            "urgent_count": "긴급요청 수",
            "stock_qty": "재고수량",
            "stock_shortage_qty": "재고부족수량",
            "request_pack": "요청 PACK",
            "yongma_in_pack": "용마입고 PACK",
            "yongma_wait_pack": "용마입고대기 PACK",
            "packable_stock_pcs": "포장가능재고(PCS)",
            "production_shortage_pcs": "생산부족 PCS",
        }
    )
    out["_daily_item_code_base"] = out["품목코드"]
    out = apply_daily_item_standard(out, product_col="대표 제품명")
    out = out.sort_values(
        [
            "_daily_status_sort",
            "_daily_negative_sort",
            "재고부족수량",
            "용마입고대기 PACK",
            "_daily_expected_date_sort",
            "품목코드",
        ],
        ascending=[True, False, False, False, True, True],
        na_position="last",
        kind="stable",
    )
    return out[
        visible_columns
        + [
            "_daily_item_code_base",
            "_daily_status_sort",
            "_daily_negative_sort",
            "_daily_expected_date_sort",
        ]
    ].copy()


def daily_inventory_detail_column_order(df: pd.DataFrame) -> list[str]:
    columns = [
        "대응상태",
        "품목코드",
        "기간구분",
        "제품명",
        "PACK",
        "POWER",
        "CP",
        "재고수량",
        "긴급요청",
        "요청 PACK",
        "용마입고 PACK",
        "용마입고대기 PACK",
        "포장가능재고(PCS)",
        "생산부족 PCS",
        "생산진도율",
        "생산완료예상일",
    ]
    return visible_columns(df, columns)


def daily_inventory_search_variants(token: str) -> list[str]:
    normalized = clean_str(token).replace("−", "-").replace("–", "-").replace("—", "-")
    variants = [normalized]
    variants.extend(expand_product_query_terms(normalized))

    pack_label = extract_daily_pack_label(normalized)
    if pack_label:
        variants.append(pack_label)

    power_label = daily_power_label(normalized)
    if power_label:
        variants.append(power_label)

    cp_label = normalize_toric_cp_label(normalized) if not normalized.startswith("-") else ""
    if cp_label:
        variants.append(cp_label)

    return list(dict.fromkeys([variant for variant in variants if clean_str(variant)]))


def is_power_query_token(token: str) -> bool:
    normalized = clean_str(token).replace("−", "-").replace("–", "-").replace("—", "-")
    if normalized.upper() == "PL":
        return True
    return bool(re.fullmatch(r"[+-]?\d+(?:\.\d+)?", normalized) and (normalized.startswith(("+", "-")) or "." in normalized))


def is_item_code_query_token(token: str) -> bool:
    normalized = clean_str(token).replace("−", "-").replace("–", "-").replace("—", "-").upper()
    return bool(re.fullmatch(r"[A-Z]\d+[+-]\d+(?:\.\d+)?", normalized))


def daily_inventory_query_mask(view: pd.DataFrame, query: str) -> pd.Series:
    tokens = split_lookup_query_terms(query)
    if not tokens:
        return pd.Series(True, index=view.index)

    text_cols = [
        col
        for col in ["품목코드", "제품명", "재고표 제품명", "제품코드", "요청제품명", "대상품목", "CP"]
        if col in view.columns
    ]
    mask = pd.Series(True, index=view.index)
    for token in tokens:
        normalized = clean_str(token).replace("−", "-").replace("–", "-").replace("—", "-")
        token_mask = pd.Series(False, index=view.index)
        power_label = daily_power_label(normalized) if is_power_query_token(normalized) else ""
        pack_label = extract_daily_pack_label(normalized)
        cp_label = normalize_toric_cp_label(normalized) if not normalized.startswith("-") else ""
        if is_item_code_query_token(normalized) and "품목코드" in view.columns:
            token_mask = contains_any_query_term(view["품목코드"].fillna(""), [normalized.upper()])
        elif cp_label and "CP" in view.columns:
            token_mask = contains_any_query_term(view["CP"].fillna(""), [cp_label])
        elif power_label and "POWER" in view.columns:
            token_mask = contains_any_query_term(view["POWER"].fillna(""), [power_label])
        elif pack_label and "PACK" in view.columns:
            token_mask = contains_any_query_term(view["PACK"].fillna(""), [pack_label])
        else:
            variants = daily_inventory_search_variants(normalized)
            for col in text_cols:
                token_mask = token_mask | contains_any_query_term(view[col].fillna(""), variants)
        mask = mask & token_mask
    return mask


def build_sales_pack_detail_view(code_summary: pd.DataFrame) -> pd.DataFrame:
    if code_summary.empty:
        return pd.DataFrame(columns=["판매코드", "PACK", "요청", "용마입고", "미입고", "생산완료예상일"])
    work = with_operational_columns(code_summary)
    out = (
        work.groupby(["sales_code", "_pack_label"], dropna=False)
        .agg(
            request_pack=("request_pack", "sum"),
            packing_pack=("packing_recognized_pack", "sum"),
            yongma_in_pack=("yongma_recognized_pack", "sum"),
            expected_date=("production_due_date", max_datetime),
        )
        .reset_index()
        .rename(columns={"sales_code": "판매코드", "_pack_label": "PACK", "request_pack": "요청", "yongma_in_pack": "용마입고"})
    )
    out["미입고"] = (out["요청"] - out["용마입고"]).clip(lower=0.0)
    out["생산완료예상일"] = out["expected_date"].map(display_date_or_dash)
    out["_pack_sort"] = out["PACK"].map(pack_sort_key)
    return out.sort_values(["미입고", "_pack_sort", "요청"], ascending=[False, True, False], kind="stable")[
        ["판매코드", "PACK", "요청", "용마입고", "미입고", "생산완료예상일"]
    ]


def build_lot_receipt_status_view(
    packing_df: pd.DataFrame,
    yongma_df: pd.DataFrame,
    code_summary: pd.DataFrame,
) -> pd.DataFrame:
    columns = [
        "기간구분",
        "판매코드",
        "제품명",
        "LOTNO",
        "포장일",
        "포장실적수량",
        "용마입고수량",
        "입고대기수량",
        "상태",
    ]
    if packing_df.empty or code_summary.empty:
        return pd.DataFrame(columns=columns)

    domestic_code_keys = set(code_summary.get("sales_code_key", pd.Series(dtype=str)).map(clean_str)) - {""}
    if not domestic_code_keys:
        return pd.DataFrame(columns=columns)

    pack = packing_df.copy()
    for col in ["sales_code_key", "packing_lot_key", "packing_barcode_key"]:
        if col not in pack.columns:
            pack[col] = ""
    pack["sales_code_key"] = pack["sales_code_key"].map(clean_str)
    pack = pack[pack["sales_code_key"].isin(domestic_code_keys)].copy()
    if pack.empty:
        return pd.DataFrame(columns=columns)

    for col in ["packing_product_name", "packing_lot", "packing_barcode"]:
        if col not in pack.columns:
            pack[col] = ""
    if "packing_date" not in pack.columns:
        pack["packing_date"] = pd.NaT

    grouped = (
        pack.groupby(
            [
                "sales_code_key",
                "sales_code",
                "packing_product_name",
                "packing_lot",
                "packing_lot_key",
                "packing_barcode",
                "packing_barcode_key",
            ],
            dropna=False,
        )
        .agg(
            packing_pack=("packing_pack", "sum"),
            packing_date=("packing_date", min_datetime),
        )
        .reset_index()
    )
    grouped["용마입고수량"] = 0.0

    if yongma_df is not None and not yongma_df.empty:
        yongma = yongma_df.copy()
        yongma["sales_code_key"] = yongma["sales_code_key"].map(clean_str)
        yongma = yongma[yongma["sales_code_key"].isin(domestic_code_keys)].copy()
        yongma["yongma_lot_key"] = yongma["yongma_lot_key"].map(clean_str)
        code_meta = code_summary.copy()
        code_meta["sales_code_key"] = code_meta.get("sales_code_key", pd.Series("", index=code_meta.index)).map(clean_str)
        code_sales_by_key = build_first_value_map(code_meta, "sales_code_key", "sales_code")
        code_product_by_key = build_first_value_map(code_meta, "sales_code_key", "product_name")
        if "factory_group" not in code_meta.columns:
            code_meta["factory_group"] = "(미기재)"
        code_factory_by_key = build_first_value_map(code_meta, "sales_code_key", "factory_group")
        code_meta = add_period_group_columns(code_meta)
        code_period_by_key = build_first_value_map(code_meta, "sales_code_key", "period_group")
        receipt_only_rows: list[dict[str, Any]] = []

        def add_receipt_to_indices(indices: list[int], qty: float) -> None:
            remaining = float(qty)
            for idx in indices:
                if remaining <= 0:
                    break
                packed = pd.to_numeric(grouped.at[idx, "packing_pack"], errors="coerce")
                received = pd.to_numeric(grouped.at[idx, "용마입고수량"], errors="coerce")
                capacity = max(0.0, float(packed if not pd.isna(packed) else 0.0) - float(received if not pd.isna(received) else 0.0))
                add_qty = min(remaining, capacity) if capacity > 0 else 0.0
                if add_qty > 0:
                    grouped.at[idx, "용마입고수량"] += add_qty
                    remaining -= add_qty
            if remaining > 0 and indices:
                grouped.at[indices[0], "용마입고수량"] += remaining

        for _, receipt in yongma.iterrows():
            code_key = clean_str(receipt.get("sales_code_key", ""))
            lot_key = clean_str(receipt.get("yongma_lot_key", ""))
            qty_value = pd.to_numeric(receipt.get("yongma_in_pack", 0.0), errors="coerce")
            qty = 0.0 if pd.isna(qty_value) else float(qty_value)
            if not code_key or not lot_key or qty <= 0:
                continue

            candidates = grouped[grouped["sales_code_key"] == code_key]
            if candidates.empty:
                receipt_only_rows.append(
                    {
                        "sales_code_key": code_key,
                        "sales_code": clean_str(receipt.get("sales_code", "")) or code_sales_by_key.get(code_key, ""),
                        "packing_product_name": clean_str(receipt.get("yongma_product_name", ""))
                        or code_product_by_key.get(code_key, ""),
                        "factory_group": code_factory_by_key.get(code_key, "(미기재)"),
                        "period_group": code_period_by_key.get(code_key, "FRP"),
                        "packing_lot": clean_str(receipt.get("yongma_lot", "")) or "(용마 LOT 미기재)",
                        "packing_lot_key": lot_key,
                        "packing_barcode": "",
                        "packing_barcode_key": "",
                        "packing_pack": 0.0,
                        "packing_date": pd.NaT,
                        "용마입고수량": qty,
                    }
                )
                continue
            exact = candidates[candidates["packing_lot_key"] == lot_key]
            target = exact
            if target.empty:
                barcode_match = candidates[
                    candidates["packing_barcode_key"].astype(str).str.contains(lot_key, regex=False, na=False)
                ]
                target = barcode_match
            if target.empty:
                target = candidates.sort_values(["packing_date", "packing_lot"], na_position="last", kind="stable")
                add_receipt_to_indices(target.index.tolist(), qty)
                continue
            add_receipt_to_indices(target.index.tolist(), qty)

        if receipt_only_rows:
            receipt_only = (
                pd.DataFrame(receipt_only_rows)
                .groupby(
                    [
                        "sales_code_key",
                        "sales_code",
                        "packing_product_name",
                        "packing_lot",
                        "packing_lot_key",
                        "packing_barcode",
                        "packing_barcode_key",
                    ],
                    dropna=False,
                )
                .agg(
                    packing_pack=("packing_pack", "sum"),
                    packing_date=("packing_date", min_datetime),
                    용마입고수량=("용마입고수량", "sum"),
                )
                .reset_index()
            )
            grouped = pd.concat([grouped, receipt_only], ignore_index=True)

    grouped["포장실적수량"] = pd.to_numeric(grouped["packing_pack"], errors="coerce").fillna(0.0)
    grouped["용마입고수량"] = pd.to_numeric(grouped["용마입고수량"], errors="coerce").fillna(0.0)
    grouped["입고대기수량"] = (grouped["포장실적수량"] - grouped["용마입고수량"]).clip(lower=0.0)
    grouped["상태"] = np.select(
        [
            (grouped["포장실적수량"] <= 0) & (grouped["용마입고수량"] > 0),
            (grouped["입고대기수량"] > 0) & (grouped["용마입고수량"] > 0),
            grouped["입고대기수량"] > 0,
        ],
        ["용마입고만", "부분입고", "입고대기"],
        default="입고완료",
    )
    grouped["_status_sort"] = grouped["상태"].map({"입고대기": 0, "부분입고": 1, "용마입고만": 2}).fillna(3)
    grouped["포장일"] = grouped["packing_date"].map(display_date_or_dash)
    code_factory = with_operational_columns(code_summary)[["sales_code_key", "factory_group"]].copy()
    code_period = with_operational_columns(code_summary)[["sales_code_key", "period_group"]].copy()
    code_factory["sales_code_key"] = code_factory["sales_code_key"].map(clean_str)
    code_period["sales_code_key"] = code_period["sales_code_key"].map(clean_str)
    code_factory = code_factory.drop_duplicates("sales_code_key", keep="first")
    code_period = code_period.drop_duplicates("sales_code_key", keep="first")
    grouped = grouped.merge(code_factory, on="sales_code_key", how="left")
    grouped = grouped.merge(code_period, on="sales_code_key", how="left")
    if "factory_group_x" in grouped.columns or "factory_group_y" in grouped.columns:
        grouped["factory_group"] = grouped.get("factory_group_x", "").where(
            grouped.get("factory_group_x", "").map(clean_str) != "",
            grouped.get("factory_group_y", ""),
        )
        grouped = grouped.drop(columns=["factory_group_x", "factory_group_y"], errors="ignore")
    if "period_group_x" in grouped.columns or "period_group_y" in grouped.columns:
        grouped["period_group"] = grouped.get("period_group_x", "").where(
            grouped.get("period_group_x", "").map(clean_str) != "",
            grouped.get("period_group_y", ""),
        )
        grouped = grouped.drop(columns=["period_group_x", "period_group_y"], errors="ignore")
    grouped["factory_group"] = grouped["factory_group"].map(clean_str).replace("", "(미기재)").fillna("(미기재)")
    grouped["period_group"] = grouped["period_group"].map(clean_str).replace("", "FRP").fillna("FRP")
    grouped = grouped.rename(
        columns={
            "factory_group": "공장구분",
            "period_group": "기간구분",
            "sales_code": "판매코드",
            "packing_product_name": "제품명",
            "packing_lot": "LOTNO",
        }
    )
    grouped["제품명"] = grouped["제품명"].replace("", "(미기재)")
    return grouped.sort_values(
        ["_status_sort", "입고대기수량", "포장실적수량"],
        ascending=[True, False, False],
        kind="stable",
    )[columns].copy()


def build_production_progress_main_view(code_summary: pd.DataFrame, pack_labels: list[str]) -> pd.DataFrame:
    if code_summary.empty:
        return pd.DataFrame(
            columns=[
                "생산코드",
                "제품명",
                *pack_labels,
                "요청합계",
                "생산부족",
                "포장부족",
                "진도율",
                "판매코드수",
            ]
        )
    work = with_operational_columns(code_summary)
    base = build_production_code_view(work).rename(
        columns={
            "요청 PACK": "요청합계",
            "생산부족수량": "생산부족",
            "포장부족수량": "포장부족",
            "포장진도율": "진도율",
            "연결 판매코드 수": "판매코드수",
        }
    )
    pivot = build_pack_pivot(work, ["production_code_display"], pack_labels).rename(
        columns={"production_code_display": "생산코드"}
    )
    out = base.merge(pivot, on="생산코드", how="left")
    for label in pack_labels:
        out[label] = out[label].fillna(0.0)
    return out[
        ["생산코드", "제품명", *pack_labels, "요청합계", "생산부족", "포장부족", "진도율", "판매코드수", "상태"]
    ].sort_values(["포장부족", "생산부족", "요청합계"], ascending=[False, False, False], kind="stable")


def prepare_production_power_rows(code_summary: pd.DataFrame) -> pd.DataFrame:
    work = with_operational_columns(code_summary)
    work = work[work["production_code_display"].map(is_p_production_code)].copy()
    work = add_allocated_production_basis(work)
    work["_production_shortage_pcs"] = pd.to_numeric(
        work["_allocated_production_shortage_qty"],
        errors="coerce",
    ).fillna(0.0)
    work["_basis_difference_pcs"] = pd.to_numeric(
        work.get("_basis_difference_pcs", pd.Series(0.0, index=work.index)),
        errors="coerce",
    ).fillna(0.0)
    work["_packing_shortage_pack"] = pd.to_numeric(
        work.get("code_packing_shortage_pack", pd.Series(0.0, index=work.index)),
        errors="coerce",
    ).fillna(0.0)
    work["_packing_recognized_pcs"] = recognized_packing_pcs(work)
    work["_power_sort"] = pd.to_numeric(work["power_value"], errors="coerce").fillna(999999.0)
    return work


def attach_deduped_sample_available_pcs(rows: pd.DataFrame) -> pd.DataFrame:
    work = rows.copy()
    sample_available = pd.to_numeric(
        work.get("sample_available_pcs", pd.Series(0.0, index=work.index)),
        errors="coerce",
    ).fillna(0.0)
    production_key = work.get("production_code_key", pd.Series("", index=work.index)).map(clean_str)
    display_key = work.get("production_code_display", pd.Series("", index=work.index)).map(normalize_match_key)
    sample_key = production_key.where(production_key != "", display_key)
    duplicated_sample = (sample_key.map(clean_str) != "") & sample_key.duplicated(keep="first")
    work["_dedup_sample_available_pcs"] = sample_available.where(~duplicated_sample, 0.0)
    return work


def refresh_scope_production_shortage(rows: pd.DataFrame) -> pd.DataFrame:
    work = add_allocated_production_basis(rows)
    work["_production_shortage_pcs"] = pd.to_numeric(
        work["_allocated_production_shortage_qty"],
        errors="coerce",
    ).fillna(0.0)
    work["_basis_difference_pcs"] = pd.to_numeric(
        work.get("_basis_difference_pcs", pd.Series(0.0, index=work.index)),
        errors="coerce",
    ).fillna(0.0)
    return work


def available_production_power_options(code_summary: pd.DataFrame) -> list[str]:
    work = prepare_production_power_rows(code_summary)
    source = work[["POWER", "_power_sort"]].drop_duplicates().sort_values("_power_sort", ascending=True, kind="stable")
    values = source["POWER"].astype(str).tolist()
    return ["전체"] + values


def filter_production_power_rows(
    code_summary: pd.DataFrame,
    product_query: str,
    production_query: str,
    power_label: str,
    pack_label: str,
    sample_scope: str,
    product_group: str,
    factory_group: str,
    period_group: str = "전체",
) -> pd.DataFrame:
    out = prepare_production_power_rows(code_summary)
    product_q = product_query.strip()
    if product_q:
        name_match = out["product_name"].astype(str).str.contains(product_q, case=False, na=False)
        base_match = out["base_product_name"].astype(str).str.contains(product_q, case=False, na=False)
        out = out[name_match | base_match]
    production_q = production_query.strip()
    if production_q:
        out = out[out["production_code_display"].astype(str).str.contains(production_q, case=False, na=False)]
    if power_label != "전체":
        out = out[out["POWER"] == power_label]
    if pack_label != "전체":
        out = out[out["_pack_label"] == pack_label]
    if sample_scope == "본품":
        out = out[out["본품/샘플"] == "본품"]
    elif sample_scope == "샘플":
        out = out[out["본품/샘플"] == "샘플"]
    if product_group != "전체":
        out = out[out["제품분류"] == product_group]
    if factory_group != "전체":
        out = out[out["factory_group"] == factory_group]
    if period_group != "전체":
        out = out[out["period_group"] == period_group]
    return refresh_scope_production_shortage(out).copy()


def is_p_production_code(value: Any) -> bool:
    return clean_str(value).upper().startswith("P")


def bottleneck_status(production_progress: Any, packing_progress: Any) -> str:
    production = float(pd.to_numeric(production_progress, errors="coerce") or 0.0)
    packing = float(pd.to_numeric(packing_progress, errors="coerce") or 0.0)
    if production < 20.0 and packing < 20.0:
        return "미착수 ⚫"
    if production < packing - 20.0:
        return "생산 병목 🔴"
    if packing < production - 20.0:
        return "포장 병목 🟠"
    return "정상 🟢"


def status_from_progress(packing_pack: Any, packing_progress: Any) -> str:
    packing = float(pd.to_numeric(packing_pack, errors="coerce") or 0.0)
    progress = float(pd.to_numeric(packing_progress, errors="coerce") or 0.0)
    return classify_status(packing, progress)


def production_code_prefix(value: Any) -> str:
    text = clean_str(value)
    if not text or text == "(생산코드 미기재)":
        return "(생산코드 미기재)"
    return text[:5].upper()


def build_production_power_main_view(
    rows: pd.DataFrame,
    pack_labels: list[str],
    shortage_only: bool = False,
) -> pd.DataFrame:
    visible_columns = [
        "생산코드",
        "기간구분",
        "대표 제품명",
        *pack_labels,
        "요청합계(PACK)",
        "포장부족(PACK)",
        "포장가능재고(PCS)",
        "생산부족수량(PCS)",
        "기준차이",
        "생산진도율",
        "포장진도율",
        "생산완료예상일",
    ]
    if rows.empty:
        return pd.DataFrame(
            columns=visible_columns
            + [
                "요청합계(PCS)",
                "생산부족수량",
                "기준차이(PCS)",
                "포장부족수량",
                "포장가능재고(PCS)",
                "병목 상태",
                "_production_code_prefix",
                "_expected_date_sort",
            ]
        )

    work = rows.copy()
    if "_packing_recognized_pcs" not in work.columns:
        work["_packing_recognized_pcs"] = recognized_packing_pcs(work)
    work = attach_deduped_sample_available_pcs(work)
    work["_production_code_prefix"] = work["production_code_display"].map(production_code_prefix)
    group_cols = ["_production_code_prefix"]
    base = (
        work.groupby(group_cols, dropna=False)
        .agg(
            representative_product=("base_product_name", first_nonempty),
            factory_group=("factory_group", join_unique),
            period_group=("period_group", first_nonempty),
            request_pack=("request_pack", "sum"),
            request_pcs=("request_pcs", "sum"),
            packing_pack=("packing_recognized_pack", "sum"),
            packing_pcs=("_packing_recognized_pcs", "sum"),
            production_shortage_pcs=("_production_shortage_pcs", "sum"),
            basis_difference_pcs=("_basis_difference_pcs", "sum"),
            packing_shortage_pack=("_packing_shortage_pack", "sum"),
            sample_available_pcs=("_dedup_sample_available_pcs", "sum"),
            expected_date=("production_due_date", max_datetime),
        )
        .reset_index()
    )

    pack_pivot = (
        work.pivot_table(
            index=group_cols,
            columns="_pack_label",
            values="request_pack",
            aggfunc="sum",
            dropna=True,
        )
        .fillna(0.0)
        .reset_index()
        .rename_axis(None, axis=1)
    )
    for label in pack_labels:
        if label not in pack_pivot.columns:
            pack_pivot[label] = 0.0

    grouped = base.merge(pack_pivot[group_cols + pack_labels], on=group_cols, how="left")
    for label in pack_labels:
        grouped[label] = grouped[label].fillna(0.0)
    grouped["생산진도율"] = calc_production_progress_pct(grouped["request_pcs"], grouped["production_shortage_pcs"])
    grouped["포장진도율"] = np.where(
        grouped["request_pcs"] > 0,
        grouped["packing_pcs"] / grouped["request_pcs"] * 100.0,
        0.0,
    )
    grouped["포장진도율"] = np.clip(grouped["포장진도율"], 0.0, 100.0)
    grouped["포장가능재고(PCS)"] = (
        grouped["request_pcs"] - grouped["production_shortage_pcs"] + grouped["sample_available_pcs"]
    ).clip(lower=0.0).round(0)
    grouped["병목 상태"] = [
        bottleneck_status(prod, pack)
        for prod, pack in zip(grouped["생산진도율"], grouped["포장진도율"])
    ]
    grouped["상태"] = [
        status_from_progress(packing, progress)
        for packing, progress in zip(grouped["packing_pack"], grouped["포장진도율"])
    ]
    grouped["_expected_date_sort"] = pd.to_datetime(grouped["expected_date"], errors="coerce")
    grouped["생산완료예상일"] = grouped["expected_date"].map(display_date_or_dash)

    out = grouped.rename(
        columns={
            "_production_code_prefix": "생산코드",
            "factory_group": "공장구분",
            "period_group": "기간구분",
            "representative_product": "대표 제품명",
            "request_pack": "요청합계(PACK)",
            "request_pcs": "요청합계(PCS)",
            "packing_pcs": "포장실적(PCS)",
            "production_shortage_pcs": "생산부족수량",
            "basis_difference_pcs": "기준차이(PCS)",
            "packing_shortage_pack": "포장부족수량",
        }
    )
    out["생산부족수량(PCS)"] = out["생산부족수량"]
    out["기준차이"] = np.where(pd.to_numeric(out["기준차이(PCS)"], errors="coerce").fillna(0.0) > 0, "기준차이", "")
    out["포장부족(PACK)"] = out["포장부족수량"]
    out["_production_code_prefix"] = out["생산코드"]
    if shortage_only:
        out = out[(out["생산부족수량"] > 0) | (out["포장부족수량"] > 0)].copy()

    out = out.sort_values(
        ["_expected_date_sort", "포장부족수량", "생산부족수량"],
        ascending=[True, False, False],
        na_position="last",
        kind="stable",
    )
    return_columns = list(
        dict.fromkeys(
            visible_columns
            + [
                "요청합계(PCS)",
                "포장실적(PCS)",
                "생산부족수량",
                "기준차이(PCS)",
                "포장부족수량",
                "포장가능재고(PCS)",
                "병목 상태",
                "_production_code_prefix",
                "_expected_date_sort",
            ]
        )
    )
    for col in return_columns:
        if col not in out.columns:
            out[col] = ""
    return out[return_columns].copy()


def build_production_power_detail_view(
    rows: pd.DataFrame,
    pack_labels: list[str],
    production_prefix: str | None = None,
    wip_df: pd.DataFrame | None = None,
) -> pd.DataFrame:
    visible_columns = [
        "생산코드 전체",
        "기간구분",
        "대표 제품명",
        "POWER",
        *pack_labels,
        "요청합계(PACK)",
        "포장부족(PACK)",
        "포장가능재고(PCS)",
        "생산부족수량(PCS)",
        "기준차이",
        "생산진도율",
        "포장진도율",
        "생산완료예상일",
    ]
    if rows.empty:
        return pd.DataFrame(
            columns=visible_columns
            + [
                "요청합계(PCS)",
                "생산부족수량",
                "기준차이(PCS)",
                "포장부족수량",
                "포장가능재고(PCS)",
                *WIP_PROCESS_COLUMNS,
                "_production_code_prefix",
                "_expected_date_sort",
                "_power_sort",
            ]
        )

    work = rows.copy()
    if "_packing_recognized_pcs" not in work.columns:
        work["_packing_recognized_pcs"] = recognized_packing_pcs(work)
    work = attach_deduped_sample_available_pcs(work)
    work["_production_code_prefix"] = work["production_code_display"].map(production_code_prefix)
    if production_prefix is not None:
        work = work[work["_production_code_prefix"] == production_prefix].copy()
    if work.empty:
        return pd.DataFrame(
            columns=visible_columns
            + [
                "요청합계(PCS)",
                "생산부족수량",
                "포장부족수량",
                "포장가능재고(PCS)",
                *WIP_PROCESS_COLUMNS,
                "_production_code_prefix",
                "_expected_date_sort",
                "_power_sort",
            ]
        )

    group_cols = ["_production_code_prefix", "production_code_display", "POWER", "_power_sort"]
    base = (
        work.groupby(group_cols, dropna=False)
        .agg(
            representative_product=("base_product_name", first_nonempty),
            factory_group=("factory_group", join_unique),
            period_group=("period_group", first_nonempty),
            request_pack=("request_pack", "sum"),
            request_pcs=("request_pcs", "sum"),
            packing_pack=("packing_recognized_pack", "sum"),
            packing_pcs=("_packing_recognized_pcs", "sum"),
            production_shortage_pcs=("_production_shortage_pcs", "sum"),
            basis_difference_pcs=("_basis_difference_pcs", "sum"),
            packing_shortage_pack=("_packing_shortage_pack", "sum"),
            sample_available_pcs=("_dedup_sample_available_pcs", "sum"),
            expected_date=("production_due_date", max_datetime),
        )
        .reset_index()
    )
    pack_pivot = (
        work.pivot_table(
            index=group_cols,
            columns="_pack_label",
            values="request_pack",
            aggfunc="sum",
            dropna=True,
        )
        .fillna(0.0)
        .reset_index()
        .rename_axis(None, axis=1)
    )
    for label in pack_labels:
        if label not in pack_pivot.columns:
            pack_pivot[label] = 0.0

    grouped = base.merge(pack_pivot[group_cols + pack_labels], on=group_cols, how="left")
    for label in pack_labels:
        grouped[label] = grouped[label].fillna(0.0)
    grouped["생산진도율"] = calc_production_progress_pct(grouped["request_pcs"], grouped["production_shortage_pcs"])
    grouped["포장진도율"] = np.where(
        grouped["request_pcs"] > 0,
        grouped["packing_pcs"] / grouped["request_pcs"] * 100.0,
        0.0,
    )
    grouped["포장진도율"] = np.clip(grouped["포장진도율"], 0.0, 100.0)
    grouped["포장가능재고(PCS)"] = (
        grouped["request_pcs"] - grouped["production_shortage_pcs"] + grouped["sample_available_pcs"]
    ).clip(lower=0.0).round(0)
    grouped["상태"] = [
        status_from_progress(packing, progress)
        for packing, progress in zip(grouped["packing_pack"], grouped["포장진도율"])
    ]
    grouped["_expected_date_sort"] = pd.to_datetime(grouped["expected_date"], errors="coerce")
    grouped["생산완료예상일"] = grouped["expected_date"].map(display_date_or_dash)

    out = grouped.rename(
        columns={
            "production_code_display": "생산코드 전체",
            "factory_group": "공장구분",
            "period_group": "기간구분",
            "representative_product": "대표 제품명",
            "request_pack": "요청합계(PACK)",
            "request_pcs": "요청합계(PCS)",
            "packing_pcs": "포장실적(PCS)",
            "production_shortage_pcs": "생산부족수량",
            "basis_difference_pcs": "기준차이(PCS)",
            "packing_shortage_pack": "포장부족수량",
        }
    )
    out["생산부족수량(PCS)"] = out["생산부족수량"]
    out["기준차이"] = np.where(pd.to_numeric(out["기준차이(PCS)"], errors="coerce").fillna(0.0) > 0, "기준차이", "")
    out["포장부족(PACK)"] = out["포장부족수량"]
    out["_wip_production_code_key"] = out["생산코드 전체"].map(normalize_match_key)
    if wip_df is not None and not wip_df.empty:
        wip_grouped = (
            wip_df.groupby("production_code_key", dropna=False)[WIP_PROCESS_COLUMNS]
            .sum()
            .reset_index()
        )
        out = out.merge(
            wip_grouped,
            left_on="_wip_production_code_key",
            right_on="production_code_key",
            how="left",
        )
        out = out.drop(columns=["production_code_key"], errors="ignore")
    for col in WIP_PROCESS_COLUMNS:
        if col not in out.columns:
            out[col] = 0.0
        out[col] = pd.to_numeric(out[col], errors="coerce").fillna(0.0)
    out = sort_power_detail_default(
        out,
        extra_cols=["_expected_date_sort", "포장부족수량", "생산부족수량"],
        extra_ascending=[True, False, False],
    )
    return_columns = list(
        dict.fromkeys(
            visible_columns
            + [
                "요청합계(PCS)",
                "포장실적(PCS)",
                *WIP_PROCESS_COLUMNS,
                "생산부족수량",
                "기준차이(PCS)",
                "포장부족수량",
                "포장가능재고(PCS)",
                "_production_code_prefix",
                "_expected_date_sort",
                "_power_sort",
                "_wip_production_code_key",
            ]
        )
    )
    for col in return_columns:
        if col not in out.columns:
            out[col] = ""
    return out[return_columns].copy()


def calc_production_power_kpis(view: pd.DataFrame) -> dict[str, float]:
    if view.empty:
        return {
            "production_code_count": 0.0,
            "request_pack": 0.0,
            "request_pcs": 0.0,
            "production_shortage_pcs": 0.0,
            "packing_shortage_pack": 0.0,
            "production_progress_pct": 0.0,
            "packing_progress_pct": 0.0,
            "production_bottleneck_count": 0.0,
            "packing_bottleneck_count": 0.0,
        }
    request_pack = float(view["요청합계(PACK)"].sum())
    request_pcs = float(view["요청합계(PCS)"].sum())
    production_shortage_pcs = float(view["생산부족수량"].sum())
    packing_shortage_pack = float(view["포장부족수량"].sum())
    production_progress = (
        (request_pcs - production_shortage_pcs) / request_pcs * 100.0
        if request_pcs > 0
        else 0.0
    )
    packing_done_pack = max(0.0, request_pack - packing_shortage_pack)
    packing_progress = (packing_done_pack / request_pack * 100.0) if request_pack > 0 else 0.0
    return {
        "production_code_count": float(view["생산코드"].nunique()),
        "request_pack": request_pack,
        "request_pcs": request_pcs,
        "production_shortage_pcs": production_shortage_pcs,
        "packing_shortage_pack": packing_shortage_pack,
        "production_progress_pct": min(100.0, max(0.0, production_progress)),
        "packing_progress_pct": min(100.0, max(0.0, packing_progress)),
        "production_bottleneck_count": float(view["병목 상태"].astype(str).str.contains("생산 병목", na=False).sum()),
        "packing_bottleneck_count": float(view["병목 상태"].astype(str).str.contains("포장 병목", na=False).sum()),
    }


def render_production_power_kpis(view: pd.DataFrame, unit_mode: str = UNIT_PACK) -> None:
    kpi = calc_production_power_kpis(view)
    if unit_mode == UNIT_PCS:
        items = [
            ("생산코드 수", f"{int(kpi['production_code_count']):,}", "normal"),
            ("총 요청 PCS", format_int(kpi["request_pcs"]), "normal"),
            ("총 생산부족수량(PCS)", format_int(kpi["production_shortage_pcs"]), "risk"),
            ("생산진도율", f"{kpi['production_progress_pct']:.1f}%", "normal"),
            ("포장진도율", f"{kpi['packing_progress_pct']:.1f}%", "normal"),
        ]
    else:
        items = [
            ("생산코드 수", f"{int(kpi['production_code_count']):,}", "normal"),
            ("총 요청 PACK", format_int(kpi["request_pack"]), "normal"),
            ("총 포장부족(PACK)", format_int(kpi["packing_shortage_pack"]), "warn"),
            ("총 생산부족수량(PCS)", format_int(kpi["production_shortage_pcs"]), "risk"),
        ]
    cards = "".join(
        "<div class='mini-kpi-card'>"
        f"<div class='metric-label'>{escape(label)}</div>"
        f"<div class='metric-value {tone}'>{value}</div>"
        "</div>"
        for label, value, tone in items
    )
    st.markdown(f"<div class='mini-kpi-grid'>{cards}</div>", unsafe_allow_html=True)


def due_d_day_label(value: Any) -> str:
    due = pd.to_datetime(value, errors="coerce")
    if pd.isna(due):
        return "-"
    today = pd.Timestamp.now(tz="Asia/Seoul").tz_localize(None).normalize()
    days = int((due.normalize() - today).days)
    if days <= 0:
        return "D-0 🔴"
    if days <= 1:
        return f"D-{days} 🟠"
    if days <= 3:
        return f"D-{days} 🟡"
    return f"D-{days} 🟢"


def render_pack_composition_chart(selected_row: pd.Series, pack_labels: list[str]) -> None:
    chart_df = pd.DataFrame(
        {
            "PACK": pack_labels,
            "필요팩": [float(selected_row.get(label, 0.0)) for label in pack_labels],
        }
    )
    fig = px.bar(
        chart_df,
        x="PACK",
        y="필요팩",
        text="필요팩",
        title="PACK 구성",
        color_discrete_sequence=[NAVY],
    )
    fig.update_traces(texttemplate="%{text:,.0f}", textposition="outside")
    fig.update_layout(
        paper_bgcolor=WHITE,
        plot_bgcolor=WHITE,
        margin=dict(l=8, r=8, t=48, b=8),
        yaxis_title="필요팩",
        xaxis_title="",
        showlegend=False,
    )
    st.plotly_chart(fig, width="stretch")


def render_production_progress_panel(selected_row: pd.Series) -> None:
    production_progress = float(selected_row.get("생산진도율", 0.0))
    packing_progress = float(selected_row.get("포장진도율", 0.0))
    due_label = due_d_day_label(selected_row.get("_expected_date_sort", pd.NaT))
    panel = (
        "<div class='progress-summary-panel'>"
        "<div>"
        "<div class='section-sub'>생산/포장 Progress Bar</div>"
        f"{progress_cell_html(production_progress, '생산')}"
        f"{progress_cell_html(packing_progress, '포장')}"
        "</div>"
        "<div class='dday-box'>"
        "<div class='metric-label'>생산완료예상 D-Day</div>"
        f"<div class='dday-value'>{escape(due_label)}</div>"
        "</div>"
        "</div>"
    )
    st.markdown(panel, unsafe_allow_html=True)


def build_production_sales_detail_view(rows: pd.DataFrame, production_code: str, power_label: str) -> pd.DataFrame:
    scope = rows[(rows["production_code_display"] == production_code) & (rows["POWER"] == power_label)].copy()
    columns = [
        "판매코드",
        "제품명",
        "PACK 단위",
        "필요팩",
        "요청PCS",
        "포장완료PACK",
        "포장부족PACK",
        "생산필요수량(PCS)",
        "생산부족수량(PCS)",
        "생산진도율",
        "포장진도율",
        "생산완료예상일",
    ]
    if scope.empty:
        return pd.DataFrame(columns=columns + ["_pack_sort"])

    grouped = (
        scope.groupby(["sales_code", "product_name", "_pack_label", "_pack_sort"], dropna=False)
        .agg(
            request_pack=("request_pack", "sum"),
            request_pcs=("request_pcs", "sum"),
            packing_pack=("packing_recognized_pack", "sum"),
            production_shortage_pcs=("production_shortage_qty", "sum"),
            expected_date=("production_due_date", max_datetime),
        )
        .reset_index()
        .rename(
            columns={
                "sales_code": "판매코드",
                "product_name": "제품명",
                "_pack_label": "PACK 단위",
                "request_pack": "필요팩",
                "request_pcs": "요청PCS",
                "packing_pack": "포장완료PACK",
            }
        )
    )
    grouped["포장부족PACK"] = (grouped["필요팩"] - grouped["포장완료PACK"]).clip(lower=0.0)
    grouped["생산필요수량(PCS)"] = grouped["production_shortage_pcs"]
    grouped["생산부족수량(PCS)"] = grouped["production_shortage_pcs"]
    grouped["생산진도율"] = calc_production_progress_pct(grouped["요청PCS"], grouped["production_shortage_pcs"])
    grouped["포장진도율"] = np.where(
        grouped["필요팩"] > 0,
        grouped["포장완료PACK"] / grouped["필요팩"] * 100.0,
        0.0,
    )
    grouped["포장진도율"] = np.clip(grouped["포장진도율"], 0.0, 100.0)
    grouped["생산완료예상일"] = grouped["expected_date"].map(display_date_or_dash)
    grouped = grouped.sort_values(
        ["_pack_sort", "expected_date", "포장부족PACK", "필요팩"],
        ascending=[True, True, False, False],
        na_position="last",
        kind="stable",
    )
    return grouped[columns + ["_pack_sort"]].copy()


def production_scope_from_row(code_summary: pd.DataFrame, production_code: str) -> pd.DataFrame:
    work = with_operational_columns(code_summary)
    return work[work["production_code_display"] == production_code].copy()


def sales_status_label(row: pd.Series) -> str:
    shortage = float(row.get("포장부족", 0.0))
    due = pd.to_datetime(row.get("production_due_date", row.get("request_due_date", pd.NaT)), errors="coerce")
    today = pd.Timestamp.now(tz="Asia/Seoul").tz_localize(None).normalize()
    if shortage <= 0:
        return "완료"
    if pd.notna(due) and due <= today + pd.Timedelta(days=7):
        return "긴급"
    return "부족"


def sales_code_base(value: Any) -> str:
    prefix = extract_sales_prefix(value)
    if re.fullmatch(r"S\d{3}", prefix.upper()):
        return prefix.upper()
    return prefix.upper() if prefix else clean_str(value)


@st.cache_data(show_spinner=False, max_entries=24)
def build_sales_order_main_view(
    code_summary: pd.DataFrame,
    stock_threshold_pack: float = INVENTORY_STOCK_THRESHOLD_DEFAULT,
    today_key: str | None = None,
) -> pd.DataFrame:
    if code_summary.empty:
        return pd.DataFrame(
            columns=[
                "우선등급",
                "판매코드",
                "생산코드",
                "기간구분",
                "제품분류",
                "제품명",
                "PACK",
                "POWER",
                "요청PACK",
                "요청PCS",
                "생산요청물량",
                "생산요청물량(PACK)",
                "생산요청물량(PCS)",
                "용마입고수량",
                "용마입고수량(PACK)",
                "용마입고수량(PCS)",
                "용마입고대기수량",
                "용마입고대기수량(PACK)",
                "용마입고대기수량(PCS)",
                "포장가능재고(PCS)",
                "용마창고재고 (PACK)",
                "재고기준(PACK)",
                "재고부족(PACK)",
                "생산부족",
                "생산부족(PCS)",
                "포장부족",
                "포장부족(PCS)",
                "생산진도율",
                "용마입고율",
                "생산완료예상일",
                "상태",
            ]
        )
    work = add_allocated_production_basis(with_operational_columns(code_summary))
    pack_unit = pd.to_numeric(work.get("pack_unit", pd.Series(np.nan, index=work.index)), errors="coerce")
    request_pack = pd.to_numeric(work.get("request_pack", pd.Series(0.0, index=work.index)), errors="coerce").fillna(0.0)
    request_pcs = pd.to_numeric(work.get("request_pcs", pd.Series(0.0, index=work.index)), errors="coerce").fillna(0.0)
    implied_unit = np.where(request_pack > 0, request_pcs / request_pack, np.nan)
    pcs_per_pack = pack_unit.where(pack_unit > 0, implied_unit)
    pcs_per_pack = pd.Series(pcs_per_pack, index=work.index).replace([np.inf, -np.inf], np.nan).fillna(1.0)
    pcs_per_pack = pcs_per_pack.where(pcs_per_pack > 0, 1.0)
    yongma_in_pack = pd.to_numeric(
        work.get("yongma_recognized_pack", pd.Series(0.0, index=work.index)),
        errors="coerce",
    ).fillna(0.0)
    packing_pack = pd.to_numeric(
        work.get("packing_recognized_pack", pd.Series(0.0, index=work.index)),
        errors="coerce",
    ).fillna(0.0)
    work["_yongma_in_pcs"] = (yongma_in_pack * pcs_per_pack).clip(lower=0.0)
    work["_yongma_wait_pcs"] = ((packing_pack - yongma_in_pack).clip(lower=0.0) * pcs_per_pack).clip(lower=0.0)
    work["_packing_shortage_pcs"] = ((request_pack - yongma_in_pack).clip(lower=0.0) * pcs_per_pack).clip(lower=0.0)
    grouped = (
        work.groupby("sales_code", dropna=False)
        .agg(
            production_code=("production_code_display", join_unique),
            factory_group=("factory_group", join_unique),
            period_group=("period_group", first_nonempty),
            product_group=("제품분류", first_nonempty),
            product_name=("product_name", join_unique),
            pack_label=("_pack_label", join_unique),
            power=("POWER", first_nonempty),
            power_value=("power_value", "min"),
            request_pack=("request_pack", "sum"),
            request_pcs=("request_pcs", "sum"),
            packing_pack=("packing_recognized_pack", "sum"),
            yongma_in_pack=("yongma_recognized_pack", "sum"),
            yongma_in_pcs=("_yongma_in_pcs", "sum"),
            yongma_wait_pcs=("_yongma_wait_pcs", "sum"),
            packing_shortage_pcs=("_packing_shortage_pcs", "sum"),
            available_stock_pack=("available_stock_pack", sum_numeric_or_nan),
            production_shortage=("production_shortage_qty", "sum"),
            sample_available_pcs=("_allocated_sample_available_pcs", "sum"),
            production_due_date=("production_due_date", max_datetime),
        )
        .reset_index()
        .rename(
            columns={
                "sales_code": "판매코드",
                "production_code": "생산코드",
                "factory_group": "공장구분",
                "period_group": "기간구분",
                "product_group": "제품분류",
                "product_name": "제품명",
                "pack_label": "PACK",
                "power": "POWER",
                "request_pack": "요청PACK",
                "request_pcs": "요청PCS",
                "yongma_in_pack": "용마입고수량",
                "yongma_in_pcs": "용마입고수량(PCS)",
                "yongma_wait_pcs": "용마입고대기수량(PCS)",
                "packing_shortage_pcs": "포장부족(PCS)",
                "available_stock_pack": "용마창고재고 (PACK)",
                "production_shortage": "생산부족",
                "sample_available_pcs": "샘플신청가능수량",
            }
        )
    )
    grouped["용마입고대기수량"] = (grouped["packing_pack"] - grouped["용마입고수량"]).clip(lower=0.0)
    grouped["용마입고수량(PACK)"] = grouped["용마입고수량"]
    grouped["용마입고대기수량(PACK)"] = grouped["용마입고대기수량"]
    grouped["포장가능재고(PCS)"] = (
        grouped["요청PCS"] - grouped["생산부족"] + grouped["샘플신청가능수량"]
    ).clip(lower=0.0)
    grouped["포장부족"] = (grouped["요청PACK"] - grouped["용마입고수량"]).clip(lower=0.0)
    grouped["생산진도율"] = calc_production_progress_pct(grouped["요청PCS"], grouped["생산부족"])
    grouped["용마입고율"] = np.where(
        grouped["요청PACK"] > 0,
        grouped["용마입고수량"] / grouped["요청PACK"] * 100.0,
        0.0,
    )
    grouped["용마입고율"] = np.clip(grouped["용마입고율"], 0.0, 100.0)
    grouped["생산완료예상일"] = grouped["production_due_date"].map(display_date_or_dash)
    grouped["상태"] = grouped.apply(sales_status_label, axis=1)
    grouped = add_priority_columns(
        grouped,
        stock_threshold_pack,
        shortage_col="포장부족",
        due_col="production_due_date",
        stock_col="용마창고재고 (PACK)",
        request_col="요청PACK",
    )
    grouped["요청합계(PACK)"] = grouped["요청PACK"]
    grouped["요청합계(PCS)"] = grouped["요청PCS"]
    grouped["생산요청물량"] = grouped["요청PCS"]
    grouped["생산요청물량(PACK)"] = grouped["요청PACK"]
    grouped["생산요청물량(PCS)"] = grouped["요청PCS"]
    grouped["생산필요수량(PCS)"] = grouped["생산부족"]
    grouped["생산부족수량(PCS)"] = grouped["생산부족"]
    grouped["생산부족(PCS)"] = grouped["생산부족"]
    grouped["포장부족(PACK)"] = grouped["포장부족"]
    return grouped[
        [
            "우선등급",
            "판매코드",
            "생산코드",
            "기간구분",
            "제품분류",
            "제품명",
            "PACK",
            "POWER",
            "요청PACK",
            "요청PCS",
            "요청합계(PACK)",
            "요청합계(PCS)",
            "생산요청물량",
            "생산요청물량(PACK)",
            "생산요청물량(PCS)",
            "용마입고수량",
            "용마입고수량(PACK)",
            "용마입고수량(PCS)",
            "용마입고대기수량",
            "용마입고대기수량(PACK)",
            "용마입고대기수량(PCS)",
            "포장가능재고(PCS)",
            "샘플신청가능수량",
            "용마창고재고 (PACK)",
            "재고기준(PACK)",
            "재고부족(PACK)",
            "생산부족",
            "생산필요수량(PCS)",
            "생산부족수량(PCS)",
            "생산부족(PCS)",
            "포장부족",
            "포장부족(PACK)",
            "포장부족(PCS)",
            "생산진도율",
            "용마입고율",
            "생산완료예상일",
            "상태",
            "power_value",
            "_priority_sort",
            "_request_due_date_sort",
        ]
    ].sort_values(
        ["_priority_sort", "_request_due_date_sort", "재고부족(PACK)", "포장부족", "생산부족"],
        ascending=[True, True, False, False, False],
        na_position="last",
        kind="stable",
    )


def build_sales_code_group_main_view(
    sales_detail_view: pd.DataFrame,
    stock_threshold_pack: float = INVENTORY_STOCK_THRESHOLD_DEFAULT,
) -> pd.DataFrame:
    columns = [
        "우선등급",
        "판매코드",
        "기간구분",
        "제품분류",
        "대표 제품명",
        "생산코드",
        "PACK",
        "POWER 수",
        "생산요청물량(PACK)",
        "생산요청물량(PCS)",
        "용마입고수량(PACK)",
        "용마입고수량(PCS)",
        "용마입고대기수량(PACK)",
        "용마입고대기수량(PCS)",
        "포장가능재고(PCS)",
        "샘플신청가능수량",
        "용마창고재고 (PACK)",
        "재고기준(PACK)",
        "재고부족(PACK)",
        "생산부족(PCS)",
        "포장부족(PACK)",
        "포장부족(PCS)",
        "생산진도율",
        "용마입고율",
        "생산완료예상일",
        "상태",
        "_sales_code_base",
        "_priority_sort",
        "_request_due_date_sort",
    ]
    if sales_detail_view.empty:
        return pd.DataFrame(columns=columns)

    work = sales_detail_view.copy()
    work["_sales_code_base"] = work["판매코드"].map(sales_code_base)
    for col in [
        "요청PACK",
        "요청PCS",
        "용마입고수량(PACK)",
        "용마입고수량(PCS)",
        "용마입고대기수량(PACK)",
        "용마입고대기수량(PCS)",
        "포장가능재고(PCS)",
        "샘플신청가능수량",
        "용마창고재고 (PACK)",
        "생산부족(PCS)",
        "포장부족(PACK)",
        "포장부족(PCS)",
    ]:
        if col not in work.columns:
            work[col] = 0.0
        work[col] = pd.to_numeric(work[col], errors="coerce").fillna(0.0)
    if "_request_due_date_sort" not in work.columns:
        work["_request_due_date_sort"] = pd.to_datetime(work.get("생산완료예상일", pd.NaT), errors="coerce")
    work["_request_due_date_sort"] = pd.to_datetime(work["_request_due_date_sort"], errors="coerce")

    grouped = (
        work.groupby("_sales_code_base", dropna=False)
        .agg(
            period_group=("기간구분", first_nonempty),
            product_group=("제품분류", first_nonempty),
            product_name=("제품명", first_nonempty),
            production_code=("생산코드", join_unique),
            pack_label=("PACK", join_unique),
            power_count=("POWER", "nunique"),
            request_pack=("요청PACK", "sum"),
            request_pcs=("요청PCS", "sum"),
            yongma_in_pack=("용마입고수량(PACK)", "sum"),
            yongma_in_pcs=("용마입고수량(PCS)", "sum"),
            yongma_wait_pack=("용마입고대기수량(PACK)", "sum"),
            yongma_wait_pcs=("용마입고대기수량(PCS)", "sum"),
            packable_pcs=("포장가능재고(PCS)", "sum"),
            sample_available_pcs=("샘플신청가능수량", "sum"),
            available_stock_pack=("용마창고재고 (PACK)", sum_numeric_or_nan),
            production_shortage_pcs=("생산부족(PCS)", "sum"),
            packing_shortage_pack=("포장부족(PACK)", "sum"),
            packing_shortage_pcs=("포장부족(PCS)", "sum"),
            production_due_date=("_request_due_date_sort", max_datetime),
        )
        .reset_index()
        .rename(
            columns={
                "_sales_code_base": "판매코드",
                "period_group": "기간구분",
                "product_group": "제품분류",
                "product_name": "대표 제품명",
                "production_code": "생산코드",
                "pack_label": "PACK",
                "power_count": "POWER 수",
                "request_pack": "생산요청물량(PACK)",
                "request_pcs": "생산요청물량(PCS)",
                "yongma_in_pack": "용마입고수량(PACK)",
                "yongma_in_pcs": "용마입고수량(PCS)",
                "yongma_wait_pack": "용마입고대기수량(PACK)",
                "yongma_wait_pcs": "용마입고대기수량(PCS)",
                "packable_pcs": "포장가능재고(PCS)",
                "sample_available_pcs": "샘플신청가능수량",
                "available_stock_pack": "용마창고재고 (PACK)",
                "production_shortage_pcs": "생산부족(PCS)",
                "packing_shortage_pack": "포장부족(PACK)",
                "packing_shortage_pcs": "포장부족(PCS)",
            }
        )
    )
    grouped["생산진도율"] = calc_production_progress_pct(grouped["생산요청물량(PCS)"], grouped["생산부족(PCS)"])
    grouped["용마입고율"] = np.where(
        grouped["생산요청물량(PACK)"] > 0,
        grouped["용마입고수량(PACK)"] / grouped["생산요청물량(PACK)"] * 100.0,
        0.0,
    )
    grouped["용마입고율"] = np.clip(grouped["용마입고율"], 0.0, 100.0)
    grouped["생산완료예상일"] = grouped["production_due_date"].map(display_date_or_dash)
    grouped["상태"] = grouped.apply(
        lambda row: sales_status_label(
            pd.Series(
                {
                    "포장부족": row["포장부족(PACK)"],
                    "production_due_date": row["production_due_date"],
                }
            )
        ),
        axis=1,
    )
    grouped = add_priority_columns(
        grouped,
        stock_threshold_pack,
        shortage_col="포장부족(PACK)",
        due_col="production_due_date",
        stock_col="용마창고재고 (PACK)",
        request_col="생산요청물량(PACK)",
    )
    grouped["_sales_code_base"] = grouped["판매코드"]
    grouped = grouped.sort_values(
        ["_priority_sort", "_request_due_date_sort", "재고부족(PACK)", "포장부족(PACK)", "생산부족(PCS)"],
        ascending=[True, True, False, False, False],
        na_position="last",
        kind="stable",
    )
    return grouped[columns].copy()


PRODUCT_COMPLETION_STATUS_FILTERS = ["전체", "생산중", "생산완료", "미계획"]
PRODUCT_COMPLETION_STATUS_LABELS = {
    "생산완료": "🟢 생산완료",
    "생산중": "🟠 생산중",
    "미계획": "🔴 미계획",
}
PRODUCT_COMPLETION_MAIN_COLUMNS = [
    "판매코드",
    "제품명",
    "POWER수",
    "생산요청물량 (PCS)",
    "용마입고수량 (PCS)",
    "용마입고대기 (PCS)",
    "포장부족수량 (PCS)",
    "포장가능수량 (PCS)",
    "생산부족수량 (PCS)",
    "생산완료예상일",
    "생산상태",
]
PRODUCT_COMPLETION_DETAIL_COLUMNS = [
    "POWER",
    "생산요청물량 (PCS)",
    "용마입고수량 (PCS)",
    "용마입고대기 (PCS)",
    "포장부족수량 (PCS)",
    "포장가능수량 (PCS)",
    "생산부족수량 (PCS)",
    "생산완료예상일",
]


def row_pcs_per_pack(work: pd.DataFrame) -> pd.Series:
    pack_unit = pd.to_numeric(work.get("pack_unit", pd.Series(np.nan, index=work.index)), errors="coerce")
    request_pack = pd.to_numeric(work.get("request_pack", pd.Series(0.0, index=work.index)), errors="coerce").fillna(0.0)
    request_pcs = pd.to_numeric(work.get("request_pcs", pd.Series(0.0, index=work.index)), errors="coerce").fillna(0.0)
    implied_unit = np.where(request_pack > 0, request_pcs / request_pack, np.nan)
    unit = pack_unit.where(pack_unit > 0, implied_unit)
    unit = pd.Series(unit, index=work.index).replace([np.inf, -np.inf], np.nan).fillna(1.0)
    return unit.where(unit > 0, 1.0)


def production_completion_status_key(plan_date: Any, production_shortage_pcs: Any) -> str:
    plan = pd.to_datetime(plan_date, errors="coerce")
    shortage = to_number_value(production_shortage_pcs)
    if shortage <= 0:
        return "생산완료"
    if pd.isna(plan):
        return "미계획"
    return "생산중"


def product_completion_display_date(value: Any) -> str:
    text = format_date(value)
    return text if text else "-"


def product_completion_status_label(status_key: Any) -> str:
    return PRODUCT_COMPLETION_STATUS_LABELS.get(clean_str(status_key), clean_str(status_key))


@st.cache_data(show_spinner=False, max_entries=24)
def build_product_completion_power_view(code_summary: pd.DataFrame) -> pd.DataFrame:
    columns = PRODUCT_COMPLETION_DETAIL_COLUMNS + [
        "판매코드",
        "제품명",
        "기간구분",
        "제품분류",
        "_sales_code_base",
        "_section",
        "_status_key",
        "_expected_date_sort",
        "power_value",
    ]
    if code_summary.empty:
        return pd.DataFrame(columns=columns)

    work = add_allocated_production_basis(with_operational_columns(code_summary))
    request_pack = pd.to_numeric(work.get("request_pack", pd.Series(0.0, index=work.index)), errors="coerce").fillna(0.0)
    request_pcs = pd.to_numeric(work.get("request_pcs", pd.Series(0.0, index=work.index)), errors="coerce").fillna(0.0)
    work = work[(request_pack > 0) | (request_pcs > 0)].copy()
    if work.empty:
        return pd.DataFrame(columns=columns)

    pcs_per_pack = row_pcs_per_pack(work)
    packing_pack = pd.to_numeric(
        work.get("packing_recognized_pack", pd.Series(0.0, index=work.index)),
        errors="coerce",
    ).fillna(0.0)
    yongma_pack = pd.to_numeric(
        work.get("yongma_recognized_pack", pd.Series(0.0, index=work.index)),
        errors="coerce",
    ).fillna(0.0)
    production_shortage = pd.to_numeric(
        work.get("_allocated_production_shortage_qty", pd.Series(0.0, index=work.index)),
        errors="coerce",
    ).fillna(0.0)
    sample_available = pd.to_numeric(
        work.get("_allocated_sample_available_pcs", pd.Series(0.0, index=work.index)),
        errors="coerce",
    ).fillna(0.0)

    work["_sales_code_base"] = work["sales_code"].map(sales_code_base)
    work["_request_pcs"] = request_pcs
    work["_yongma_in_pcs"] = (yongma_pack * pcs_per_pack).clip(lower=0.0)
    work["_yongma_wait_pcs"] = ((packing_pack - yongma_pack).clip(lower=0.0) * pcs_per_pack).clip(lower=0.0)
    work["_packing_pcs"] = (packing_pack * pcs_per_pack).clip(lower=0.0)
    work["_packing_shortage_pcs"] = (request_pcs - work["_packing_pcs"]).clip(lower=0.0)
    work["_packable_pcs"] = (request_pcs - production_shortage + sample_available).clip(lower=0.0)
    work["_production_shortage_pcs"] = production_shortage.clip(lower=0.0)
    work["_production_plan_date"] = pd.to_datetime(
        work.get("production_plan_date", pd.Series(pd.NaT, index=work.index)),
        errors="coerce",
    )
    work["_expected_date_sort"] = pd.to_datetime(
        work.get("production_complete_expected_date", pd.Series(pd.NaT, index=work.index)),
        errors="coerce",
    )
    missing_expected = work["_expected_date_sort"].isna() & work["_production_plan_date"].notna()
    work.loc[missing_expected, "_expected_date_sort"] = (
        work.loc[missing_expected, "_production_plan_date"] + pd.Timedelta(days=5)
    )
    production_code_source = work.get("production_code", pd.Series("", index=work.index))
    p_code_source = work.get("p_code", pd.Series("", index=work.index))
    work["_base_p_code"] = [
        first_nonempty(
            [
                extract_base_p_code_key(production_code),
                extract_base_p_code_key(p_code),
            ]
        )
        for production_code, p_code in zip(production_code_source, p_code_source)
    ]
    has_base_p_code = work["_base_p_code"].map(clean_str) != ""
    if has_base_p_code.any():
        latest_plan_by_p = work.loc[has_base_p_code].groupby("_base_p_code")["_production_plan_date"].transform("max")
        latest_expected_by_p = work.loc[has_base_p_code].groupby("_base_p_code")["_expected_date_sort"].transform("max")
        work.loc[has_base_p_code, "_production_plan_date"] = latest_plan_by_p.where(
            latest_plan_by_p.notna(),
            work.loc[has_base_p_code, "_production_plan_date"],
        )
        work.loc[has_base_p_code, "_expected_date_sort"] = latest_expected_by_p.where(
            latest_expected_by_p.notna(),
            work.loc[has_base_p_code, "_expected_date_sort"],
        )
    work["_section"] = work["본품분류"].map(family_card_section)

    grouped = (
        work.groupby(["_sales_code_base", "POWER", "power_value"], dropna=False)
        .agg(
            product_name=("product_name", first_nonempty),
            period_group=("period_group", first_nonempty),
            product_group=("제품분류", first_nonempty),
            section=("_section", first_nonempty),
            request_pcs=("_request_pcs", "sum"),
            yongma_in_pcs=("_yongma_in_pcs", "sum"),
            yongma_wait_pcs=("_yongma_wait_pcs", "sum"),
            packing_shortage_pcs=("_packing_shortage_pcs", "sum"),
            packable_pcs=("_packable_pcs", "sum"),
            production_shortage_pcs=("_production_shortage_pcs", "sum"),
            production_plan_date=("_production_plan_date", min_datetime),
            expected_date=("_expected_date_sort", max_datetime),
        )
        .reset_index()
        .rename(
            columns={
                "_sales_code_base": "판매코드",
                "product_name": "제품명",
                "period_group": "기간구분",
                "product_group": "제품분류",
                "section": "_section",
                "request_pcs": "생산요청물량 (PCS)",
                "yongma_in_pcs": "용마입고수량 (PCS)",
                "yongma_wait_pcs": "용마입고대기 (PCS)",
                "packing_shortage_pcs": "포장부족수량 (PCS)",
                "packable_pcs": "포장가능수량 (PCS)",
                "production_shortage_pcs": "생산부족수량 (PCS)",
            }
        )
    )
    for col in [
        "생산요청물량 (PCS)",
        "용마입고수량 (PCS)",
        "용마입고대기 (PCS)",
        "포장부족수량 (PCS)",
        "포장가능수량 (PCS)",
        "생산부족수량 (PCS)",
    ]:
        grouped[col] = pd.to_numeric(grouped[col], errors="coerce").fillna(0.0).round(0).astype("int64")
    grouped["_status_key"] = grouped.apply(
        lambda row: production_completion_status_key(row["production_plan_date"], row["생산부족수량 (PCS)"]),
        axis=1,
    )
    grouped["생산상태"] = grouped["_status_key"].map(product_completion_status_label)
    grouped["생산완료예상일"] = grouped["expected_date"].map(product_completion_display_date)
    grouped.loc[grouped["_status_key"] == "생산완료", "생산완료예상일"] = "-"
    grouped["_expected_date_sort"] = pd.to_datetime(grouped["expected_date"], errors="coerce")
    grouped["_sales_code_base"] = grouped["판매코드"]
    return sort_power_detail_default(
        grouped[columns + ["production_plan_date", "expected_date", "생산상태"]],
        extra_cols=["_expected_date_sort", "생산부족수량 (PCS)", "판매코드"],
        extra_ascending=[True, False, True],
    )


@st.cache_data(show_spinner=False, max_entries=24)
def build_product_completion_main_view(power_view: pd.DataFrame) -> pd.DataFrame:
    columns = PRODUCT_COMPLETION_MAIN_COLUMNS + [
        "_sales_code_base",
        "_section",
        "_status_key",
        "_expected_date_sort",
    ]
    if power_view.empty:
        return pd.DataFrame(columns=columns)

    work = power_view.copy()
    grouped = (
        work.groupby("_sales_code_base", dropna=False)
        .agg(
            product_name=("제품명", first_nonempty),
            period_group=("기간구분", first_nonempty),
            product_group=("제품분류", first_nonempty),
            section=("_section", first_nonempty),
            power_count=("POWER", "nunique"),
            request_pcs=("생산요청물량 (PCS)", "sum"),
            yongma_in_pcs=("용마입고수량 (PCS)", "sum"),
            yongma_wait_pcs=("용마입고대기 (PCS)", "sum"),
            packing_shortage_pcs=("포장부족수량 (PCS)", "sum"),
            packable_pcs=("포장가능수량 (PCS)", "sum"),
            production_shortage_pcs=("생산부족수량 (PCS)", "sum"),
            production_plan_date=("production_plan_date", min_datetime),
            expected_date=("_expected_date_sort", max_datetime),
        )
        .reset_index()
        .rename(
            columns={
                "_sales_code_base": "판매코드",
                "product_name": "제품명",
                "period_group": "기간구분",
                "product_group": "제품분류",
                "section": "_section",
                "power_count": "POWER수",
                "request_pcs": "생산요청물량 (PCS)",
                "yongma_in_pcs": "용마입고수량 (PCS)",
                "yongma_wait_pcs": "용마입고대기 (PCS)",
                "packing_shortage_pcs": "포장부족수량 (PCS)",
                "packable_pcs": "포장가능수량 (PCS)",
                "production_shortage_pcs": "생산부족수량 (PCS)",
            }
        )
    )
    for col in [
        "POWER수",
        "생산요청물량 (PCS)",
        "용마입고수량 (PCS)",
        "용마입고대기 (PCS)",
        "포장부족수량 (PCS)",
        "포장가능수량 (PCS)",
        "생산부족수량 (PCS)",
    ]:
        grouped[col] = pd.to_numeric(grouped[col], errors="coerce").fillna(0.0).round(0).astype("int64")
    grouped["_status_key"] = grouped.apply(
        lambda row: production_completion_status_key(row["production_plan_date"], row["생산부족수량 (PCS)"]),
        axis=1,
    )
    grouped["생산상태"] = grouped["_status_key"].map(product_completion_status_label)
    grouped["생산완료예상일"] = grouped["expected_date"].map(product_completion_display_date)
    grouped.loc[grouped["_status_key"] == "생산완료", "생산완료예상일"] = "-"
    grouped["_expected_date_sort"] = pd.to_datetime(grouped["expected_date"], errors="coerce")
    grouped["_sales_code_base"] = grouped["판매코드"]
    return grouped[columns].sort_values(
        ["_expected_date_sort", "생산부족수량 (PCS)", "판매코드"],
        ascending=[True, False, True],
        na_position="last",
        kind="stable",
    )


def filter_product_completion_view(
    view: pd.DataFrame,
    period_filter: str,
    unified_query: str,
) -> pd.DataFrame:
    if view.empty:
        return view.copy()
    out = view.copy()
    if period_filter != "전체":
        section_filter = "1DAY" if period_filter == "1-DAY" else period_filter
        out = out[out["_section"].astype(str) == section_filter]
    if unified_query.strip():
        search_columns = list(out.columns[:2])
        out = filter_dataframe_by_terms(out, unified_query, search_columns)
    return out.copy()


def product_completion_column_config() -> dict[str, Any]:
    numeric_format = "%,.0f"
    column_config = drilldown_column_config()
    column_config.update(
        {
            "판매코드": st.column_config.TextColumn("판매코드", width=82),
            "제품명": st.column_config.TextColumn("제품명", width=270),
            "POWER수": st.column_config.NumberColumn("POWER", format=numeric_format, width=74),
            "생산요청물량 (PCS)": st.column_config.NumberColumn("요청PCS", format=numeric_format, width=118),
            "용마입고수량 (PCS)": st.column_config.NumberColumn("입고PCS", format=numeric_format, width=116),
            "용마입고대기 (PCS)": st.column_config.NumberColumn("대기PCS", format=numeric_format, width=116),
            "포장부족수량 (PCS)": st.column_config.NumberColumn("포장부족", format=numeric_format, width=116),
            "포장가능수량 (PCS)": st.column_config.NumberColumn("포장가능", format=numeric_format, width=116),
            "생산부족수량 (PCS)": st.column_config.NumberColumn("생산부족", format=numeric_format, width=116),
            "생산완료예상일": st.column_config.TextColumn("예상일", width=92),
            "생산상태": st.column_config.TextColumn("상태", width=92),
        }
    )
    return column_config


def product_completion_summary_html(view: pd.DataFrame) -> str:
    status = view.get("_status_key", pd.Series(dtype="object")).map(clean_str)
    total = len(view)
    done = int((status == "생산완료").sum())
    active = int((status == "생산중").sum())
    unplanned = int((status == "미계획").sum())
    shortage_total = pd.to_numeric(
        view.get("생산부족수량 (PCS)", pd.Series(dtype="float64")),
        errors="coerce",
    ).fillna(0.0).sum()
    items = [
        ("총 판매코드", format_int(total), "neutral"),
        ("생산완료", format_int(done), "done"),
        ("생산중", format_int(active), "active"),
        ("미계획", format_int(unplanned), "risk"),
        ("생산부족 PCS", format_int(float(shortage_total)), "risk" if shortage_total > 0 else "neutral"),
    ]
    cards = "".join(
        "<div class='completion-summary-item {tone}'>"
        f"<span>{escape(label)}</span>"
        f"<strong>{escape(value)}</strong>"
        "</div>".format(tone=tone)
        for label, value, tone in items
    )
    return f"<div class='completion-summary-card'>{cards}</div>"


def render_product_completion_main_table(
    title: str,
    sub: str,
    df: pd.DataFrame,
    key: str,
    height: int,
) -> pd.Series | None:
    if title or sub:
        render_panel_title(title, sub)
    if df.empty:
        st.warning("조건에 맞는 데이터가 없습니다.")
        return None

    st.markdown(product_completion_summary_html(df), unsafe_allow_html=True)
    display_df = dataframe_for_streamlit(df)
    event = st.dataframe(
        display_df,
        hide_index=True,
        height=dataframe_auto_height(len(display_df), height, row_height=48),
        width="stretch",
        column_config=product_completion_column_config(),
        column_order=visible_columns(display_df, PRODUCT_COMPLETION_MAIN_COLUMNS),
        on_select="rerun",
        selection_mode="single-row",
        key=key,
    )
    return get_selected_row(event, df)


def render_product_completion_detail_dialog(
    selected_row: pd.Series,
    detail_view: pd.DataFrame,
    table_nonce_key: str,
) -> None:
    sales_code = clean_str(selected_row.get("_sales_code_base", selected_row.get("판매코드", "")))
    product_name = clean_str(selected_row.get("제품명", selected_row.get("대표 제품명", "")))
    title = f"판매코드 {sales_code} POWER 상세 - {product_name}"

    @st.dialog(title, width="large")
    def _dialog() -> None:
        detail_display = sort_power_detail_default(
            detail_view,
            extra_cols=["_expected_date_sort", "생산부족수량 (PCS)"],
            extra_ascending=[True, False],
        )
        st.caption(f"{sales_code}에 해당하는 POWER 기준 생산 완료 현황 | 표시 건수: {len(detail_display):,}")
        if detail_display.empty:
            st.warning("상세 데이터가 없습니다.")
        else:
            st.dataframe(
                dataframe_for_streamlit(detail_display),
                hide_index=True,
                height=dataframe_auto_height(len(detail_display), 520),
                width="stretch",
                column_config=drilldown_column_config(),
                column_order=visible_columns(detail_display, PRODUCT_COMPLETION_DETAIL_COLUMNS),
            )
        if st.button("닫기", key="close_product_completion_detail_dialog", width="stretch"):
            st.session_state[table_nonce_key] = int(st.session_state.get(table_nonce_key, 0)) + 1
            st.rerun()

    _dialog()


def render_product_completion_section(code_summary: pd.DataFrame) -> None:
    st.markdown("<div class='section-gap'></div>", unsafe_allow_html=True)
    render_panel_title(
        "제품별 생산 완료 현황",
        "판매코드 기준 생산 진행 현황과 누수규격검사 계획일 기준 생산완료예상일을 확인합니다.",
    )
    power_view = build_product_completion_power_view(code_summary)
    main_view = build_product_completion_main_view(power_view)
    if main_view.empty:
        st.info("표시할 제품별 생산 완료 현황 데이터가 없습니다.")
        return

    if st.session_state.get("product_completion_period_group_filter") not in PERIOD_GROUP_ORDER:
        st.session_state["product_completion_period_group_filter"] = "전체"

    f1, f2 = st.columns([1.2, 3.2], gap="small")
    with f1:
        period_filter = st.segmented_control(
            "기간구분",
            options=PERIOD_GROUP_ORDER,
            default=st.session_state.get("product_completion_period_group_filter", "전체"),
            key="product_completion_period_group_filter",
        )
    period_filter = str(period_filter or "전체")
    with f2:
        unified_query = st.text_input(
            "통합검색",
            value="",
            placeholder="예: S120, 소울브라운",
            key="product_completion_unified_query",
        )

    filtered = filter_product_completion_view(
        main_view,
        period_filter,
        unified_query,
    )

    table_nonce_key = "product_completion_table_nonce"
    table_nonce = int(st.session_state.get(table_nonce_key, 0))
    selected_row = render_product_completion_main_table(
        "",
        "",
        filtered,
        key=f"product_completion_table_{table_nonce}",
        height=620,
    )
    if selected_row is not None:
        selected_sales_code = clean_str(selected_row.get("_sales_code_base", selected_row.get("판매코드", "")))
        detail_scope = power_view[power_view["_sales_code_base"] == selected_sales_code].copy()
        detail_scope = sort_power_detail_default(
            detail_scope,
            extra_cols=["_expected_date_sort", "생산부족수량 (PCS)"],
            extra_ascending=[True, False],
        )
        render_product_completion_detail_dialog(selected_row, detail_scope, table_nonce_key)


def build_urgent_sales_packing_view(sales_view: pd.DataFrame, max_rows: int = 20) -> pd.DataFrame:
    columns = [
        "우선등급",
        "기간구분",
        "판매코드",
        "제품명",
        "POWER",
        "PACK",
        "생산요청물량(PACK)",
        "포장부족(PACK)",
        "생산완료예상일",
    ]
    if sales_view.empty:
        return pd.DataFrame(columns=columns)

    out = sales_view[
        (pd.to_numeric(sales_view["포장부족"], errors="coerce").fillna(0.0) > 0)
        & (sales_view["우선등급"].isin(["A 긴급", "B 주의"]))
    ].copy()
    if out.empty:
        return pd.DataFrame(columns=columns)

    out = out.sort_values(
        ["_priority_sort", "_request_due_date_sort", "재고부족(PACK)", "포장부족", "생산부족"],
        ascending=[True, True, False, False, False],
        na_position="last",
        kind="stable",
    )
    return out[columns].head(max_rows).copy()


def filter_sales_order_view(
    sales_view: pd.DataFrame,
    product_query: str = "",
    production_query: str = "",
    sales_query: str = "",
    pack_label: str = "전체",
    product_group: str = "전체",
    power_label: str = "전체",
) -> pd.DataFrame:
    if sales_view.empty:
        return sales_view.copy()

    out = sales_view
    product_q = product_query.strip()
    if product_q and "제품명" in out.columns:
        out = out[out["제품명"].astype(str).str.contains(product_q, case=False, na=False, regex=False)]
    production_q = production_query.strip()
    if production_q and "생산코드" in out.columns:
        out = out[out["생산코드"].astype(str).str.contains(production_q, case=False, na=False, regex=False)]
    sales_q = sales_query.strip()
    if sales_q and "판매코드" in out.columns:
        out = out[out["판매코드"].astype(str).str.contains(sales_q, case=False, na=False, regex=False)]
    if pack_label != "전체" and "PACK" in out.columns:
        out = out[out["PACK"] == pack_label]
    if product_group != "전체" and "제품분류" in out.columns:
        out = out[out["제품분류"] == product_group]
    if power_label != "전체" and "POWER" in out.columns:
        out = out[out["POWER"] == power_label]
    return out.copy()


def split_search_terms(query: str) -> list[str]:
    text = clean_str(query)
    if not text:
        return []
    return [term for term in [clean_str(part) for part in text.split(",")] if term]


def matching_product_alias_groups(term: str) -> list[list[str]]:
    compact = compact_search_text(term)
    if not compact:
        return []

    groups: list[list[str]] = []
    for alias, values in PRODUCT_QUERY_ALIASES.items():
        alias_compact = compact_search_text(alias)
        if alias_compact and (alias_compact in compact or compact in alias_compact):
            candidates = [alias, *values]
            groups.append(list(dict.fromkeys([candidate for candidate in candidates if clean_str(candidate)])))
    return groups


def search_candidate_mask(raw_haystack: pd.Series, compact_haystack: pd.Series, candidate: str) -> pd.Series:
    raw_term = clean_str(candidate).lower()
    compact_term = compact_search_text(candidate)
    mask = pd.Series(False, index=raw_haystack.index)
    if raw_term:
        mask = mask | raw_haystack.str.contains(raw_term, regex=False, na=False)
    if compact_term:
        mask = mask | compact_haystack.str.contains(compact_term, regex=False, na=False)
    return mask


def filter_dataframe_by_terms(
    df: pd.DataFrame,
    query: str,
    columns: list[str] | None = None,
) -> pd.DataFrame:
    terms = split_search_terms(query)
    if df.empty or not terms:
        return df.copy()

    searchable_columns = [col for col in (columns or list(df.columns)) if col in df.columns]
    searchable_columns = [col for col in searchable_columns if not str(col).startswith("_")]
    if not searchable_columns:
        return df.copy()

    source = df[searchable_columns].fillna("").astype(str)
    raw_haystack = source.apply(lambda row: " ".join(row.values.tolist()).lower(), axis=1)
    compact_haystack = raw_haystack.map(compact_search_text)
    mask = pd.Series(False, index=df.index)
    for term in terms:
        term_mask = search_candidate_mask(raw_haystack, compact_haystack, term)
        alias_groups = matching_product_alias_groups(term)
        if alias_groups:
            alias_mask = pd.Series(True, index=df.index)
            for group in alias_groups:
                group_mask = pd.Series(False, index=df.index)
                for candidate in group:
                    group_mask = group_mask | search_candidate_mask(raw_haystack, compact_haystack, candidate)
                alias_mask = alias_mask & group_mask
            term_mask = term_mask | alias_mask
        mask = mask | term_mask
    return df[mask].copy()


def render_urgent_sales_packing_list(sales_view: pd.DataFrame) -> None:
    urgent_view = build_urgent_sales_packing_view(sales_view)
    render_panel_title(
        "긴급 포장 리스트",
        "용마 보유 재고는 긴급도 판단에만 사용하고, 표에는 PACK 기준 요청·부족 수량만 표시합니다.",
    )
    if urgent_view.empty:
        st.info("현재 기준에 해당하는 긴급 포장 판매코드가 없습니다.")
    else:
        st.dataframe(
            urgent_view,
            hide_index=True,
            height=dataframe_auto_height(len(urgent_view), 260),
            width="stretch",
            column_config=drilldown_column_config(),
        )


def sales_scope_from_row(code_summary: pd.DataFrame, sales_code: str) -> pd.DataFrame:
    work = with_operational_columns(code_summary)
    return work[work["sales_code"] == sales_code].copy()


def production_shortage_pack_equivalent(work: pd.DataFrame) -> pd.Series:
    if "_allocated_production_shortage_qty" not in work.columns:
        work = add_allocated_production_basis(work)
    pack_unit = pd.to_numeric(work.get("pack_unit", pd.Series(np.nan, index=work.index)), errors="coerce")
    implied_unit = np.where(work["request_pack"] > 0, work["request_pcs"] / work["request_pack"], np.nan)
    unit = pack_unit.where(pack_unit > 0, implied_unit)
    unit = pd.Series(unit, index=work.index).replace([np.inf, -np.inf], np.nan).fillna(1.0)
    unit = unit.where(unit > 0, 1.0)
    return (work["_allocated_production_shortage_qty"] / unit).clip(lower=0.0)


def build_power_summary_view(code_summary: pd.DataFrame) -> pd.DataFrame:
    work = with_operational_columns(code_summary)
    work = work[work["power_value"].notna()].copy()
    if work.empty:
        return pd.DataFrame(
            columns=[
                "POWER",
                "기간구분",
                "요청합계(PACK)",
                "요청합계(PCS)",
                "포장 PACK",
                "포장부족(PACK)",
                "생산필요수량(PCS)",
                "생산부족수량(PCS)",
                "생산진도율",
                "포장진도율",
                "power_value",
            ]
        )
    work = add_allocated_production_basis(work)
    grouped = (
        work.groupby(["power_value", "POWER"], dropna=False)
        .agg(
            factory_group=("factory_group", join_unique),
            period_group=("period_group", first_nonempty),
            request_pack=("request_pack", "sum"),
            request_pcs=("request_pcs", "sum"),
            packing_pack=("packing_recognized_pack", "sum"),
            production_shortage_pcs=("_allocated_production_shortage_qty", "sum"),
        )
        .reset_index()
        .rename(
            columns={
                "request_pack": "요청합계(PACK)",
                "factory_group": "공장구분",
                "period_group": "기간구분",
                "request_pcs": "요청합계(PCS)",
                "packing_pack": "포장 PACK",
                "production_shortage_pcs": "생산부족수량(PCS)",
            }
        )
    )
    grouped["포장부족(PACK)"] = (grouped["요청합계(PACK)"] - grouped["포장 PACK"]).clip(lower=0.0)
    grouped["생산필요수량(PCS)"] = grouped["생산부족수량(PCS)"]
    grouped["생산진도율"] = calc_production_progress_pct(grouped["요청합계(PCS)"], grouped["생산부족수량(PCS)"])
    grouped["포장진도율"] = np.where(
        grouped["요청합계(PACK)"] > 0,
        grouped["포장 PACK"] / grouped["요청합계(PACK)"] * 100.0,
        0.0,
    )
    grouped["포장진도율"] = np.clip(grouped["포장진도율"], 0.0, 100.0)
    return grouped[
        [
            "POWER",
            "기간구분",
            "요청합계(PACK)",
            "요청합계(PCS)",
            "포장 PACK",
            "포장부족(PACK)",
            "생산필요수량(PCS)",
            "생산부족수량(PCS)",
            "생산진도율",
            "포장진도율",
            "power_value",
        ]
    ].sort_values("power_value", ascending=True, kind="stable")


def build_power_sku_detail_view(code_summary: pd.DataFrame, power_label: str) -> pd.DataFrame:
    work = add_allocated_production_basis(with_operational_columns(code_summary))
    scope = work[work["POWER"] == power_label].copy()
    if scope.empty:
        return pd.DataFrame(
            columns=[
                "생산코드",
                "판매코드",
                "기간구분",
                "제품명",
                "PACK",
                "요청합계(PACK)",
                "요청합계(PCS)",
                "포장부족(PACK)",
                "생산필요수량(PCS)",
                "생산부족수량(PCS)",
                "생산완료예상일",
            ]
        )
    out = (
        scope.groupby(["production_code_display", "sales_code", "period_group", "product_name", "_pack_label"], dropna=False)
        .agg(
            request_pack=("request_pack", "sum"),
            request_pcs=("request_pcs", "sum"),
            packing_pack=("packing_recognized_pack", "sum"),
            production_shortage_pcs=("production_shortage_qty", "sum"),
            expected_date=("production_due_date", max_datetime),
        )
        .reset_index()
        .rename(
            columns={
                "production_code_display": "생산코드",
                "sales_code": "판매코드",
                "period_group": "기간구분",
                "product_name": "제품명",
                "_pack_label": "PACK",
                "request_pack": "요청합계(PACK)",
                "request_pcs": "요청합계(PCS)",
                "production_shortage_pcs": "생산부족수량(PCS)",
            }
        )
    )
    out["포장부족(PACK)"] = (out["요청합계(PACK)"] - out["packing_pack"]).clip(lower=0.0)
    out["생산필요수량(PCS)"] = out["생산부족수량(PCS)"]
    out["생산완료예상일"] = out["expected_date"].map(display_date_or_dash)
    return out[
        [
            "생산코드",
            "판매코드",
            "기간구분",
            "제품명",
            "PACK",
            "요청합계(PACK)",
            "요청합계(PCS)",
            "포장부족(PACK)",
            "생산필요수량(PCS)",
            "생산부족수량(PCS)",
            "생산완료예상일",
        ]
    ].sort_values(
        ["포장부족(PACK)", "요청합계(PACK)"], ascending=[False, False], kind="stable"
    )


def empty_inventory_detail_view() -> pd.DataFrame:
    return pd.DataFrame(
        columns=["판매코드", "기간구분", "WMS제품명", "용마창고재고 (PACK)", "총수량(PACK)", "제품규격", "전송일자", "매칭여부"]
    )


def build_inventory_detail_view(code_summary: pd.DataFrame, sales_code: str) -> pd.DataFrame:
    if code_summary.empty:
        return empty_inventory_detail_view()
    work = with_operational_columns(code_summary)
    scope = work[work["sales_code"].astype(str) == str(sales_code)].copy()
    if scope.empty:
        return empty_inventory_detail_view()

    grouped = (
        scope.groupby("sales_code", dropna=False)
        .agg(
            factory_group=("factory_group", join_unique),
            period_group=("period_group", first_nonempty),
            inventory_product_name=("inventory_product_name", first_nonempty),
            available_stock_pack=("available_stock_pack", sum_numeric_or_nan),
            inventory_total_stock_pack=("inventory_total_stock_pack", sum_numeric_or_nan),
            inventory_product_spec=("inventory_product_spec", first_nonempty),
            inventory_updated_at=("inventory_updated_at", max_datetime),
            inventory_matched=("inventory_matched", "max"),
        )
        .reset_index()
        .rename(
            columns={
                "sales_code": "판매코드",
                "factory_group": "공장구분",
                "period_group": "기간구분",
                "inventory_product_name": "WMS제품명",
                "available_stock_pack": "용마창고재고 (PACK)",
                "inventory_total_stock_pack": "총수량(PACK)",
                "inventory_product_spec": "제품규격",
                "inventory_updated_at": "전송일자",
                "inventory_matched": "매칭여부",
            }
        )
    )
    grouped["전송일자"] = grouped["전송일자"].map(display_date_or_dash)
    grouped["매칭여부"] = np.where(grouped["매칭여부"], "매칭", "미매칭")
    return grouped[["판매코드", "기간구분", "WMS제품명", "용마창고재고 (PACK)", "총수량(PACK)", "제품규격", "전송일자", "매칭여부"]]


def build_inventory_prefix_detail_view(code_summary: pd.DataFrame, sales_code_prefix: str) -> pd.DataFrame:
    if code_summary.empty:
        return empty_inventory_detail_view()
    work = with_operational_columns(code_summary)
    scope = work[work["sales_code"].map(sales_code_base) == sales_code_prefix].copy()
    if scope.empty:
        return empty_inventory_detail_view()

    grouped = (
        scope.groupby("sales_code", dropna=False)
        .agg(
            factory_group=("factory_group", join_unique),
            period_group=("period_group", first_nonempty),
            inventory_product_name=("inventory_product_name", first_nonempty),
            available_stock_pack=("available_stock_pack", sum_numeric_or_nan),
            inventory_total_stock_pack=("inventory_total_stock_pack", sum_numeric_or_nan),
            inventory_product_spec=("inventory_product_spec", first_nonempty),
            inventory_updated_at=("inventory_updated_at", max_datetime),
            inventory_matched=("inventory_matched", "max"),
        )
        .reset_index()
        .rename(
            columns={
                "sales_code": "판매코드",
                "factory_group": "공장구분",
                "period_group": "기간구분",
                "inventory_product_name": "WMS제품명",
                "available_stock_pack": "용마창고재고 (PACK)",
                "inventory_total_stock_pack": "총수량(PACK)",
                "inventory_product_spec": "제품규격",
                "inventory_updated_at": "전송일자",
                "inventory_matched": "매칭여부",
            }
        )
    )
    grouped["전송일자"] = grouped["전송일자"].map(display_date_or_dash)
    grouped["매칭여부"] = np.where(grouped["매칭여부"], "매칭", "미매칭")
    grouped = grouped.sort_values(["판매코드"], kind="stable")
    return grouped[["판매코드", "기간구분", "WMS제품명", "용마창고재고 (PACK)", "총수량(PACK)", "제품규격", "전송일자", "매칭여부"]]


def ppt_rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.replace("#", ""))


def apply_ppt_font(
    run: Any,
    size: int | float,
    bold: bool = False,
    color: str = TEXT_DARK,
) -> None:
    run.font.name = PPT_FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = ppt_rgb(color)

    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        font_el = r_pr.find(qn(tag))
        if font_el is None:
            font_el = OxmlElement(tag)
            r_pr.append(font_el)
        font_el.set("typeface", PPT_FONT_NAME)


def set_cell_text(
    cell: Any,
    text: str,
    size: int = 9,
    bold: bool = False,
    color: str = TEXT_DARK,
    align: PP_ALIGN = PP_ALIGN.CENTER,
) -> None:
    cell.text = ""
    cell.margin_left = Inches(0.04)
    cell.margin_right = Inches(0.04)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.text_frame.word_wrap = True
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    paragraph.space_before = Pt(0)
    run = paragraph.add_run()
    run.text = text
    apply_ppt_font(run, size=size, bold=bold, color=color)


def add_textbox(
    slide: Any,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    size: int,
    bold: bool = False,
    color: str = TEXT_DARK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    vertical_anchor: Any | None = None,
) -> None:
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = True
    if vertical_anchor is not None:
        frame.vertical_anchor = vertical_anchor
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    paragraph.space_before = Pt(0)
    run = paragraph.add_run()
    run.text = text
    apply_ppt_font(run, size=size, bold=bold, color=color)


def truncate_report_text(value: Any, max_chars: int = 34) -> str:
    text = clean_str(value)
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "..."


def add_report_rule(
    slide: Any,
    left: float,
    top: float,
    width: float,
    color: str = MID_GRAY,
    vertical: bool = False,
) -> None:
    shape_width = 0.01 if vertical else width
    shape_height = width if vertical else 0.01
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(shape_width),
        Inches(shape_height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ppt_rgb(color)
    shape.line.fill.background()


def add_report_shape(
    slide: Any,
    shape_type: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: str,
    line_color: str | None = None,
    line_width: float = 0.5,
) -> Any:
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ppt_rgb(fill_color)
    if line_color:
        shape.line.color.rgb = ppt_rgb(line_color)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def report_progress_color(value: Any) -> str:
    num = pd.to_numeric(value, errors="coerce")
    if pd.isna(num):
        return REPORT_HEADER
    return REPORT_HEADER


def add_kpi_card(
    slide: Any,
    title: str,
    kpi: dict[str, float],
    dot_color: str,
    left: float,
    top: float,
    width: float,
    height: float,
    emphasis: bool = False,
) -> None:
    add_report_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
        REPORT_PANEL,
        REPORT_PANEL_LINE,
        0.5,
    )
    add_textbox(
        slide,
        title,
        left + 0.18,
        top + 0.16,
        width - 0.36,
        0.24,
        10 if emphasis else 9.2,
        True,
        dot_color if emphasis else REPORT_HEADER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    if emphasis:
        add_textbox(
            slide,
            "요청 대비 진행 현황",
            left + width - 1.35,
            top + 0.18,
            1.12,
            0.18,
            6.5,
            True,
            REPORT_MUTED,
            PP_ALIGN.RIGHT,
            MSO_ANCHOR.MIDDLE,
        )

    receipt_progress = kpi.get("progress_pct", 0.0)
    packing_progress = kpi.get("packing_progress_pct", 0.0)
    production_progress = kpi.get("production_progress_pct", 0.0)
    metric_font = 14.5 if emphasis else 10.5
    metric_label_size = 6.8 if emphasis else 6.2
    progress_color = "#64748B"

    if emphasis:
        metrics = [
            ("요청 PACK", format_report_value(kpi.get("request_pack", 0.0)), REPORT_HEADER),
            ("생산진도", format_report_value(production_progress, True), COLOR_BLUE),
            ("포장진도", format_report_value(packing_progress, True), COLOR_ORANGE),
            ("용마입고율", format_report_value(receipt_progress, True), COLOR_AMBER),
            ("미입고 PACK", format_report_value(kpi.get("shortage_pack", 0.0)), COLOR_DANGER),
        ]
        metric_top = top + 0.62
        col_width = (width - 0.42) / len(metrics)
        for idx, (label, value, value_color) in enumerate(metrics):
            metric_left = left + 0.2 + idx * col_width
            if idx:
                add_report_rule(slide, metric_left - 0.08, metric_top - 0.03, 0.68, "#D1D5DB", vertical=True)
            add_textbox(
                slide,
                label,
                metric_left,
                metric_top,
                col_width - 0.1,
                0.14,
                metric_label_size,
                True,
                REPORT_MUTED,
                PP_ALIGN.LEFT,
                MSO_ANCHOR.MIDDLE,
            )
            add_textbox(
                slide,
                value,
                metric_left,
                metric_top + 0.22,
                col_width - 0.1,
                0.27,
                metric_font,
                True,
                value_color,
                PP_ALIGN.LEFT,
                MSO_ANCHOR.MIDDLE,
            )
    else:
        metrics = [
            ("요청 PACK", format_report_value(kpi.get("request_pack", 0.0))),
            ("생산진도율", format_report_value(production_progress, True)),
            ("포장진도율", format_report_value(packing_progress, True)),
            ("용마입고율", format_report_value(receipt_progress, True)),
            ("미입고 PACK", format_report_value(kpi.get("shortage_pack", 0.0))),
            ("생산부족 PCS", format_report_value(kpi.get("production_shortage_pcs", 0.0))),
        ]
        metric_top = top + 0.48
        col_width = (width - 0.5) / 2
        row_gap = 0.43
        for idx, (label, value) in enumerate(metrics):
            metric_left = left + 0.2 + (idx % 2) * col_width
            metric_row_top = metric_top + (idx // 2) * row_gap
            if idx % 2 == 1:
                add_report_rule(slide, metric_left - 0.1, metric_row_top - 0.02, 0.34, REPORT_PANEL_LINE, vertical=True)
            add_textbox(
                slide,
                label,
                metric_left,
                metric_row_top,
                col_width - 0.12,
                0.13,
                metric_label_size,
                True,
                REPORT_MUTED,
                PP_ALIGN.LEFT,
                MSO_ANCHOR.MIDDLE,
            )
            add_textbox(
                slide,
                value,
                metric_left,
                metric_row_top + 0.18,
                col_width - 0.12,
                0.18,
                metric_font,
                True,
                REPORT_HEADER,
                PP_ALIGN.LEFT,
                MSO_ANCHOR.MIDDLE,
            )

    def add_card_progress_row(row_top: float, label: str, value: Any, color: str) -> None:
        pct = max(0.0, min(100.0, to_report_float(value)))
        label_width = 0.5
        pct_width = 0.44
        bar_left = left + 0.2 + label_width
        bar_width = max(0.0, width - label_width - pct_width - 0.58)
        add_textbox(
            slide,
            label,
            left + 0.18,
            row_top - 0.01,
            label_width,
            0.18,
            6.7,
            True,
            REPORT_HEADER,
            PP_ALIGN.LEFT,
            MSO_ANCHOR.MIDDLE,
        )
        add_report_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, row_top + 0.07, bar_width, 0.04, REPORT_FAINT)
        fill_width = bar_width * pct / 100.0
        if fill_width > 0:
            add_report_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, row_top + 0.07, fill_width, 0.04, color, color)
        add_textbox(
            slide,
            f"{pct:.1f}%",
            left + width - pct_width - 0.2,
            row_top - 0.01,
            pct_width,
            0.18,
            7 if emphasis else 6.4,
            True,
            REPORT_HEADER if emphasis else REPORT_MUTED,
            PP_ALIGN.RIGHT,
            MSO_ANCHOR.MIDDLE,
        )

    progress_top = top + (1.34 if emphasis else 1.66)
    add_card_progress_row(progress_top, "생산", production_progress, COLOR_BLUE if emphasis else progress_color)
    add_card_progress_row(progress_top + 0.24, "포장", packing_progress, COLOR_ORANGE if emphasis else progress_color)
    add_card_progress_row(progress_top + 0.48, "용마입고", receipt_progress, COLOR_AMBER if emphasis else progress_color)

    if emphasis:
        packing_pack = to_report_float(kpi.get("packing_pack", 0.0))
        yongma_in_pack = to_report_float(kpi.get("yongma_in_pack", 0.0))
        shortage_pack = to_report_float(kpi.get("shortage_pack", 0.0))
        production_shortage = to_report_float(kpi.get("production_shortage_pcs", 0.0))
        details = (
            f"용마입고 {format_report_value(yongma_in_pack)} PACK     "
            f"포장대기 {format_report_value(max(0.0, packing_pack - yongma_in_pack))} PACK     "
            f"포장필요 {format_report_value(shortage_pack)} PACK     "
            f"생산부족 {format_report_value(production_shortage)} PCS"
        )
        add_textbox(
            slide,
            details,
            left + 0.2,
            top + height - 0.28,
            width - 0.4,
            0.18,
            7.1,
            True,
            REPORT_MUTED,
            PP_ALIGN.LEFT,
            MSO_ANCHOR.MIDDLE,
        )


def format_report_value(value: Any, is_percent: bool = False) -> str:
    num = pd.to_numeric(value, errors="coerce")
    if pd.isna(num):
        return "0.0%" if is_percent else "0"
    return f"{float(num):.1f}%" if is_percent else format_int(float(num))


def sanitize_excel_sheet_name(name: str) -> str:
    cleaned = re.sub(r"[\[\]\:\*\?\/\\]", "_", clean_str(name))
    return (cleaned or "Sheet")[:31]


def dataframe_for_excel(df: pd.DataFrame) -> pd.DataFrame:
    if df is None or df.empty:
        return pd.DataFrame()
    out = df.copy()
    drop_cols = [
        col
        for col in out.columns
        if str(col).startswith("_") or str(col) in {"power_value"}
    ]
    if drop_cols:
        out = out.drop(columns=drop_cols, errors="ignore")
    for col in out.columns:
        if pd.api.types.is_datetime64tz_dtype(out[col]):
            out[col] = out[col].dt.tz_localize(None)
    return out


def apply_code_display_terms(df: pd.DataFrame) -> pd.DataFrame:
    if df is None or df.empty:
        return pd.DataFrame() if df is None else df.copy()

    out = df.copy()
    rename_map: dict[str, str] = {}
    if "품목코드" in out.columns and "판매코드" not in out.columns:
        rename_map["품목코드"] = "판매코드"
    if "S코드" in out.columns and "판매코드" not in out.columns:
        rename_map["S코드"] = "판매코드"
    if "제품코드" in out.columns:
        rename_map["제품코드"] = "판매코드(기본)" if "판매코드" in out.columns or "품목코드" in out.columns or "S코드" in out.columns else "판매코드"
    if rename_map:
        out = out.rename(columns=rename_map)
    return out


def excel_text_length(value: Any) -> int:
    if value is None:
        return 0
    try:
        if pd.isna(value):
            return 0
    except (TypeError, ValueError):
        pass
    return len(str(value))


def make_unique_excel_columns(columns: pd.Index) -> list[str]:
    used_counts: dict[str, int] = {}
    unique_columns: list[str] = []
    for idx, column in enumerate(columns, start=1):
        base_name = clean_str(column) or f"컬럼{idx}"
        current_count = used_counts.get(base_name, 0)
        used_counts[base_name] = current_count + 1
        if current_count:
            unique_columns.append(f"{base_name}_{current_count + 1}")
        else:
            unique_columns.append(base_name)
    return unique_columns


@st.cache_data(show_spinner=False, max_entries=24)
def build_excel_download_bytes(sheets: dict[str, pd.DataFrame]) -> bytes:
    output = BytesIO()
    used_names: set[str] = set()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        for raw_name, df in sheets.items():
            base_name = sanitize_excel_sheet_name(raw_name)
            sheet_name = base_name
            suffix = 1
            while sheet_name in used_names:
                suffix += 1
                sheet_name = f"{base_name[:28]}_{suffix}"
            used_names.add(sheet_name)

            excel_df = apply_code_display_terms(dataframe_for_excel(df))
            if excel_df.empty:
                excel_df = pd.DataFrame({"내용": ["조건에 맞는 데이터가 없습니다."]})
            excel_df = excel_df.copy()
            excel_df.columns = make_unique_excel_columns(excel_df.columns)
            excel_df.to_excel(writer, sheet_name=sheet_name, index=False)

            worksheet = writer.sheets[sheet_name]
            for col_idx, col_name in enumerate(excel_df.columns, start=1):
                column_values = excel_df.iloc[:, col_idx - 1].head(300).tolist()
                value_lengths = [excel_text_length(value) for value in column_values]
                max_len = max([len(str(col_name)), *(value_lengths if value_lengths else [0])])
                worksheet.column_dimensions[worksheet.cell(row=1, column=col_idx).column_letter].width = min(
                    max(max_len + 2, 10),
                    45,
                )
    return output.getvalue()


def dataframe_light_signature(df: pd.DataFrame) -> tuple[Any, ...]:
    if df is None or df.empty:
        return ("empty",)

    columns = tuple(str(col) for col in df.columns)
    sample = pd.concat([df.head(5), df.tail(5)], ignore_index=False)
    sample_hash = tuple(int(value) for value in pd.util.hash_pandas_object(sample.astype(str), index=True).to_numpy())
    numeric_sums: list[tuple[str, float]] = []
    for col in df.columns[:20]:
        series = pd.to_numeric(df[col], errors="coerce")
        if series.notna().any():
            numeric_sums.append((str(col), round(float(series.fillna(0.0).sum()), 6)))
    return (tuple(df.shape), columns, tuple(numeric_sums), sample_hash)


def sheet_collection_signature(sheets: dict[str, pd.DataFrame]) -> tuple[Any, ...]:
    return tuple((clean_str(name), dataframe_light_signature(df)) for name, df in sheets.items())


def render_lazy_binary_download(
    label: str,
    prepare_label: str,
    file_name: str,
    mime: str,
    build_bytes: Callable[[], bytes],
    signature: tuple[Any, ...],
    key: str,
    width: str = "stretch",
) -> None:
    bytes_key = f"{key}_bytes"
    signature_key = f"{key}_signature"
    if st.session_state.get(signature_key) != signature:
        st.session_state.pop(bytes_key, None)
        st.session_state[signature_key] = signature

    if st.button(prepare_label, key=f"{key}_prepare", width=width):
        with st.spinner(f"{label} 생성 중..."):
            st.session_state[bytes_key] = build_bytes()
            st.session_state[signature_key] = signature

    if bytes_key in st.session_state:
        st.download_button(
            label,
            data=st.session_state[bytes_key],
            file_name=file_name,
            mime=mime,
            width=width,
            key=f"{key}_download",
            on_click="ignore",
        )


def render_excel_download(
    label: str,
    file_prefix: str,
    sheets: dict[str, pd.DataFrame],
    key: str,
    width: str = "stretch",
) -> None:
    timestamp = pd.Timestamp.now(tz="Asia/Seoul").strftime("%Y%m%d_%H%M")
    render_lazy_binary_download(
        label,
        f"{label} 준비",
        f"{file_prefix}_{timestamp}.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        build_bytes=lambda: build_excel_download_bytes(sheets),
        signature=sheet_collection_signature(sheets),
        key=key,
        width=width,
    )


def build_priority_report_view(product_view: pd.DataFrame, max_rows: int = 6) -> pd.DataFrame:
    columns = ["제품명", "요청 PACK", "생산진도율", "용마입고율", "생산부족수량", "미입고수량", "상태"]
    if product_view.empty:
        return pd.DataFrame(columns=columns)

    view = product_view.copy()
    if "용마입고율" not in view.columns:
        view["용마입고율"] = 0.0
    if "미입고수량" not in view.columns:
        view["미입고수량"] = view.get("포장부족수량", 0.0)
    for col in columns:
        if col not in view.columns:
            view[col] = ""
    return view.sort_values(
        ["미입고수량", "생산부족수량", "요청 PACK"],
        ascending=[False, False, False],
        kind="stable",
    ).head(max_rows)[columns].copy()


def build_daily_exception_report_view(
    daily_inventory_df: pd.DataFrame | None,
    code_summary: pd.DataFrame,
    sample_available_df: pd.DataFrame | None = None,
    max_rows: int = 5,
    response_view: pd.DataFrame | None = None,
) -> tuple[dict[str, float], pd.DataFrame]:
    columns = ["품목코드", "제품명", "현재 재고수량", "부족수량", "포장가능재고(PCS)", "대응가능 여부"]
    empty_kpis = {"request_out_count": 0.0, "negative_count": 0.0, "waiting_pcs": 0.0}
    if response_view is None and (daily_inventory_df is None or daily_inventory_df.empty):
        return empty_kpis, pd.DataFrame(columns=columns)

    view = (
        response_view
        if response_view is not None
        else build_daily_inventory_response_view(daily_inventory_df, code_summary, sample_available_df)
    )
    if view.empty or "대응상태" not in view.columns:
        return empty_kpis, pd.DataFrame(columns=columns)

    out = view[view["대응상태"] == "요청외 긴급"].copy()
    if out.empty:
        return empty_kpis, pd.DataFrame(columns=columns)

    out["재고수량"] = pd.to_numeric(out["재고수량"], errors="coerce")
    out["재고부족수량"] = pd.to_numeric(out.get("재고부족수량", 0.0), errors="coerce").fillna(0.0)
    pack_units = out["PACK"].map(pack_unit_from_label)
    out["부족수량"] = (out["재고부족수량"] * pack_units).clip(lower=0.0).round(0)
    out["포장가능재고(PCS)"] = pd.to_numeric(out["포장가능재고(PCS)"], errors="coerce").fillna(0.0)
    out["대응가능 여부"] = [
        classify_exception_response(available, shortage)
        for available, shortage in zip(out["포장가능재고(PCS)"], out["부족수량"])
    ]
    out["현재 재고수량"] = out["재고수량"]
    out["_stock_missing_sort"] = out["재고수량"].isna().astype(int)
    out["_stock_sort"] = out["재고수량"].fillna(0.0)
    out["_response_sort"] = out["대응가능 여부"].map({"대응 필요": 0, "일부 가능": 1, "충당 가능": 2}).fillna(3)
    kpis = {
        "request_out_count": float(len(out)),
        "negative_count": float((out["재고수량"] < 0).sum()),
        "waiting_pcs": float(out["포장가능재고(PCS)"].sum()),
    }
    detail = out.sort_values(
        ["_response_sort", "_stock_missing_sort", "_stock_sort", "포장가능재고(PCS)", "품목코드"],
        ascending=[True, True, True, False, True],
        kind="stable",
    ).head(max_rows)
    for col in columns:
        if col not in detail.columns:
            detail[col] = ""
    return kpis, detail[columns].copy()


def build_urgent_request_summary_view(
    daily_inventory_df: pd.DataFrame | None,
    code_summary: pd.DataFrame,
    sample_available_df: pd.DataFrame | None = None,
    response_view: pd.DataFrame | None = None,
) -> pd.DataFrame:
    columns = ["판매코드", "요청구분", "제품명", "SKU 수"]
    if response_view is None and (daily_inventory_df is None or daily_inventory_df.empty):
        return pd.DataFrame(columns=columns)

    view = (
        response_view
        if response_view is not None
        else build_daily_inventory_response_view(daily_inventory_df, code_summary, sample_available_df)
    )
    if view.empty or "긴급요청" not in view.columns:
        return pd.DataFrame(columns=columns)

    urgent = view[view["긴급요청"].fillna(False).astype(bool)].copy()
    if urgent.empty:
        return pd.DataFrame(columns=columns)

    urgent["판매코드"] = [
        extract_sales_prefix(product_code) or extract_sales_prefix(item_code)
        for product_code, item_code in zip(
            urgent.get("제품코드", pd.Series("", index=urgent.index)),
            urgent.get("품목코드", pd.Series("", index=urgent.index)),
        )
    ]
    urgent = urgent[urgent["판매코드"].map(lambda value: bool(re.fullmatch(r"S\d{3}", clean_str(value))))].copy()
    if urgent.empty:
        return pd.DataFrame(columns=columns)

    urgent["_sku_key"] = [
        clean_str(item_code)
        or "|".join([clean_str(product_code), clean_str(pack), clean_str(power)])
        for item_code, product_code, pack, power in zip(
            urgent.get("품목코드", pd.Series("", index=urgent.index)),
            urgent["판매코드"],
            urgent.get("PACK", pd.Series("", index=urgent.index)),
            urgent.get("POWER", pd.Series("", index=urgent.index)),
        )
    ]
    urgent["_request_pack"] = pd.to_numeric(
        urgent.get("요청 PACK", pd.Series(0.0, index=urgent.index)),
        errors="coerce",
    ).fillna(0.0)
    urgent["_request_scope"] = np.where(urgent["_request_pack"] > 0, "요청내", "요청외")
    urgent["_request_in_sku_key"] = urgent["_sku_key"].where(urgent["_request_scope"] == "요청내", "")
    urgent["_request_out_sku_key"] = urgent["_sku_key"].where(urgent["_request_scope"] == "요청외", "")
    grouped = (
        urgent.groupby("판매코드", dropna=False)
        .agg(
            제품명=("제품명", join_unique),
            request_in_count=("_request_in_sku_key", lambda series: len({clean_str(value) for value in series if clean_str(value)})),
            request_out_count=("_request_out_sku_key", lambda series: len({clean_str(value) for value in series if clean_str(value)})),
            sku_count=("_sku_key", "nunique"),
        )
        .reset_index()
    )
    grouped = grouped.rename(
        columns={
            "request_in_count": "요청내 SKU",
            "request_out_count": "요청외 SKU",
            "sku_count": "SKU 수",
        }
    )
    for col in ["요청내 SKU", "요청외 SKU"]:
        grouped[col] = pd.to_numeric(grouped[col], errors="coerce").fillna(0).astype(int)
    grouped["요청구분"] = np.where(grouped["요청외 SKU"] > 0, "요청외", "요청내")
    grouped["SKU 수"] = pd.to_numeric(grouped["SKU 수"], errors="coerce").fillna(0).astype(int)
    return grouped.sort_values(["SKU 수", "판매코드"], ascending=[False, True], kind="stable")[columns].copy()


def classify_exception_response(available_pcs: Any, shortage_pcs: Any) -> str:
    available_num = pd.to_numeric(available_pcs, errors="coerce")
    shortage_num = pd.to_numeric(shortage_pcs, errors="coerce")
    available = 0.0 if pd.isna(available_num) else float(available_num)
    shortage = 0.0 if pd.isna(shortage_num) else float(shortage_num)
    if available <= 0:
        return "대응 필요"
    if available >= shortage:
        return "충당 가능"
    return "일부 가능"


def render_exception_summary_table(exception_detail: pd.DataFrame) -> None:
    if exception_detail.empty:
        st.warning("요청외 긴급 품목이 없습니다.")
        return

    headers = ["판매코드", "제품명", "현재 재고수량", "부족수량", "포장가능재고(PCS)", "대응가능 여부"]
    rows: list[str] = []
    for _, row in exception_detail.iterrows():
        stock = pd.to_numeric(row.get("현재 재고수량", np.nan), errors="coerce")
        shortage = pd.to_numeric(row.get("부족수량", 0.0), errors="coerce")
        available = pd.to_numeric(row.get("포장가능재고(PCS)", 0.0), errors="coerce")
        response = clean_str(row.get("대응가능 여부", ""))
        response_class = {
            "충당 가능": "ok",
            "일부 가능": "partial",
            "대응 필요": "need",
        }.get(response, "need")
        stock_text = "-" if pd.isna(stock) else format_int(float(stock))
        stock_class = "num negative" if not pd.isna(stock) and float(stock) < 0 else "num"
        rows.append(
            "<tr>"
            f"<td>{escape(str(row.get('품목코드', '')))}</td>"
            f"<td class='left'>{escape(str(row.get('제품명', '')))}</td>"
            f"<td class='{stock_class}'>{stock_text}</td>"
            f"<td class='num shortage'>{format_int(float(shortage) if not pd.isna(shortage) else 0.0)}</td>"
            f"<td class='num'>{format_int(float(available) if not pd.isna(available) else 0.0)}</td>"
            f"<td><span class='response-badge {response_class}'>{escape(response)}</span></td>"
            "</tr>"
        )
    header_html = "".join(f"<th class='{'left' if header == '제품명' else ''}'>{escape(header)}</th>" for header in headers)
    st.markdown(
        "<div class='table-wrap compact-table'>"
        "<table class='ops-table urgent-summary-table'>"
        f"<thead><tr>{header_html}</tr></thead>"
        f"<tbody>{''.join(rows)}</tbody>"
        "</table>"
        "</div>",
        unsafe_allow_html=True,
    )


def render_urgent_request_summary_table(summary_view: pd.DataFrame) -> None:
    if summary_view.empty:
        st.warning("긴급요청 품목이 없습니다.")
        return

    rows: list[str] = []
    for _, row in summary_view.iterrows():
        sku_num = pd.to_numeric(row.get("SKU 수", 0), errors="coerce")
        sku_count = 0.0 if pd.isna(sku_num) else float(sku_num)
        scope = clean_str(row.get("요청구분", ""))
        scope_class = "in" if scope == "요청내" else "out"
        rows.append(
            "<tr>"
            f"<td class='left code-cell'>{escape(clean_str(row.get('판매코드', '')))}</td>"
            f"<td class='left'><span class='request-scope-badge {scope_class}'>{escape(scope)}</span></td>"
            f"<td class='left'>{escape(clean_str(row.get('제품명', '')))}</td>"
            f"<td class='num shortage'>{format_int(sku_count)}</td>"
            "</tr>"
        )
    st.markdown(
        "<div class='table-wrap compact-table'>"
        "<table class='ops-table urgent-summary-table main-summary-table urgent-request-summary-table'>"
        "<colgroup>"
        "<col class='summary-code-col'>"
        "<col class='summary-scope-col'>"
        "<col class='summary-product-col'>"
        "<col class='summary-number-col'>"
        "</colgroup>"
        "<thead><tr>"
        "<th class='left'>판매코드</th>"
        "<th class='left'>요청구분</th>"
        "<th class='left'>제품명</th>"
        "<th class='num'>SKU 수</th>"
        "</tr></thead>"
        f"<tbody>{''.join(rows)}</tbody>"
        "</table>"
        "</div>",
        unsafe_allow_html=True,
    )


def render_urgent_request_compact(summary_view: pd.DataFrame) -> None:
    if summary_view.empty:
        st.warning("긴급요청 품목이 없습니다.")
        return
    rows: list[str] = []
    for _, row in summary_view.head(8).iterrows():
        sku_num = pd.to_numeric(row.get("SKU 수", 0), errors="coerce")
        sku_count = 0.0 if pd.isna(sku_num) else float(sku_num)
        scope = clean_str(row.get("요청구분", ""))
        scope_class = "in" if scope == "요청내" else "out"
        rows.append(
            "<div class='urgent-list-row'>"
            f"<span class='urgent-code'>{escape(clean_str(row.get('판매코드', '')))}</span>"
            f"<span class='request-scope-badge {scope_class}'>{escape(scope)}</span>"
            f"<span class='urgent-product'>{escape(clean_str(row.get('제품명', '')))}</span>"
            f"<span class='urgent-sku'>SKU {format_int(sku_count)}</span>"
            "</div>"
        )
    st.markdown(
        f"<div class='panel-box dashboard-card'><div class='urgent-list'>{''.join(rows)}</div></div>",
        unsafe_allow_html=True,
    )


def to_report_float(value: Any) -> float:
    num = pd.to_numeric(value, errors="coerce")
    if pd.isna(num):
        return 0.0
    return float(num)


def add_report_status_badge(slide: Any, status: str, left: float, top: float, width: float, height: float) -> None:
    status_text = clean_str(status) or "-"
    if status_text in {"완료", "입고완료"}:
        fill_color = "#E8F5F0"
        line_color = "#9ED8C5"
        text_color = COLOR_TEAL
    elif status_text in {"진행중"}:
        fill_color = "#FFF4DE"
        line_color = "#E4B968"
        text_color = COLOR_AMBER
    else:
        fill_color = COLOR_ALERT_BG
        line_color = COLOR_ALERT_BD
        text_color = COLOR_ORANGE

    add_report_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
        fill_color,
        line_color,
        0.5,
    )
    add_textbox(
        slide,
        status_text,
        left,
        top,
        width,
        height,
        7.5,
        True,
        text_color,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )


def add_priority_report_table(
    slide: Any,
    priority_view: pd.DataFrame,
    left: float = 0.45,
    top: float = 3.62,
    width: float = 8.35,
    height: float = 3.28,
) -> None:
    add_report_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
        REPORT_PANEL,
        REPORT_PANEL_LINE,
        0.5,
    )
    add_report_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, 0.38, REPORT_NAVY, REPORT_NAVY, 0.5)

    headers = ["순위", "제품명", "요청 PACK", "용마입고율", "미입고 PACK", "생산진도율"]
    col_widths = [0.58, 3.18, 1.08, 1.14, 1.12, 1.08]
    col_lefts = [left + 0.16]
    for width in col_widths[:-1]:
        col_lefts.append(col_lefts[-1] + width)

    for idx, header in enumerate(headers):
        add_textbox(
            slide,
            header,
            col_lefts[idx],
            top + 0.04,
            col_widths[idx],
            0.3,
            8.2,
            True,
            "#FFFFFF",
            PP_ALIGN.LEFT if idx == 1 else PP_ALIGN.CENTER if idx == 0 else PP_ALIGN.RIGHT,
            MSO_ANCHOR.MIDDLE,
        )

    if priority_view.empty:
        add_textbox(
            slide,
            "조건에 맞는 제품 데이터가 없습니다.",
            left + 0.2,
            top + 0.55,
            width - 0.4,
            0.35,
            8.4,
            False,
            REPORT_MUTED,
            PP_ALIGN.LEFT,
            MSO_ANCHOR.MIDDLE,
        )
        return

    row_height = 0.48
    for row_idx, (_, row) in enumerate(priority_view.iterrows(), start=1):
        row_top = top + 0.38 + (row_idx - 1) * row_height
        production_progress = to_report_float(row["생산진도율"])
        receipt_progress = to_report_float(row["용마입고율"])
        receipt_shortage = to_report_float(row["미입고수량"])

        if row_idx % 2 == 0:
            add_report_shape(slide, MSO_SHAPE.RECTANGLE, left + 0.06, row_top, width - 0.12, row_height, REPORT_ROW_ALT)

        cell_top = row_top + 0.01
        cell_height = row_height - 0.02
        values = [
            str(row_idx),
            truncate_report_text(row["제품명"], max_chars=32),
            format_report_value(row["요청 PACK"]),
            format_report_value(receipt_progress, True),
            format_report_value(receipt_shortage),
            format_report_value(production_progress, True),
        ]
        colors = [
            REPORT_MUTED,
            REPORT_HEADER,
            REPORT_HEADER,
            REPORT_HEADER,
            REPORT_MUTED if receipt_shortage <= 0 else REPORT_ACCENT,
            REPORT_HEADER,
        ]
        bolds = [False, False, False, False, receipt_shortage > 0, False]
        aligns = [PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT]

        for col_idx, value in enumerate(values):
            add_textbox(
                slide,
                value,
                col_lefts[col_idx],
                cell_top,
                col_widths[col_idx],
                cell_height,
                8.6 if col_idx != 1 else 8.8,
                bolds[col_idx],
                colors[col_idx],
                aligns[col_idx],
                MSO_ANCHOR.MIDDLE,
            )

        add_report_rule(slide, left + 0.1, row_top + row_height, width - 0.2, REPORT_FAINT)


def add_daily_exception_report_panel(
    slide: Any,
    exception_kpis: dict[str, float],
    exception_view: pd.DataFrame,
    left: float = 8.95,
    top: float = 3.62,
    width: float = 3.95,
    height: float = 3.28,
) -> None:
    add_report_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
        REPORT_PANEL,
        REPORT_PANEL_LINE,
        0.5,
    )
    add_report_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, 0.38, REPORT_NAVY, REPORT_NAVY, 0.5)

    headers = ["판매코드", "제품명", "재고", "가용 PCS"]
    col_widths = [0.9, 1.44, 0.56, 0.78]
    col_lefts = [left + 0.12]
    for col_width in col_widths[:-1]:
        col_lefts.append(col_lefts[-1] + col_width)
    for idx, header in enumerate(headers):
        add_textbox(
            slide,
            header,
            col_lefts[idx],
            top + 0.04,
            col_widths[idx],
            0.3,
            7.8,
            True,
            "#FFFFFF",
            PP_ALIGN.LEFT if idx in {0, 1} else PP_ALIGN.RIGHT,
            MSO_ANCHOR.MIDDLE,
        )

    if exception_view.empty:
        add_textbox(
            slide,
            "요청물량 외 긴급 대응 품목이 없습니다.",
            left + 0.18,
            top + 0.72,
            width - 0.36,
            0.28,
            8.2,
            False,
            REPORT_MUTED,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )
        return

    row_height = 0.48
    for row_idx, (_, row) in enumerate(exception_view.iterrows(), start=1):
        row_top = top + 0.38 + (row_idx - 1) * row_height
        stock = pd.to_numeric(row.get("현재 재고수량", row.get("재고수량", np.nan)), errors="coerce")
        waiting_pcs = pd.to_numeric(row.get("포장가능재고(PCS)", np.nan), errors="coerce")
        stock_text = "-" if pd.isna(stock) else format_report_value(stock)
        waiting_text = "-" if pd.isna(waiting_pcs) else format_report_value(waiting_pcs)
        stock_color = REPORT_ACCENT if pd.notna(stock) and float(stock) < 0 else REPORT_HEADER
        waiting_color = REPORT_HEADER
        if row_idx % 2 == 0:
            add_report_shape(slide, MSO_SHAPE.RECTANGLE, left + 0.03, row_top, width - 0.06, row_height, REPORT_ROW_ALT)
        values = [
            truncate_report_text(row.get("품목코드", ""), 12),
            truncate_report_text(row.get("제품명", ""), 20),
            stock_text,
            waiting_text,
        ]
        colors = [REPORT_HEADER, REPORT_HEADER, stock_color, waiting_color]
        aligns = [PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT]
        for col_idx, value in enumerate(values):
            add_textbox(
                slide,
                value,
                col_lefts[col_idx],
                row_top + 0.03,
                col_widths[col_idx],
                row_height - 0.06,
                7.7 if col_idx != 1 else 7.9,
                col_idx == 2 and stock_color == REPORT_ACCENT,
                colors[col_idx],
                aligns[col_idx],
                MSO_ANCHOR.MIDDLE,
            )
        add_report_rule(slide, left + 0.08, row_top + row_height, width - 0.16, REPORT_PANEL_LINE)


def add_urgent_request_summary_panel(
    slide: Any,
    urgent_summary_view: pd.DataFrame,
    left: float = 8.95,
    top: float = 3.56,
    width: float = 3.95,
    height: float = 3.47,
) -> None:
    add_report_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
        REPORT_PANEL,
        REPORT_PANEL_LINE,
        0.5,
    )
    add_report_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, 0.36, REPORT_NAVY, REPORT_NAVY, 0.5)

    headers = ["판매코드", "구분", "제품명", "SKU"]
    content_width = max(0.0, width - 0.24)
    if width >= 8.0:
        col_widths = [1.0, 1.1, max(1.0, content_width - 2.85), 0.75]
        product_max_chars = 68
        body_font_size = 8.0
        product_font_size = 7.8
    else:
        col_widths = [0.58, 0.58, 2.12, 0.38]
        product_max_chars = 22
        body_font_size = 6.5
        product_font_size = 6.1
    col_lefts = [left + 0.12]
    for col_width in col_widths[:-1]:
        col_lefts.append(col_lefts[-1] + col_width)

    for idx, header in enumerate(headers):
        add_textbox(
            slide,
            header,
            col_lefts[idx],
            top + 0.04,
            col_widths[idx],
            0.27,
            7.2,
            True,
            "#FFFFFF",
            PP_ALIGN.RIGHT if idx == 3 else PP_ALIGN.LEFT,
            MSO_ANCHOR.MIDDLE,
        )

    if urgent_summary_view.empty:
        add_textbox(
            slide,
            "긴급요청 품목이 없습니다.",
            left + 0.18,
            top + 0.7,
            width - 0.36,
            0.3,
            8.0,
            False,
            REPORT_MUTED,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )
        return

    max_rows = 11
    rows = urgent_summary_view.head(max_rows).copy()
    hidden_count = max(0, len(urgent_summary_view) - len(rows))
    row_height = min(0.27, (height - 0.58) / max(len(rows), 1))
    for row_idx, (_, row) in enumerate(rows.iterrows(), start=1):
        row_top = top + 0.36 + (row_idx - 1) * row_height
        if row_idx % 2 == 0:
            add_report_shape(slide, MSO_SHAPE.RECTANGLE, left + 0.04, row_top, width - 0.08, row_height, REPORT_ROW_ALT)

        scope = clean_str(row.get("요청구분", ""))
        scope_color = REPORT_ACCENT if scope == "요청외" else COLOR_TEAL
        values = [
            clean_str(row.get("판매코드", "")),
            scope,
            truncate_report_text(row.get("제품명", ""), product_max_chars),
            format_report_value(row.get("SKU 수", 0.0)),
        ]
        colors = [COLOR_BLUE, scope_color, REPORT_HEADER, REPORT_ACCENT]
        bolds = [True, True, False, True]
        aligns = [PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.RIGHT]
        for col_idx, value in enumerate(values):
            add_textbox(
                slide,
                value,
                col_lefts[col_idx],
                row_top + 0.01,
                col_widths[col_idx],
                row_height - 0.02,
                body_font_size if col_idx != 2 else product_font_size,
                bolds[col_idx],
                colors[col_idx],
                aligns[col_idx],
                MSO_ANCHOR.MIDDLE,
            )
        add_report_rule(slide, left + 0.08, row_top + row_height, width - 0.16, REPORT_FAINT)

    note = "요청외 SKU 포함 시 판매코드 전체 요청외"
    if hidden_count:
        note = f"{note} / 외 {hidden_count:,}개"
    add_textbox(
        slide,
        note,
        left + 0.12,
        top + height - 0.19,
        width - 0.24,
        0.16,
        5.9,
        False,
        REPORT_MUTED,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )


def add_report_progress_bar(
    slide: Any,
    value: Any,
    left: float,
    top: float,
    width: float,
    color: str,
    text_color: str = REPORT_HEADER,
) -> None:
    pct = max(0.0, min(100.0, to_report_float(value)))
    bar_width = max(0.0, width - 0.55)
    add_report_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top + 0.09, bar_width, 0.04, REPORT_FAINT)
    fill_width = bar_width * pct / 100.0
    if fill_width > 0:
        add_report_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top + 0.09, fill_width, 0.04, color, color)
    add_textbox(
        slide,
        f"{pct:.1f}%",
        left + bar_width + 0.07,
        top,
        0.48,
        0.22,
        7.1,
        True,
        text_color,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )


def add_family_progress_summary_panel(
    slide: Any,
    family_view: pd.DataFrame,
    left: float = 0.45,
    top: float = 3.56,
    width: float = 12.45,
    height: float = 3.47,
    max_rows: int = 10,
) -> None:
    add_report_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
        REPORT_PANEL,
        REPORT_PANEL_LINE,
        0.5,
    )
    add_report_shape(slide, MSO_SHAPE.RECTANGLE, left + 0.02, top + 0.02, width - 0.04, 0.38, REPORT_TABLE_HEADER, REPORT_TABLE_HEADER, 0.5)

    headers = ["제품 분류", "요청 PACK", "생산진도", "포장진도", "용마입고", "생산부족 PCS"]
    col_widths = [2.05, 1.25, 2.25, 2.25, 2.25, 1.65]
    col_lefts = [left + 0.16]
    for col_width in col_widths[:-1]:
        col_lefts.append(col_lefts[-1] + col_width)

    for idx, header in enumerate(headers):
        add_textbox(
            slide,
            header,
            col_lefts[idx],
            top + 0.06,
            col_widths[idx],
            0.26,
            8.0,
            True,
            REPORT_HEADER,
            PP_ALIGN.RIGHT if idx in {1, 5} else PP_ALIGN.LEFT,
            MSO_ANCHOR.MIDDLE,
        )
    add_report_rule(slide, left + 0.1, top + 0.4, width - 0.2, REPORT_PANEL_LINE)

    if family_view.empty:
        add_textbox(
            slide,
            "표시할 제품 분류별 진도 데이터가 없습니다.",
            left + 0.2,
            top + 0.76,
            width - 0.4,
            0.3,
            8.4,
            False,
            REPORT_MUTED,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )
        return

    rows = family_view.head(max_rows).copy()
    hidden_count = max(0, len(family_view) - len(rows))
    row_height = min(0.34, (height - 0.7) / max(len(rows), 1))
    for row_idx, (_, row) in enumerate(rows.iterrows(), start=1):
        row_top = top + 0.43 + (row_idx - 1) * row_height

        family = truncate_report_text(row.get("본품분류", ""), 24)
        request_pack = to_report_float(row.get("요청 PACK", 0.0))
        production_progress = to_report_float(row.get("생산진도율", 0.0))
        packing_progress = to_report_float(row.get("포장진도율", 0.0))
        receipt_progress = to_report_float(row.get("용마입고율", 0.0))
        production_shortage = to_report_float(row.get("생산부족수량", 0.0))
        shortage_color = COLOR_DANGER if production_shortage > 0 else REPORT_HEADER

        add_textbox(
            slide,
            family,
            col_lefts[0],
            row_top + 0.03,
            col_widths[0],
            row_height - 0.06,
            8.1,
            True,
            REPORT_HEADER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide,
            format_report_value(request_pack),
            col_lefts[1],
            row_top + 0.03,
            col_widths[1],
            row_height - 0.06,
            8.0,
            True,
            REPORT_HEADER,
            PP_ALIGN.RIGHT,
            MSO_ANCHOR.MIDDLE,
        )
        add_report_progress_bar(slide, production_progress, col_lefts[2], row_top + 0.03, col_widths[2] - 0.12, COLOR_BLUE)
        add_report_progress_bar(slide, packing_progress, col_lefts[3], row_top + 0.03, col_widths[3] - 0.12, COLOR_ORANGE)
        add_report_progress_bar(slide, receipt_progress, col_lefts[4], row_top + 0.03, col_widths[4] - 0.12, COLOR_AMBER)
        add_textbox(
            slide,
            format_report_value(production_shortage),
            col_lefts[5],
            row_top + 0.03,
            col_widths[5],
            row_height - 0.06,
            8.0,
            True,
            shortage_color,
            PP_ALIGN.RIGHT,
            MSO_ANCHOR.MIDDLE,
        )
        add_report_rule(slide, left + 0.1, row_top + row_height, width - 0.2, REPORT_FAINT)

    note = "제품군별 생산지시 PACK, 생산·포장·용마입고율, 생산부족 PCS 기준"
    if hidden_count:
        note = f"{note} / 외 {hidden_count:,}개 분류"
    add_textbox(
        slide,
        note,
        left + 0.16,
        top + height - 0.2,
        width - 0.32,
        0.16,
        6.2,
        False,
        REPORT_MUTED,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )


def add_urgent_request_summary_slide(
    prs: Presentation,
    urgent_summary_view: pd.DataFrame,
    scope_label: str,
    generated_at: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = ppt_rgb(REPORT_BG)

    add_report_shape(slide, MSO_SHAPE.RECTANGLE, 0.0, 0.0, 13.333, 0.88, REPORT_NAVY)
    add_report_shape(slide, MSO_SHAPE.RECTANGLE, 0.0, 0.86, 13.333, 0.03, REPORT_ACCENT, REPORT_ACCENT)
    add_textbox(
        slide,
        "요청 긴급 판매코드 요약",
        0.45,
        0.16,
        6.8,
        0.32,
        18,
        True,
        "#FFFFFF",
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide,
        "일일 재고표 기준 긴급요청 품목을 판매코드 단위로 집계",
        0.45,
        0.51,
        7.2,
        0.18,
        8.3,
        False,
        "#CBD5E1",
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide,
        f"기준: {scope_label} / 산출시각: {generated_at}",
        8.1,
        0.33,
        4.75,
        0.2,
        7.8,
        False,
        "#E5E7EB",
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )

    total_sku = int(pd.to_numeric(urgent_summary_view.get("SKU 수", pd.Series(dtype=float)), errors="coerce").fillna(0).sum())
    request_out_count = int((urgent_summary_view.get("요청구분", pd.Series(dtype=str)).astype(str) == "요청외").sum())
    add_textbox(
        slide,
        f"판매코드 {len(urgent_summary_view):,}개 / SKU {total_sku:,}개 / 요청외 판매코드 {request_out_count:,}개",
        0.62,
        1.12,
        8.0,
        0.24,
        10,
        True,
        REPORT_HEADER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide,
        "요청외는 해당 판매코드 안에 요청외 긴급 SKU가 1개 이상 포함된 경우로 분류합니다.",
        0.62,
        6.98,
        11.8,
        0.2,
        7.5,
        False,
        REPORT_MUTED,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )

    left = 0.62
    top = 1.5
    width = 12.1
    height = 5.3
    add_report_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, REPORT_PANEL, REPORT_PANEL_LINE, 0.5)
    add_report_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, 0.42, REPORT_NAVY, REPORT_NAVY, 0.5)

    headers = ["판매코드", "요청구분", "제품명", "SKU 수"]
    col_widths = [1.05, 1.25, 8.55, 0.95]
    col_lefts = [left + 0.16]
    for col_width in col_widths[:-1]:
        col_lefts.append(col_lefts[-1] + col_width)

    for idx, header in enumerate(headers):
        add_textbox(
            slide,
            header,
            col_lefts[idx],
            top + 0.06,
            col_widths[idx],
            0.3,
            8.5,
            True,
            "#FFFFFF",
            PP_ALIGN.RIGHT if idx == 3 else PP_ALIGN.LEFT,
            MSO_ANCHOR.MIDDLE,
        )

    if urgent_summary_view.empty:
        add_textbox(
            slide,
            "긴급요청 품목이 없습니다.",
            left + 0.2,
            top + 0.72,
            width - 0.4,
            0.35,
            9,
            False,
            REPORT_MUTED,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )
        return

    rows = urgent_summary_view.head(12).copy()
    row_height = 0.39
    for row_idx, (_, row) in enumerate(rows.iterrows(), start=1):
        row_top = top + 0.42 + (row_idx - 1) * row_height
        if row_idx % 2 == 0:
            add_report_shape(slide, MSO_SHAPE.RECTANGLE, left + 0.08, row_top, width - 0.16, row_height, REPORT_ROW_ALT)
        scope = clean_str(row.get("요청구분", ""))
        scope_color = REPORT_ACCENT if scope == "요청외" else COLOR_TEAL
        values = [
            clean_str(row.get("판매코드", "")),
            scope,
            truncate_report_text(row.get("제품명", ""), 64),
            format_report_value(row.get("SKU 수", 0.0)),
        ]
        colors = [COLOR_BLUE, scope_color, REPORT_HEADER, REPORT_ACCENT]
        aligns = [PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.RIGHT]
        for col_idx, value in enumerate(values):
            add_textbox(
                slide,
                value,
                col_lefts[col_idx],
                row_top + 0.02,
                col_widths[col_idx],
                row_height - 0.04,
                8.7 if col_idx != 2 else 8.4,
                col_idx in {0, 1, 3},
                colors[col_idx],
                aligns[col_idx],
                MSO_ANCHOR.MIDDLE,
            )
        add_report_rule(slide, left + 0.1, row_top + row_height, width - 0.2, REPORT_FAINT)

    hidden_count = max(0, len(urgent_summary_view) - len(rows))
    if hidden_count:
        add_textbox(
            slide,
            f"외 {hidden_count:,}개 판매코드는 대시보드와 엑셀 다운로드에서 확인 가능합니다.",
            left + 0.2,
            top + height - 0.36,
            width - 0.4,
            0.2,
            7.5,
            False,
            REPORT_MUTED,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )


def add_report_legend(slide: Any) -> None:
    legend_items = [
        (COLOR_TEAL, "진도율 정상 (>=80%)"),
        (COLOR_ORANGE, "진도율 미달 / 미입고 / 경고"),
        (COLOR_AMBER, "생산부족 발생"),
        ("#AAAAAA", "부족 없음"),
    ]
    for idx, (color, label) in enumerate(legend_items):
        left = 0.3 + idx * 3.1
        add_report_shape(slide, MSO_SHAPE.RECTANGLE, left, 7.28, 0.18, 0.1, color, color)
        add_textbox(
            slide,
            label,
            left + 0.22,
            7.22,
            2.8,
            0.22,
            7.5,
            False,
            TEXT_SECONDARY,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )


@st.cache_data(show_spinner=False, max_entries=8)
def build_ppt_report(
    product_view: pd.DataFrame,
    code_summary: pd.DataFrame,
    product_names: pd.Series,
    scope_label: str,
    daily_inventory_df: pd.DataFrame | None = None,
    sample_available_df: pd.DataFrame | None = None,
) -> bytes:
    work = add_allocated_production_basis(code_summary)
    work = code_summary_for_products(work, product_names)
    scope_kpis = build_scope_kpis(work)
    main_products, _ = split_main_sample(product_view)
    family_view = build_family_progress_view(main_products)
    urgent_summary_view = build_urgent_request_summary_view(
        daily_inventory_df,
        code_summary,
        sample_available_df,
    )
    urgent_sku_count = int(
        pd.to_numeric(urgent_summary_view.get("SKU 수", pd.Series(dtype=float)), errors="coerce")
        .fillna(0)
        .sum()
    )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = ppt_rgb(REPORT_BG)

    add_report_shape(slide, MSO_SHAPE.RECTANGLE, 0.0, 0.0, 13.333, 0.78, REPORT_SOFT_BG, REPORT_SOFT_BG)
    add_report_rule(slide, 0.45, 0.78, 12.45, REPORT_PANEL_LINE)
    add_textbox(
        slide,
        "국내 생산·포장 운영현황",
        0.45,
        0.1,
        6.5,
        0.38,
        28,
        True,
        REPORT_HEADER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide,
        "국내 요청 물량의 생산·포장·용마 입고 현황",
        0.45,
        0.5,
        6.6,
        0.18,
        9.2,
        False,
        REPORT_MUTED,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    generated_at = pd.Timestamp.now(tz="Asia/Seoul").strftime("%Y-%m-%d %H:%M")
    add_textbox(
        slide,
        f"기준 : {scope_label}",
        8.2,
        0.19,
        4.7,
        0.18,
        8.2,
        False,
        REPORT_HEADER,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide,
        f"산출 : {generated_at}",
        8.2,
        0.45,
        4.7,
        0.18,
        8.2,
        False,
        REPORT_MUTED,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )

    kpi_map = {name: kpi for name, kpi in scope_kpis}
    total_kpi = kpi_map.get("전체", {})
    total_progress = to_report_float(total_kpi.get("progress_pct", 0.0))
    total_packing_progress = to_report_float(total_kpi.get("packing_progress_pct", 0.0))
    total_shortage = to_report_float(total_kpi.get("shortage_pack", 0.0))
    exception_count = float(urgent_sku_count)
    if total_shortage > 0 or exception_count > 0:
        banner_fill = "#FFFFFF"
        banner_color = REPORT_ACCENT
        status_label = "주의"
    else:
        banner_fill = "#FFFFFF"
        banner_color = COLOR_BLUE
        status_label = "정상"
    banner_text = (
        f"포장률 {total_packing_progress:.1f}%     "
        f"용마입고율 {total_progress:.1f}%     "
        f"미입고 {format_report_value(total_shortage)} PACK     "
        f"긴급 SKU {format_report_value(exception_count)}"
    )

    add_report_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.45, 0.94, 12.45, 0.38, banner_fill, REPORT_PANEL_LINE, 0.5)
    add_report_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.62, 1.03, 0.62, 0.2, banner_color, banner_color, 0.4)
    add_textbox(
        slide,
        status_label,
        0.62,
        1.035,
        0.62,
        0.17,
        7.0,
        True,
        "#FFFFFF",
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide,
        banner_text,
        1.38,
        1.035,
        11.1,
        0.2,
        8.7,
        True,
        REPORT_HEADER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )

    add_textbox(slide, "KPI", 0.45, 1.46, 2.0, 0.24, 12.5, True, REPORT_HEADER)
    add_kpi_card(slide, "전체 KPI", total_kpi, COLOR_BLUE, 0.45, 1.74, 5.7, 2.28, emphasis=True)
    add_kpi_card(slide, "본품 KPI", kpi_map.get("본품", {}), COLOR_TEAL, 6.35, 1.74, 3.1, 2.28)
    add_kpi_card(slide, "샘플 KPI", kpi_map.get("샘플", {}), COLOR_TEAL, 9.75, 1.74, 3.15, 2.28)

    add_textbox(slide, "제품 분류별 진도현황", 0.45, 4.24, 12.45, 0.24, 12.5, True, REPORT_HEADER)

    add_family_progress_summary_panel(slide, family_view, left=0.45, top=4.55, width=12.45, height=2.48, max_rows=7)
    add_textbox(
        slide,
        "진도율은 생산요청 기준이며, 제품 분류별 현황은 본품 제품군 기준으로 산출됩니다.",
        0.45,
        7.18,
        12.4,
        0.22,
        7.5,
        False,
        REPORT_MUTED,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )

    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def build_status_count_chart(main_df: pd.DataFrame, sample_df: pd.DataFrame) -> px.bar:
    parts: list[pd.DataFrame] = []
    for category, frame in [("본품", main_df), ("샘플", sample_df)]:
        counts = frame["상태"].value_counts().reindex(STATUS_ORDER, fill_value=0).reset_index()
        counts.columns = ["상태", "제품 수"]
        counts["구분"] = category
        parts.append(counts)
    status_df = pd.concat(parts, ignore_index=True)

    fig = px.bar(
        status_df,
        x="상태",
        y="제품 수",
        color="구분",
        barmode="group",
        title="상태별 제품 수",
        text="제품 수",
        color_discrete_map={"본품": NAVY, "샘플": SOFT_NAVY},
    )
    fig.update_traces(texttemplate="%{text:,.0f}")
    fig.update_layout(
        paper_bgcolor=WHITE,
        plot_bgcolor=WHITE,
        margin=dict(l=8, r=8, t=52, b=8),
        legend_title_text="",
    )
    return fig


def build_progress_compare_chart(
    total_kpi: dict[str, float],
    main_kpi: dict[str, float],
    sample_kpi: dict[str, float],
) -> px.bar:
    compare_df = pd.DataFrame(
        {
            "구분": ["전체", "본품", "샘플"],
            "진도율(%)": [total_kpi["progress_pct"], main_kpi["progress_pct"], sample_kpi["progress_pct"]],
        }
    )
    fig = px.bar(
        compare_df,
        x="구분",
        y="진도율(%)",
        color="구분",
        title="전체/본품/샘플 진도율 비교",
        text="진도율(%)",
        color_discrete_map={"전체": COLOR_BLUE, "본품": COLOR_TEAL, "샘플": COLOR_AMBER},
    )
    fig.update_traces(texttemplate="%{text:.1f}%")
    fig.update_layout(
        paper_bgcolor=WHITE,
        plot_bgcolor=WHITE,
        margin=dict(l=8, r=8, t=52, b=8),
        legend_title_text="",
        showlegend=False,
    )
    fig.update_yaxes(rangemode="tozero")
    return fig


def build_shortage_top_chart(main_df: pd.DataFrame, sample_df: pd.DataFrame, top_n: int = 10) -> px.bar | None:
    source = pd.concat(
        [
            main_df.assign(구분="본품"),
            sample_df.assign(구분="샘플"),
        ],
        ignore_index=True,
    )
    source = source[source["포장부족수량"] > 0].copy()
    if source.empty:
        return None
    source = source.nlargest(top_n, "포장부족수량").sort_values("포장부족수량", ascending=True)
    source["라벨"] = source.apply(lambda r: f"[{r['구분']}] {r['제품명']}", axis=1)

    fig = px.bar(
        source,
        x="포장부족수량",
        y="라벨",
        color="구분",
        orientation="h",
        title=f"포장부족 TOP {min(top_n, len(source))}",
        text="포장부족수량",
        color_discrete_map={"본품": COLOR_ORANGE, "샘플": COLOR_AMBER},
    )
    fig.update_traces(texttemplate="%{text:,.0f}")
    fig.update_layout(
        paper_bgcolor=WHITE,
        plot_bgcolor=WHITE,
        margin=dict(l=8, r=8, t=52, b=8),
        legend_title_text="",
        yaxis_title="",
    )
    return fig


def build_product_progress_gap_chart(df: pd.DataFrame, top_n: int = 18) -> px.bar | None:
    if df.empty:
        return None
    source = df.copy()
    source = source.sort_values(
        ["포장부족수량", "생산부족수량", "요청 PACK"],
        ascending=[False, False, False],
        kind="stable",
    ).head(top_n)
    if source.empty:
        return None
    if "용마입고율" not in source.columns:
        source["용마입고율"] = 0.0
    chart_source = source[["제품명", "생산진도율", "용마입고율"]]
    chart_df = chart_source.melt(
        id_vars="제품명",
        value_vars=["생산진도율", "용마입고율"],
        var_name="지표",
        value_name="진도율",
    )
    product_order = source["제품명"].tolist()[::-1]
    fig = px.bar(
        chart_df,
        x="진도율",
        y="제품명",
        color="지표",
        barmode="group",
        orientation="h",
        category_orders={"제품명": product_order, "지표": ["생산진도율", "용마입고율"]},
        title="제품별 생산진도율 vs 용마입고율",
        text="진도율",
        color_discrete_map={"생산진도율": COLOR_BLUE, "용마입고율": COLOR_TEAL},
    )
    fig.update_traces(texttemplate="%{text:.1f}%")
    fig.update_layout(
        paper_bgcolor=WHITE,
        plot_bgcolor=WHITE,
        margin=dict(l=8, r=8, t=52, b=8),
        legend_title_text="",
        yaxis_title="",
    )
    fig.update_xaxes(range=[0, 100])
    return fig


def render_style() -> None:
    st.markdown(
        f"""
        <style>
        @import url('https://cdn.jsdelivr.net/gh/orioncactus/pretendard/dist/web/static/pretendard.css');

        html, body, [class*="css"], .stApp {{
            font-family: 'Pretendard', 'Apple SD Gothic Neo', sans-serif !important;
            color: {TEXT_PRIMARY};
        }}
        .stApp {{
            background: {BG_PAGE};
            color: {TEXT_PRIMARY};
        }}
        :root {{
            --ui-gap: 12px;
            --color-blue: {COLOR_BLUE};
            --color-teal: {COLOR_TEAL};
            --color-orange: {COLOR_ORANGE};
            --color-amber: {COLOR_AMBER};
            --color-purple: {COLOR_AMBER};
            --color-danger: {COLOR_DANGER};
            --bg-page: {BG_PAGE};
            --bg-card: {BG_CARD};
            --bg-section: {BG_SECTION};
            --text-primary: {TEXT_PRIMARY};
            --text-secondary: {TEXT_SECONDARY};
            --text-tertiary: {TEXT_TERTIARY};
            --border-default: {BORDER_DEFAULT};
            --border-light: {BORDER_LIGHT};
            --kpi-card-height: 320px;
        }}
        .block-container {{
            padding: 24px 32px 32px !important;
            max-width: 100% !important;
        }}
        h1 {{
            font-size: 22px !important;
            font-weight: 500 !important;
            color: {TEXT_PRIMARY} !important;
            margin-bottom: 4px !important;
            letter-spacing: 0 !important;
        }}
        h2 {{
            font-size: 15px !important;
            font-weight: 500 !important;
            color: {TEXT_PRIMARY} !important;
            margin-bottom: 2px !important;
            letter-spacing: 0 !important;
        }}
        h3 {{
            font-size: 13px !important;
            font-weight: 500 !important;
            color: {TEXT_PRIMARY} !important;
            letter-spacing: 0 !important;
        }}
        [data-baseweb="tab-list"] {{
            gap: 0 !important;
            border-bottom: 1.5px solid {BORDER_DEFAULT} !important;
            background: transparent !important;
            margin-bottom: 20px !important;
        }}
        [data-baseweb="tab"] {{
            font-size: 13px !important;
            padding: 10px 20px !important;
            color: {TEXT_SECONDARY} !important;
            background: transparent !important;
            letter-spacing: 0 !important;
        }}
        [aria-selected="true"][data-baseweb="tab"] {{
            color: {COLOR_ORANGE} !important;
            font-weight: 500 !important;
        }}
        [data-baseweb="tab-highlight"] {{
            background-color: {COLOR_ORANGE} !important;
            height: 2px !important;
        }}
        [data-baseweb="tab-border"] {{
            display: none !important;
        }}
        [data-testid="stSegmentedControl"] {{
            margin: 2px 0 14px 0 !important;
        }}
        [data-testid="stSegmentedControl"] button {{
            border: 1px solid {BORDER_DEFAULT} !important;
            border-radius: 8px !important;
            background: {BG_CARD} !important;
            color: {TEXT_PRIMARY} !important;
            font-size: 13px !important;
            font-weight: 500 !important;
            padding: 10px 20px !important;
            box-shadow: none !important;
        }}
        [data-testid="stSegmentedControl"] button[aria-pressed="true"] {{
            background: {BG_CARD} !important;
            color: {COLOR_DANGER} !important;
            font-weight: 600 !important;
            border-color: {COLOR_DANGER} !important;
            border-bottom-color: {COLOR_DANGER} !important;
        }}
        .dashboard-nav-divider {{
            height: 1.5px;
            background: {BORDER_DEFAULT};
            margin: -15px 0 20px 0;
        }}
        [data-testid="stRadio"] label,
        [data-testid="stCheckbox"] label {{
            font-size: 13px !important;
            color: {TEXT_SECONDARY} !important;
        }}
        [data-testid="stRadio"] [data-testid="stMarkdownContainer"] p,
        [data-testid="stCheckbox"] [data-testid="stMarkdownContainer"] p {{
            font-size: 13px !important;
            color: {TEXT_SECONDARY} !important;
        }}
        [data-testid="metric-container"] {{
            background: {BG_CARD};
            border: 0.5px solid {BORDER_DEFAULT};
            border-radius: 12px;
            padding: 14px 18px;
            box-shadow: none;
        }}
        [data-testid="stMetricLabel"] {{
            font-size: 11px !important;
            color: {TEXT_TERTIARY} !important;
            font-weight: 400 !important;
        }}
        [data-testid="stMetricValue"] {{
            font-size: 20px !important;
            line-height: 1.1 !important;
            font-weight: 500 !important;
            color: {TEXT_PRIMARY} !important;
        }}
        [data-testid="stMetricDelta"] {{
            display: none !important;
        }}
        .completion-summary-card {{
            display: grid;
            grid-template-columns: repeat(5, minmax(120px, 1fr));
            gap: 0;
            margin: 12px 0 14px;
            border: 1px solid {BORDER_DEFAULT};
            border-radius: 12px;
            background: {BG_CARD};
            box-shadow: 0 2px 8px rgba(15, 23, 42, 0.04);
            overflow: hidden;
        }}
        .completion-summary-item {{
            padding: 14px 18px;
            border-right: 1px solid {BORDER_DEFAULT};
        }}
        .completion-summary-item:last-child {{
            border-right: 0;
        }}
        .completion-summary-item span {{
            display: block;
            font-size: 11px;
            line-height: 1.2;
            color: {TEXT_TERTIARY};
            font-weight: 600;
            margin-bottom: 6px;
        }}
        .completion-summary-item strong {{
            display: block;
            font-size: 20px;
            line-height: 1.15;
            color: {TEXT_PRIMARY};
            font-weight: 700;
            letter-spacing: 0;
        }}
        .completion-summary-item.done strong {{
            color: #16A34A;
        }}
        .completion-summary-item.active strong {{
            color: {COLOR_ORANGE};
        }}
        .completion-summary-item.risk strong {{
            color: {COLOR_DANGER};
        }}
        [data-testid="stDataFrame"] {{
            border: 1px solid {BORDER_DEFAULT} !important;
            border-radius: 12px !important;
            overflow: hidden;
            background: {BG_CARD} !important;
            box-shadow: 0 2px 8px rgba(15, 23, 42, 0.04) !important;
        }}
        [data-testid="stDataFrame"] th {{
            background: #FFFFFF !important;
            font-size: 12px !important;
            font-weight: 600 !important;
            color: {TEXT_SECONDARY} !important;
            padding: 12px 12px !important;
            border-bottom: 1px solid {BORDER_DEFAULT} !important;
            white-space: nowrap;
            text-align: center !important;
        }}
        [data-testid="stDataFrame"] td {{
            font-size: 13px !important;
            color: {TEXT_PRIMARY} !important;
            padding: 10px 12px !important;
            border-bottom: 0.5px solid {BORDER_LIGHT} !important;
        }}
        [data-testid="stDataFrame"] table,
        [data-testid="stDataFrame"] thead,
        [data-testid="stDataFrame"] tbody,
        [data-testid="stDataFrame"] tr,
        [data-testid="stDataFrame"] td,
        [data-testid="stDataFrame"] [role="grid"],
        [data-testid="stDataFrame"] [role="rowgroup"],
        [data-testid="stDataFrame"] [role="row"],
        [data-testid="stDataFrame"] [role="gridcell"] {{
            background: {BG_CARD} !important;
        }}
        [data-testid="stDataFrame"] [role="columnheader"] {{
            background: #FFFFFF !important;
            font-weight: 600 !important;
            color: #374151 !important;
            border-color: {BORDER_DEFAULT} !important;
            justify-content: center !important;
        }}
        [data-testid="stDataFrame"] tr:hover td {{
            background: #F8FAFC !important;
            cursor: pointer !important;
        }}
        [data-testid="stDataFrame"] [role="row"] {{
            min-height: 48px !important;
        }}
        [data-testid="stDataFrame"] [role="columnheader"],
        [data-testid="stDataFrame"] [role="gridcell"] {{
            min-height: 48px !important;
            align-items: center !important;
            border-color: {BORDER_DEFAULT} !important;
        }}
        [data-testid="stDataFrame"] [role="row"]:hover [role="gridcell"] {{
            background: #F8FAFC !important;
            cursor: pointer !important;
        }}
        [data-testid="stDataFrame"] [role="row"]:nth-child(even) [role="gridcell"] {{
            background: #FCFCFD !important;
        }}
        [data-testid="stDataFrame"] [role="row"]:hover [role="gridcell"]:first-child {{
            box-shadow: inset 4px 0 0 {COLOR_BLUE} !important;
        }}
        [data-testid="stDataFrame"] [aria-selected="true"] [role="gridcell"],
        [data-testid="stDataFrame"] [role="row"][aria-selected="true"] [role="gridcell"] {{
            background: #EFF6FF !important;
        }}
        [data-testid="stDataFrame"] [role="progressbar"] {{
            height: 4px !important;
            border-radius: 999px !important;
            background: {BORDER_DEFAULT} !important;
            overflow: hidden !important;
        }}
        [data-testid="stDataFrame"] [role="progressbar"] > div {{
            border-radius: 999px !important;
        }}
        [data-testid="stTextInput"] input,
        [data-testid="stNumberInput"] input,
        [data-testid="stSelectbox"] > div > div,
        [data-testid="stMultiSelect"] > div > div {{
            font-size: 13px !important;
            border-radius: 8px !important;
            border: 0.5px solid rgba(0,0,0,0.15) !important;
            background: {BG_CARD} !important;
            color: {TEXT_PRIMARY} !important;
        }}
        [data-testid="stTextInput"] label,
        [data-testid="stNumberInput"] label,
        [data-testid="stSelectbox"] label,
        [data-testid="stMultiSelect"] label {{
            font-size: 12px !important;
            font-weight: 500 !important;
            color: {TEXT_SECONDARY} !important;
        }}
        [data-testid="stMultiSelect"] [data-baseweb="tag"] {{
            background: {BG_CARD} !important;
            color: #993C1D !important;
            border-radius: 20px !important;
            border: 1px solid {COLOR_ALERT_BD} !important;
            font-size: 12px !important;
            font-weight: 500 !important;
        }}
        [data-testid="stButton"] button,
        [data-testid="stDownloadButton"] button {{
            border-radius: 10px !important;
            font-size: 13px !important;
            font-weight: 500 !important;
            border: 0.5px solid rgba(0,0,0,0.15) !important;
            background: {BG_CARD} !important;
            color: {TEXT_PRIMARY} !important;
            box-shadow: none !important;
        }}
        [data-testid="stButton"] button:hover,
        [data-testid="stDownloadButton"] button:hover {{
            background: {BG_PAGE} !important;
            border-color: rgba(0,0,0,0.25) !important;
            color: {TEXT_PRIMARY} !important;
        }}
        hr {{
            border-color: rgba(0,0,0,0.08) !important;
            margin: 20px 0 !important;
        }}
        ::-webkit-scrollbar {{
            width: 6px;
            height: 6px;
        }}
        ::-webkit-scrollbar-track {{
            background: transparent;
        }}
        ::-webkit-scrollbar-thumb {{
            background: rgba(0,0,0,0.15);
            border-radius: 3px;
        }}
        .kpi-panel {{
            background: {BG_CARD};
            border: 0.5px solid {BORDER_DEFAULT};
            border-radius: 12px;
            padding: 14px 18px;
            box-shadow: none;
            margin-bottom: 0;
            height: 100%;
        }}
        .drill-kpi {{
            margin-bottom: 12px;
        }}
        .kpi-title {{
            font-size: 15px;
            font-weight: 500;
            color: {TEXT_PRIMARY};
            margin-bottom: var(--ui-gap);
        }}
        .kpi-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(120px, 1fr));
            gap: var(--ui-gap);
        }}
        .scope-kpi .kpi-grid {{
            grid-template-columns: repeat(3, minmax(0, 1fr));
            gap: 8px;
        }}
        .kpi-card {{
            border: 0.5px solid {BORDER_DEFAULT};
            border-radius: 12px;
            padding: 14px 18px;
            background: {BG_CARD};
            min-height: 72px;
            display: flex;
            flex-direction: column;
            justify-content: space-between;
        }}
        .metric-label {{
            font-size: 11px;
            font-weight: 400;
            color: {TEXT_TERTIARY};
            margin-bottom: 6px;
        }}
        .metric-value {{
            font-size: 20px;
            line-height: 1.1;
            font-weight: 500;
            color: {TEXT_PRIMARY};
            white-space: nowrap;
            overflow-wrap: normal;
            word-break: normal;
        }}
        .scope-kpi .metric-value {{
            font-size: 20px;
        }}
        @media (max-width: 1100px) {{
            .scope-kpi .kpi-grid {{
                grid-template-columns: repeat(2, minmax(0, 1fr));
            }}
        }}
        .metric-value.warn {{
            color: {COLOR_ORANGE};
        }}
        .metric-value.risk {{
            color: {COLOR_ORANGE};
        }}
        .metric-value.normal {{
            color: {TEXT_PRIMARY};
        }}
        .metric-value.good {{
            color: {COLOR_TEAL};
        }}
        .metric-value.mid {{
            color: {COLOR_AMBER};
        }}
        .metric-value.muted {{
            color: {TEXT_TERTIARY};
        }}
        .metric-note {{
            color: {TEXT_TERTIARY};
            font-size: 11px;
            line-height: 1.2;
            margin-top: 4px;
        }}
        .mini-kpi-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(150px, 1fr));
            gap: 10px;
            margin-bottom: 12px;
        }}
        .mini-kpi-card {{
            background: {BG_CARD};
            border: 0.5px solid {BORDER_DEFAULT};
            border-radius: 12px;
            padding: 14px 18px;
            min-height: 76px;
            display: flex;
            flex-direction: column;
            justify-content: space-between;
            box-shadow: none;
        }}
        @media (max-width: 1100px) {{
            .mini-kpi-grid {{
                grid-template-columns: repeat(2, minmax(0, 1fr));
            }}
        }}
        .shortage-card {{
            border-color: {COLOR_ALERT_BD};
            background: {BG_CARD};
        }}
        .status-board {{
            display: grid;
            grid-template-columns: minmax(320px, 1.1fr) minmax(420px, 1.9fr);
            gap: 10px;
            margin: 2px 0 10px;
            align-items: stretch;
        }}
        .status-main,
        .status-tile {{
            background: {BG_CARD};
            border: 0.5px solid {BORDER_DEFAULT};
            border-radius: 8px;
            box-shadow: none;
        }}
        .status-main {{
            padding: 16px 18px;
            border-left: 4px solid {TEXT_TERTIARY};
            box-shadow: none;
        }}
        .status-board.warn .status-main {{
            border-left-color: {COLOR_AMBER};
        }}
        .status-board.risk .status-main {{
            border-left-color: {COLOR_ORANGE};
        }}
        .status-head {{
            display: flex;
            align-items: center;
            justify-content: space-between;
            gap: 12px;
            margin-bottom: 10px;
        }}
        .status-head strong {{
            color: {TEXT_PRIMARY};
            font-size: 14px;
            font-weight: 700;
        }}
        .status-pill {{
            display: inline-flex;
            align-items: center;
            border-radius: 999px;
            padding: 4px 9px;
            font-size: 11px;
            font-weight: 700;
            color: {TEXT_SECONDARY};
            background: {BG_CARD};
            border: 1px solid {BORDER_DEFAULT};
        }}
        .status-pill.warn {{
            color: {COLOR_AMBER};
            background: {BG_CARD};
            border-color: {COLOR_AMBER};
        }}
        .status-pill.risk {{
            color: {COLOR_ORANGE};
            background: {BG_CARD};
            border-color: {COLOR_ALERT_BD};
        }}
        .status-main-value {{
            color: {TEXT_PRIMARY};
            font-size: 34px;
            line-height: 1;
            font-weight: 700;
            font-variant-numeric: tabular-nums;
            margin-bottom: 12px;
        }}
        .status-flow {{
            display: flex;
            width: 100%;
            height: 10px;
            border-radius: 999px;
            background: {BG_SECTION};
            overflow: hidden;
        }}
        .status-flow-fill.receipt {{
            background: {COLOR_TEAL};
        }}
        .status-flow-fill.shortage {{
            background: {COLOR_ORANGE};
        }}
        .status-flow-legend {{
            display: flex;
            justify-content: space-between;
            gap: 10px;
            flex-wrap: wrap;
            margin-top: 9px;
            color: {TEXT_SECONDARY};
            font-size: 11px;
            font-variant-numeric: tabular-nums;
        }}
        .status-tile-grid {{
            display: grid;
            grid-template-columns: repeat(3, minmax(0, 1fr));
            gap: 8px;
        }}
        .status-tile {{
            min-height: 72px;
            padding: 11px 13px;
            display: flex;
            flex-direction: column;
            justify-content: space-between;
        }}
        .status-tile .metric-value {{
            font-size: 20px;
        }}
        @media (max-width: 1100px) {{
            .status-board {{
                grid-template-columns: 1fr;
            }}
            .status-tile-grid {{
                grid-template-columns: repeat(2, minmax(0, 1fr));
            }}
        }}
        @media (max-width: 640px) {{
            .status-tile-grid {{
                grid-template-columns: 1fr;
            }}
            .status-flow-legend {{
                flex-direction: column;
                gap: 4px;
            }}
        }}
        .panel-box {{
            background: {BG_CARD};
            border: 0.5px solid {BORDER_DEFAULT};
            border-radius: 8px;
            padding: 12px 14px;
            box-shadow: none;
        }}
        .family-section {{
            display: flex;
            flex-direction: column;
            gap: 8px;
        }}
        .family-section + .family-section {{
            margin-top: 16px;
        }}
        .family-section-title {{
            display: inline-block;
            width: fit-content;
            color: #444441;
            background: {BG_SECTION};
            font-size: 13px;
            font-weight: 500;
            line-height: 1.2;
            padding: 5px 12px;
            border-radius: 8px;
        }}
        .family-grid {{
            display: grid;
            grid-template-columns: repeat(5, minmax(0, 1fr));
            gap: 8px;
        }}
        .family-card {{
            border: 0.5px solid {BORDER_DEFAULT};
            border-radius: 8px;
            background: {BG_CARD};
            padding: 11px 13px;
            min-height: 106px;
            display: flex;
            flex-direction: column;
            gap: 8px;
        }}
        .family-card:has(.progress-fill.production.risk),
        .family-card:has(.progress-fill.production.warn) {{
            border-color: {COLOR_ALERT_BD};
        }}
        .family-head {{
            display: flex;
            align-items: baseline;
            justify-content: space-between;
            gap: 10px;
        }}
        .family-head span {{
            color: {TEXT_PRIMARY};
            font-size: 12px;
            font-weight: 500;
            white-space: nowrap;
            overflow: hidden;
            text-overflow: ellipsis;
        }}
        .family-head b {{
            color: {TEXT_TERTIARY};
            font-size: 11px;
            font-weight: 400;
            font-variant-numeric: tabular-nums;
            white-space: nowrap;
        }}
        .family-shortages {{
            display: flex;
            justify-content: flex-start;
            gap: 8px;
            color: {TEXT_SECONDARY};
            font-size: 10px;
        }}
        .family-shortages b {{
            color: {COLOR_ORANGE};
            font-variant-numeric: tabular-nums;
        }}
        .top-list {{
            display: flex;
            flex-direction: column;
            gap: 7px;
        }}
        .top-row {{
            display: grid;
            grid-template-columns: 32px minmax(220px, 1fr) 120px minmax(240px, 0.9fr);
            gap: 12px;
            align-items: center;
            border-bottom: 0.5px solid {BORDER_LIGHT};
            padding: 10px 16px;
            background: {BG_CARD};
        }}
        .top-rank {{
            font-size: 12px;
            font-weight: 500;
            color: {TEXT_TERTIARY};
            text-align: center;
        }}
        .top-name {{
            color: {TEXT_PRIMARY};
            font-size: 13px;
            font-weight: 400;
            overflow-wrap: anywhere;
        }}
        .top-shortage {{
            color: {COLOR_ORANGE};
            font-size: 13px;
            font-weight: 500;
            text-align: right;
            font-variant-numeric: tabular-nums;
        }}
        .top-progress {{
            min-width: 0;
        }}
        .gap-list {{
            display: flex;
            flex-direction: column;
            gap: 7px;
        }}
        .gap-row {{
            display: grid;
            grid-template-columns: 32px minmax(220px, 1fr) minmax(220px, 0.78fr) minmax(220px, 0.78fr) 82px;
            gap: 12px;
            align-items: center;
            border-bottom: 0.5px solid {BORDER_LIGHT};
            padding: 10px 16px;
            background: {BG_CARD};
        }}
        .gap-progress {{
            min-width: 0;
        }}
        .gap-value {{
            color: {COLOR_ORANGE};
            font-size: 13px;
            font-weight: 500;
            text-align: right;
            font-variant-numeric: tabular-nums;
        }}
        @media (max-width: 900px) {{
            .family-grid {{
                grid-template-columns: repeat(auto-fit, minmax(240px, 1fr));
            }}
            .top-row {{
                grid-template-columns: 28px 1fr;
            }}
            .top-shortage {{
                text-align: left;
            }}
            .top-progress {{
                grid-column: 2;
            }}
            .gap-row {{
                grid-template-columns: 28px 1fr;
            }}
            .gap-progress, .gap-value {{
                grid-column: 2;
            }}
            .gap-value {{
                text-align: left;
            }}
        }}
        .drill-panel {{
            margin-bottom: 12px;
        }}
        .table-wrap {{
            max-height: 640px;
            overflow: auto;
            border: 0.5px solid {BORDER_DEFAULT};
            border-radius: 8px;
            background: {BG_CARD};
        }}
        .compact-table {{
            max-height: 360px;
        }}
        .main-summary-table {{
            width: 100%;
            min-width: 980px;
        }}
        .main-summary-table .summary-rank-col {{
            width: 6%;
        }}
        .main-summary-table .summary-product-col {{
            width: 30%;
        }}
        .main-summary-table .summary-number-col {{
            width: 13%;
        }}
        .main-summary-table .summary-progress-col {{
            width: 17%;
        }}
        .gap-summary-table .summary-product-col {{
            width: 35%;
        }}
        .gap-summary-table .summary-progress-col {{
            width: 22%;
        }}
        .gap-summary-table .summary-number-col {{
            width: 15%;
        }}
        .urgent-request-summary-table .summary-code-col {{
            width: 12%;
        }}
        .urgent-request-summary-table .summary-scope-col {{
            width: 14%;
        }}
        .urgent-request-summary-table .summary-product-col {{
            width: 60%;
        }}
        .urgent-request-summary-table .summary-number-col {{
            width: 14%;
        }}
        .ops-table {{
            width: 100%;
            border-collapse: collapse;
            table-layout: fixed;
        }}
        .ops-table th {{
            position: sticky;
            top: 0;
            background: #FFFFFF;
            color: {TEXT_SECONDARY};
            font-size: 11px;
            font-weight: 600;
            border-bottom: 1px solid {BORDER_DEFAULT};
            padding: 8px 12px;
            z-index: 1;
        }}
        .ops-table td {{
            border-bottom: 1px solid {BORDER_DEFAULT};
            padding: 7px 12px;
            font-size: 12px;
            color: {TEXT_PRIMARY};
            vertical-align: middle;
            background: {BG_CARD};
            height: 48px;
        }}
        .ops-table tbody tr:hover td {{
            background: #F8FAFC;
        }}
        .ops-table td.left, .ops-table th.left {{
            text-align: left;
        }}
        .ops-table td.num, .ops-table th.num {{
            text-align: right;
            font-variant-numeric: tabular-nums;
        }}
        .ops-table td.num.shortage {{
            color: {COLOR_DANGER};
            font-weight: 700;
        }}
        .ops-table td.num.muted {{
            color: {TEXT_SECONDARY};
            font-weight: 500;
        }}
        .ops-table td.num.negative {{
            color: {COLOR_ORANGE};
            font-weight: 700;
        }}
        .ops-table td.code-cell {{
            color: {COLOR_BLUE};
            font-weight: 700;
            font-variant-numeric: tabular-nums;
        }}
        .request-scope-badge {{
            display: inline-flex;
            align-items: center;
            justify-content: center;
            border-radius: 999px;
            padding: 3px 8px;
            font-size: 11px;
            font-weight: 700;
            border: 0.5px solid {BORDER_DEFAULT};
            background: {BG_CARD};
            color: {TEXT_SECONDARY};
            white-space: nowrap;
        }}
        .request-scope-badge.in {{
            color: {COLOR_TEAL};
            border-color: #B9E3D4;
        }}
        .request-scope-badge.out {{
            color: {COLOR_ORANGE};
            border-color: {COLOR_ALERT_BD};
        }}
        .request-scope-badge.mixed {{
            color: {COLOR_AMBER};
            border-color: #E4B968;
        }}
        .response-badge {{
            display: inline-flex;
            align-items: center;
            min-width: 72px;
            justify-content: center;
            border-radius: 999px;
            padding: 3px 8px;
            font-size: 11px;
            font-weight: 600;
            background: {BG_CARD};
            color: {TEXT_SECONDARY};
            border: 1px solid {BORDER_DEFAULT};
        }}
        .response-badge.partial {{
            color: {COLOR_AMBER};
            border-color: #E4B968;
        }}
        .response-badge.need {{
            color: {COLOR_ORANGE};
            border-color: {COLOR_ALERT_BD};
        }}
        .ops-table td.power-cell {{
            text-align: center;
            font-variant-numeric: tabular-nums;
            font-weight: 500;
            color: {TEXT_PRIMARY};
        }}
        .ops-table td.power-cell.high {{
            color: {TEXT_PRIMARY};
        }}
        .progress-cell {{
            display: flex;
            align-items: center;
            gap: 8px;
        }}
        .progress-name {{
            min-width: 28px;
            color: {TEXT_SECONDARY};
            font-size: 11px;
            font-weight: 500;
        }}
        .progress-track {{
            flex: 1;
            min-width: 80px;
            height: 4px;
            border-radius: 999px;
            background: {BORDER_DEFAULT};
            overflow: hidden;
        }}
        .progress-fill {{
            height: 100%;
            border-radius: 999px;
            background: {COLOR_ORANGE};
        }}
        .progress-fill.done {{
            background: {COLOR_ORANGE};
        }}
        .progress-fill.active {{
            background: {COLOR_ORANGE};
        }}
        .progress-fill.warn {{
            background: {COLOR_ORANGE};
        }}
        .progress-fill.risk {{
            background: {COLOR_ORANGE};
        }}
        .progress-fill.production {{
            background: {COLOR_BLUE};
        }}
        .progress-fill.packing {{
            background: {COLOR_ORANGE};
        }}
        .progress-fill.receipt {{
            background: {COLOR_AMBER};
        }}
        .progress-fill.risk.receipt {{
            background: {COLOR_AMBER};
        }}
        .progress-text {{
            min-width: 52px;
            text-align: right;
            font-size: 11px;
            color: {TEXT_TERTIARY};
            font-variant-numeric: tabular-nums;
        }}
        .status-badge {{
            display: inline-block;
            padding: 3px 8px;
            border-radius: 6px;
            border: 1px solid transparent;
            font-size: 11px;
            font-weight: 400;
            line-height: 1.2;
        }}
        .status-badge.done {{
            background: {BG_CARD};
            color: #3B6D11;
            border-color: #CBD5E1;
        }}
        .status-badge.active {{
            background: {BG_CARD};
            color: {COLOR_BLUE};
            border-color: #BFDBFE;
        }}
        .status-badge.warn {{
            background: {BG_CARD};
            color: #993C1D;
            font-weight: 500;
            border-color: {COLOR_ALERT_BD};
        }}
        .status-badge.waiting {{
            background: {BG_CARD};
            color: #5F5E5A;
            border-color: {BORDER_DEFAULT};
        }}
        .status-badge.risk {{
            background: {BG_CARD};
            color: #993C1D;
            border-color: {COLOR_ALERT_BD};
        }}
        .section-title {{
            color: {TEXT_PRIMARY};
            font-weight: 500;
            font-size: 15px;
            margin-bottom: 4px;
        }}
        .section-sub {{
            color: {TEXT_SECONDARY};
            font-size: 12px;
            font-weight: 400;
            margin-bottom: 10px;
        }}
        .breadcrumb {{
            display: flex;
            gap: 8px;
            align-items: center;
            color: {TEXT_SECONDARY};
            font-size: 12px;
            margin: 2px 0 10px 0;
        }}
        .breadcrumb span {{
            color: {COLOR_BLUE};
            font-weight: 500;
        }}
        .breadcrumb b {{
            color: {TEXT_SECONDARY};
        }}
        .progress-summary-panel {{
            display: grid;
            grid-template-columns: minmax(0, 1fr) 160px;
            gap: 14px;
            align-items: stretch;
            background: {BG_CARD};
            border: 0.5px solid {BORDER_DEFAULT};
            border-radius: 12px;
            padding: 16px 20px;
            margin-bottom: 12px;
        }}
        .progress-summary-panel .progress-cell {{
            margin: 10px 0;
        }}
        .dday-box {{
            border: 0.5px solid {BORDER_DEFAULT};
            border-radius: 12px;
            padding: 12px;
            display: flex;
            flex-direction: column;
            justify-content: center;
            background: {BG_PAGE};
        }}
        .dday-value {{
            color: {COLOR_ORANGE};
            font-size: 20px;
            font-weight: 500;
            font-variant-numeric: tabular-nums;
        }}

        /* Enterprise dashboard refresh: layout and visual treatment only. */
        .stApp {{
            background: {BG_PAGE} !important;
        }}
        [data-testid="stHeader"],
        .stApp > header {{
            background: {BG_PAGE} !important;
            border-bottom: 0 !important;
            box-shadow: none !important;
        }}
        [data-testid="stToolbar"],
        [data-testid="stDecoration"] {{
            background: transparent !important;
        }}
        .block-container {{
            padding: 72px 28px 40px !important;
            max-width: 1720px !important;
        }}
        [data-testid="stSidebar"] {{
            background: #FFFFFF !important;
            border-right: 1px solid {BORDER_DEFAULT};
        }}
        [data-testid="stSidebar"] > div:first-child {{
            padding: 22px 14px 24px !important;
        }}
        .sidebar-brand {{
            display: flex;
            align-items: center;
            gap: 10px;
            padding: 0 8px 22px;
            border-bottom: 1px solid {BORDER_LIGHT};
            margin-bottom: 20px;
        }}
        .sidebar-logo-dot {{
            width: 24px;
            height: 24px;
            border-radius: 999px;
            background: {COLOR_BLUE};
            box-shadow: none;
        }}
        .sidebar-brand-title {{
            color: {TEXT_PRIMARY};
            font-size: 18px;
            line-height: 1;
            font-weight: 700;
            letter-spacing: 0.02em;
        }}
        .sidebar-brand-sub {{
            color: #64748B;
            font-size: 11px;
            line-height: 1.3;
            font-weight: 600;
            margin-top: 4px;
        }}
        .sidebar-nav {{
            display: grid;
            gap: 4px;
            padding: 0 4px;
        }}
        .sidebar-nav-item {{
            position: relative;
            display: flex;
            align-items: center;
            height: 48px;
            padding: 0 16px 0 20px;
            border-radius: 10px;
            color: #374151 !important;
            text-decoration: none !important;
            font-size: 14px;
            line-height: 1;
            font-weight: 500;
            background: #FFFFFF;
            transition: background-color 0.15s ease;
        }}
        .sidebar-nav-item::before {{
            content: "";
            position: absolute;
            left: 0;
            top: 10px;
            bottom: 10px;
            width: 4px;
            border-radius: 999px;
            background: transparent;
        }}
        .sidebar-nav-item:hover {{
            background: #F8FAFC;
            color: #374151 !important;
            text-decoration: none !important;
        }}
        .sidebar-nav-item.active {{
            background: #EFF6FF;
            color: {COLOR_BLUE} !important;
            font-weight: 600;
        }}
        .sidebar-nav-item.active::before {{
            background: {COLOR_BLUE};
        }}
        .sidebar-nav-item span {{
            color: inherit;
            white-space: nowrap;
        }}
        .sidebar-section-title {{
            color: {COLOR_BLUE};
            font-size: 11px;
            line-height: 1.2;
            font-weight: 700;
            padding: 14px 8px 8px;
            border-top: 1px solid {BORDER_LIGHT};
            margin-top: 12px;
        }}
        .sidebar-muted-menu {{
            display: grid;
            gap: 4px;
            padding: 0 0 10px;
        }}
        .sidebar-muted-menu span {{
            display: flex;
            align-items: center;
            min-height: 36px;
            padding: 0 10px;
            border-radius: 8px;
            color: #64748B;
            font-size: 13px;
            font-weight: 700;
        }}
        .sidebar-filter-note {{
            margin: 0 8px;
            padding: 12px;
            border-radius: 8px;
            background: {BG_CARD};
            color: #64748B;
            font-size: 12px;
            line-height: 1.45;
            font-weight: 600;
            border: 1px solid {BORDER_LIGHT};
        }}
        .app-header {{
            margin-bottom: 18px;
            padding-top: 4px;
            overflow: visible;
        }}
        .app-title {{
            color: {TEXT_PRIMARY};
            font-size: 34px;
            line-height: 1.45;
            font-weight: 700;
            letter-spacing: 0;
            margin: 0 0 6px;
            min-height: 50px;
            overflow: visible;
        }}
        .app-basis {{
            color: {TEXT_SECONDARY};
            font-size: 13px;
            line-height: 1.5;
            font-weight: 500;
        }}
        .header-date {{
            color: {TEXT_SECONDARY};
            font-size: 12px;
            font-weight: 600;
            text-align: right;
            margin-bottom: 8px;
        }}
        .section-title {{
            color: {TEXT_PRIMARY};
            font-size: 18px;
            line-height: 1.35;
            font-weight: 700;
            margin: 0 0 6px;
            letter-spacing: 0;
        }}
        .section-title-row {{
            display: flex;
            align-items: baseline;
            justify-content: space-between;
            gap: 12px;
            margin: 0 0 4px;
        }}
        .section-title-row .section-title {{
            margin: 0;
        }}
        .section-view-link {{
            color: {COLOR_BLUE};
            font-size: 11px;
            line-height: 1.2;
            font-weight: 700;
            white-space: nowrap;
        }}
        .section-gap {{
            height: 32px;
        }}
        .section-sub {{
            color: {TEXT_SECONDARY};
            font-size: 12px;
            line-height: 1.4;
            font-weight: 500;
            margin-bottom: 16px;
            min-height: 17px;
            max-height: 17px;
            overflow: hidden;
            white-space: nowrap;
            text-overflow: ellipsis;
        }}
        [data-testid="stSegmentedControl"] {{
            margin: 0 0 24px !important;
        }}
        [data-testid="stSegmentedControl"] label {{
            color: {TEXT_SECONDARY} !important;
            font-size: 12px !important;
            font-weight: 600 !important;
            margin-bottom: 8px !important;
        }}
        [data-testid="stSegmentedControl"] div[role="radiogroup"] {{
            gap: 8px !important;
            background: {BG_CARD} !important;
        }}
        [data-testid="stSegmentedControl"] button {{
            min-height: 40px !important;
            border: 1px solid {BORDER_DEFAULT} !important;
            border-radius: 8px !important;
            background: {BG_CARD} !important;
            color: {TEXT_PRIMARY} !important;
            font-size: 14px !important;
            font-weight: 600 !important;
            padding: 9px 18px !important;
            box-shadow: none !important;
            transition: background 0.2s ease, border-color 0.2s ease, transform 0.2s ease !important;
        }}
        [data-testid="stSegmentedControl"] button:hover {{
            transform: translateY(-1px);
            background: {BG_CARD} !important;
            border-color: #CBD5E1 !important;
            color: {TEXT_PRIMARY} !important;
        }}
        [data-testid="stSegmentedControl"] button[aria-pressed="true"],
        [data-testid="stSegmentedControl"] button[aria-selected="true"],
        [data-testid="stSegmentedControl"] button[aria-checked="true"],
        [data-testid="stSegmentedControl"] [role="radio"][aria-checked="true"],
        [data-testid="stSegmentedControl"] [data-checked="true"],
        [data-testid="stSegmentedControl"] [data-selected="true"],
        [data-testid="stSegmentedControl"] label:has(input:checked),
        [data-testid="stSegmentedControl"] input:checked + div,
        [data-testid="stSegmentedControl"] input:checked ~ div {{
            background: {BG_CARD} !important;
            border-color: {COLOR_DANGER} !important;
            color: {COLOR_DANGER} !important;
            border-bottom-color: {COLOR_DANGER} !important;
            box-shadow: none !important;
        }}
        [data-testid="stSegmentedControl"] button[aria-pressed="true"] *,
        [data-testid="stSegmentedControl"] button[aria-selected="true"] *,
        [data-testid="stSegmentedControl"] button[aria-checked="true"] *,
        [data-testid="stSegmentedControl"] [role="radio"][aria-checked="true"] *,
        [data-testid="stSegmentedControl"] [data-checked="true"] *,
        [data-testid="stSegmentedControl"] [data-selected="true"] *,
        [data-testid="stSegmentedControl"] label:has(input:checked) *,
        [data-testid="stSegmentedControl"] input:checked + div *,
        [data-testid="stSegmentedControl"] input:checked ~ div * {{
            color: {COLOR_DANGER} !important;
        }}
        [data-testid="stSegmentedControl"] [data-baseweb="button-group"] button,
        [data-testid="stSegmentedControl"] [role="radio"],
        [data-testid="stSegmentedControl"] label[data-baseweb="radio"],
        [data-testid="stSegmentedControl"] label[data-baseweb="radio"] > div,
        [data-testid="stSegmentedControl"] label:has(input) > div {{
            background: {BG_CARD} !important;
            border-color: {BORDER_DEFAULT} !important;
            color: {TEXT_PRIMARY} !important;
            box-shadow: none !important;
        }}
        [data-testid="stSegmentedControl"] [data-baseweb="button-group"] button[aria-pressed="true"],
        [data-testid="stSegmentedControl"] [role="radio"][aria-checked="true"],
        [data-testid="stSegmentedControl"] label[data-baseweb="radio"]:has(input:checked),
        [data-testid="stSegmentedControl"] label[data-baseweb="radio"]:has(input:checked) > div,
        [data-testid="stSegmentedControl"] label:has(input:checked) > div {{
            background: {BG_CARD} !important;
            border-color: {COLOR_DANGER} !important;
            color: {COLOR_DANGER} !important;
        }}
        [data-testid="stSegmentedControl"] [data-baseweb="button-group"] button *,
        [data-testid="stSegmentedControl"] [role="radio"] *,
        [data-testid="stSegmentedControl"] label[data-baseweb="radio"] *,
        [data-testid="stSegmentedControl"] label:has(input) * {{
            color: {TEXT_PRIMARY} !important;
        }}
        [data-testid="stSegmentedControl"] [data-baseweb="button-group"] button[aria-pressed="true"] *,
        [data-testid="stSegmentedControl"] [role="radio"][aria-checked="true"] *,
        [data-testid="stSegmentedControl"] label[data-baseweb="radio"]:has(input:checked) *,
        [data-testid="stSegmentedControl"] label:has(input:checked) *,
        [data-testid="stSegmentedControl"] input:checked + div *,
        [data-testid="stSegmentedControl"] input:checked ~ div * {{
            color: {COLOR_DANGER} !important;
        }}
        [data-testid="stSegmentedControl"] [role="radio"][aria-checked="false"],
        [data-testid="stSegmentedControl"] [role="radio"]:not([aria-checked="true"]) {{
            border-color: {BORDER_DEFAULT} !important;
            background: {BG_CARD} !important;
            color: {TEXT_PRIMARY} !important;
        }}
        .dashboard-nav-divider {{
            height: 1px;
            background: {BORDER_DEFAULT};
            margin: -10px 0 24px;
        }}
        [data-testid="stTextInput"] input,
        [data-testid="stNumberInput"] input,
        [data-testid="stSelectbox"] > div > div,
        [data-testid="stMultiSelect"] > div > div {{
            min-height: 40px !important;
            border-radius: 8px !important;
            border: 1px solid {BORDER_DEFAULT} !important;
            background: {BG_CARD} !important;
            box-shadow: none !important;
        }}
        [data-testid="stMultiSelect"] [data-baseweb="tag"] {{
            background: {BG_CARD} !important;
            border: 1px solid {COLOR_ALERT_BD} !important;
            color: #993C1D !important;
        }}
        [data-testid="stTextInput"] label,
        [data-testid="stNumberInput"] label,
        [data-testid="stSelectbox"] label,
        [data-testid="stMultiSelect"] label,
        [data-testid="stCheckbox"] label,
        [data-testid="stRadio"] label {{
            color: {TEXT_SECONDARY} !important;
            font-size: 12px !important;
            font-weight: 600 !important;
        }}
        [data-testid="stButton"] button,
        [data-testid="stDownloadButton"] button {{
            height: 40px !important;
            min-height: 40px !important;
            padding: 0 18px !important;
            border-radius: 10px !important;
            border: 1px solid {BORDER_DEFAULT} !important;
            background: {BG_CARD} !important;
            color: {TEXT_PRIMARY} !important;
            font-size: 14px !important;
            font-weight: 600 !important;
            box-shadow: none !important;
            transition: transform 0.2s ease, background 0.2s ease, border-color 0.2s ease !important;
        }}
        [data-testid="stButton"] button:hover,
        [data-testid="stDownloadButton"] button:hover {{
            transform: translateY(-2px);
            background: #FFFFFF !important;
            border-color: #CBD5E1 !important;
            color: {TEXT_PRIMARY} !important;
        }}
        .kpi-dashboard-block {{
            margin: 0;
        }}
        .kpi-dashboard-label {{
            color: {TEXT_PRIMARY};
            font-size: 14px;
            font-weight: 700;
            margin: 4px 0 10px;
        }}
        .kpi-dashboard-block.status-board {{
            display: block;
            grid-template-columns: none;
            gap: 0;
            align-items: initial;
        }}
        .kpi-dashboard-head {{
            display: flex;
            align-items: flex-start;
            justify-content: space-between;
            gap: 16px;
            margin-bottom: 12px;
        }}
        .overall-kpi-card,
        .kpi-panel,
        .panel-box,
        .mini-kpi-card,
        [data-testid="metric-container"],
        [data-testid="stDataFrame"] {{
            background: {BG_CARD} !important;
            border: 1px solid {BORDER_DEFAULT} !important;
            border-radius: 12px !important;
            box-shadow: 0 2px 8px rgba(0,0,0,0.04) !important;
        }}
        .overall-kpi-card {{
            padding: 20px 24px;
            height: var(--kpi-card-height);
            min-height: var(--kpi-card-height);
            max-height: var(--kpi-card-height);
            border-color: #93B5FF !important;
            transition: transform 0.2s ease;
        }}
        .overall-kpi-card:hover,
        .kpi-panel:hover,
        .panel-box:hover {{
            transform: translateY(-2px);
        }}
        .overall-kpi-title,
        .kpi-panel-head {{
            display: flex;
            align-items: baseline;
            justify-content: space-between;
            gap: 12px;
            margin-bottom: 14px;
        }}
        .overall-kpi-title span,
        .kpi-title {{
            color: {COLOR_BLUE};
            font-size: 15px;
            line-height: 1.35;
            font-weight: 700;
            margin: 0;
        }}
        .scope-kpi .kpi-title {{
            color: #4B5563;
        }}
        .scope-kpi.main-kpi .kpi-title {{
            color: #4B5563;
        }}
        .scope-kpi.sample-kpi .kpi-title {{
            color: #4B5563;
        }}
        .overall-kpi-title b {{
            color: {TEXT_SECONDARY};
            font-size: 11px;
            font-weight: 600;
        }}
        .scope-kpi .metric-value,
        .scope-kpi .metric-value.primary,
        .scope-kpi .metric-value.warning,
        .scope-kpi .metric-value.warn,
        .scope-kpi .metric-value.mid,
        .scope-kpi .metric-value.purple,
        .scope-kpi .metric-value.danger,
        .scope-kpi .metric-value.risk {{
            color: {TEXT_PRIMARY};
        }}
        .kpi-rating.large {{
            font-size: 13px;
        }}
        .kpi-divider-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(130px, 1fr));
            gap: 0;
            border-top: 0;
            border-bottom: 0;
        }}
        .kpi-divider-grid.overall {{
            grid-template-columns: repeat(5, minmax(0, 1fr));
        }}
        .kpi-metric {{
            min-width: 0;
            padding: 12px 16px;
            border-left: 1px solid #D1D5DB;
        }}
        .kpi-metric:first-child {{
            border-left: 0;
        }}
        .metric-label {{
            color: {TEXT_TERTIARY};
            font-size: 11px;
            line-height: 1.35;
            font-weight: 700;
            margin-bottom: 8px;
        }}
        .metric-value {{
            color: {TEXT_PRIMARY};
            font-size: 22px;
            line-height: 1.1;
            font-weight: 700;
            font-variant-numeric: tabular-nums;
            white-space: nowrap;
        }}
        .metric-value.primary {{
            color: {COLOR_BLUE};
        }}
        .metric-value.good {{
            color: {COLOR_TEAL};
        }}
        .metric-value.warning,
        .metric-value.warn,
        .metric-value.mid {{
            color: {COLOR_ORANGE};
        }}
        .metric-value.purple {{
            color: {COLOR_AMBER};
        }}
        .metric-value.danger,
        .metric-value.risk {{
            color: {COLOR_DANGER};
        }}
        .overall-kpi-card .metric-value {{
            font-size: 26px;
        }}
        .metric-value.quantity {{
            display: flex;
            align-items: baseline;
            gap: 4px;
            min-width: 0;
            white-space: nowrap;
        }}
        .metric-value.quantity .metric-number {{
            min-width: 0;
        }}
        .metric-value.quantity .metric-unit {{
            color: inherit;
            font-size: 0.68em;
            font-weight: 700;
            letter-spacing: 0;
        }}
        .metric-subvalue {{
            color: {TEXT_TERTIARY};
            font-size: 12px;
            line-height: 1.25;
            font-weight: 700;
            margin-top: 6px;
            white-space: nowrap;
            font-variant-numeric: tabular-nums;
        }}
        .overall-kpi-card .metric-value.quantity {{
            font-size: clamp(18px, 1.3vw, 22px);
        }}
        .scope-kpi .metric-value.quantity {{
            font-size: clamp(15px, 1vw, 18px);
        }}
        .scope-kpi .metric-subvalue {{
            font-size: 11px;
        }}
        .kpi-progress-stack {{
            display: grid;
            gap: 8px;
            margin-top: 14px;
        }}
        .kpi-progress-row {{
            display: grid;
            grid-template-columns: 44px minmax(100px, 1fr) 54px;
            align-items: center;
            gap: 8px;
        }}
        .kpi-progress-row span,
        .kpi-progress-row b {{
            color: {TEXT_SECONDARY};
            font-size: 11px;
            font-weight: 700;
            font-variant-numeric: tabular-nums;
        }}
        .kpi-progress-row b {{
            text-align: right;
        }}
        .overall-kpi-card .kpi-progress-row b {{
            color: {TEXT_PRIMARY};
            font-size: 15px;
            font-weight: 700;
        }}
        .kpi-progress-track {{
            height: 4px;
            border-radius: 999px;
            background: {BORDER_DEFAULT};
            overflow: hidden;
        }}
        .kpi-progress-fill {{
            height: 100%;
            border-radius: 999px;
        }}
        .kpi-progress-fill.production,
        .progress-fill.production {{
            background: {COLOR_BLUE};
        }}
        .kpi-progress-fill.packing,
        .progress-fill.packing,
        .progress-fill.warn,
        .progress-fill.risk {{
            background: {COLOR_ORANGE};
        }}
        .kpi-progress-fill.receipt,
        .progress-fill.receipt {{
            background: {COLOR_AMBER};
        }}
        .progress-fill.good.production,
        .progress-fill.mid.production,
        .progress-fill.warn.production,
        .progress-fill.risk.production {{
            background: {COLOR_BLUE};
        }}
        .progress-fill.good.packing,
        .progress-fill.mid.packing,
        .progress-fill.warn.packing,
        .progress-fill.risk.packing {{
            background: {COLOR_ORANGE};
        }}
        .progress-fill.good.receipt,
        .progress-fill.mid.receipt,
        .progress-fill.warn.receipt,
        .progress-fill.risk.receipt {{
            background: {COLOR_AMBER};
        }}
        .progress-fill.done,
        .progress-fill.active {{
            background: {COLOR_ORANGE};
        }}
        .scope-kpi .kpi-progress-fill.production,
        .scope-kpi .kpi-progress-fill.packing,
        .scope-kpi .kpi-progress-fill.receipt {{
            background: #64748B;
        }}
        .overall-kpi-foot {{
            display: flex;
            flex-wrap: wrap;
            gap: 10px 18px;
            margin-top: 12px;
            color: {TEXT_SECONDARY};
            font-size: 12px;
            font-weight: 700;
            font-variant-numeric: tabular-nums;
        }}
        .scope-kpi {{
            padding: 16px 20px;
            margin-bottom: 0;
            height: var(--kpi-card-height);
            min-height: var(--kpi-card-height);
            max-height: var(--kpi-card-height);
            transition: transform 0.2s ease;
        }}
        .scope-kpi .kpi-divider-grid {{
            grid-template-columns: repeat(2, minmax(0, 1fr));
        }}
        .scope-kpi .kpi-metric {{
            padding: 6px 14px;
        }}
        .scope-kpi .metric-label {{
            margin-bottom: 5px;
        }}
        .scope-kpi .metric-value {{
            font-size: 20px;
        }}
        .scope-kpi .kpi-progress-stack {{
            gap: 6px;
            margin-top: 10px;
        }}
        .scope-kpi .kpi-progress-row {{
            grid-template-columns: 44px minmax(80px, 1fr) 44px;
            gap: 8px;
        }}
        .metric-strip {{
            padding: 0;
            overflow: hidden;
            margin-bottom: 24px;
        }}
        .metric-strip .kpi-divider-grid {{
            border-top: 0;
            border-bottom: 0;
        }}
        .mini-kpi-grid {{
            gap: 16px;
            margin-bottom: 24px;
        }}
        .mini-kpi-card {{
            padding: 20px;
            min-height: 92px;
        }}
        .production-dialog-summary {{
            border: 1px solid {BORDER_DEFAULT};
            border-radius: 12px;
            background: #FFFFFF;
            padding: 18px;
            margin: 2px 0 14px;
            box-shadow: 0 2px 8px rgba(15, 23, 42, 0.04);
        }}
        .production-dialog-head {{
            display: flex;
            align-items: flex-start;
            justify-content: space-between;
            gap: 16px;
            margin-bottom: 6px;
        }}
        .production-dialog-code {{
            color: {COLOR_BLUE};
            font-size: 18px;
            line-height: 1.2;
            font-weight: 700;
            letter-spacing: 0;
        }}
        .production-dialog-product {{
            color: {TEXT_PRIMARY};
            font-size: 14px;
            line-height: 1.35;
            font-weight: 700;
            margin-top: 4px;
        }}
        .production-dialog-count {{
            color: {TEXT_SECONDARY};
            font-size: 12px;
            font-weight: 600;
            margin-bottom: 14px;
        }}
        .production-dialog-metrics {{
            display: grid;
            grid-template-columns: repeat(6, minmax(0, 1fr));
            border-top: 1px solid {BORDER_LIGHT};
        }}
        .production-dialog-metric {{
            padding: 14px 14px 4px;
            border-right: 1px solid {BORDER_LIGHT};
            min-width: 0;
        }}
        .production-dialog-metric:last-child {{
            border-right: 0;
        }}
        .production-dialog-metric p {{
            margin: 0 0 7px;
            color: {TEXT_SECONDARY};
            font-size: 12px;
            line-height: 1.2;
            font-weight: 700;
        }}
        .production-dialog-metric strong {{
            display: block;
            min-width: 0;
            color: {TEXT_PRIMARY};
            font-size: 16px;
            line-height: 1.25;
            font-weight: 700;
            letter-spacing: 0;
            font-variant-numeric: tabular-nums;
            white-space: nowrap;
            overflow: hidden;
            text-overflow: ellipsis;
        }}
        .production-dialog-metric strong.primary {{
            color: {COLOR_BLUE};
        }}
        .production-dialog-metric strong.warning {{
            color: {COLOR_ORANGE};
        }}
        .production-dialog-metric strong.danger {{
            color: {COLOR_DANGER};
        }}
        .production-dialog-metric span {{
            display: block;
            margin-top: 5px;
            color: {TEXT_TERTIARY};
            font-size: 11px;
            line-height: 1.2;
            font-weight: 700;
            white-space: nowrap;
            font-variant-numeric: tabular-nums;
        }}
        .production-dialog-section-title {{
            color: {TEXT_PRIMARY};
            font-size: 14px;
            font-weight: 700;
            margin: 6px 0 8px;
        }}
        .panel-box {{
            padding: 18px 20px;
            margin-bottom: 24px;
        }}
        .dashboard-card {{
            height: auto;
            min-height: 0;
            max-height: none;
            margin-bottom: 0;
            overflow: hidden;
        }}
        .lower-card-control-spacer {{
            height: 64px;
            min-height: 64px;
            max-height: 64px;
        }}
        .family-progress-panel {{
            height: auto;
            min-height: 0;
            max-height: none;
            padding: 24px;
        }}
        .family-table {{
            display: grid;
            gap: 0;
        }}
        .family-table-row {{
            display: grid;
            grid-template-columns: minmax(118px, 1.16fr) minmax(126px, 0.94fr) minmax(112px, 1fr) minmax(112px, 1fr) minmax(112px, 1fr) minmax(86px, 0.82fr);
            align-items: center;
            gap: 14px;
            height: 56px;
            min-height: 56px;
            max-height: 56px;
            border-bottom: 1px solid {BORDER_LIGHT};
            overflow: hidden;
        }}
        .family-table-row:last-child {{
            border-bottom: 0;
        }}
        .family-table-head {{
            height: 40px;
            min-height: 40px;
            max-height: 40px;
            color: #64748B;
            font-size: 13px;
            line-height: 1.2;
            font-weight: 700;
        }}
        .family-name {{
            display: flex;
            align-items: center;
            gap: 10px;
            min-width: 0;
            color: {TEXT_PRIMARY};
            font-size: 14px;
            font-weight: 700;
        }}
        .family-name b {{
            min-width: 0;
            overflow: hidden;
            white-space: nowrap;
            text-overflow: ellipsis;
        }}
        .family-dot {{
            width: 9px;
            height: 9px;
            border-radius: 999px;
            flex: 0 0 auto;
            background: {COLOR_BLUE};
        }}
        .family-dot.dot-1 {{ background: {COLOR_DANGER}; }}
        .family-dot.dot-2 {{ background: {COLOR_TEAL}; }}
        .family-dot.dot-3 {{ background: #334155; }}
        .family-dot.dot-4 {{ background: {COLOR_DANGER}; }}
        .family-dot.dot-5 {{ background: {COLOR_ORANGE}; }}
        .family-dot.dot-6 {{ background: {COLOR_AMBER}; }}
        .family-dot.dot-7 {{ background: {COLOR_AMBER}; }}
        .family-num {{
            color: {TEXT_PRIMARY};
            font-size: 14px;
            font-weight: 700;
            text-align: right;
            font-variant-numeric: tabular-nums;
            white-space: nowrap;
        }}
        .family-request {{
            display: flex;
            flex-direction: column;
            align-items: flex-end;
            gap: 3px;
            min-width: 0;
            text-align: right;
            font-variant-numeric: tabular-nums;
        }}
        .family-request strong {{
            color: {TEXT_PRIMARY};
            font-size: 14px;
            line-height: 1.05;
            font-weight: 700;
            white-space: nowrap;
        }}
        .family-request span {{
            color: {TEXT_MUTED};
            font-size: 11px;
            line-height: 1.05;
            font-weight: 600;
            white-space: nowrap;
        }}
        .family-request-head {{
            color: #64748B;
            font-size: 13px;
            font-weight: 700;
            text-align: right;
            white-space: nowrap;
        }}
        .family-num.shortage.normal {{
            color: {TEXT_PRIMARY};
        }}
        .family-num.shortage.danger {{
            color: {COLOR_DANGER};
        }}
        .family-progress-metric {{
            display: grid;
            grid-template-columns: minmax(52px, 1fr) 46px;
            align-items: center;
            gap: 9px;
            min-width: 0;
        }}
        .family-progress-track {{
            height: 4px;
            border-radius: 999px;
            background: {BORDER_DEFAULT};
            overflow: hidden;
        }}
        .family-progress-fill {{
            height: 100%;
            border-radius: 999px;
        }}
        .family-progress-fill.production {{
            background: {COLOR_BLUE};
        }}
        .family-progress-fill.packing {{
            background: {COLOR_ORANGE};
        }}
        .family-progress-fill.receipt {{
            background: {COLOR_AMBER};
        }}
        .family-progress-metric span {{
            color: {TEXT_SECONDARY};
            font-size: 13px;
            font-weight: 700;
            text-align: right;
            font-variant-numeric: tabular-nums;
        }}
        .family-section {{
            gap: 12px;
        }}
        .family-section + .family-section {{
            margin-top: 24px;
        }}
        .family-section-title {{
            background: #EEF2FF;
            color: {COLOR_BLUE};
            border-radius: 999px;
            font-size: 12px;
            font-weight: 700;
            padding: 6px 12px;
        }}
        .family-grid {{
            gap: 12px;
        }}
        .family-card {{
            border: 1px solid {BORDER_DEFAULT};
            border-radius: 12px;
            padding: 16px;
            min-height: 128px;
            gap: 12px;
            box-shadow: none;
            transition: transform 0.2s ease, border-color 0.2s ease;
        }}
        .family-card:hover {{
            transform: translateY(-2px);
            border-color: #CBD5E1;
        }}
        .family-shortages b,
        .ops-table td.num.shortage,
        .ops-table td.num.negative {{
            color: {COLOR_DANGER};
            font-weight: 700;
        }}
        .table-wrap {{
            border: 1px solid {BORDER_DEFAULT};
            border-radius: 12px;
            background: {BG_CARD};
            box-shadow: 0 2px 8px rgba(0,0,0,0.04);
        }}
        .ops-table th {{
            position: sticky;
            top: 0;
            z-index: 2;
            background: #FFFFFF;
            color: {TEXT_SECONDARY};
            font-size: 12px;
            font-weight: 600;
            padding: 12px 14px;
            border-bottom: 1px solid {BORDER_DEFAULT};
        }}
        .ops-table td {{
            color: {TEXT_PRIMARY};
            font-size: 13px;
            font-weight: 500;
            padding: 12px 14px;
            border-bottom: 1px solid {BORDER_DEFAULT};
            height: 48px;
        }}
        .ops-table tbody tr:hover td {{
            background: #F8FAFC;
        }}
        .progress-track {{
            height: 4px;
            border-radius: 999px;
            background: {BORDER_DEFAULT};
        }}
        .progress-fill {{
            border-radius: 999px;
            background: {COLOR_ORANGE};
        }}
        .progress-text {{
            color: {TEXT_SECONDARY};
            font-size: 11px;
            font-weight: 600;
        }}
        .top-row,
        .gap-row {{
            border: 1px solid {BORDER_DEFAULT};
            border-radius: 10px;
            padding: 12px 16px;
            margin-bottom: 8px;
            background: {BG_CARD};
        }}
        .top-shortage,
        .gap-value {{
            color: {COLOR_DANGER};
            font-weight: 700;
        }}
        .rank-list {{
            display: grid;
            gap: 10px;
        }}
        .rank-list-row {{
            display: grid;
            grid-template-columns: 22px minmax(92px, 1fr) minmax(52px, 0.75fr) 58px;
            align-items: center;
            gap: 10px;
            height: 34px;
            min-height: 34px;
            max-height: 34px;
            overflow: hidden;
        }}
        .rank-num {{
            color: #64748B;
            font-size: 12px;
            font-weight: 700;
            text-align: right;
            font-variant-numeric: tabular-nums;
        }}
        .rank-name {{
            color: {TEXT_PRIMARY};
            font-size: 12px;
            line-height: 1.25;
            font-weight: 700;
            overflow: hidden;
            white-space: nowrap;
            text-overflow: ellipsis;
        }}
        .rank-bar {{
            height: 4px;
            border-radius: 999px;
            background: {BORDER_DEFAULT};
            overflow: hidden;
        }}
        .rank-bar i {{
            display: block;
            height: 100%;
            border-radius: 999px;
            background: #64748B;
        }}
        .rank-value {{
            color: {COLOR_DANGER};
            font-size: 12px;
            font-weight: 700;
            text-align: right;
            font-variant-numeric: tabular-nums;
        }}
        .urgent-list {{
            display: grid;
            gap: 0;
        }}
        .urgent-list-row {{
            display: grid;
            grid-template-columns: 42px 58px minmax(0, 1fr) 44px;
            align-items: center;
            gap: 8px;
            height: 42px;
            min-height: 42px;
            max-height: 42px;
            padding: 0 2px;
            border-bottom: 1px solid {BORDER_LIGHT};
            overflow: hidden;
        }}
        .urgent-list-row:last-child {{
            border-bottom: 0;
        }}
        .urgent-code {{
            color: {COLOR_BLUE};
            font-size: 12px;
            font-weight: 700;
            font-variant-numeric: tabular-nums;
        }}
        .urgent-sku {{
            color: #64748B;
            font-size: 11px;
            font-weight: 700;
            text-align: right;
            white-space: nowrap;
        }}
        .urgent-product {{
            color: #64748B;
            font-size: 11px;
            line-height: 1.2;
            font-weight: 700;
            overflow: hidden;
            white-space: nowrap;
            text-overflow: ellipsis;
        }}
        .request-scope-badge,
        .response-badge,
        .status-badge {{
            border-radius: 999px;
            padding: 3px 9px;
            font-size: 11px;
            font-weight: 700;
            border: 1px solid transparent;
        }}
        .response-badge {{
            background: {BG_CARD};
            color: #475569;
            min-width: 56px;
            border-color: {BORDER_DEFAULT};
        }}
        .response-badge.partial,
        .request-scope-badge.mixed {{
            background: {BG_CARD};
            color: #9A3412;
            border-color: #FED7AA;
        }}
        .response-badge.need,
        .request-scope-badge.out,
        .status-badge.warn,
        .status-badge.risk {{
            background: {BG_CARD};
            color: {COLOR_DANGER};
            border-color: #FECACA;
        }}
        .request-scope-badge.in,
        .status-badge.done {{
            background: {BG_CARD};
            color: #475569;
            border-color: #CBD5E1;
        }}
        .status-badge.active {{
            background: {BG_CARD};
            color: {COLOR_BLUE};
            border-color: #BFDBFE;
        }}
        .breadcrumb {{
            margin: 4px 0 16px;
        }}
        hr {{
            border-color: {BORDER_DEFAULT} !important;
            margin: 24px 0 !important;
        }}
        @media (max-width: 1200px) {{
            .kpi-divider-grid.overall {{
                grid-template-columns: repeat(3, minmax(0, 1fr));
            }}
            .kpi-metric:nth-child(4) {{
                border-left: 0;
            }}
        }}
        @media (max-width: 760px) {{
            .block-container {{
                padding: 20px 16px 32px !important;
            }}
            .app-title {{
                font-size: 28px;
            }}
            .kpi-divider-grid,
            .kpi-divider-grid.overall {{
                grid-template-columns: 1fr;
            }}
            .kpi-metric,
            .kpi-metric:nth-child(4) {{
                border-left: 0;
                border-top: 1px solid {BORDER_LIGHT};
            }}
            .kpi-metric:first-child {{
                border-top: 0;
            }}
            .kpi-progress-row {{
                grid-template-columns: 56px minmax(100px, 1fr) 52px;
            }}
        }}
        /* Enterprise table reset: remove legacy tinted cell and empty-area backgrounds. */
        .table-wrap,
        .ops-table,
        .ops-table thead,
        .ops-table tbody,
        .ops-table tr,
        .ops-table th,
        .ops-table td,
        .ops-table tbody tr:nth-child(odd) td,
        .ops-table tbody tr:nth-child(even) td {{
            background: #FFFFFF !important;
            background-color: #FFFFFF !important;
        }}
        .ops-table th {{
            background: #FFFFFF !important;
            background-color: #FFFFFF !important;
            font-weight: 600 !important;
            border-bottom: 1px solid #E5E7EB !important;
        }}
        .ops-table td {{
            height: 48px !important;
            border-bottom: 1px solid #E5E7EB !important;
        }}
        .ops-table tbody tr:hover td {{
            background: #F8FAFC !important;
        }}
        [data-testid="stDataFrame"],
        [data-testid="stDataFrame"] > div,
        [data-testid="stDataFrame"] div[data-testid],
        [data-testid="stDataFrame"] [class*="empty"],
        [data-testid="stDataFrame"] [class*="Empty"],
        [data-testid="stDataFrame"] [class*="blank"],
        [data-testid="stDataFrame"] [class*="Blank"],
        [data-testid="stDataFrame"] [data-testid="stElementToolbar"],
        [data-testid="stDataFrame"] [data-testid="stTable"] {{
            background: #FFFFFF !important;
            background-color: #FFFFFF !important;
        }}
        [data-testid="stDataFrame"] th,
        [data-testid="stDataFrame"] [role="columnheader"] {{
            background: #FFFFFF !important;
            background-color: #FFFFFF !important;
            font-weight: 600 !important;
            color: #374151 !important;
            border-bottom: 1px solid #E5E7EB !important;
            justify-content: center !important;
        }}
        [data-testid="stDataFrame"] td,
        [data-testid="stDataFrame"] [role="gridcell"] {{
            min-height: 48px !important;
            border-bottom: 1px solid #E5E7EB !important;
        }}
        [data-testid="stDataFrame"] [role="row"]:nth-child(even) [role="gridcell"] {{
            background: #FCFCFD !important;
            background-color: #FCFCFD !important;
        }}
        [data-testid="stDataFrame"] tr:hover td,
        [data-testid="stDataFrame"] [role="row"]:hover [role="gridcell"] {{
            background: #F8FAFC !important;
            background-color: #F8FAFC !important;
            cursor: pointer !important;
        }}
        [data-testid="stDataFrame"] [role="row"]:hover [role="gridcell"]:first-child {{
            box-shadow: inset 4px 0 0 #2563EB !important;
        }}
        [data-testid="stDataFrame"] [aria-selected="true"] [role="gridcell"],
        [data-testid="stDataFrame"] [role="row"][aria-selected="true"] [role="gridcell"] {{
            background: #EFF6FF !important;
            background-color: #EFF6FF !important;
        }}
        </style>
        """,
        unsafe_allow_html=True,
    )


def render_panel_title(title: str, sub: str, link_label: str = "") -> None:
    link_html = f"<div class='section-view-link'>{escape(link_label)}</div>" if link_label else ""
    st.markdown(
        "<div class='section-title-row'>"
        f"<div class='section-title'>{escape(title)}</div>"
        f"{link_html}"
        "</div>"
        f"<div class='section-sub'>{escape(sub)}</div>",
        unsafe_allow_html=True,
    )


def build_sales_code_view(code_summary: pd.DataFrame) -> pd.DataFrame:
    view = code_summary.rename(
        columns={
            "sales_code": "판매코드",
            "product_name": "제품명",
            "request_pack": "요청 PACK",
            "request_pcs": "요청 PCS",
            "packing_pack": "포장 PACK",
            "production_code": "생산코드",
            "q_code": "분리코드",
            "r_code": "사출코드",
            "production_basis_qty": "누수규격검사 생산수량",
            "production_shortage_qty": "생산부족수량",
            "production_progress_pct": "생산진도율",
        }
    )
    view = finalize_summary(view)
    return view


def format_sales_code_view(view: pd.DataFrame) -> pd.DataFrame:
    out = view.copy()
    for col in ["요청 PACK", "포장 PACK", "부족 PACK"]:
        out[col] = out[col].map(format_int)
    out["생산진도율"] = out["생산진도율"].map(lambda x: f"{float(x):.1f}%")
    out["용마입고율"] = out["용마입고율"].map(lambda x: f"{float(x):.1f}%")
    return out[
        [
            "판매코드",
            "요청 PACK",
            "포장 PACK",
            "부족 PACK",
            "생산진도율",
            "용마입고율",
        ]
    ]


def format_production_code_view(view: pd.DataFrame) -> pd.DataFrame:
    out = view.copy()
    for col in ["요청 PACK", "생산부족수량", "포장부족수량"]:
        out[col] = out[col].map(format_int)
    out["생산진도율"] = out["생산진도율"].map(lambda x: f"{float(x):.1f}%")
    out["포장진도율"] = out["포장진도율"].map(lambda x: f"{float(x):.1f}%")
    if "생산완료예상일" not in out.columns and "납기일" in out.columns:
        out["생산완료예상일"] = out["납기일"]
    out["생산완료예상일"] = out["생산완료예상일"].map(format_date)
    return out[
        [
            "생산코드",
            "연결 판매코드 수",
            "제품명",
            "요청 PACK",
            "생산진도율",
            "생산부족수량",
            "포장진도율",
            "포장부족수량",
            "생산완료예상일",
        ]
    ]


def calc_drilldown_kpi(df: pd.DataFrame) -> dict[str, float]:
    if df.empty:
        return {
            "request_qty": 0.0,
            "request_pcs": 0.0,
            "production_shortage_qty": 0.0,
            "production_progress_pct": 0.0,
            "packing_qty": 0.0,
            "shortage_qty": 0.0,
        }
    request_qty = float(df["요청 PACK"].sum()) if "요청 PACK" in df.columns else float(df["요청수량"].sum())
    if "요청 PCS" in df.columns:
        request_pcs = float(df["요청 PCS"].sum())
    elif "요청PCS" in df.columns:
        request_pcs = float(df["요청PCS"].sum())
    else:
        request_pcs = request_qty
    production_shortage_qty = float(df["생산부족수량"].sum()) if "생산부족수량" in df.columns else 0.0
    production_progress_pct = (
        (request_pcs - production_shortage_qty) / request_pcs * 100.0
        if request_pcs > 0
        else 0.0
    )
    return {
        "request_qty": request_qty,
        "request_pcs": request_pcs,
        "production_shortage_qty": production_shortage_qty,
        "production_progress_pct": max(0.0, min(100.0, production_progress_pct)),
        "packing_qty": float(df["포장 PACK"].sum()) if "포장 PACK" in df.columns else float(df["포장수량"].sum()),
        "shortage_qty": float(df["부족 PACK"].sum()) if "부족 PACK" in df.columns else float(df["부족수량"].sum()),
    }


def render_drilldown_kpi(kpi: dict[str, float]) -> None:
    metrics = [
        ("요청 PACK", format_int(kpi["request_qty"]), "normal"),
        ("생산부족수량", format_int(kpi["production_shortage_qty"]), "danger" if kpi["production_shortage_qty"] > 0 else "normal"),
        ("생산진도율", f"{kpi['production_progress_pct']:.1f}%", "primary"),
        ("포장수량", format_int(kpi["packing_qty"]), "normal"),
        ("포장부족수량", format_int(kpi["shortage_qty"]), "danger" if kpi["shortage_qty"] > 0 else "normal"),
    ]
    metric_html = "".join(kpi_metric_item_html(label, value, tone) for label, value, tone in metrics)
    panel_html = f"""
    <div class='kpi-panel drill-kpi'>
      <div class='kpi-divider-grid'>{metric_html}</div>
      <div class='kpi-progress-stack'>
        {kpi_progress_line_html("생산", kpi["production_progress_pct"], "production")}
      </div>
    </div>
    """
    st.markdown(panel_html, unsafe_allow_html=True)


def build_product_drilldown_view(product_summary: pd.DataFrame) -> pd.DataFrame:
    out = product_summary.copy()
    out = out.rename(
        columns={
            "요청 PACK": "요청수량",
            "포장 PACK": "포장수량",
            "포장부족수량": "포장부족수량",
        }
    )
    return out[
        ["제품명", "요청수량", "생산진도율", "포장진도율", "생산부족수량", "포장부족수량", "상태"]
    ]


def build_pack_unit_view(code_summary: pd.DataFrame, product_name: str) -> pd.DataFrame:
    if code_summary.empty:
        return pd.DataFrame(columns=["팩 단위", "요청수량", "포장수량", "부족수량", "진도율", "_sort"])

    selected_base = strip_pack_unit_suffix(product_name)
    work = code_summary.copy()
    if "base_product_name" not in work.columns:
        work["base_product_name"] = work["product_name"].map(strip_pack_unit_suffix)
    if "pack_unit" not in work.columns:
        work["pack_unit"] = work["product_name"].map(extract_pack_unit)
    if "pack_unit_label" not in work.columns:
        work["pack_unit_label"] = [
            format_pack_unit_label(unit, name)
            for unit, name in zip(work["pack_unit"], work["product_name"])
        ]

    scope = work[work["base_product_name"] == selected_base].copy()
    if scope.empty:
        scope = work[work["product_name"] == product_name].copy()
    if scope.empty:
        return pd.DataFrame(columns=["팩 단위", "요청수량", "포장수량", "부족수량", "진도율", "_sort"])

    grouped = (
        scope.groupby(["pack_unit", "pack_unit_label"], dropna=False)
        .agg(
            request_qty=("request_pack", "sum"),
            packing_qty=("packing_pack", "sum"),
        )
        .reset_index()
    )
    grouped["shortage_qty"] = (grouped["request_qty"] - grouped["packing_qty"]).clip(lower=0.0)
    grouped["progress_pct"] = np.where(
        grouped["request_qty"] > 0,
        grouped["packing_qty"] / grouped["request_qty"] * 100.0,
        0.0,
    )
    grouped["progress_pct"] = np.clip(grouped["progress_pct"], 0.0, 100.0)
    grouped["_sort"] = pd.to_numeric(grouped["pack_unit"], errors="coerce").fillna(999999.0)
    grouped = grouped.sort_values("_sort", kind="stable")

    out = grouped.rename(
        columns={
            "pack_unit_label": "팩 단위",
            "request_qty": "요청수량",
            "packing_qty": "포장수량",
            "shortage_qty": "부족수량",
            "progress_pct": "진도율",
        }
    )[["팩 단위", "요청수량", "포장수량", "부족수량", "진도율", "_sort"]]

    total_request = float(out["요청수량"].sum())
    total_packing = float(out["포장수량"].sum())
    total_shortage = max(0.0, total_request - total_packing)
    total_progress = (total_packing / total_request * 100.0) if total_request > 0 else 0.0
    total_row = pd.DataFrame(
        [
            {
                "팩 단위": "전체",
                "요청수량": total_request,
                "포장수량": total_packing,
                "부족수량": total_shortage,
                "진도율": min(100.0, max(0.0, total_progress)),
                "_sort": 1000000.0,
            }
        ]
    )
    return pd.concat([out, total_row], ignore_index=True)


def pack_unit_column_config() -> dict[str, Any]:
    numeric_format = "%,.0f"
    return {
        "요청수량": st.column_config.NumberColumn("요청수량", format=numeric_format),
        "포장수량": st.column_config.NumberColumn("포장수량", format=numeric_format),
        "부족수량": st.column_config.NumberColumn("부족수량", format=numeric_format),
        "진도율": st.column_config.ProgressColumn(
            "진도율",
            min_value=0,
            max_value=100,
            format="%.2f%%",
            color=COLOR_ORANGE,
        ),
        "_sort": None,
    }


def build_production_drilldown_view(code_summary: pd.DataFrame) -> pd.DataFrame:
    out = build_production_code_view(code_summary).rename(
        columns={
            "요청 PACK": "요청수량",
            "포장 PACK": "포장수량",
        }
    )
    return out[
        [
            "생산코드",
            "요청수량",
            "생산부족수량",
            "생산진도율",
            "포장부족수량",
            "연결 판매코드 수",
            "포장진도율",
            "상태",
        ]
    ]


def build_sales_drilldown_view(code_summary: pd.DataFrame) -> pd.DataFrame:
    out = build_sales_code_view(code_summary).rename(
        columns={
            "요청 PACK": "요청수량",
            "포장 PACK": "포장수량",
            "부족 PACK": "부족수량",
        }
    )
    return out[
        [
            "판매코드",
            "요청수량",
            "포장수량",
            "부족수량",
            "생산진도율",
            "용마입고율",
            "생산코드",
            "상태",
        ]
    ]


def build_power_drilldown_view(code_summary: pd.DataFrame) -> pd.DataFrame:
    out = build_power_detail(code_summary).rename(
        columns={
            "요청수량": "요청수량",
            "포장수량": "포장수량",
            "부족수량": "부족수량",
            "진도율": "포장진도율",
        }
    )
    if out.empty:
        return pd.DataFrame(
            columns=["POWER", "요청수량", "포장수량", "부족수량", "생산진도율", "포장진도율", "power_value"]
        )
    out = out.sort_values("power_value", ascending=True, kind="stable")
    return out[
        ["POWER", "요청수량", "포장수량", "부족수량", "생산진도율", "포장진도율", "power_value"]
    ]


def drilldown_column_config() -> dict[str, Any]:
    numeric_format = "%,.0f"
    return {
        "품목코드": st.column_config.TextColumn("판매코드"),
        "S코드": st.column_config.TextColumn("판매코드"),
        "CP": st.column_config.TextColumn("CP"),
        "요청합계(PACK)": st.column_config.NumberColumn("요청합계(PACK)", format=numeric_format),
        "요청합계(PCS)": st.column_config.NumberColumn("요청합계(PCS)", format=numeric_format),
        "생산요청물량": st.column_config.NumberColumn("생산요청물량", format=numeric_format),
        "생산요청물량(PACK)": st.column_config.NumberColumn("생산요청물량(PACK)", format=numeric_format),
        "생산요청물량(PCS)": st.column_config.NumberColumn("생산요청물량(PCS)", format=numeric_format),
        "용마창고재고 (PACK)": st.column_config.NumberColumn("용마창고재고 (PACK)", format=numeric_format),
        "총수량(PACK)": st.column_config.NumberColumn("총수량(PACK)", format=numeric_format),
        "재고기준(PACK)": st.column_config.NumberColumn("재고기준(PACK)", format=numeric_format),
        "재고부족(PACK)": st.column_config.NumberColumn("재고부족(PACK)", format=numeric_format),
        "용마입고 PACK": st.column_config.NumberColumn("용마입고 PACK", format=numeric_format),
        "용마입고": st.column_config.NumberColumn("용마입고", format=numeric_format),
        "미입고": st.column_config.NumberColumn("미입고", format=numeric_format),
        "미입고 PACK": st.column_config.NumberColumn("미입고 PACK", format=numeric_format),
        "용마입고수량": st.column_config.NumberColumn("용마입고수량", format=numeric_format),
        "용마입고수량(PACK)": st.column_config.NumberColumn("용마입고수량(PACK)", format=numeric_format),
        "용마입고수량(PCS)": st.column_config.NumberColumn("용마입고수량(PCS)", format=numeric_format),
        "용마입고대기 PACK": st.column_config.NumberColumn("용마입고대기 PACK", format=numeric_format),
        "용마입고대기수량": st.column_config.NumberColumn("용마입고대기수량", format=numeric_format),
        "용마입고대기수량(PACK)": st.column_config.NumberColumn("용마입고대기수량(PACK)", format=numeric_format),
        "용마입고대기수량(PCS)": st.column_config.NumberColumn("용마입고대기수량(PCS)", format=numeric_format),
        "포장가능재고(PCS)": st.column_config.NumberColumn("포장가능재고(PCS)", format=numeric_format),
        "샘플신청가능수량": st.column_config.NumberColumn("샘플신청가능수량", format=numeric_format),
        "순위": st.column_config.NumberColumn("순위", format="%d", width="small"),
        "현재 재고수량": st.column_config.NumberColumn("현재 재고수량", format=numeric_format),
        "부족수량": st.column_config.NumberColumn("부족수량", format=numeric_format),
        "상세 건수": st.column_config.NumberColumn("상세 건수", format=numeric_format),
        "긴급요청 수": st.column_config.NumberColumn("긴급요청 수", format=numeric_format),
        "미입고(PACK)": st.column_config.NumberColumn("미입고(PACK)", format=numeric_format),
        "미입고수량": st.column_config.NumberColumn("미입고수량", format=numeric_format),
        "입고대기수량": st.column_config.NumberColumn("입고대기수량", format=numeric_format),
        "제품필요수량": st.column_config.NumberColumn("제품필요수량", format=numeric_format),
        "생산필요수량(PCS)": st.column_config.NumberColumn("생산필요수량(PCS)", format=numeric_format),
        "생산부족 PCS": st.column_config.NumberColumn("생산부족 PCS", format=numeric_format),
        "생산부족수량(PCS)": st.column_config.NumberColumn("생산부족수량(PCS)", format=numeric_format),
        "기준차이": st.column_config.TextColumn("기준차이", width="small"),
        "기준차이(PCS)": st.column_config.NumberColumn("기준차이(PCS)", format=numeric_format),
        "포장부족(PACK)": st.column_config.NumberColumn("포장부족(PACK)", format=numeric_format),
        "포장부족(PCS)": st.column_config.NumberColumn("포장부족(PCS)", format=numeric_format),
        "포장부족(재고 PCS)": st.column_config.NumberColumn("포장부족(재고 PCS)", format=numeric_format),
        "검사접착": st.column_config.NumberColumn("검사접착", format=numeric_format),
        "누수규격검사": st.column_config.NumberColumn("누수규격검사", format=numeric_format),
        "5P 필요팩": st.column_config.NumberColumn("5P 필요팩", format=numeric_format),
        "10P 필요팩": st.column_config.NumberColumn("10P 필요팩", format=numeric_format),
        "30P 필요팩": st.column_config.NumberColumn("30P 필요팩", format=numeric_format),
        "80P 필요팩": st.column_config.NumberColumn("80P 필요팩", format=numeric_format),
        "90P 필요팩": st.column_config.NumberColumn("90P 필요팩", format=numeric_format),
        "기타팩 필요팩": st.column_config.NumberColumn("기타팩 필요팩", format=numeric_format),
        "진도율": st.column_config.ProgressColumn(
            "진도율",
            min_value=0,
            max_value=100,
            format="%.1f%%",
            color=COLOR_ORANGE,
        ),
        "전체진도율": st.column_config.ProgressColumn(
            "전체진도율",
            min_value=0,
            max_value=100,
            format="%.1f%%",
            color=COLOR_ORANGE,
        ),
        "요청합계": st.column_config.NumberColumn("요청합계", format=numeric_format),
        "생산부족": st.column_config.NumberColumn("생산부족", format=numeric_format),
        "포장부족": st.column_config.NumberColumn("포장부족", format=numeric_format),
        "판매코드수": st.column_config.NumberColumn("판매코드수", format=numeric_format),
        "판매코드 수": st.column_config.NumberColumn("판매코드 수", format=numeric_format),
        "POWER 수": st.column_config.NumberColumn("POWER 수", format=numeric_format),
        "POWER수": st.column_config.NumberColumn("POWER수", format=numeric_format),
        "생산요청물량 (PCS)": st.column_config.NumberColumn("생산요청물량 (PCS)", format=numeric_format),
        "용마입고수량 (PCS)": st.column_config.NumberColumn("용마입고수량 (PCS)", format=numeric_format),
        "용마입고대기 (PCS)": st.column_config.NumberColumn("용마입고대기 (PCS)", format=numeric_format),
        "포장부족수량 (PCS)": st.column_config.NumberColumn("포장부족수량 (PCS)", format=numeric_format),
        "포장가능수량 (PCS)": st.column_config.NumberColumn("포장가능수량 (PCS)", format=numeric_format),
        "생산부족수량 (PCS)": st.column_config.NumberColumn("생산부족수량 (PCS)", format=numeric_format),
        "생산완료예상일": st.column_config.TextColumn("생산완료예상일"),
        "생산상태": st.column_config.TextColumn("생산상태"),
        "요청 PACK": st.column_config.NumberColumn("요청 PACK", format=numeric_format),
        "포장 PACK": st.column_config.NumberColumn("포장 PACK", format=numeric_format),
        "부족 PACK": st.column_config.NumberColumn("부족 PACK", format=numeric_format),
        "요청": st.column_config.NumberColumn("요청", format=numeric_format),
        "포장": st.column_config.NumberColumn("포장", format=numeric_format),
        "부족": st.column_config.NumberColumn("부족", format=numeric_format),
        "요청PACK": st.column_config.NumberColumn("요청PACK", format=numeric_format),
        "요청PCS": st.column_config.NumberColumn("요청PCS", format=numeric_format),
        "생산": st.column_config.NumberColumn("생산", format=numeric_format),
        "필요팩": st.column_config.NumberColumn("필요팩", format=numeric_format),
        "포장완료PACK": st.column_config.NumberColumn("포장완료PACK", format=numeric_format),
        "포장부족PACK": st.column_config.NumberColumn("포장부족PACK", format=numeric_format),
        "요청수량": st.column_config.NumberColumn("요청수량", format=numeric_format),
        "생산부족(PCS)": st.column_config.NumberColumn("생산부족(PCS)", format=numeric_format),
        "생산부족수량": st.column_config.NumberColumn("생산부족수량", format=numeric_format),
        "포장부족수량": st.column_config.NumberColumn("포장부족수량", format=numeric_format),
        "생산진도율": st.column_config.ProgressColumn(
            "생산진도율",
            min_value=0,
            max_value=100,
            format="%.1f%%",
            color=COLOR_BLUE,
        ),
        "용마입고율": st.column_config.ProgressColumn(
            "용마입고율",
            min_value=0,
            max_value=100,
            format="%.1f%%",
            color=COLOR_AMBER,
        ),
        "포장수량": st.column_config.NumberColumn("포장수량", format=numeric_format),
        "부족수량": st.column_config.NumberColumn("부족수량", format=numeric_format),
        "포장진도율": st.column_config.ProgressColumn(
            "포장진도율",
            min_value=0,
            max_value=100,
            format="%.1f%%",
            color=COLOR_ORANGE,
        ),
        "PACK 지시율": st.column_config.NumberColumn("PACK 지시율", format="%.1f%%"),
        "PCS 지시율": st.column_config.NumberColumn("PCS 지시율", format="%.1f%%"),
        "GAP": st.column_config.NumberColumn("GAP", format="%.1f"),
        "power_value": None,
        "_power_sort": None,
        "_production_code_prefix": None,
        "_min_due_date_sort": None,
        "_priority_sort": None,
        "_request_due_date_sort": None,
        "_pack_sort": None,
        "_daily_item_code_base": None,
        "_daily_status_sort": None,
        "_daily_negative_sort": None,
        "_daily_min_due_date_sort": None,
    }


def get_selected_row(selection_event: Any, df: pd.DataFrame) -> pd.Series | None:
    rows: list[int] = []
    if hasattr(selection_event, "selection"):
        rows = list(getattr(selection_event.selection, "rows", []) or [])
    elif isinstance(selection_event, dict):
        rows = list(selection_event.get("selection", {}).get("rows", []) or [])
    if not rows:
        return None
    row_idx = int(rows[0])
    if row_idx < 0 or row_idx >= len(df):
        return None
    return df.iloc[row_idx]


def dataframe_auto_height(row_count: int, max_height: int, min_height: int = 92, row_height: int = 36) -> int:
    rows = max(int(row_count), 1)
    return int(min(max_height, max(min_height, 50 + row_height * rows)))


def render_selectable_table(
    title: str,
    sub: str,
    df: pd.DataFrame,
    key: str,
    height: int,
    column_order: list[str] | None = None,
) -> pd.Series | None:
    render_panel_title(title, sub)
    if df.empty:
        st.warning("조건에 맞는 데이터가 없습니다.")
        return None
    display_df = dataframe_for_streamlit(df)
    column_config = drilldown_column_config()
    for col in display_df.columns:
        if re.match(r"^\d+(?:\.\d+)?P(?:\(PCS\))?$", str(col)):
            column_config[col] = st.column_config.NumberColumn(str(col), format="%,.0f")
    event = st.dataframe(
        display_df,
        hide_index=True,
        height=dataframe_auto_height(len(display_df), height),
        width="stretch",
        column_config=column_config,
        column_order=visible_columns(display_df, column_order) if column_order is not None else None,
        on_select="rerun",
        selection_mode="single-row",
        key=key,
    )
    return get_selected_row(event, df)


def production_dialog_status_info(selected_row: pd.Series) -> tuple[str, str]:
    shortage = pd.to_numeric(
        pd.Series([selected_row.get("생산부족수량(PCS)", selected_row.get("생산부족수량", 0.0))]),
        errors="coerce",
    ).fillna(0.0).iloc[0]
    return ("생산완료", "done") if float(shortage) <= 0 else ("생산중", "warn")


def production_dialog_metric_html(label: str, value: str, tone: str = "normal", note: str = "") -> str:
    note_html = f"<span>{escape(note)}</span>" if note else ""
    return (
        "<div class='production-dialog-metric'>"
        f"<p>{escape(label)}</p>"
        f"<strong class='{escape(tone)}'>{escape(value)}</strong>"
        f"{note_html}"
        "</div>"
    )


def production_dialog_summary_html(selected_row: pd.Series, detail_count: int) -> str:
    production_code = clean_str(selected_row.get("생산코드", ""))
    product_name = clean_str(selected_row.get("대표 제품명", ""))
    request_pcs = float(pd.to_numeric(pd.Series([selected_row.get("요청합계(PCS)", 0.0)]), errors="coerce").fillna(0.0).iloc[0])
    request_pack = float(pd.to_numeric(pd.Series([selected_row.get("요청합계(PACK)", 0.0)]), errors="coerce").fillna(0.0).iloc[0])
    production_shortage = float(
        pd.to_numeric(
            pd.Series([selected_row.get("생산부족수량(PCS)", selected_row.get("생산부족수량", 0.0))]),
            errors="coerce",
        ).fillna(0.0).iloc[0]
    )
    packing_shortage = float(pd.to_numeric(pd.Series([selected_row.get("포장부족(PACK)", 0.0)]), errors="coerce").fillna(0.0).iloc[0])
    production_progress = float(pd.to_numeric(pd.Series([selected_row.get("생산진도율", 0.0)]), errors="coerce").fillna(0.0).iloc[0])
    packing_progress = float(pd.to_numeric(pd.Series([selected_row.get("포장진도율", 0.0)]), errors="coerce").fillna(0.0).iloc[0])
    expected_date = clean_str(selected_row.get("생산완료예상일", "-")) or "-"
    status_label, status_tone = production_dialog_status_info(selected_row)

    metrics = "".join(
        [
            production_dialog_metric_html(
                "생산요청",
                f"{format_int(request_pcs)} PCS",
                "normal",
                f"{format_int(request_pack)} PACK",
            ),
            production_dialog_metric_html("생산부족", f"{format_int(production_shortage)} PCS", "danger" if production_shortage > 0 else "normal"),
            production_dialog_metric_html("생산진도율", f"{production_progress:.1f}%", "primary"),
            production_dialog_metric_html("포장진도율", f"{packing_progress:.1f}%", "warning"),
            production_dialog_metric_html("포장부족", f"{format_int(packing_shortage)} PACK", "warning" if packing_shortage > 0 else "normal"),
            production_dialog_metric_html("완료예정일", expected_date, "normal"),
        ]
    )
    return (
        "<div class='production-dialog-summary'>"
        "<div class='production-dialog-head'>"
        "<div>"
        f"<div class='production-dialog-code'>{escape(production_code)}</div>"
        f"<div class='production-dialog-product'>{escape(product_name)}</div>"
        "</div>"
        f"<span class='status-badge {status_tone}'>{escape(status_label)}</span>"
        "</div>"
        f"<div class='production-dialog-count'>POWER 상세 {detail_count:,}건</div>"
        f"<div class='production-dialog-metrics'>{metrics}</div>"
        "</div>"
    )


def build_production_pack_dialog_view(selected_row: pd.Series, pack_labels: list[str]) -> pd.DataFrame:
    rows: list[dict[str, Any]] = []
    for label in pack_labels:
        request_pack = float(pd.to_numeric(pd.Series([selected_row.get(label, 0.0)]), errors="coerce").fillna(0.0).iloc[0])
        rows.append(
            {
                "PACK": label,
                "요청 PACK": request_pack,
                "요청 PCS": request_pack * pack_unit_from_label(label),
            }
        )
    out = pd.DataFrame(rows)
    total_pack = float(out["요청 PACK"].sum()) if not out.empty else 0.0
    total_pcs = float(out["요청 PCS"].sum()) if not out.empty else 0.0
    total_row = pd.DataFrame(
        [
            {
                "PACK": "합계",
                "요청 PACK": total_pack,
                "요청 PCS": total_pcs,
            }
        ]
    )
    return pd.concat([out, total_row], ignore_index=True)


def production_pack_dialog_column_config() -> dict[str, Any]:
    column_config = drilldown_column_config()
    column_config["요청 PCS"] = st.column_config.NumberColumn("요청 PCS", format="%,.0f")
    return column_config


def build_production_power_dialog_view(detail_view: pd.DataFrame) -> pd.DataFrame:
    columns = [
        "POWER",
        "요청합계(PCS)",
        "포장실적(PCS)",
        "포장가능재고(PCS)",
        "검사접착",
        "누수규격검사",
        "생산부족수량(PCS)",
        "포장부족(PACK)",
        "생산진도율",
        "포장진도율",
        "생산완료예상일",
    ]
    if detail_view.empty:
        return pd.DataFrame(columns=columns)
    out = detail_view.copy()
    for col in columns:
        if col not in out.columns:
            out[col] = 0.0 if col not in {"POWER", "생산완료예상일"} else ""
    return out[columns].copy()


def render_production_power_detail_dialog(
    selected_row: pd.Series,
    detail_view: pd.DataFrame,
    pack_labels: list[str],
    table_nonce_key: str,
) -> None:
    production_code = clean_str(selected_row.get("생산코드", ""))
    product_name = clean_str(selected_row.get("대표 제품명", ""))
    title = f"생산코드 {production_code} 상세 - {product_name}"

    @st.dialog(title, width="large")
    def _dialog() -> None:
        detail_display = sort_power_detail_default(
            detail_view,
            extra_cols=["_expected_date_sort", "포장부족수량", "생산부족수량"],
            extra_ascending=[True, False, False],
        )
        st.markdown(production_dialog_summary_html(selected_row, len(detail_display)), unsafe_allow_html=True)
        if detail_display.empty:
            st.warning("상세 데이터가 없습니다.")
        else:
            pack_view = build_production_pack_dialog_view(selected_row, pack_labels)
            power_view = build_production_power_dialog_view(detail_display)
            pack_col, power_col = st.columns([1.0, 2.2], gap="small")
            with pack_col:
                st.markdown("<div class='production-dialog-section-title'>PACK 구성 현황</div>", unsafe_allow_html=True)
                st.dataframe(
                    dataframe_for_streamlit(pack_view),
                    hide_index=True,
                    height=dataframe_auto_height(len(pack_view), 390, row_height=34),
                    width="stretch",
                    column_config=production_pack_dialog_column_config(),
                )
            with power_col:
                st.markdown("<div class='production-dialog-section-title'>POWER별 상세 현황</div>", unsafe_allow_html=True)
                st.dataframe(
                    dataframe_for_streamlit(power_view),
                    hide_index=True,
                    height=dataframe_auto_height(len(power_view), 390, row_height=34),
                    width="stretch",
                    column_config=drilldown_column_config(),
                )
        if st.button("닫기", key="close_production_power_detail_dialog", width="stretch"):
            st.session_state[table_nonce_key] = int(st.session_state.get(table_nonce_key, 0)) + 1
            st.rerun()

    _dialog()


def render_daily_inventory_detail_dialog(
    selected_row: pd.Series,
    detail_view: pd.DataFrame,
    table_nonce_key: str,
) -> None:
    item_code = clean_str(selected_row.get("_daily_item_code_base", selected_row.get("품목코드", "")))
    product_name = clean_str(selected_row.get("대표 제품명", selected_row.get("제품명", "")))
    title = f"판매코드 {item_code} 상세 - {product_name}"

    @st.dialog(title, width="large")
    def _dialog() -> None:
        detail_display = sort_power_detail_default(detail_view)
        st.caption(f"{item_code}에 해당하는 PACK/POWER별 재고 대응 상세 | 표시 건수: {len(detail_display):,}")
        if detail_display.empty:
            st.warning("상세 데이터가 없습니다.")
        else:
            st.dataframe(
                dataframe_for_streamlit(detail_display),
                hide_index=True,
                height=dataframe_auto_height(len(detail_display), 520),
                width="stretch",
                column_config=drilldown_column_config(),
                column_order=daily_inventory_detail_column_order(detail_display),
            )
        if st.button("닫기", key="close_daily_inventory_detail_dialog", width="stretch"):
            st.session_state[table_nonce_key] = int(st.session_state.get(table_nonce_key, 0)) + 1
            st.rerun()

    _dialog()


def render_sales_code_detail_dialog(
    selected_row: pd.Series,
    detail_view: pd.DataFrame,
    inventory_view: pd.DataFrame,
    unit_mode: str,
    table_nonce_key: str,
) -> None:
    sales_code = clean_str(selected_row.get("_sales_code_base", selected_row.get("판매코드", "")))
    product_name = clean_str(selected_row.get("대표 제품명", selected_row.get("제품명", "")))
    title = f"판매코드 {sales_code} 상세 - {product_name}"

    @st.dialog(title, width="large")
    def _dialog() -> None:
        detail_display = sort_power_detail_default(detail_view)
        st.caption(f"{sales_code}에 해당하는 POWER/PACK별 출고·오더 상세 | 표시 건수: {len(detail_display):,}")
        if detail_display.empty:
            st.warning("상세 데이터가 없습니다.")
        else:
            st.dataframe(
                dataframe_for_streamlit(detail_display.drop(columns=["power_value"], errors="ignore")),
                hide_index=True,
                height=dataframe_auto_height(len(detail_display), 430),
                width="stretch",
                column_config=drilldown_column_config(),
                column_order=sales_progress_column_order(detail_display, unit_mode),
            )

        st.markdown("<div style='height:10px'></div>", unsafe_allow_html=True)
        st.caption("용마WMS재고현황 기준 판매코드별 PACK 재고")
        if inventory_view.empty or set(inventory_view.get("매칭여부", pd.Series(dtype=str)).astype(str)) == {"미매칭"}:
            st.info("선택한 판매코드와 매칭되는 WMS 재고가 없습니다.")
        else:
            st.dataframe(
                dataframe_for_streamlit(inventory_view),
                hide_index=True,
                height=dataframe_auto_height(len(inventory_view), 160),
                width="stretch",
                column_config=drilldown_column_config(),
            )
        if st.button("닫기", key="close_sales_code_detail_dialog", width="stretch"):
            st.session_state[table_nonce_key] = int(st.session_state.get(table_nonce_key, 0)) + 1
            st.rerun()

    _dialog()


def render_daily_inventory_tab(
    daily_inventory_df: pd.DataFrame,
    code_summary: pd.DataFrame,
    sample_available_df: pd.DataFrame | None = None,
    lot_status_df: pd.DataFrame | None = None,
    selected_period: str = "전체",
) -> None:
    render_panel_title(
        "일일 재고 대응",
        "매일 공유되는 재고현황표 기준으로 요청물량 외 긴급 품목과 재고부족을 확인합니다.",
    )
    if daily_inventory_df.empty:
        st.warning("일일 재고현황표를 찾지 못했거나 처리할 데이터가 없습니다.")
        return

    scoped_code_summary = filter_operational_code_summary(
        code_summary,
        period_group=selected_period,
    )
    scoped_lot_status_df = lot_status_df
    if lot_status_df is not None and not lot_status_df.empty:
        scoped_lot_status_df = filter_by_period_group(lot_status_df, selected_period)

    response_view = build_daily_inventory_response_view(
        daily_inventory_df,
        scoped_code_summary,
        sample_available_df,
        scoped_lot_status_df,
    )
    response_view = filter_by_period_group(response_view, selected_period)
    if response_view.empty:
        st.warning("표시할 일일 재고 대응 데이터가 없습니다.")
        return

    urgent_count = int(response_view["긴급요청"].sum())
    negative_count = int((response_view["재고수량"] < 0).sum())
    request_out_count = int((response_view["대응상태"] == "요청외 긴급").sum())
    request_in_count = int(response_view["대응상태"].isin(["요청내 긴급", "요청내 재고부족"]).sum())
    shortage_qty = float(response_view["재고부족수량"].sum())
    render_metric_card_grid(
        [
            ("긴급요청 품목", f"{urgent_count:,}", "warn" if urgent_count else "normal"),
            ("요청외 긴급", f"{request_out_count:,}", "warn" if request_out_count else "normal"),
            ("요청내 부족/긴급", f"{request_in_count:,}", "warn" if request_in_count else "normal"),
            (
                "재고부족수량",
                format_int(shortage_qty),
                "warn" if shortage_qty > 0 else "normal",
                f"음수 {negative_count:,}품목" if negative_count else "",
            ),
        ]
    )

    f1, f2, f3 = st.columns([2.4, 1.7, 1.2], gap="small")
    with f1:
        query = st.text_input(
            "제품명/판매코드/POWER 검색",
            value="",
            placeholder="예: 소울브라운, 40P, -06.50",
            key="daily_inventory_query",
        )
    with f2:
        statuses = st.multiselect(
            "대응상태",
            sorted(response_view["대응상태"].dropna().astype(str).unique().tolist()),
            default=sorted(response_view["대응상태"].dropna().astype(str).unique().tolist()),
            key="daily_inventory_status",
        )
    with f3:
        important_only = st.checkbox("긴급/부족만 보기", value=True, key="daily_inventory_important_only")

    view = response_view.copy()
    if query.strip():
        view = view[daily_inventory_query_mask(view, query)].copy()
    if statuses:
        view = view[view["대응상태"].isin(statuses)].copy()
    if important_only:
        view = view[(view["긴급요청"]) | (view["재고수량"] < 0) | (view["재고부족수량"] > 0)].copy()

    hidden_daily_inventory_cols = [
        "재고표 제품명",
        "전일재고",
        "재고증감",
        "재고부족수량",
        "요청제품명",
        "판매코드 수",
        "대상품목",
        "포장부족(재고 PCS)",
        "포장 PACK",
        "미입고 PACK",
    ]
    detail_view = view.drop(columns=hidden_daily_inventory_cols, errors="ignore")
    main_view = build_daily_inventory_main_view(view)
    full_export_view = response_view.drop(columns=hidden_daily_inventory_cols, errors="ignore")

    dl_col, _ = st.columns([1.2, 4.8], gap="small")
    with dl_col:
        render_excel_download(
            "엑셀 다운로드",
            "일일_재고_대응",
            {
                "일일 재고 대응": main_view,
                "일일 재고 상세": detail_view,
                "일일 재고 전체": full_export_view,
            },
            key="download_daily_inventory_excel",
        )

    table_nonce_key = "daily_inventory_main_table_nonce"
    table_nonce = int(st.session_state.get(table_nonce_key, 0))
    selected_daily_row = render_selectable_table(
        "일일 재고 대응 테이블",
        f"판매코드 기준 집계 | 표시 건수: {len(main_view):,} | 상세 건수: {len(view):,}",
        main_view,
        key=f"daily_inventory_table_{table_nonce}",
        height=560,
        column_order=[
            "대응상태",
            "품목코드",
            "기간구분",
            "대표 제품명",
            "긴급요청 수",
            "재고수량",
            "요청 PACK",
            "용마입고 PACK",
            "용마입고대기 PACK",
            "포장가능재고(PCS)",
            "생산부족 PCS",
            "생산진도율",
            "생산완료예상일",
        ],
    )
    if selected_daily_row is not None:
        selected_item_code = clean_str(
            selected_daily_row.get("_daily_item_code_base", selected_daily_row.get("품목코드", ""))
        )
        detail_scope = detail_view[
            detail_view["품목코드"].map(daily_item_code_base) == selected_item_code
        ].copy()
        detail_scope["_daily_status_sort"] = detail_scope["대응상태"].map(daily_inventory_status_rank)
        detail_scope["_daily_pack_sort"] = detail_scope["PACK"].map(pack_sort_key)
        detail_scope["_daily_power_sort"] = pd.to_numeric(
            detail_scope["POWER"].astype(str).str.replace("-00.00", "0", regex=False),
            errors="coerce",
        ).fillna(0.0)
        detail_scope["_daily_stock_shortage_sort"] = pd.to_numeric(
            detail_scope.get("재고수량", pd.Series(0.0, index=detail_scope.index)),
            errors="coerce",
        ).fillna(0.0)
        detail_scope = sort_power_detail_default(
            detail_scope,
            extra_cols=["_daily_pack_sort", "_daily_status_sort", "_daily_stock_shortage_sort"],
            extra_ascending=[True, True, True],
        ).drop(
            columns=[
                "_daily_status_sort",
                "_daily_pack_sort",
                "_daily_power_sort",
                "_daily_stock_shortage_sort",
            ],
            errors="ignore",
        )
        render_daily_inventory_detail_dialog(selected_daily_row, detail_scope, table_nonce_key)


def render_product_summary_tab(
    product_summary: pd.DataFrame,
    code_summary: pd.DataFrame,
    request_df: pd.DataFrame | None = None,
    instruction_df: pd.DataFrame | None = None,
    daily_inventory_df: pd.DataFrame | None = None,
    sample_available_df: pd.DataFrame | None = None,
    selected_period: str = "전체",
) -> None:
    product_summary = filter_by_period_group(product_summary, selected_period)
    code_summary = filter_operational_code_summary(code_summary, period_group=selected_period)
    if request_df is not None and not request_df.empty:
        request_df = filter_by_period_group(request_df, selected_period)
    if instruction_df is not None and not instruction_df.empty:
        instruction_df = filter_by_period_group(instruction_df, selected_period)
    main_products, _ = split_main_sample(product_summary)
    stock_threshold_pack = float(INVENTORY_STOCK_THRESHOLD_DEFAULT)

    family_view = build_family_progress_view(main_products)
    request_level_source = request_df if request_df is not None and not request_df.empty else code_summary
    category_request_view = build_category_request_summary_view(request_level_source, instruction_df)
    top_shortage_view = build_top_shortage_view(product_summary, top_n=10)
    gap_top_view = build_gap_top_view(product_summary, top_n=10)
    daily_response_view = (
        build_daily_inventory_response_view(daily_inventory_df, code_summary, sample_available_df)
        if daily_inventory_df is not None and not daily_inventory_df.empty
        else None
    )
    exception_kpis, exception_detail = build_daily_exception_report_view(
        daily_inventory_df,
        code_summary,
        sample_available_df,
        max_rows=10,
        response_view=daily_response_view,
    )
    urgent_summary_view = build_urgent_request_summary_view(
        daily_inventory_df,
        code_summary,
        sample_available_df,
        response_view=daily_response_view,
    )

    kpi_head_period_col, kpi_head_spacer_col, kpi_head_action_col = st.columns(
        [1.35, 1.1, 2.15],
        gap="small",
        vertical_alignment="center",
    )
    with kpi_head_period_col:
        st.segmented_control(
            "기간구분",
            options=PERIOD_GROUP_ORDER,
            default=selected_period if selected_period in PERIOD_GROUP_ORDER else "전체",
            label_visibility="visible",
            key="product_summary_period_group_filter",
        )
    with kpi_head_spacer_col:
        st.empty()
    with kpi_head_action_col:
        action_cols = st.columns([1, 1], gap="small")
        with action_cols[0]:
            render_lazy_binary_download(
                "PPT 다운로드",
                "PPT 보고서 준비",
                f"국내_제품_포장현황_운영보고서_{pd.Timestamp.now(tz='Asia/Seoul').strftime('%Y%m%d_%H%M')}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                build_bytes=lambda: build_ppt_report(
                    product_view=product_summary,
                    code_summary=code_summary,
                    product_names=product_summary["제품명"],
                    scope_label="전체",
                    daily_inventory_df=daily_inventory_df,
                    sample_available_df=sample_available_df,
                ),
                signature=(
                    dataframe_light_signature(product_summary),
                    dataframe_light_signature(code_summary),
                    dataframe_light_signature(daily_inventory_df) if daily_inventory_df is not None else ("none",),
                    dataframe_light_signature(sample_available_df) if sample_available_df is not None else ("none",),
                ),
                key="download_ppt_report",
                width="stretch",
            )
        with action_cols[1]:
            render_excel_download(
                "엑셀 다운로드",
                "제품_진도_현황",
                {
                    "제품 요약": product_summary,
                    "신규분류요약별 요청지시율": category_request_view,
                    "미입고 TOP10": top_shortage_view,
                    "본품 분류별 진도": family_view,
                    "생산완료 후 미입고 TOP10": gap_top_view,
                    "요청 긴급 요약": urgent_summary_view,
                    "요청 긴급 상세": exception_detail,
                },
                key="download_product_progress_excel",
            )
    scope_kpis = {
        name: kpi
        for name, kpi in build_scope_kpis(add_allocated_production_basis(code_summary))
        if name in {"본품", "샘플"}
    }
    kpi_cols = st.columns([5.5, 2.25, 2.25], gap="large")
    with kpi_cols[0]:
        render_status_board(
            product_summary,
            code_summary,
            daily_inventory_df,
            sample_available_df,
            stock_threshold_pack,
            exception_kpis=exception_kpis,
        )
    with kpi_cols[1]:
        render_kpi_panel("본품 KPI", scope_kpis.get("본품", calc_kpi_from_code_summary(pd.DataFrame())))
    with kpi_cols[2]:
        render_kpi_panel("샘플 KPI", scope_kpis.get("샘플", calc_kpi_from_code_summary(pd.DataFrame())))

    st.markdown("<div class='section-gap'></div>", unsafe_allow_html=True)
    render_panel_title(
        "제품 분류별 진도 현황",
        "제품군별 생산지시 PCS와 PACK, 생산진도율, 용마입고율, 생산부족 PCS를 비교합니다.",
    )
    render_family_progress_cards(family_view)
    render_product_completion_section(code_summary)

    st.markdown("<div style='height:10px'></div>", unsafe_allow_html=True)
    with st.expander("신규분류요약별 요청 대비 지시 수준", expanded=False):
        st.caption("3Q전체물량은 요청량, 생산지시물량은 지시량으로 보고 신규분류요약별 지시율과 미지시 PCS를 집계합니다.")
        render_request_instruction_level_cards(category_request_view)
        render_category_request_summary_table(category_request_view)


def render_production_code_tab(
    code_summary: pd.DataFrame,
    wip_df: pd.DataFrame | None = None,
    selected_period: str = "전체",
) -> None:
    render_panel_title(
        "생산코드 상세",
        "생산코드 기준으로 제품군 위험도를 확인하고, 선택 시 POWER별 상세를 팝업으로 확인합니다.",
    )
    production_unit_mode = UNIT_PACK
    pack_options = available_pack_options(code_summary)
    pack_labels = PRODUCTION_CODE_PACK_LABELS
    power_options = available_production_power_options(code_summary)
    group_options = available_product_group_options(code_summary)

    pc1, pc2, pc3, pc4, pc5 = st.columns([3.7, 1.55, 1.1, 1.1, 0.95], gap="small")
    with pc1:
        integrated_query = st.text_input(
            "통합검색",
            value="",
            placeholder="예: 소울브라운, S145, P0019",
            key="tab_production_integrated_query",
        )
    with pc2:
        selected_group = st.selectbox(
            "분류 선택",
            options=group_options,
            index=0,
            key="tab_production_group",
        )
    with pc3:
        selected_pack = st.selectbox(
            "PACK 선택",
            options=pack_options,
            index=0,
            key="tab_production_pack",
        )
    with pc4:
        selected_power = st.selectbox(
            "POWER 선택",
            options=power_options,
            index=0,
            key="tab_production_power",
        )
    with pc5:
        shortage_only = st.checkbox("부족품만 보기", value=False, key="tab_production_shortage_only")

    production_source = filter_production_power_rows(
        code_summary,
        product_query="",
        production_query="",
        power_label=selected_power,
        pack_label=selected_pack,
        sample_scope="전체",
        product_group=selected_group,
        factory_group="전체",
        period_group=selected_period,
    )
    production_source = filter_dataframe_by_terms(
        production_source,
        integrated_query,
        columns=[
            "period_group",
            "production_code_display",
            "sales_code",
            "product_name",
            "base_product_name",
            "product_name_code",
            "POWER",
            "_pack_label",
            "제품분류",
            "본품/샘플",
        ],
    )
    production_view = build_production_power_main_view(
        production_source,
        pack_labels=pack_labels,
        shortage_only=shortage_only,
    )
    production_detail_view = build_production_power_detail_view(
        production_source,
        pack_labels=pack_labels,
        wip_df=wip_df,
    )
    render_production_power_kpis(production_view, unit_mode=production_unit_mode)
    production_main_export = production_view[
        production_progress_column_order(production_view, pack_labels, production_unit_mode)
    ].copy()
    production_detail_export = production_detail_view[
        production_power_detail_column_order(production_detail_view, pack_labels)
    ].copy()
    dl_col, _ = st.columns([1.2, 4.8], gap="small")
    with dl_col:
        render_excel_download(
            "엑셀 다운로드",
            "생산코드_상세",
            {
                "생산코드 집계": production_main_export,
                "POWER 상세": production_detail_export,
            },
            key="download_production_code_excel",
        )

    table_nonce_key = "production_code_main_table_nonce"
    table_nonce = int(st.session_state.get(table_nonce_key, 0))
    selected_production_row = render_selectable_table(
        "생산코드 메인 테이블",
        f"생산코드 기준 집계 | 생산완료예상일, 포장부족, 생산부족 순 정렬 | 표시 건수: {len(production_view):,}",
        production_view,
        key=f"production_code_main_table_{table_nonce}",
        height=620,
        column_order=production_progress_column_order(production_view, pack_labels, production_unit_mode),
    )
    if selected_production_row is None:
        return

    selected_production = clean_str(selected_production_row.get("_production_code_prefix", selected_production_row.get("생산코드", "")))
    detail_view = build_production_power_detail_view(
        production_source,
        pack_labels=pack_labels,
        production_prefix=selected_production,
        wip_df=wip_df,
    )
    render_production_power_detail_dialog(
        selected_production_row,
        detail_view,
        pack_labels,
        table_nonce_key,
    )


def render_sales_code_tab(code_summary: pd.DataFrame, selected_period: str = "전체") -> None:
    render_panel_title(
        "판매코드 상세",
        "출고/오더 관점에서 판매코드별 생산·포장 진도와 생산완료예상일 상태를 확인합니다.",
    )
    pack_options = available_pack_options(code_summary)
    power_options = available_power_options(code_summary)

    fc1, fc2, _ = st.columns([1.55, 1.25, 3.2], gap="small")
    with fc1:
        sales_unit_mode = st.radio(
            "조회 단위 선택",
            UNIT_OPTIONS,
            index=0,
            horizontal=True,
            key="sales_progress_unit_mode",
        )
        if sales_unit_mode == UNIT_PCS:
            st.caption("포장가능재고·생산부족 기준 조회")
        else:
            st.caption("용마입고·포장부족 기준 조회")
    with fc2:
        stock_threshold_pack = st.number_input(
            "긴급 재고 기준(PACK)",
            min_value=0,
            value=INVENTORY_STOCK_THRESHOLD_DEFAULT,
            step=10,
            key="sales_inventory_stock_threshold_pack",
        )

    period_scoped_code_summary = filter_operational_code_summary(
        code_summary,
        period_group=selected_period,
    )
    sales_base = build_sales_order_main_view(
        period_scoped_code_summary,
        stock_threshold_pack=float(stock_threshold_pack),
        today_key=pd.Timestamp.now(tz="Asia/Seoul").strftime("%Y-%m-%d"),
    )
    priority_tab, detail_tab = st.tabs(["포장 우선순위", "판매코드 세부리스트"])

    with priority_tab:
        render_urgent_sales_packing_list(sales_base)

    with detail_tab:
        sf1, sf2, sf3, sf4 = st.columns([4.5, 1.25, 1.25, 0.95], gap="small")
        with sf1:
            integrated_query = st.text_input(
                "통합검색",
                value="",
                placeholder="예: 소울브라운, S145, P0019",
                key="tab_sales_integrated_query",
            )
        with sf2:
            selected_pack = st.selectbox("PACK 선택", options=pack_options, index=0, key="tab_sales_pack")
        with sf3:
            selected_power = st.selectbox("POWER 선택", options=power_options, index=0, key="tab_sales_power")
        with sf4:
            shortage_only = st.checkbox("부족품만 보기", value=False, key="tab_sales_shortage_only")

        sales_detail_view = filter_sales_order_view(
            sales_base,
            pack_label=selected_pack,
            power_label=selected_power,
        )
        sales_main_unfiltered = build_sales_code_group_main_view(
            sales_detail_view,
            stock_threshold_pack=float(stock_threshold_pack),
        )
        sales_search_columns = list(
            dict.fromkeys(
                sales_group_column_order(sales_main_unfiltered, sales_unit_mode)
                + ["제품분류", "생산코드", "생산요청물량(PCS)", "용마입고수량(PCS)", "포장부족(PCS)"]
            )
        )
        sales_main_view = filter_dataframe_by_terms(
            sales_main_unfiltered,
            integrated_query,
            columns=sales_search_columns,
        )
        if shortage_only and "포장부족(PACK)" in sales_main_view.columns:
            sales_main_view = sales_main_view[
                pd.to_numeric(sales_main_view["포장부족(PACK)"], errors="coerce").fillna(0.0) > 0
            ].copy()
        visible_sales_codes = set(sales_main_view.get("_sales_code_base", sales_main_view.get("판매코드", pd.Series(dtype=str))).map(clean_str))
        if integrated_query.strip() or shortage_only:
            if visible_sales_codes:
                sales_detail_export_view = sales_detail_view[
                    sales_detail_view["판매코드"].map(sales_code_base).isin(visible_sales_codes)
                ].copy()
            else:
                sales_detail_export_view = sales_detail_view.iloc[0:0].copy()
        else:
            sales_detail_export_view = sales_detail_view
        dl_col, _ = st.columns([1.2, 4.8], gap="small")
        with dl_col:
            render_excel_download(
                "엑셀 다운로드",
                "판매코드_상세",
                {
                    "긴급 포장 리스트": build_urgent_sales_packing_view(sales_base),
                    "판매코드 집계": sales_main_view,
                    "POWER 상세": sales_detail_export_view,
                },
                key="download_sales_code_excel",
            )
        table_nonce_key = "sales_code_main_table_nonce"
        table_nonce = int(st.session_state.get(table_nonce_key, 0))
        selected_sales_row = render_selectable_table(
            "판매코드",
            f"판매코드 기준 집계 | 표시 건수: {len(sales_main_view):,} | 상세 건수: {len(sales_detail_export_view):,}",
            sales_main_view,
            key=f"sales_code_main_table_{table_nonce}",
            height=620,
            column_order=sales_group_column_order(sales_main_view, sales_unit_mode),
        )
        if selected_sales_row is None:
            return

        selected_sales = clean_str(selected_sales_row.get("_sales_code_base", selected_sales_row.get("판매코드", "")))
        detail_scope = sales_detail_view[sales_detail_view["판매코드"].map(sales_code_base) == selected_sales].copy()
        detail_scope["_pack_sort"] = detail_scope["PACK"].map(pack_sort_rank) if "PACK" in detail_scope.columns else 0.0
        if "power_value" not in detail_scope.columns:
            detail_scope["power_value"] = pd.to_numeric(
                detail_scope.get("POWER", pd.Series(0.0, index=detail_scope.index)),
                errors="coerce",
            ).fillna(999999.0)
        detail_scope = sort_power_detail_default(
            detail_scope,
            extra_cols=["_pack_sort", "_priority_sort", "_request_due_date_sort", "포장부족(PACK)"],
            extra_ascending=[True, True, True, False],
        ).drop(columns=["_pack_sort"], errors="ignore")
        inventory_source = filter_operational_code_summary(
            period_scoped_code_summary,
            pack_label=selected_pack,
            power_label=selected_power,
        )
        inventory_view = build_inventory_prefix_detail_view(inventory_source, selected_sales)
        render_sales_code_detail_dialog(
            selected_sales_row,
            detail_scope,
            inventory_view,
            sales_unit_mode,
            table_nonce_key,
        )


def render_power_tab(code_summary: pd.DataFrame, selected_period: str = "전체") -> None:
    render_panel_title(
        "POWER 상세",
        "렌즈 POWER 기준 요청/생산/포장/부족 현황과 하위 생산·판매코드를 확인합니다.",
    )
    power_unit_mode = render_unit_selector("power_progress_unit_mode")
    power_options = available_power_options(code_summary)

    pf1, pf2, pf3, pf4 = st.columns([2.0, 1.7, 1.7, 1.2], gap="small")
    with pf1:
        product_query = st.text_input("제품명", value="", placeholder="제품명/SKU 일부 입력", key="tab_power_product_query")
    with pf2:
        production_query = st.text_input("생산코드", value="", placeholder="예: P3015", key="tab_power_production_query")
    with pf3:
        sales_query = st.text_input("판매코드", value="", placeholder="예: S309", key="tab_power_sales_query")
    with pf4:
        selected_power = st.selectbox("POWER", options=power_options, index=0, key="tab_power_power")

    power_source = filter_operational_code_summary(
        code_summary,
        product_query=product_query,
        production_query=production_query,
        sales_query=sales_query,
        power_label=selected_power,
        period_group=selected_period,
    )
    power_detail_for_heatmap = build_power_detail(power_source)
    heatmap = build_power_heatmap(power_detail_for_heatmap)
    if heatmap is not None:
        st.plotly_chart(heatmap, width="stretch")

    power_summary = build_power_summary_view(power_source)
    ops_kpi = calc_power_ops_kpi(power_detail_for_heatmap)
    if power_unit_mode == UNIT_PCS:
        request_pcs = float(power_summary["요청합계(PCS)"].sum()) if not power_summary.empty else 0.0
        production_shortage_pcs = (
            float(power_summary["생산부족수량(PCS)"].sum()) if not power_summary.empty else 0.0
        )
        request_pack = float(power_summary["요청합계(PACK)"].sum()) if not power_summary.empty else 0.0
        packing_shortage_pack = float(power_summary["포장부족(PACK)"].sum()) if not power_summary.empty else 0.0
        production_progress = (
            (request_pcs - production_shortage_pcs) / request_pcs * 100.0
            if request_pcs > 0
            else 0.0
        )
        packing_progress = (
            (request_pack - packing_shortage_pack) / request_pack * 100.0
            if request_pack > 0
            else 0.0
        )
        production_progress = min(100.0, max(0.0, production_progress))
        packing_progress = min(100.0, max(0.0, packing_progress))
        render_metric_card_grid(
            [
                ("대상 도수", f"{ops_kpi['rows']:,}", "normal"),
                ("요청합계(PCS)", format_int(request_pcs), "normal"),
                ("생산부족수량(PCS)", format_int(production_shortage_pcs), "warn"),
                ("생산진도율", f"{production_progress:.1f}%", metric_progress_tone(production_progress)),
                ("포장진도율", f"{packing_progress:.1f}%", metric_progress_tone(packing_progress)),
            ]
        )
    else:
        render_metric_card_grid(
            [
                ("대상 도수", f"{ops_kpi['rows']:,}", "normal"),
                ("포장부족 도수", f"{ops_kpi['shortage_rows']:,}", "warn"),
                ("미착수 도수", f"{ops_kpi['not_started_rows']:,}", "warn"),
                ("하이파워 부족", f"{ops_kpi['high_power_shortage_rows']:,}", "warn"),
                ("포장부족(PACK) 합계", format_int(ops_kpi["shortage_qty"]), "warn"),
            ]
        )

    dl_col, _ = st.columns([1.2, 4.8], gap="small")
    with dl_col:
        render_excel_download(
            "엑셀 다운로드",
            "POWER_상세",
            {
                "POWER 요약": power_summary,
                "POWER 상세": power_detail_for_heatmap,
            },
            key="download_power_excel",
        )

    selected_power_row = render_selectable_table(
        "POWER 히트맵 상세",
        f"POWER 기준 요청/생산/포장/부족 | 표시 건수: {len(power_summary):,}",
        power_summary.drop(columns=["power_value"], errors="ignore"),
        key="power_summary_table",
        height=430,
        column_order=power_progress_column_order(power_summary, power_unit_mode),
    )
    if selected_power_row is None:
        return

    selected_power_detail = str(selected_power_row["POWER"])
    st.markdown(f"<div class='breadcrumb'>POWER <span>{escape(selected_power_detail)}</span></div>", unsafe_allow_html=True)
    sku_detail = build_power_sku_detail_view(power_source, selected_power_detail)
    render_selectable_table(
        "SKU 상세",
        f"{selected_power_detail} 기준 생산코드/판매코드 상세 | 표시 건수: {len(sku_detail):,}",
        sku_detail,
        key="power_sku_detail_table",
        height=320,
        column_order=visible_columns(
            sku_detail,
            [
                "생산코드",
                "판매코드",
                "기간구분",
                "제품명",
                "PACK",
                "요청합계(PACK)",
                "포장부족(PACK)",
                "생산부족수량(PCS)",
                "생산완료예상일",
            ]
            if power_unit_mode == UNIT_PACK
            else [
                "생산코드",
                "판매코드",
                "기간구분",
                "제품명",
                "PACK",
                "요청합계(PCS)",
                "생산필요수량(PCS)",
                "생산부족수량(PCS)",
                "생산완료예상일",
            ],
        ),
    )


def render_drilldown_tab(product_summary: pd.DataFrame, code_summary: pd.DataFrame) -> None:
    render_drilldown_kpi(calc_drilldown_kpi(product_summary))

    filter_col1, filter_col2 = st.columns([3, 2], gap="small")
    with filter_col1:
        query = st.text_input("제품명 검색", value="", placeholder="제품명 일부 입력", key="drill_product_query")
    with filter_col2:
        statuses = st.multiselect("상태 필터", STATUS_ORDER, default=STATUS_ORDER, key="drill_product_status")

    product_filtered = apply_filters(product_summary, query=query, statuses=statuses)
    product_filtered = product_filtered.sort_values(
        ["포장부족수량", "생산부족수량", "요청 PACK"],
        ascending=[False, False, False],
        kind="stable",
    )
    product_view = build_product_drilldown_view(product_filtered)
    selected_product_row = render_selectable_table(
        "제품 진도현황",
        f"제품 기준 요청/생산부족/포장 현황 | 표시 건수: {len(product_view):,}",
        product_view,
        key="drill_product_table",
        height=430,
    )
    if selected_product_row is None:
        return

    selected_product = str(selected_product_row["제품명"])
    if "base_product_name" in code_summary.columns:
        product_scope = code_summary[code_summary["base_product_name"] == selected_product].copy()
    else:
        product_scope = code_summary[code_summary["product_name"].map(strip_pack_unit_suffix) == selected_product].copy()
    if product_scope.empty:
        product_scope = code_summary[code_summary["product_name"] == selected_product].copy()
    st.markdown(f"<div class='breadcrumb'>제품 <span>{escape(selected_product)}</span></div>", unsafe_allow_html=True)

    pack_unit_view = build_pack_unit_view(code_summary, selected_product)
    render_panel_title(
        "팩 단위 포장 진도",
        f"{strip_pack_unit_suffix(selected_product)} 기준 팩 단위 요청/포장/부족/진도율",
    )
    if pack_unit_view.empty:
        st.warning("팩 단위 상세 데이터가 없습니다.")
    else:
        st.dataframe(
            pack_unit_view,
            hide_index=True,
            height=dataframe_auto_height(len(pack_unit_view), 260),
            width="stretch",
            column_config=pack_unit_column_config(),
        )

    production_view = build_production_drilldown_view(product_scope)
    production_view = production_view.sort_values(
        ["포장부족수량", "생산부족수량", "요청수량"],
        ascending=[False, False, False],
        kind="stable",
    )
    selected_production_row = render_selectable_table(
        "생산코드",
        f"{selected_product} 기준 생산코드 현황 | 표시 건수: {len(production_view):,}",
        production_view,
        key="drill_production_table",
        height=300,
    )
    if selected_production_row is None:
        return

    selected_production = str(selected_production_row["생산코드"])
    production_key = selected_production if selected_production != "(생산코드 미기재)" else ""
    production_scope = product_scope[product_scope["production_code"].replace("", "(생산코드 미기재)") == selected_production].copy()
    if production_scope.empty and production_key:
        production_scope = product_scope[product_scope["production_code"] == production_key].copy()
    st.markdown(
        "<div class='breadcrumb'>"
        f"제품 <span>{escape(selected_product)}</span>"
        f"<b>›</b> 생산코드 <span>{escape(selected_production)}</span>"
        "</div>",
        unsafe_allow_html=True,
    )

    sales_view = build_sales_drilldown_view(production_scope)
    sales_view = sales_view.sort_values(["부족수량", "요청수량"], ascending=[False, False], kind="stable")
    selected_sales_row = render_selectable_table(
        "판매코드",
        f"{selected_production} 기준 판매코드 현황 | 표시 건수: {len(sales_view):,}",
        sales_view,
        key="drill_sales_table",
        height=280,
    )
    if selected_sales_row is None:
        return

    selected_sales = str(selected_sales_row["판매코드"])
    sales_scope = production_scope[production_scope["sales_code"] == selected_sales].copy()
    st.markdown(
        "<div class='breadcrumb'>"
        f"제품 <span>{escape(selected_product)}</span>"
        f"<b>›</b> 생산코드 <span>{escape(selected_production)}</span>"
        f"<b>›</b> 판매코드 <span>{escape(selected_sales)}</span>"
        "</div>",
        unsafe_allow_html=True,
    )

    power_view = build_power_drilldown_view(sales_scope)
    render_selectable_table(
        "POWER",
        f"{selected_sales} 기준 POWER 상세 운영현황 | 표시 건수: {len(power_view):,}",
        power_view,
        key="drill_power_table",
        height=240,
    )


def file_fingerprint(path: Path | None) -> tuple[str, int, int] | None:
    if path is None:
        return None
    stat = path.stat()
    return (str(path.resolve()), int(stat.st_mtime_ns), int(stat.st_size))


def dashboard_cache_fingerprint() -> tuple[int, str, str]:
    try:
        path = Path(__file__).resolve()
        digest = hashlib.sha256(path.read_bytes()).hexdigest()
        return (DATA_CACHE_VERSION, str(path), digest)
    except Exception:
        return (DATA_CACHE_VERSION, "dashboard.py", "")


@st.cache_data(show_spinner="데이터 파일을 읽는 중입니다. 잠시만 기다려 주세요.", max_entries=8)
def load_dashboard_data(
    request_fingerprint: tuple[str, int, int],
    packing_fingerprint: tuple[str, int, int],
    progress_fingerprint: tuple[str, int, int] | None,
    inventory_fingerprint: tuple[str, int, int] | None,
    daily_inventory_fingerprint: tuple[str, int, int] | None,
    product_master_fingerprint: tuple[str, int, int] | None,
    wip_fingerprint: tuple[str, int, int] | None,
    cache_version: tuple[int, str, str] | int,
) -> tuple[pd.DataFrame, pd.DataFrame, pd.DataFrame, pd.DataFrame, pd.DataFrame, pd.DataFrame, pd.DataFrame, pd.DataFrame, pd.DataFrame]:
    request_file = Path(request_fingerprint[0])
    packing_file = Path(packing_fingerprint[0])
    progress_file = Path(progress_fingerprint[0]) if progress_fingerprint is not None else None
    inventory_file = Path(inventory_fingerprint[0]) if inventory_fingerprint is not None else None
    daily_inventory_file = Path(daily_inventory_fingerprint[0]) if daily_inventory_fingerprint is not None else None
    product_master_file = Path(product_master_fingerprint[0]) if product_master_fingerprint is not None else None
    wip_file = Path(wip_fingerprint[0]) if wip_fingerprint is not None else None

    product_master_df = normalize_product_code_master(product_master_file)
    request_df = normalize_request(request_file, product_master_file, product_master_df=product_master_df)
    instruction_df = normalize_instruction_request(request_file, product_master_file, product_master_df=product_master_df)
    progress_basis_df = instruction_df if not instruction_df.empty else request_df
    packing_df, yongma_df, sample_available_df = normalize_packing_workbook(packing_file)
    wip_df = normalize_wip(wip_file)
    inventory_df = normalize_inventory(inventory_file)
    daily_inventory_df = normalize_daily_inventory_file(daily_inventory_file)
    daily_inventory_df = enrich_daily_inventory_from_wms(daily_inventory_df, inventory_df)
    product_summary, _unmatched_packing_total, code_summary = build_summaries(
        progress_basis_df,
        packing_df,
        yongma_df,
        product_master_df,
    )
    code_summary = attach_inventory_to_code_summary(code_summary, inventory_df)
    progress_df, _progress_info = normalize_progress(progress_file, progress_basis_df)
    production_progress_df = filter_progress_for_production_month(progress_df)
    code_summary = attach_progress_to_code_summary(code_summary, production_progress_df)
    product_summary = enrich_product_summary_from_code_summary(product_summary, code_summary)
    code_summary = attach_sample_available_to_code_summary(code_summary, sample_available_df)
    code_summary = with_operational_columns(code_summary)
    product_summary = attach_inventory_to_product_summary(product_summary, code_summary)
    return product_summary, code_summary, packing_df, yongma_df, daily_inventory_df, sample_available_df, instruction_df, request_df, wip_df


def get_sidebar_tab_from_query() -> str | None:
    try:
        value = st.query_params.get("tab")
    except Exception:
        return None
    if isinstance(value, list):
        value = value[0] if value else None
    key = str(value or "")
    return SIDEBAR_NAV_KEY_TO_TAB.get(key)


def sidebar_nav_item_html(tab: str, nav_key: str, active_tab: str) -> str:
    active_class = " active" if tab == active_tab else ""
    return (
        f'<a class="sidebar-nav-item{active_class}" href="?tab={escape(nav_key)}" target="_self">'
        f'<span>{escape(tab)}</span>'
        "</a>"
    )


def render_dashboard_nav() -> str:
    with st.sidebar:
        st.markdown(
            """
            <div class="sidebar-brand">
              <div class="sidebar-logo-dot"></div>
              <div>
                <div class="sidebar-brand-title">INTEROJO</div>
                <div class="sidebar-brand-sub">Domestic Dashboard</div>
              </div>
            </div>
            """,
            unsafe_allow_html=True,
        )
        selected = get_sidebar_tab_from_query() or str(
            st.session_state.get("dashboard_active_tab_sidebar", DASHBOARD_TABS[0])
        )
        if selected not in DASHBOARD_TABS:
            selected = DASHBOARD_TABS[0]
        st.session_state["dashboard_active_tab_sidebar"] = selected
        nav_html = "".join(
            sidebar_nav_item_html(tab, nav_key, selected)
            for tab, nav_key in SIDEBAR_NAV_ITEMS
        )
        st.markdown(f"<nav class='sidebar-nav'>{nav_html}</nav>", unsafe_allow_html=True)
    return str(selected or DASHBOARD_TABS[0])


def render_period_group_filter(active_tab: str) -> str:
    filter_key_by_tab = {
        "제품 진도 현황": "product_summary",
        "일일 재고 대응": "daily_inventory",
        "생산코드 상세": "production_code",
        "판매코드 상세": "sales_code",
    }
    if active_tab == "제품 진도 현황":
        key = f"{filter_key_by_tab[active_tab]}_period_group_filter"
        selected = st.session_state.get(key, "전체")
        return str(selected if selected in PERIOD_GROUP_ORDER else "전체")
    filter_col, _ = st.columns([1.6, 4.4], gap="small")
    with filter_col:
        selected = st.segmented_control(
            "기간구분",
            options=PERIOD_GROUP_ORDER,
            default="전체",
            key=f"{filter_key_by_tab.get(active_tab, 'common')}_period_group_filter",
        )
    return str(selected or "전체")


def main() -> None:
    render_style()
    today_label = pd.Timestamp.now(tz="Asia/Seoul").strftime("%Y-%m-%d")
    header_left, header_right = st.columns([5.2, 1.2], gap="large", vertical_alignment="center")
    with header_left:
        st.markdown(
            "<div class='app-header'>"
            "<div class='app-title'>국내 생산·포장 현황</div>"
            f"<div class='app-basis'>기준일 {today_label} · 생산지시 기준 {REQUEST_DUE_MONTH_LABEL} 생산완료예상일 · "
            f"용마입고 기준 {PACKING_RECEIPT_BASE_DATE_LABEL}부터 · 지시수준 3Q전체물량 대비 생산지시물량</div>"
            "</div>",
            unsafe_allow_html=True,
        )
    with header_right:
        st.markdown(f"<div class='header-date'>{today_label}</div>", unsafe_allow_html=True)
        if st.button("새로고침", key="refresh_dashboard_header", width="stretch"):
            st.cache_data.clear()
            st.rerun()
    active_tab = render_dashboard_nav()

    base_dir = Path.cwd()
    try:
        files = discover_source_files(base_dir)
        (
            product_summary,
            code_summary,
            packing_df,
            yongma_df,
            daily_inventory_df,
            sample_available_df,
            instruction_df,
            request_df,
            wip_df,
        ) = load_dashboard_data(
            file_fingerprint(files.request_file),
            file_fingerprint(files.packing_file),
            file_fingerprint(files.progress_file),
            file_fingerprint(files.inventory_file),
            file_fingerprint(files.daily_inventory_file),
            file_fingerprint(files.product_master_file),
            file_fingerprint(files.wip_file),
            dashboard_cache_fingerprint(),
        )
        lot_status_df = pd.DataFrame()
    except DashboardConfigError as exc:
        st.error("데이터 설정 오류")
        for msg in exc.messages:
            st.write(f"- {msg}")
        st.stop()
    except Exception as exc:
        st.error(f"처리 중 오류가 발생했습니다: {exc}")
        st.stop()

    selected_period = render_period_group_filter(active_tab)

    if active_tab == "제품 진도 현황":
        render_product_summary_tab(
            product_summary,
            code_summary,
            request_df,
            instruction_df,
            daily_inventory_df,
            sample_available_df,
            selected_period,
        )
    elif active_tab == "일일 재고 대응":
        render_daily_inventory_tab(daily_inventory_df, code_summary, sample_available_df, lot_status_df, selected_period)
    elif active_tab == "생산코드 상세":
        render_production_code_tab(code_summary, wip_df, selected_period)
    elif active_tab == "판매코드 상세":
        render_sales_code_tab(code_summary, selected_period)


if __name__ == "__main__":
    main()
